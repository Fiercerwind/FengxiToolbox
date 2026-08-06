import ast
import concurrent.futures
import ctypes
import dis
import inspect
import io
import json
import marshal
import importlib
import importlib.util
import os
import re
import shutil
import subprocess
import sys
import tempfile
import threading
import types
import time
import platform
import customtkinter
import pypdf
import pywinstyles
import pythoncom
import tkinter
import tkinter.colorchooser
import tkinter.filedialog
import tkinter.font
import tkinter.messagebox
import tkinter.ttk
import windnd
import win32com.client
import win32gui
from win32com.shell import shell, shellcon
from PIL import Image as PILImage, ImageDraw, ImageFont
from pathlib import Path
from reportlab.lib import pagesizes
from reportlab.pdfbase import ttfonts
from reportlab.pdfgen import canvas
from tools.fx_audio_task import (
    AUDIO_VALID_AUDIO_EXTS,
    AUDIO_VALID_VIDEO_EXTS,
    AudioTaskCallbacks,
    build_audio_output_path as _audio_task_build_output_path,
    collect_audio_files as _audio_task_collect_files,
    get_audio_task_args as _audio_task_get_args,
    get_audio_transcribe_args as _audio_task_get_transcribe_args,
    process_one_audio_file as _audio_task_process_one_file,
    run_audio_task_core,
)
from tools.fx_speech_to_text import (
    SPEECH_LANGUAGE_OPTIONS,
    SPEECH_MODEL_OPTIONS,
    SPEECH_OUTPUT_FORMATS,
    format_srt_timestamp as _speech_format_timestamp,
    transcribe_media_file as _speech_transcribe_media_file,
)
from tools.fx_convert_core import (
    CONVERT_IMAGE_EXTS,
    CONVERT_MODE_SPECS,
    collect_convert_files as _convert_core_collect_files,
    describe_convert_mode as _convert_core_describe_mode,
    normalize_convert_mode as _convert_core_normalize_mode,
    plan_convert_output_path as _convert_core_plan_output_path,
)
from tools.fx_convert_task import (
    ConvertFileContext,
    ConvertImgsToPdfCallbacks,
    convert_md_to_pdf_file,
    convert_pdf_to_md_file,
    convert_txt_to_word_file,
    process_convert_file,
    run_convert_imgs_to_pdf_task_core,
)
from tools.fx_image_pdf_task import (
    IMAGE_TO_PDF_EXTS,
    ImagePdfTaskCallbacks,
    ImagePdfTaskOptions,
    build_image_pdf_output_path as _image_pdf_task_build_output_path,
    collect_image_to_pdf_files as _image_pdf_task_collect_files,
    image_file_to_pdf as _image_pdf_task_image_file_to_pdf,
    run_image_pdf_task_core,
)
from tools.fx_performance import FxPerformanceRecorder, load_performance_entries
from tools.fx_pdf_compress_core import (
    PDF_COMPRESS_LEVELS,
    PDF_IMAGE_COMPRESS_LEVELS,
    build_pdf_compress_profile_stamp as _pdf_compress_core_build_profile_stamp,
    build_pdf_compress_output_path as _pdf_compress_core_build_output_path,
    compress_pdf_file as _pdf_compress_core_compress_pdf_file,
    pdf_compress_meta_matches as _pdf_compress_core_meta_matches,
    write_pdf_compress_meta as _pdf_compress_core_write_meta,
)
from tools.fx_queue_history import (
    QueueHistoryContext,
    build_queue_history_search_blob,
    filter_queue_history_entries,
    load_queue_history,
    normalize_queue_history_entry,
    prune_queue_history_entries,
    queue_history_entry_timestamp,
    queue_status_text,
    save_queue_history,
)
from tools.fx_resume import (
    format_resume_skip_message,
    is_nonempty_file,
)
from tools.fx_watermark_checkpoint import (
    build_checkpoint_identity as _watermark_checkpoint_identity,
    build_checkpoint_path as _watermark_checkpoint_path,
    checkpoint_item_reusable as _watermark_checkpoint_item_reusable,
    clear_checkpoint as _watermark_checkpoint_clear,
    load_checkpoint as _watermark_checkpoint_load,
    mark_item as _watermark_checkpoint_mark_item,
    mark_state as _watermark_checkpoint_mark_state,
    prepare_checkpoint as _watermark_checkpoint_prepare,
    split_resumable_files as _watermark_checkpoint_split_files,
)
from tools.fx_runtime_patches import wrap_callable
from tools.fx_startup_patches import StartupPatchContext, install_startup_performance_patch
from tools.fx_file_manager_core import (
    apply_rename_to_file as _file_core_apply_rename_to_file,
    deduplicate_files as _file_core_deduplicate_files,
    normalize_file_rename_spec as _file_core_normalize_rename_spec,
    plan_renamed_output_path as _file_core_plan_renamed_output_path,
    rename_file_name as _file_core_rename_file_name,
    run_file_dedup_task as _file_core_run_file_dedup_task,
)
from tools.fx_file_manager_task import run_file_dedup_task_core as _run_file_dedup_task_core
from tools.fx_meta_core import (
    build_meta_output_path as _meta_core_build_output_path,
    modify_file_timestamp as _meta_core_modify_file_timestamp,
    modify_office_meta as _meta_core_modify_office_meta,
    modify_pdf_author as _meta_core_modify_pdf_author,
    process_meta_file as _meta_core_process_meta_file,
)
from tools.fx_task_history_exports import (
    TaskHistoryExportContext,
    build_diagnostic_summary,
    build_recent_history_diagnostic_snapshot,
    build_task_history_diagnostic_filename,
    build_task_history_export_filename,
    build_task_history_log_export_filename,
    build_task_history_log_export_text,
    build_task_history_report_export_filename,
    build_task_history_report_text,
    diagnostic_path_replacements,
    diagnostic_write_json,
    diagnostic_write_text,
    export_task_history_diagnostic_package,
    export_task_history_entry,
    export_task_history_log,
    export_task_history_report,
    redact_diagnostic_payload,
    redact_diagnostic_text,
)
from tools.fx_user_prefs import (
    UserPrefsContext,
    get_saved_output_strategy as _prefs_get_saved_output_strategy,
    get_saved_remove_wm_mode as _prefs_get_saved_remove_wm_mode,
    get_saved_watermark_filename_rule_settings as _prefs_get_saved_watermark_filename_rule_settings,
    get_saved_watermark_text as _prefs_get_saved_watermark_text,
    get_active_last_settings_category as _prefs_get_active_last_settings_category,
    delete_preset_entry as _prefs_delete_preset_entry,
    find_preset_entry as _prefs_find_preset_entry,
    load_presets as _prefs_load_presets,
    load_last_settings as _prefs_load_last_settings,
    load_user_prefs as _prefs_load_user_prefs,
    make_preset_id as _prefs_make_preset_id,
    save_preset_entry as _prefs_save_preset_entry,
    save_presets as _prefs_save_presets,
    save_last_settings_entry as _prefs_save_last_settings_entry,
    save_output_strategy as _prefs_save_output_strategy,
    save_remove_wm_mode as _prefs_save_remove_wm_mode,
    save_user_prefs as _prefs_save_user_prefs,
    save_watermark_filename_rule_settings as _prefs_save_watermark_filename_rule_settings,
    save_watermark_text as _prefs_save_watermark_text,
)
from tools.fx_watermark_core import (
    add_watermark_to_pdf as _watermark_core_add_watermark_to_pdf,
    add_watermark_to_word as _watermark_core_add_watermark_to_word,
    create_watermark_packet as _watermark_core_create_watermark_packet,
    normalize_watermark_page_range as _watermark_core_normalize_page_range,
    open_word_document_safely as _watermark_core_open_word_document_safely,
    watermark_color_to_hex as _watermark_core_color_to_hex,
)
from tools.fx_zip_core import (
    estimate_zip_progress_units,
    normalize_zip_archive_policy,
    normalize_zip_depth_range,
    normalize_zip_max_depth,
    normalize_zip_mode,
    plan_zip_archives,
    run_zip_task as _zip_core_run_zip_task,
)


_FX_ORIGINAL_WIN32COM_DISPATCH = win32com.client.Dispatch
_FX_ORIGINAL_WIN32COM_DISPATCH_EX = win32com.client.DispatchEx

BOOTSTRAP_STARTED_AT = time.perf_counter()
RUNTIME_BIN = Path(__file__).with_name("fengxi_runtime.bin")
DEBUG_LOG = Path(tempfile.gettempdir()) / "fx_toolbox_loader.log"
DEFAULT_STARTUP_TAB = "watermark"
STARTUP_LAYOUT_FIRST_DELAY_MS = 450
STARTUP_LAYOUT_STAGE_DELAY_MS = 90
STARTUP_DEFAULT_TAB_INIT_DELAY_MS = 30
LAZY_TAB_SPECS = {
    "watermark": {"init": "init_watermark_ui"},
    "remove_wm": {"init": "init_remove_wm_ui"},
    "convert": {"init": "init_convert_ui"},
    "audio": {"init": "init_audio_ui"},
    "zip": {"init": "init_zip_ui"},
    "pdf": {"init": "init_pdf_ui"},
    "image": {"init": "init_img_ui"},
    "meta": {"init": "init_meta_ui"},
    "file": {"init": "init_file_ui"},
}
TAB_LAYOUT_ATTRS = {
    "watermark": "tab_wm",
    "remove_wm": "tab_rm_wm",
    "convert": "tab_cv",
    "audio": "tab_audio",
    "zip": "tab_zip",
    "pdf": "tab_pdf",
    "image": "tab_img",
    "meta": "tab_meta",
    "file": "tab_file",
}
LAZY_ATTR_PREFIXES = (
    ("pdf_", "pdf"),
    ("file_", "file"),
    ("zip_", "zip"),
    ("cv_", "convert"),
    ("rm_wm_", "remove_wm"),
    ("audio_", "audio"),
    ("img_", "image"),
    ("meta_", "meta"),
    ("wm_", "watermark"),
)
SIDEBAR_BUTTON_SPECS = {
    "btn_nav_wm": {"label": "批量水印", "icon": "shield"},
    "btn_nav_rm_wm": {"label": "去除水印", "icon": "eraser"},
    "btn_nav_cv": {"label": "格式转换", "icon": "swap"},
    "btn_nav_audio": {"label": "音频工具", "icon": "music"},
    "btn_nav_zip": {"label": "批量压缩", "icon": "box"},
    "btn_nav_pdf": {"label": "PDF 工具", "icon": "document"},
    "btn_nav_img": {"label": "图片工厂", "icon": "image"},
    "btn_nav_meta": {"label": "属性隐私", "icon": "lock"},
    "btn_nav_file": {"label": "文件管家", "icon": "folder"},
}
SIDEBAR_AUX_BUTTON_SPECS = {
    "btn_help_proxy": {"label": "使用教程", "icon": "book"},
    "btn_donate": {"label": "赞助作者", "icon": "coffee"},
}
SIDEBAR_ICON_NAV = (213, 228, 245, 255)
SIDEBAR_ICON_HELP = (232, 222, 208, 255)
SIDEBAR_ICON_DONATE = (239, 206, 123, 255)
CONTENT_ICON_PRIMARY = (173, 200, 228, 255)
CONTENT_ICON_SECONDARY = (194, 214, 236, 255)
SIDEBAR_AUX_STYLES = {
    "help": {
        "fg_color": "#161B21",
        "hover_color": "#222A34",
        "border_color": "#566274",
        "text_color": "#E8EDF5",
    },
    "donate": {
        "fg_color": "#241A10",
        "hover_color": "#322517",
        "border_color": "#B89352",
        "text_color": "#F6E2B2",
    },
}
FAST_SIDEBAR_BUILD_FONT = ("Microsoft YaHei UI", 12)
APP_ICON_PNG = "fengxi_app_icon.png"
APP_ICON_ICO = "fengxi_app_icon.ico"
DONATE_QR_PNG = "donate_qr.png"
APP_RELEASE_VERSION = "4.0.0"
APP_DISPLAY_VERSION = "4.0"
ZIP_MODE_LABEL_TEXTS = {
    "仅压缩总文件 (Total Zip)": "仅压缩总文件 (Total Zip)",
    "全层级递归压缩 (Total + Recursive)": "全层级递归压缩 (Total + Recursive)",
    "智能混合模式 (Smart Recursive) [推荐]": "智能混合模式 (Smart Recursive) [推荐]",
}
ZIP_IMPLEMENTATION_HELP_LINES = (
    "1. 三种模式：总文件模式只生成根目录总包；全层级递归为选定层级内每个文件夹生成包；智能混合在递归规划基础上按内容边界停止下钻。",
    "2. 智能扫描顺序：从输入根目录开始，根目录计为第 1 层，子文件夹依次为第 2、3 层；目录和文件均按名称忽略大小写排序，结果顺序稳定。",
    "3. 智能建包规则：每个实际访问到且处于层级范围内的目录都会计划一个 ZIP；叶子目录和空目录也会建包，空目录结构会保留。",
    "4. 混合边界：目录同时包含普通文件和子文件夹时，当前目录是边界；若当前层在范围内则压当前层，之后不再为其子文件夹单独建包。",
    "5. 起始层之前的边界：即使混合边界早于所选起始层，也会立即停止该分支；因为当前层不在范围内，所以当前层和更深子目录都不会建包。",
    "6. 继续下钻：目录只有子文件夹，或只有已有压缩包文件加子文件夹时，会在符合层级范围时压当前层，然后继续扫描子文件夹。",
    "7. 已有压缩包不算普通文件：zip、rar、7z、tar、gz、bz2、xz、zst 不会触发混合边界；它们仍可作为普通内容进入父包，但本次计划生成的 ZIP 会自动排除，避免压缩包互相套入。",
    "8. .DS_Store 例外：匹配时忽略大小写；它不参与智能边界判断、不写入任何生成的 ZIP、单独选中时直接跳过，但不会从源目录删除或修改。",
    "9. 层级范围：两个输入框都留空表示不限；只填结束 4 表示第 1-4 层；起始 2、结束 4 表示只建第 2-4 层；达到结束层后不再下钻。",
    "10. 输出位置：根目录递归/智能包写在根目录内并命名为“根目录名.zip”；子目录包写在该目录的上一级并命名为“子目录名.zip”。",
    "11. 单文件与总包：单文件输出“原文件名_Backup.zip”；总文件模式输出“根目录名_Backup.zip”；目录包会包含正常文件与空目录，但排除 .DS_Store 和本次任务的计划输出包。",
    "12. 已有同名包：复用模式仅跳过结构有效的 ZIP，适合断点续跑；旧包损坏时会重建；重建模式会先删除同名旧输出，再按当前规则重新压缩。",
    "13. 压缩与执行：使用 ZIP Deflate、压缩级别 6；批量压缩按计划单线程逐包执行，进度以计划包数量计算，每个包都会记录源目录和输出位置。",
    "14. 停止与失败：停止请求在两个压缩包之间响应，不会中断正在写入的单个大包；某个包失败后会记录失败项并继续后续任务，最终状态会标记失败。",
    "15. 数据安全：压缩不会删除普通源文件或源目录；只有选择重建策略时，才会删除本次计划输出位置上的同名旧 ZIP。",
)
ZIP_MODE_DESCRIPTION_TEXT = "功能说明：新版智能压缩完整实现规则\n" + "\n".join(ZIP_IMPLEMENTATION_HELP_LINES)
ZIP_ARCHIVE_POLICY_REUSE_LABEL = "复用已有压缩包（断点续跑）"
ZIP_ARCHIVE_POLICY_REBUILD_LABEL = "删除旧包并重新压缩"
ZIP_ARCHIVE_POLICY_LABELS = (ZIP_ARCHIVE_POLICY_REUSE_LABEL, ZIP_ARCHIVE_POLICY_REBUILD_LABEL)
ZIP_ARCHIVE_POLICY_LABEL_TO_VALUE = {
    ZIP_ARCHIVE_POLICY_REUSE_LABEL: "reuse_existing",
    ZIP_ARCHIVE_POLICY_REBUILD_LABEL: "rebuild_existing",
}
ZIP_ARCHIVE_POLICY_VALUE_TO_LABEL = {
    "reuse_existing": ZIP_ARCHIVE_POLICY_REUSE_LABEL,
    "rebuild_existing": ZIP_ARCHIVE_POLICY_REBUILD_LABEL,
}
HELP_TAB_TITLE = "使用教程"
INLINE_HELP_SECTIONS = (
    (
        "三步上手",
        (
            "1. 在顶部输入框拖入文件/文件夹，或点击选择入口。拖拽单个文件时会保留真实文件路径，不会降级成父文件夹。",
            "2. 在左侧选择功能，再在右侧页面配置参数。PDF 和图片模块会先显示功能入口，点进具体功能后再显示细项。",
            "3. 点击底部开始执行。正式执行前会弹出任务预览，确认本次将处理的文件数、跳过数量、输出策略和高风险选项。",
            "4. 运行中看进度条旁的文字：当前文件、当前阶段、总进度和预计剩余时间会尽量同步显示。",
            "5. 常用设置会自动记住上一次的选择，不需要单独打开预设中心。",
        ),
    ),
    (
        "输出与安全确认",
        (
            "默认策略：文件夹任务通常在原目录生成【处理完成】结果文件夹；单文件任务通常在同目录生成新文件。",
            "覆盖原文件：只有在页面明确开启覆盖时才会尝试替换原文件，适合确认无误后的重复处理。",
            "原目录新文件：更适合试跑和排查问题，可以保留原始资料。",
            "删除源文件、覆盖原文件、文件去重都属于高风险操作，开始前预览会再次提醒。",
            "如果不确定效果，先用少量副本试跑，再把同样设置加入队列批量执行。",
        ),
    ),
    (
        "任务队列与历史记录",
        (
            "配置好当前页面后，可点击加入队列，把输入路径、功能类型和参数快照保存为待执行任务。",
            "队列历史支持按功能、状态、失败原因和关键词筛选，方便回看大批量处理结果。",
            "失败任务可以按原参数重试；成功任务也可以回放加入队列。",
            "历史详情里可以打开输出位置、导出结构化结果、导出日志或导出 Markdown 报告。",
            "过久的历史记录会自动清理，避免队列文件越来越大。",
        ),
    ),
    (
        "PDF：OCR 搜索版 PDF",
        (
            "适用场景：扫描件、图片型 PDF、无法搜索文字的资料。",
            "推荐设置：后端选择 auto，图像增强选择 auto；低质量扫描件可尝试 scan，清晰原件可改为 off。",
            "处理方式：保留原页面画面，并叠加透明文字层，生成可搜索的新 PDF。",
            "修正文本方向（OCR）：开启后由 OCR 后端判断并校正文字方向；关闭后不执行文字方向分类。",
            "质量回退：auto 会优先尝试最快路径，识别质量偏低时会继续尝试备用后端。",
            "对比报告：开启后会记录后端、图像增强候选和质量评分，便于判断哪条 OCR 路线效果更好。",
            "常见失败原因：OCR 依赖不可用、PDF 加密、页面图片损坏、路径权限不足或文档过大。",
        ),
    ),
    (
        "PDF：压缩、合并、拆分、加密",
        (
            "仅修正文本方向（不进行OCR）：这是独立于 OCR 的页面兼容功能，不需要加载 OCR 模型，也不会重新识别文字。",
            "仅修正文本方向（不进行OCR）会把页面的 90/180/270 度显示旋转写入页面内容并归零旋转值，适合修复 Noteful 等软件中的画面与文字层方向错位。",
            "处理结果会先复制到结果文件夹再修正，默认不改原 PDF；加密 PDF 请在左侧填写密码。",
            "PDF 压缩：可分别选择 PDF 压缩程度和图片压缩程度，适合分享、归档和上传限制场景。",
            "PDF 合并：按页面当前规则收集 PDF 后合成为一个文件，建议先检查文件名排序。",
            "PDF 拆分：适合把多页资料拆成独立文件。",
            "PDF 加密：请牢记密码，生成后如果忘记密码，工具箱不会替你找回。",
            "多文件压缩在开启批量并行时可提速；合并、OCR 和角度修正为了稳定性默认按单线程执行。",
        ),
    ),
    (
        "批量水印",
        (
            "支持 PDF、Word、PPT 等文档批量添加文字或图片水印。",
            "智能加水印支持文件名跳过规则：可设置按开头或结尾匹配指定字符，默认兼容跳过文件名去扩展名后以“-”结尾的文件。",
            "复制干扰层作用于 PDF、直接输出的 Word，以及先转成 PDF 的 Word/PPT。PDF 使用不可见文本；Word 使用 1pt 白色文本段落，在常规白底文档中不显色，深色/彩色页面可能看见白字。",
            "干扰只放在段落/文本块之间：局部拖选一行正文或单独复制一段正文通常不受影响；全选整页、复制整份文件或直接提取文字时会混入干扰块。",
            "干扰内容混合数字、字母、符号、乱码样式字符；Word 还会混入汉字片段。复制干扰层不能阻止截图后重新 OCR；轻度、标准、强力只控制每页干扰块数量，默认关闭。",
            "水印内容、字号、透明度、角度、跳过规则和输出策略都会自动记住上一次设置。",
            "这是稳定区功能，除非明确要求，不应改动核心加水印处理逻辑。",
        ),
    ),
    (
        "去除水印",
        (
            "适用场景：尝试清理 Word、PDF、PPT 中的明显水印对象。",
            "模式选择：默认保守（推荐），尽量避免误删正常文字和图片；标准适合常规水印；激进适合顽固水印但误删风险更高。",
            "单文件默认在同目录生成新文件；文件夹默认生成【处理完成】结果文件夹。",
            "如开启覆盖原文件，会先生成临时结果，成功后再替换原文件。",
            "去水印不是万能删除器，不同文件结构差异很大，重要文件建议先用副本试跑。",
        ),
    ),
    (
        "图片工厂",
        (
            "支持图片批量处理、图片转 PDF、以及多图合并成一个 PDF。",
            "图片转 PDF 会为每张图片生成独立 PDF；多图合并 PDF 会按文件名排序合成为一份 PDF。",
            "支持 jpg、jpeg、png、bmp、webp、tif、tiff 等常见图片格式。",
            "多张图片逐张转 PDF 可在批量并行开启时提速；多图合并为单个 PDF 会保持稳定单线程。",
            "如果图片方向、尺寸差异较大，建议先整理文件名和顺序，再合并成 PDF。",
        ),
    ),
    (
        "格式转换与音频工具",
        (
            "格式转换支持常用文档格式互转，包括 Word、PPT、PDF 等转换入口。",
            "Office/PDF 转换通常依赖本机 Office 或 WPS 环境；失败时优先检查软件安装、加密文档和文件占用。",
            "音频工具支持从视频中提取音频，以及常见音频格式转换。",
            "音视频转换依赖 ffmpeg 路线，大体积文件耗时较长，处理期间尽量不要频繁强制关闭。",
        ),
    ),
    (
        "属性隐私与文件管家",
        (
            "属性隐私支持修改 PDF/Office 作者信息，以及批量修改文件时间属性。",
            "Office 元数据依赖本机 Office COM 环境；PDF 作者信息通过 PDF 元数据写入。",
            "文件管家支持批量重命名和重复文件清理，重命名规则会自动记住上一次设置。",
            "去重基于文件内容哈希判断，删除类操作建议先用测试文件夹确认规则。",
        ),
    ),
    (
        "批量压缩",
        ZIP_IMPLEMENTATION_HELP_LINES,
    ),
    (
        "性能、进度与排障",
        (
            "批量并行只对部分多文件工作流生效，例如部分 PDF 压缩、图片逐张转 PDF、音视频逐文件转换等。",
            "稳定单线程适合 OCR、去水印、Office/PDF 转换、PDF 合并、文件去重等容易受外部依赖影响的任务。",
            "运行信息框会保留关键日志；任务历史详情可以导出日志和报告，用于复盘失败原因。",
            "如果遇到卡住或失败，先查看当前阶段、失败文件名、历史详情里的错误原因，再决定是否重试。",
        ),
    ),
    (
        "重要约束",
        (
            "批量压缩和批量水印是稳定区，非必要不改业务代码。",
            "任何重要修改前需要先做可恢复备份。",
            "不得删除项目目录外的任何文件。",
            "每次工作应先读取记忆文件，再按项目和类别渐进加载上下文。",
        ),
    ),
)
DONATE_TAB_TITLE = "赞助作者"
DONATE_SUPPORT_SENTENCE = "如果风兮工具箱帮你节省了时间，欢迎随缘赞助一杯咖啡，支持我继续把它做得更稳、更好用。"
RESULT_FOLDER_NAME = "【处理完成】结果文件夹"
OUTPUT_STRATEGY_DEFAULT = "result_folder"
OUTPUT_STRATEGY_LABELS = {
    "result_folder": "原目录新文件 / 【处理完成】结果文件夹",
    "same_dir": "同目录新文件",
    "overwrite": "覆盖原文件",
}
OUTPUT_STRATEGY_VALUES = tuple(OUTPUT_STRATEGY_LABELS.keys())
OUTPUT_STRATEGY_VALUE_TO_LABEL = dict(OUTPUT_STRATEGY_LABELS)
OUTPUT_STRATEGY_LABEL_TO_VALUE = {label: value for value, label in OUTPUT_STRATEGY_LABELS.items()}
REMOVE_WM_MODE_DEFAULT = "conservative"
REMOVE_WM_MODE_LABELS = {
    "conservative": "保守（推荐）",
    "standard": "标准",
    "aggressive": "激进",
}
REMOVE_WM_MODE_VALUES = tuple(REMOVE_WM_MODE_LABELS.keys())
REMOVE_WM_MODE_VALUE_TO_LABEL = dict(REMOVE_WM_MODE_LABELS)
REMOVE_WM_MODE_LABEL_TO_VALUE = {label: value for value, label in REMOVE_WM_MODE_LABELS.items()}
REMOVE_WM_MODE_HINTS = {
    "conservative": "默认更安全：只处理更像水印的大尺寸斜向/半透明对象，尽量避免误删正常文字和图片。",
    "standard": "平衡模式：沿用之前的识别阈值，适合常规文档水印。",
    "aggressive": "强力模式：会扩大识别范围，适合顽固水印，但更可能误删正常元素。",
}
REMOVE_WM_MODE_PROFILES = {
    "conservative": {
        "shape_width_ratio": 0.48,
        "shape_height_ratio": 0.28,
        "shape_center_ratio": 0.24,
        "shape_rotation_distance": 18.0,
        "shape_transparency": 0.25,
        "inline_width_ratio": 0.60,
        "inline_height_ratio": 0.16,
        "allow_plain_very_large": False,
        "very_large_width_ratio": 0.72,
        "very_large_height_ratio": 0.42,
    },
    "standard": {
        "shape_width_ratio": 0.35,
        "shape_height_ratio": 0.20,
        "shape_center_ratio": 0.30,
        "shape_rotation_distance": 12.0,
        "shape_transparency": 0.15,
        "inline_width_ratio": 0.45,
        "inline_height_ratio": 0.16,
        "allow_plain_very_large": False,
        "very_large_width_ratio": 0.68,
        "very_large_height_ratio": 0.36,
    },
    "aggressive": {
        "shape_width_ratio": 0.25,
        "shape_height_ratio": 0.12,
        "shape_center_ratio": 0.40,
        "shape_rotation_distance": 8.0,
        "shape_transparency": 0.08,
        "inline_width_ratio": 0.35,
        "inline_height_ratio": 0.10,
        "allow_plain_very_large": True,
        "very_large_width_ratio": 0.55,
        "very_large_height_ratio": 0.28,
    },
}
_REMOVE_WM_MODE_CONTEXT = threading.local()
QUEUE_HISTORY_LIMIT = 80
QUEUE_HISTORY_RETENTION_DAYS = 90
PROGRESS_STATUS_IDLE_TEXT = "进度：等待任务"
PARALLEL_MAX_WORKERS = 4
WATERMARK_PDF_MAX_WORKERS = 5
WATERMARK_PDF_QUEUE_FACTOR = 2
WATERMARK_UI_LOG_FLUSH_INTERVAL_SECONDS = 0.16
WATERMARK_UI_LOG_BATCH_SIZE = 12
PARALLEL_SWITCH_TEXT = "批量并行（部分生效）"
PRESET_CATEGORY_LABELS = {
    "watermark": "批量水印",
    "ocr": "OCR 搜索版 PDF",
    "pdf_compress": "PDF 压缩",
    "audio": "音频工具",
    "rename": "命名规则",
    "zip": "批量压缩",
}
PRESET_CATEGORY_TO_TASK = {
    "watermark": "watermark",
    "ocr": "pdf",
    "pdf_compress": "pdf",
    "audio": "audio",
    "rename": "file",
    "zip": "zip",
}
PRESET_LABEL_TO_CATEGORY = {label: key for key, label in PRESET_CATEGORY_LABELS.items()}
QUEUE_ERROR_MARKERS = (
    "❌",
    "🔥",
    "失败",
    "错误",
    "异常",
    "严重错误",
    "error",
    "failed",
    "traceback",
)
FEATURE_REGISTRY = {
    "watermark": {
        "label": "批量水印",
        "icon": "shield",
        "page": "watermark",
        "input": {"file": True, "folder": True, "drag_drop": True},
        "output_strategy": {"supported": True, "force_result_folder": False},
        "parallel": {"mode": "safe", "hint": "批量水印可对多文件并行处理；遇到 Word/PDF 特殊链路时会自动保护。"},
        "preview_modes": {"default": "添加水印"},
        "risk_flags": ("delete_source",),
        "stable_core": True,
    },
    "remove_wm": {
        "label": "去除水印",
        "icon": "eraser",
        "page": "remove_wm",
        "input": {"file": True, "folder": True, "drag_drop": True},
        "output_strategy": {"supported": True, "force_result_folder": False},
        "parallel": {"mode": "forced_single", "detail": {"": "去水印会调用 Office/PDF 清理链路，保持单线程更安全。"}},
        "preview_modes": {
            "conservative": "保守（推荐）",
            "standard": "标准",
            "aggressive": "激进",
        },
        "risk_flags": ("overwrite_original",),
    },
    "convert": {
        "label": "格式转换",
        "icon": "swap",
        "page": "convert",
        "input": {"file": True, "folder": True, "drag_drop": True},
        "output_strategy": {"supported": False, "force_result_folder": False},
        "parallel": {"mode": "forced_single", "detail": {"": "Office/PDF 转换依赖 COM 或重型转换器，已强制单线程。"}},
        "preview_modes": {
            "word2pdf": "Word 转 PDF",
            "pdf2word": "PDF 转 Word",
            "ppt2pdf": "PPT 转 PDF",
            "pdf2ppt": "PDF 转 PPT",
            "txt2word": "TXT 转 Word",
            "md2pdf": "Markdown 转 PDF",
            "pdf2md": "PDF 转 Markdown",
            "imgs2pdf": "多图合并 ➔ PDF电子书",
        },
    },
    "audio": {
        "label": "音频工具",
        "icon": "music",
        "page": "audio",
        "input": {"file": True, "folder": True, "drag_drop": True},
        "output_strategy": {"supported": False, "force_result_folder": False},
        "parallel": {"mode": "safe", "hint": "音视频逐文件转换可并行处理；语音转文字也可逐文件并行，但会占用较多 CPU/内存。"},
        "preview_modes": {
            "video2mp3": "视频提取音频",
            "convert": "音频格式互转",
            "transcribe": "语音转文字",
        },
        "risk_flags": ("delete_source",),
    },
    "zip": {
        "label": "批量压缩",
        "icon": "box",
        "page": "zip",
        "input": {"file": True, "folder": True, "drag_drop": True},
        "output_strategy": {"supported": True, "force_result_folder": False},
        "parallel": {"mode": "forced_single", "detail": {"": "批量压缩保持原有稳定流程，不做额外并行。"}},
        "preview_modes": {
            "total": "总体压缩",
            "recursive": "递归压缩",
            "smart_recursive": "智能混合压缩",
        },
        "stable_core": True,
    },
    "pdf": {
        "label": "PDF 工具",
        "icon": "document",
        "page": "pdf",
        "input": {"file": True, "folder": True, "drag_drop": True},
        "output_strategy": {"supported": True, "force_result_folder": True},
        "parallel": {
            "mode": "safe",
            "hint": "PDF 拆分/加密/压缩可并行处理；合并、OCR 会自动切到稳定流程。",
            "detail": {
                "merge": ("forced_single", "PDF 合并需要保持文件顺序与输出一致，已强制单线程。"),
                "ocr": ("forced_single", "OCR 会占用大量 CPU/内存，当前采用单线程稳定处理。"),
                "rotation": ("forced_single", "PDF 页面角度修正按文件顺序稳定处理。"),
            },
        },
        "preview_modes": {
            "merge": "合并",
            "split": "拆分",
            "encrypt": "加密",
            "compress": "PDF 压缩",
            "ocr": "OCR 搜索版 PDF",
            "rotation": "仅修正文本方向（不进行OCR）",
        },
        "risk_flags": ("delete_source",),
    },
    "image": {
        "label": "图片工厂",
        "icon": "image",
        "page": "image",
        "input": {"file": True, "folder": True, "drag_drop": True},
        "output_strategy": {"supported": True, "force_result_folder": True},
        "parallel": {
            "mode": "safe",
            "hint": "图片格式转换/压缩/逐张转 PDF 可并行处理；多图合并 PDF 会自动切到稳定流程。",
            "detail": {
                "merge_pdf": ("forced_single", "多图合并 PDF 需要保持图片顺序，已强制单线程。"),
            },
        },
        "preview_modes": {
            "to_pdf": "图片转 PDF",
            "merge_pdf": "多图合并 PDF",
        },
        "risk_flags": ("delete_source",),
    },
    "meta": {
        "label": "属性隐私",
        "icon": "lock",
        "page": "meta",
        "input": {"file": True, "folder": True, "drag_drop": True},
        "output_strategy": {"supported": False, "force_result_folder": False},
        "parallel": {
            "mode": "safe",
            "hint": "普通文件时间修改可并行处理；Office 元数据会自动切到稳定流程。",
            "detail": {
                "office": ("forced_single", "Office 元数据修改依赖 COM，已强制单线程。"),
                "pdf": ("forced_single", "PDF 元数据写入保持单线程，避免同名输出冲突。"),
            },
        },
        "preview_modes": {},
    },
    "file": {
        "label": "文件管家",
        "icon": "folder",
        "page": "file",
        "input": {"file": True, "folder": True, "drag_drop": True},
        "output_strategy": {"supported": False, "force_result_folder": False},
        "parallel": {
            "mode": "safe",
            "hint": "文件重命名可并行处理；文件去重会自动切到稳定流程。",
            "detail": {
                "dedup": ("forced_single", "文件去重需要全局哈希比对，已强制单线程。"),
            },
        },
        "preview_modes": {
            "rename": "批量重命名",
            "dedup": "重复文件清理",
        },
        "risk_flags": ("dedup_delete",),
    },
}
QUEUE_TASK_LABELS = {task_type: spec.get("label", task_type) for task_type, spec in FEATURE_REGISTRY.items()}


def _get_feature_spec(task_type):
    return FEATURE_REGISTRY.get(str(task_type or ""), {})


def _get_feature_label(task_type, fallback="未知功能"):
    task_key = str(task_type or "")
    spec = _get_feature_spec(task_key)
    return spec.get("label") or task_key or fallback


def _get_feature_preview_mode_label(task_type, mode_value="", fallback=""):
    spec = _get_feature_spec(task_type)
    modes = spec.get("preview_modes") if isinstance(spec.get("preview_modes"), dict) else {}
    mode_key = str(mode_value or "")
    return modes.get(mode_key) or modes.get("default") or fallback or mode_key


def _feature_supports_output_strategy(task_type):
    output = _get_feature_spec(task_type).get("output_strategy") or {}
    return bool(output.get("supported"))


def _feature_forces_result_folder(task_type):
    output = _get_feature_spec(task_type).get("output_strategy") or {}
    return bool(output.get("force_result_folder"))


def _build_parallel_sets_from_registry():
    safe_tasks = set()
    forced_tasks = set()
    forced_details = {}
    supported_hints = {}
    for task_type, spec in FEATURE_REGISTRY.items():
        parallel = spec.get("parallel") if isinstance(spec.get("parallel"), dict) else {}
        mode = parallel.get("mode", "")
        if mode == "safe":
            safe_tasks.add(task_type)
            if parallel.get("hint"):
                supported_hints[task_type] = parallel.get("hint")
        elif mode == "forced_single":
            forced_tasks.add(task_type)

        details = parallel.get("detail") if isinstance(parallel.get("detail"), dict) else {}
        for detail_key, detail_value in details.items():
            detail_mode = mode
            detail_message = detail_value
            if isinstance(detail_value, (tuple, list)) and len(detail_value) >= 2:
                detail_mode, detail_message = detail_value[0], detail_value[1]
            if detail_mode == "forced_single":
                forced_details[(task_type, str(detail_key or ""))] = str(detail_message or "当前功能使用专用流程。")
    return safe_tasks, forced_tasks, forced_details, supported_hints


def _get_feature_registry_errors():
    errors = []
    for task_type, spec in FEATURE_REGISTRY.items():
        if not spec.get("label"):
            errors.append(f"{task_type}:missing_label")
        if not isinstance(spec.get("input"), dict):
            errors.append(f"{task_type}:missing_input")
        if not isinstance(spec.get("output_strategy"), dict):
            errors.append(f"{task_type}:missing_output_strategy")
        if not isinstance(spec.get("parallel"), dict):
            errors.append(f"{task_type}:missing_parallel")
        if not isinstance(spec.get("preview_modes"), dict):
            errors.append(f"{task_type}:missing_preview_modes")
    return errors


OUTPUT_STRATEGY_SUPPORTED_TASKS = {
    task_type for task_type in FEATURE_REGISTRY if _feature_supports_output_strategy(task_type)
}
OUTPUT_STRATEGY_FORCE_RESULT_FOLDER_TASKS = {
    task_type for task_type in FEATURE_REGISTRY if _feature_forces_result_folder(task_type)
}
(
    PARALLEL_SAFE_TASKS,
    PARALLEL_FORCED_SINGLE_TASKS,
    PARALLEL_FORCED_SINGLE_DETAILS,
    PARALLEL_SUPPORTED_HINTS,
) = _build_parallel_sets_from_registry()
QUEUE_STATUS_LABELS = {
    "queued": "等待",
    "running": "执行中",
    "success": "完成",
    "failed": "失败",
    "skipped": "跳过",
    "stopped": "已停止",
}
QUEUE_HISTORY_STATUS_DEFAULT = "全部状态"
QUEUE_HISTORY_TASK_DEFAULT = "全部功能"
QUEUE_HISTORY_FAILURE_DEFAULT = "全部失败"
QUEUE_HISTORY_STATUS_OPTIONS = (
    (QUEUE_HISTORY_STATUS_DEFAULT, ""),
    ("仅完成", "success"),
    ("仅失败", "failed"),
    ("仅跳过", "skipped"),
    ("仅停止", "stopped"),
)
QUEUE_HISTORY_FAILURE_OPTIONS = (
    (QUEUE_HISTORY_FAILURE_DEFAULT, ""),
    ("路径缺失", "path_missing"),
    ("权限问题", "permission"),
    ("超时", "timeout"),
    ("依赖问题", "dependency"),
    ("部分失败", "partial_failure"),
    ("日志失败", "log_failure"),
    ("普通失败", "generic_failure"),
    ("未知失败", "unknown"),
)
QUEUE_HISTORY_TASK_OPTIONS = ((QUEUE_HISTORY_TASK_DEFAULT, ""),) + tuple(
    (label, task_type) for task_type, label in QUEUE_TASK_LABELS.items()
)
QUEUE_HISTORY_STATUS_LABEL_TO_VALUE = {label: value for label, value in QUEUE_HISTORY_STATUS_OPTIONS}
QUEUE_HISTORY_FAILURE_LABEL_TO_VALUE = {label: value for label, value in QUEUE_HISTORY_FAILURE_OPTIONS}
QUEUE_HISTORY_FAILURE_VALUE_TO_LABEL = {value: label for label, value in QUEUE_HISTORY_FAILURE_OPTIONS if value}
QUEUE_HISTORY_TASK_LABEL_TO_VALUE = {label: value for label, value in QUEUE_HISTORY_TASK_OPTIONS}
INLINE_TITLE_ICON_SPECS = {
    "批量去水印": {"icon": "eraser", "size": 20, "color": CONTENT_ICON_PRIMARY},
    "文档格式互转": {"icon": "swap", "size": 20, "color": CONTENT_ICON_PRIMARY},
    "音频/视频处理": {"icon": "music", "size": 20, "color": CONTENT_ICON_PRIMARY},
    "批量压缩": {"icon": "box", "size": 20, "color": CONTENT_ICON_PRIMARY},
    "PDF 进阶工具箱": {"icon": "document", "size": 20, "color": CONTENT_ICON_PRIMARY},
    "图片批量工厂": {"icon": "image", "size": 20, "color": CONTENT_ICON_PRIMARY},
    "属性隐私设置": {"icon": "lock", "size": 20, "color": CONTENT_ICON_PRIMARY},
    "文件管家": {"icon": "folder", "size": 20, "color": CONTENT_ICON_PRIMARY},
    "水印内容": {"icon": "document", "size": 18, "color": CONTENT_ICON_SECONDARY},
    "参数配置": {"icon": "settings", "size": 18, "color": CONTENT_ICON_SECONDARY},
    "图片处理工具": {"icon": "image", "size": 18, "color": CONTENT_ICON_SECONDARY},
    "⚡": {"icon": "box", "size": 16, "color": CONTENT_ICON_SECONDARY, "display_text": ""},
}
BIF_NEWDIALOGSTYLE = 0x00000040
BIF_USENEWUI = shellcon.BIF_EDITBOX | BIF_NEWDIALOGSTYLE
BFFM_INITIALIZED = 1
BFFM_SETSELECTIONW = 1127
LAZY_RUNTIME_IMPORT_SPECS = {
    "pdf2docx": {
        "probe": "pdf2docx",
        "symbols": {
            "Converter": ("pdf2docx", "Converter"),
        },
    },
    "moviepy": {
        "probe": "moviepy",
        "symbols": {},
    },
    "moviepy.editor": {
        "probe": "moviepy",
        "fallback_modules": ["moviepy"],
        "symbols": {
            "AudioFileClip": ("moviepy", "AudioFileClip"),
            "VideoFileClip": ("moviepy", "VideoFileClip"),
        },
    },
}


def _debug(message):
    try:
        with DEBUG_LOG.open("a", encoding="utf-8") as fh:
            fh.write(f"{message}\n")
    except Exception:
        pass


_FX_PERFORMANCE_RECORDER = None


def _get_user_pref_root():
    local_app_data = (os.environ.get("LOCALAPPDATA") or "").strip()
    if local_app_data:
        return Path(local_app_data) / "FengxiToolbox"
    try:
        return Path.home() / ".fengxi_toolbox"
    except Exception:
        return Path(__file__).resolve().parent


def _get_performance_log_file():
    return _get_user_pref_root() / "performance.jsonl"


def _get_performance_recorder():
    global _FX_PERFORMANCE_RECORDER
    if _FX_PERFORMANCE_RECORDER is None:
        _FX_PERFORMANCE_RECORDER = FxPerformanceRecorder(
            _get_performance_log_file(),
            app_version=APP_RELEASE_VERSION,
        )
    return _FX_PERFORMANCE_RECORDER


def _record_performance(event, started_at=None, task_name="", details=None):
    try:
        return _get_performance_recorder().record(
            event,
            started_at=started_at,
            task_name=task_name,
            details=details,
        )
    except Exception as exc:
        _debug(f"performance:record_error:{event}:{exc}")
        return {}


def _load_recent_performance_entries(limit=60):
    try:
        entries = load_performance_entries(_get_performance_log_file())
    except Exception as exc:
        _debug(f"performance:load_error:{exc}")
        return []
    try:
        limit = int(limit)
    except Exception:
        limit = 60
    if limit <= 0:
        return entries
    return entries[-limit:]


def _wrap_callable(owner, name, label=None):
    return wrap_callable(owner, name, label=label, debug=_debug)


LOG_DISPLAY_ICON_REPLACEMENTS = {
    "🎯": "[目标]",
    "🔎": "[检索]",
    "🚀": "[启动]",
    "✅": "[完成]",
    "⚠️": "[提示]",
    "⚠": "[提示]",
    "❌": "[错误]",
    "ℹ️": "[说明]",
    "⏹️": "[停止]",
    "⏭️": "[跳过]",
    "🎉": "[完成]",
    "🛡️": "[保护]",
    "🗑️": "[删除]",
    "📌": "[队列]",
    "↪️": "[跳过]",
    "📝": "[文本]",
    "🖼️": "[图片]",
    "⏱️": "[时间]",
    "🕵️": "[隐私]",
}
LOG_MOJIBAKE_ICON_REPLACEMENTS = {
    "\u9983\u6537": "[检索]",
    "\u9983\u6546": "[提示]",
}
LOG_MOJIBAKE_MARKERS = (
    "\u93b5\ue0a3\u5f3f",
    "\u59dd\uff45\u6e6a",
    "\u95c8\u70b4\u6309",
    "\u93b5\u5f52\u567a",
)


def _repair_mojibake_log_text(text):
    """Repair legacy GBK/UTF-8 mojibake while leaving normal Chinese untouched."""

    icon_placeholders = {}
    for source, replacement in LOG_MOJIBAKE_ICON_REPLACEMENTS.items():
        placeholder = f"__FX_LOG_ICON_{len(icon_placeholders)}__"
        if source in text:
            text = text.replace(source, placeholder)
            icon_placeholders[placeholder] = replacement
    if not any(marker in text for marker in LOG_MOJIBAKE_MARKERS):
        return _restore_log_icon_placeholders(text, icon_placeholders)

    original_marker_count = sum(text.count(marker) for marker in LOG_MOJIBAKE_MARKERS)
    for encoding in ("gb18030", "gbk"):
        try:
            repaired = text.encode(encoding).decode("utf-8")
        except (UnicodeEncodeError, UnicodeDecodeError):
            continue
        repaired_marker_count = sum(repaired.count(marker) for marker in LOG_MOJIBAKE_MARKERS)
        if repaired_marker_count < original_marker_count and "\ufffd" not in repaired:
            return _restore_log_icon_placeholders(repaired, icon_placeholders)
    return _restore_log_icon_placeholders(text, icon_placeholders)


def _restore_log_icon_placeholders(text, placeholders):
    for placeholder, replacement in placeholders.items():
        text = text.replace(placeholder, replacement)
    return text


def _normalize_log_display_text(message):
    """Replace emoji-only log decorations with stable text labels."""

    text = str(message or "")
    for source, replacement in LOG_DISPLAY_ICON_REPLACEMENTS.items():
        text = text.replace(source, replacement)
    text = _repair_mojibake_log_text(text)
    return text.replace("\ufe0f", "")


def _get_log_display_widget(app):
    box = getattr(app, "log_box", None)
    if box is None:
        return None
    return box


def _read_log_scroll_state(widget):
    try:
        top, bottom = widget.yview()
        return float(top), float(bottom), float(bottom) >= 0.995
    except Exception:
        return 0.0, 1.0, True


def _sync_log_scroll_state(app):
    widget = _get_log_display_widget(app)
    if widget is None:
        return
    top, _bottom, at_tail = _read_log_scroll_state(widget)
    try:
        app._fx_log_follow_tail = at_tail
        if not at_tail:
            app._fx_log_scroll_anchor = top
    except Exception:
        pass


def _schedule_log_scroll_state_sync(app):
    try:
        app.after_idle(lambda target=app: _sync_log_scroll_state(target))
    except Exception:
        _sync_log_scroll_state(app)


def _handle_log_scroll_interaction(app):
    try:
        app._fx_log_scroll_generation = int(getattr(app, "_fx_log_scroll_generation", 0) or 0) + 1
        app._fx_log_restore_pending = False
    except Exception:
        pass
    _sync_log_scroll_state(app)
    _schedule_log_scroll_state_sync(app)


def _bind_log_scroll_tracking(app):
    if getattr(app, "_fx_log_scroll_tracking_ready", False):
        return
    widget = _get_log_display_widget(app)
    if widget is None:
        return
    tracked_widgets = [widget]
    inner_widget = getattr(widget, "_textbox", None)
    scrollbar = getattr(widget, "_y_scrollbar", None) or getattr(widget, "_scrollbar", None)
    scrollbar_canvas = getattr(scrollbar, "_canvas", None) if scrollbar is not None else None
    if inner_widget is not None:
        tracked_widgets.append(inner_widget)
    if scrollbar is not None:
        tracked_widgets.append(scrollbar)
    if scrollbar_canvas is not None:
        tracked_widgets.append(scrollbar_canvas)
    events = (
        "<MouseWheel>",
        "<Button-4>",
        "<Button-5>",
        "<ButtonPress-1>",
        "<ButtonRelease-1>",
        "<B1-Motion>",
        "<KeyRelease>",
    )
    for target in tracked_widgets:
        for event_name in events:
            try:
                target.bind(
                    event_name,
                    lambda _event=None, owner=app: _handle_log_scroll_interaction(owner),
                    add="+",
                )
            except Exception:
                pass
    try:
        top, _bottom, at_tail = _read_log_scroll_state(widget)
        app._fx_log_follow_tail = at_tail
        app._fx_log_scroll_anchor = top
        app._fx_log_scroll_generation = int(getattr(app, "_fx_log_scroll_generation", 0) or 0)
        app._fx_log_scroll_tracking_ready = True
    except Exception:
        pass


def _restore_log_scroll_after_append(app, anchor):
    widget = _get_log_display_widget(app)
    if widget is None or getattr(app, "_fx_log_follow_tail", True):
        return
    try:
        widget.yview_moveto(max(0.0, min(1.0, float(anchor))))
    except Exception:
        try:
            inner_widget = getattr(widget, "_textbox", None)
            if inner_widget is not None:
                inner_widget.yview_moveto(max(0.0, min(1.0, float(anchor))))
        except Exception:
            pass


def _patch_log_display_behavior():
    try:
        original_log = FengxiToolboxApp.log
    except Exception as exc:
        _debug(f"log_display:missing:{exc}")
        return
    if getattr(original_log, "__fx_log_display_patch__", False):
        return

    def patched_log(self, message, *args, **kwargs):
        widget = _get_log_display_widget(self)
        _bind_log_scroll_tracking(self)
        _top, _bottom, at_tail = _read_log_scroll_state(widget) if widget is not None else (0.0, 1.0, True)
        follow_tail = bool(getattr(self, "_fx_log_follow_tail", at_tail))
        anchor = float(getattr(self, "_fx_log_scroll_anchor", _top) or _top)
        generation = int(getattr(self, "_fx_log_scroll_generation", 0) or 0)
        if not follow_tail and not getattr(self, "_fx_log_restore_pending", False):
            anchor = _top
            self._fx_log_scroll_anchor = anchor
            self._fx_log_restore_pending = True
        result = original_log(self, _normalize_log_display_text(message), *args, **kwargs)
        if not follow_tail:
            def restore(target=self, saved_anchor=anchor, saved_generation=generation):
                try:
                    if int(getattr(target, "_fx_log_scroll_generation", 0) or 0) != saved_generation:
                        return
                    if getattr(target, "_fx_log_follow_tail", True):
                        return
                    _restore_log_scroll_after_append(target, saved_anchor)
                finally:
                    if int(getattr(target, "_fx_log_scroll_generation", 0) or 0) == saved_generation:
                        target._fx_log_restore_pending = False

            try:
                self.after_idle(restore)
            except Exception:
                restore()
        return result

    patched_log.__fx_log_display_patch__ = True
    patched_log.__wrapped__ = original_log
    FengxiToolboxApp.log = patched_log
    _debug("log_display:installed")


def _draw_sidebar_icon(draw, kind, color, size=22, stroke=2):
    c = color
    viewport = 24.0
    scale = max(0.5, float(size) / viewport)
    stroke_px = max(1, int(round(stroke * scale * 0.88)))
    thin_stroke_px = max(1, int(round(stroke_px * 0.72)))

    def pt(x, y):
        return (x * scale, y * scale)

    def path(points):
        return [pt(x, y) for x, y in points]

    def line(points, width=None):
        draw.line(path(points), fill=c, width=width or stroke_px, joint="curve")

    def polygon(points, closed=True, fill=None, outline=True, width=None):
        pts = path(points)
        if fill is not None:
            draw.polygon(pts, fill=fill)
        if outline:
            outline_path = pts + ([pts[0]] if closed else [])
            draw.line(outline_path, fill=c, width=width or stroke_px, joint="curve")

    def ellipse(x1, y1, x2, y2, width=None, fill=None):
        draw.ellipse(
            (x1 * scale, y1 * scale, x2 * scale, y2 * scale),
            outline=c if width else None,
            width=width or stroke_px,
            fill=fill,
        )

    def rounded_rect(x1, y1, x2, y2, radius, width=None):
        draw.rounded_rectangle(
            (x1 * scale, y1 * scale, x2 * scale, y2 * scale),
            radius=radius * scale,
            outline=c,
            width=width or stroke_px,
        )

    if kind == "shield":
        outer = [(12.0, 3.0), (17.1, 5.2), (16.6, 11.0), (12.0, 18.5), (7.4, 11.0), (6.9, 5.2)]
        inner = [(12.0, 5.0), (15.3, 6.4), (15.0, 10.0), (12.0, 15.0), (9.0, 10.0), (8.7, 6.4)]
        polygon(outer, fill=(243, 247, 252, 255), outline=False)
        polygon(inner, fill=(86, 136, 210, 255), outline=False)
    elif kind == "eraser":
        polygon([(7.1, 14.2), (11.7, 9.6), (16.9, 14.8), (12.3, 19.4)], width=stroke_px)
        polygon([(11.7, 9.6), (13.9, 7.4), (19.1, 12.6), (16.9, 14.8)], width=stroke_px)
        line([(9.3, 16.4), (14.5, 11.2)], width=thin_stroke_px)
        line([(12.0, 19.0), (17.2, 13.8)], width=thin_stroke_px)
        line([(14.1, 8.3), (18.2, 12.4)], width=thin_stroke_px)
    elif kind == "swap":
        line([(5.8, 8.1), (17.0, 8.1)])
        polygon([(18.2, 8.1), (14.9, 5.6), (14.9, 10.6)], fill=c, outline=False)
        line([(18.2, 15.9), (7.0, 15.9)])
        polygon([(5.8, 15.9), (9.1, 13.4), (9.1, 18.4)], fill=c, outline=False)
    elif kind == "music":
        draw.ellipse((4.8 * scale, 11.0 * scale, 11.9 * scale, 18.1 * scale), fill=c)
        draw.ellipse((12.1 * scale, 11.0 * scale, 19.2 * scale, 18.1 * scale), fill=c)
        draw.rounded_rectangle(
            (9.6 * scale, 4.0 * scale, 12.4 * scale, 14.4 * scale),
            radius=1.1 * scale,
            fill=c,
        )
        draw.rounded_rectangle(
            (15.9 * scale, 4.0 * scale, 18.7 * scale, 14.4 * scale),
            radius=1.1 * scale,
            fill=c,
        )
        draw.rounded_rectangle(
            (9.6 * scale, 4.0 * scale, 18.7 * scale, 6.8 * scale),
            radius=1.1 * scale,
            fill=c,
        )
        draw.rounded_rectangle(
            (10.4 * scale, 8.4 * scale, 17.9 * scale, 10.4 * scale),
            radius=0.8 * scale,
            fill=c,
        )
    elif kind == "box":
        polygon([(12.0, 4.2), (17.6, 7.2), (12.0, 10.3), (6.4, 7.2)])
        line([(6.4, 7.2), (6.4, 15.5), (12.0, 18.7), (12.0, 10.3)])
        line([(17.6, 7.2), (17.6, 15.5), (12.0, 18.7)])
        line([(12.0, 10.3), (12.0, 18.7)], width=thin_stroke_px)
    elif kind == "document":
        rounded_rect(5.9, 3.7, 17.7, 19.1, radius=1.8)
        line([(12.8, 3.8), (12.8, 8.1), (17.4, 8.1)])
        line([(12.8, 3.8), (17.4, 8.1)])
        line([(8.4, 11.1), (14.7, 11.1)], width=thin_stroke_px)
        line([(8.4, 14.2), (14.7, 14.2)], width=thin_stroke_px)
    elif kind == "image":
        rounded_rect(5.0, 5.5, 19.0, 17.4, radius=1.2)
        ellipse(7.0, 7.6, 9.9, 10.5, width=thin_stroke_px)
        line([(6.8, 14.9), (10.0, 11.8), (12.7, 14.0), (14.2, 12.3), (17.0, 15.0)])
    elif kind == "lock":
        rounded_rect(5.8, 10.7, 18.2, 19.0, radius=1.8)
        draw.arc(
            (7.6 * scale, 4.0 * scale, 16.4 * scale, 13.5 * scale),
            start=200,
            end=340,
            fill=c,
            width=stroke_px,
        )
        ellipse(11.0, 13.1, 12.8, 14.9, width=thin_stroke_px)
        line([(11.9, 14.8), (11.9, 16.6)], width=thin_stroke_px)
    elif kind == "folder":
        line([(4.9, 8.5), (9.0, 8.5), (10.9, 6.7), (18.7, 6.7)])
        rounded_rect(4.9, 8.5, 18.9, 17.8, radius=1.5)
    elif kind == "settings":
        ellipse(8.0, 8.0, 16.0, 16.0, width=stroke_px)
        ellipse(10.5, 10.5, 13.5, 13.5, width=thin_stroke_px)
        for x1, y1, x2, y2 in (
            (12.0, 3.6, 12.0, 5.5),
            (12.0, 18.5, 12.0, 20.4),
            (3.6, 12.0, 5.5, 12.0),
            (18.5, 12.0, 20.4, 12.0),
            (6.1, 6.1, 7.5, 7.5),
            (16.5, 16.5, 17.9, 17.9),
            (16.5, 7.5, 17.9, 6.1),
            (6.1, 17.9, 7.5, 16.5),
        ):
            line([(x1, y1), (x2, y2)], width=thin_stroke_px)
    elif kind == "book":
        line([(12.0, 5.2), (12.0, 18.6)], width=thin_stroke_px)
        line([(6.4, 6.2), (10.9, 5.5), (10.9, 17.9), (6.4, 17.1), (6.4, 6.2)])
        line([(17.6, 6.2), (13.1, 5.5), (13.1, 17.9), (17.6, 17.1), (17.6, 6.2)])
    elif kind == "coffee":
        rounded_rect(5.9, 10.3, 15.8, 16.2, radius=1.6)
        draw.arc(
            (13.4 * scale, 10.8 * scale, 18.3 * scale, 15.5 * scale),
            start=270,
            end=90,
            fill=c,
            width=stroke_px,
        )
        line([(6.3, 18.2), (16.2, 18.2)], width=thin_stroke_px)
        draw.arc(
            (7.5 * scale, 4.6 * scale, 9.5 * scale, 8.4 * scale),
            start=195,
            end=355,
            fill=c,
            width=thin_stroke_px,
        )
        draw.arc(
            (10.8 * scale, 3.8 * scale, 12.8 * scale, 7.6 * scale),
            start=195,
            end=355,
            fill=c,
            width=thin_stroke_px,
        )
    else:
        ellipse(5.5, 5.5, 18.5, 18.5, width=stroke_px)


def _build_sidebar_icon_image(kind, color, size=22):
    render_scale = 4
    render_size = max(size * render_scale, 64)
    canvas = PILImage.new("RGBA", (render_size, render_size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(canvas)
    _draw_sidebar_icon(draw, kind, color, size=render_size, stroke=2)
    final_canvas = canvas.resize((size, size), PILImage.Resampling.LANCZOS)
    return customtkinter.CTkImage(light_image=final_canvas, dark_image=final_canvas, size=(size, size))


def _get_sidebar_icon_images(app):
    cache = getattr(app, "_fx_sidebar_icon_images", None)
    if cache is not None:
        return cache

    cache = {}
    for spec in SIDEBAR_BUTTON_SPECS.values():
        key = ("nav", spec["icon"])
        if key not in cache:
            cache[key] = _build_sidebar_icon_image(spec["icon"], SIDEBAR_ICON_NAV)
    help_spec = SIDEBAR_AUX_BUTTON_SPECS["btn_help_proxy"]
    donate_spec = SIDEBAR_AUX_BUTTON_SPECS["btn_donate"]
    cache[("help", help_spec["icon"])] = _build_sidebar_icon_image(help_spec["icon"], SIDEBAR_ICON_HELP)
    cache[("donate", donate_spec["icon"])] = _build_sidebar_icon_image(donate_spec["icon"], SIDEBAR_ICON_DONATE)
    app._fx_sidebar_icon_images = cache
    return cache


def _get_sidebar_brand_image(app):
    image = getattr(app, "_fx_sidebar_brand_image", None)
    if image is not None:
        return image

    png_path = _resolve_app_asset(APP_ICON_PNG)
    if not png_path.exists():
        return None

    try:
        brand_image = PILImage.open(png_path).convert("RGBA")
        image = customtkinter.CTkImage(
            light_image=brand_image,
            dark_image=brand_image,
            size=(44, 44),
        )
        app._fx_sidebar_brand_image = image
        return image
    except Exception as exc:
        _debug(f"sidebar_brand_image:error:{exc}")
        return None


def _get_inline_title_icon_image(app, kind, color, size):
    cache = getattr(app, "_fx_inline_title_icon_images", None)
    if cache is None:
        cache = {}
        app._fx_inline_title_icon_images = cache

    cache_key = (kind, color, size)
    image = cache.get(cache_key)
    if image is None:
        image = _build_sidebar_icon_image(kind, color, size=size)
        cache[cache_key] = image
    return image


def _get_inline_donate_qr_image(app):
    image = getattr(app, "_fx_inline_donate_qr_image", None)
    if image is not None:
        return image

    qr_path = _resolve_app_asset(DONATE_QR_PNG)
    if not qr_path.exists():
        return None

    try:
        qr_image = PILImage.open(qr_path).convert("RGBA")
        image = customtkinter.CTkImage(
            light_image=qr_image,
            dark_image=qr_image,
            size=(280, 280),
        )
        app._fx_inline_donate_qr_image = image
        return image
    except Exception as exc:
        _debug(f"inline_donate:qr_error:{exc}")
        return None


def _ensure_inline_help_tab(app):
    help_tab = getattr(app, "tab_help", None)
    if help_tab is not None:
        return help_tab
    try:
        help_tab = app.main_panel.add(HELP_TAB_TITLE)
    except Exception:
        help_tab = app.main_panel.tab(HELP_TAB_TITLE)
    app.tab_help = help_tab
    _build_inline_help_page(app, help_tab)
    return help_tab


def _build_inline_help_page(app, help_tab):
    if getattr(help_tab, "_fx_help_page_built", False):
        return
    try:
        help_tab.grid_rowconfigure(0, weight=1)
        help_tab.grid_columnconfigure(0, weight=1)
    except Exception:
        pass

    card = customtkinter.CTkFrame(
        help_tab,
        fg_color=globals().get("COLOR_CARD", "#2B2B2B"),
        corner_radius=18,
        border_width=1,
        border_color=globals().get("COLOR_BORDER", "#3A3A3A"),
    )
    card.grid(row=0, column=0, sticky="nsew", padx=20, pady=18)
    card.grid_rowconfigure(1, weight=1)
    card.grid_columnconfigure(0, weight=1)

    header = customtkinter.CTkFrame(card, fg_color="transparent")
    header.grid(row=0, column=0, sticky="ew", padx=28, pady=(22, 10))
    header.grid_columnconfigure(1, weight=1)

    icon = _get_inline_title_icon_image(app, "book", CONTENT_ICON_PRIMARY, 24)
    customtkinter.CTkLabel(
        header,
        text="使用教程",
        image=icon,
        compound="left",
        anchor="w",
        font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=22, weight="bold"),
        text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
    ).grid(row=0, column=0, sticky="w")
    customtkinter.CTkLabel(
        header,
        text="内置说明随功能同步更新，内容较多时可向下滚动查看。",
        anchor="e",
        font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=12),
        text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
    ).grid(row=0, column=1, sticky="e", padx=(18, 0))

    scroll = customtkinter.CTkScrollableFrame(
        card,
        fg_color=globals().get("COLOR_CARD_ALT", "#303030"),
        corner_radius=14,
        border_width=1,
        border_color=globals().get("COLOR_BORDER", "#3A3A3A"),
    )
    scroll.grid(row=1, column=0, sticky="nsew", padx=28, pady=(0, 24))
    scroll.grid_columnconfigure(0, weight=1)

    for section_index, (title, lines) in enumerate(INLINE_HELP_SECTIONS):
        section = customtkinter.CTkFrame(scroll, fg_color="transparent")
        section.grid(row=section_index, column=0, sticky="ew", padx=18, pady=(16 if section_index else 18, 2))
        section.grid_columnconfigure(0, weight=1)
        customtkinter.CTkLabel(
            section,
            text=title,
            anchor="w",
            font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=16, weight="bold"),
            text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
        ).grid(row=0, column=0, sticky="ew")
        body_text = "\n".join(f"• {line}" for line in lines)
        customtkinter.CTkLabel(
            section,
            text=body_text,
            anchor="w",
            justify="left",
            wraplength=920,
            font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=13),
            text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
        ).grid(row=1, column=0, sticky="ew", pady=(6, 0))

    help_tab._fx_help_page_built = True


def _ensure_inline_donate_tab(app):
    donate_tab = getattr(app, "tab_donate", None)
    if donate_tab is not None:
        return donate_tab
    try:
        donate_tab = app.main_panel.add(DONATE_TAB_TITLE)
    except Exception:
        donate_tab = app.main_panel.tab(DONATE_TAB_TITLE)
    app.tab_donate = donate_tab
    _build_inline_donate_page(app, donate_tab)
    return donate_tab


def _build_inline_donate_page(app, donate_tab):
    if getattr(donate_tab, "_fx_donate_page_built", False):
        return
    try:
        donate_tab.grid_rowconfigure(0, weight=1)
        donate_tab.grid_columnconfigure(0, weight=1)
    except Exception:
        pass

    card = customtkinter.CTkFrame(
        donate_tab,
        fg_color=globals().get("COLOR_CARD", "#2B2B2B"),
        corner_radius=18,
        border_width=1,
        border_color=globals().get("COLOR_BORDER", "#3A3A3A"),
    )
    card.grid(row=0, column=0, sticky="nsew", padx=20, pady=18)
    card.grid_rowconfigure(1, weight=1)
    card.grid_columnconfigure(0, weight=1)

    header = customtkinter.CTkFrame(card, fg_color="transparent")
    header.grid(row=0, column=0, sticky="ew", padx=28, pady=(22, 10))
    header.grid_columnconfigure(1, weight=1)

    icon = _get_inline_title_icon_image(app, "coffee", SIDEBAR_ICON_DONATE, 24)
    customtkinter.CTkLabel(
        header,
        text=DONATE_TAB_TITLE,
        image=icon,
        compound="left",
        anchor="w",
        font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=22, weight="bold"),
        text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
    ).grid(row=0, column=0, sticky="w")
    customtkinter.CTkLabel(
        header,
        text="感谢每一份支持，它都会变成后续维护和新功能的动力。",
        anchor="e",
        font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=12),
        text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
    ).grid(row=0, column=1, sticky="e", padx=(18, 0))

    body = customtkinter.CTkFrame(
        card,
        fg_color=globals().get("COLOR_CARD_ALT", "#303030"),
        corner_radius=14,
        border_width=1,
        border_color=globals().get("COLOR_BORDER", "#3A3A3A"),
    )
    body.grid(row=1, column=0, sticky="nsew", padx=28, pady=(0, 24))
    body.grid_columnconfigure(0, weight=1)
    body.grid_columnconfigure(1, weight=0)
    body.grid_rowconfigure(0, weight=1)

    text_panel = customtkinter.CTkFrame(body, fg_color="transparent")
    text_panel.grid(row=0, column=0, sticky="nsew", padx=(28, 24), pady=28)
    text_panel.grid_columnconfigure(0, weight=1)

    customtkinter.CTkLabel(
        text_panel,
        text="让风兮继续长风而行",
        anchor="w",
        font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=18, weight="bold"),
        text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
    ).grid(row=0, column=0, sticky="ew")
    customtkinter.CTkLabel(
        text_panel,
        text=DONATE_SUPPORT_SENTENCE,
        anchor="w",
        justify="left",
        wraplength=560,
        font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=15, weight="bold"),
        text_color="#F6E2B2",
    ).grid(row=1, column=0, sticky="ew", pady=(14, 12))
    customtkinter.CTkLabel(
        text_panel,
        text="赞助完全自愿，不影响任何功能使用。你的支持会优先用于修 bug、做兼容测试、维护 OCR/PDF/打包环境，以及继续打磨批处理体验。",
        anchor="w",
        justify="left",
        wraplength=560,
        font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=13),
        text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
    ).grid(row=2, column=0, sticky="ew", pady=(0, 16))
    customtkinter.CTkLabel(
        text_panel,
        text="谢谢你愿意把风兮工具箱留在电脑里，也谢谢你每一次反馈。",
        anchor="w",
        justify="left",
        wraplength=560,
        font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=13),
        text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
    ).grid(row=3, column=0, sticky="ew")

    qr_panel = customtkinter.CTkFrame(
        body,
        fg_color="#241A10",
        corner_radius=18,
        border_width=1,
        border_color=SIDEBAR_AUX_STYLES["donate"]["border_color"],
    )
    qr_panel.grid(row=0, column=1, sticky="n", padx=(0, 28), pady=28)
    qr_panel.grid_columnconfigure(0, weight=1)

    qr_image = _get_inline_donate_qr_image(app)
    if qr_image is not None:
        customtkinter.CTkLabel(qr_panel, text="", image=qr_image).grid(row=0, column=0, padx=18, pady=(18, 10))
    else:
        customtkinter.CTkLabel(
            qr_panel,
            text="未找到赞助二维码图片\nassets/donate_qr.png",
            justify="center",
            width=280,
            height=280,
            font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=13),
            text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
        ).grid(row=0, column=0, padx=18, pady=(18, 10))
    customtkinter.CTkLabel(
        qr_panel,
        text="扫码赞助作者",
        anchor="center",
        font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=14, weight="bold"),
        text_color=SIDEBAR_AUX_STYLES["donate"]["text_color"],
    ).grid(row=1, column=0, sticky="ew", padx=18, pady=(0, 18))

    donate_tab._fx_donate_page_built = True


def _set_help_button_selected(app, selected):
    btn = getattr(app, "btn_help_proxy", None)
    if btn is None:
        return
    try:
        if selected:
            btn.configure(
                text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
                fg_color=SIDEBAR_AUX_STYLES["help"]["hover_color"],
                border_color=globals().get("COLOR_BORDER", "#566274"),
            )
        else:
            _style_sidebar_aux_button(
                btn,
                SIDEBAR_AUX_BUTTON_SPECS["btn_help_proxy"]["label"],
                _get_sidebar_icon_images(app).get(("help", SIDEBAR_AUX_BUTTON_SPECS["btn_help_proxy"]["icon"])),
                _get_sidebar_button_font(app),
                "help",
            )
    except Exception:
        pass


def _set_donate_button_selected(app, selected):
    btn = getattr(app, "btn_donate", None)
    if btn is None:
        return
    try:
        if selected:
            btn.configure(
                text_color=SIDEBAR_AUX_STYLES["donate"]["text_color"],
                fg_color=SIDEBAR_AUX_STYLES["donate"]["hover_color"],
                border_color=SIDEBAR_AUX_STYLES["donate"]["border_color"],
            )
        else:
            _style_sidebar_aux_button(
                btn,
                SIDEBAR_AUX_BUTTON_SPECS["btn_donate"]["label"],
                _get_sidebar_icon_images(app).get(("donate", SIDEBAR_AUX_BUTTON_SPECS["btn_donate"]["icon"])),
                _get_sidebar_button_font(app),
                "donate",
            )
    except Exception:
        pass


def _set_help_action_state(app, visible, label="查看说明中"):
    if getattr(app, "is_running", False):
        return
    btn_run = getattr(app, "btn_run", None)
    btn_stop = getattr(app, "btn_stop", None)
    if btn_run is not None:
        try:
            if visible:
                btn_run.configure(state="disabled", text=label, fg_color="#455A64")
            else:
                btn_run.configure(
                    state="normal",
                    text="🚀 立即开始处理",
                    fg_color=globals().get("COLOR_ACCENT", "#8FA9B8"),
                )
        except Exception:
            pass
    if btn_stop is not None and visible:
        try:
            btn_stop.configure(state="disabled")
        except Exception:
            pass


def _show_inline_help(app):
    try:
        _ensure_inline_help_tab(app)
        app.main_panel.set(HELP_TAB_TITLE)
        app.current_task = "help"
        _set_help_button_selected(app, True)
        _set_donate_button_selected(app, False)
        _set_help_action_state(app, True, "查看使用说明中")
        for nav_name in (
            "btn_nav_wm",
            "btn_nav_rm_wm",
            "btn_nav_cv",
            "btn_nav_audio",
            "btn_nav_zip",
            "btn_nav_pdf",
            "btn_nav_img",
            "btn_nav_meta",
            "btn_nav_file",
        ):
            nav = getattr(app, nav_name, None)
            if nav is not None:
                nav.configure(
                    text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
                    fg_color="transparent",
                    border_color=globals().get("COLOR_SIDEBAR", "#202020"),
                )
        app.update_idletasks()
    except Exception as exc:
        _debug(f"inline_help:show_error:{exc}")


def _show_inline_donate(app):
    try:
        _ensure_inline_donate_tab(app)
        app.main_panel.set(DONATE_TAB_TITLE)
        app.current_task = "donate"
        _set_help_button_selected(app, False)
        _set_donate_button_selected(app, True)
        _set_help_action_state(app, True, "查看赞助作者中")
        for nav_name in (
            "btn_nav_wm",
            "btn_nav_rm_wm",
            "btn_nav_cv",
            "btn_nav_audio",
            "btn_nav_zip",
            "btn_nav_pdf",
            "btn_nav_img",
            "btn_nav_meta",
            "btn_nav_file",
        ):
            nav = getattr(app, nav_name, None)
            if nav is not None:
                nav.configure(
                    text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
                    fg_color="transparent",
                    border_color=globals().get("COLOR_SIDEBAR", "#202020"),
                )
        app.update_idletasks()
    except Exception as exc:
        _debug(f"inline_donate:show_error:{exc}")


def _normalize_icon_title_text(text):
    text = str(text or "").strip()
    if not text:
        return ""
    while text:
        first = text[0]
        if first == "[" or first.isalnum() or ("\u4e00" <= first <= "\u9fff"):
            break
        text = text[1:].lstrip()
    return text


def _resolve_inline_title_spec(raw_text):
    raw_text = str(raw_text or "")
    for canonical_text in sorted(INLINE_TITLE_ICON_SPECS, key=len, reverse=True):
        if canonical_text in raw_text:
            return canonical_text, INLINE_TITLE_ICON_SPECS[canonical_text]

    canonical_text = _normalize_icon_title_text(raw_text)
    spec = INLINE_TITLE_ICON_SPECS.get(canonical_text)
    if spec is None:
        return "", None
    return canonical_text, spec


def _apply_inline_title_icons(app, root_widget):
    if root_widget is None:
        return

    pending = [root_widget]
    while pending:
        widget = pending.pop(0)
        try:
            pending.extend(widget.winfo_children())
        except Exception:
            pass

        if not isinstance(widget, customtkinter.CTkLabel):
            continue
        try:
            raw_text = widget.cget("text")
        except Exception:
            continue

        canonical_text, spec = _resolve_inline_title_spec(raw_text)
        if spec is None:
            continue

        signature = (
            canonical_text,
            spec["icon"],
            spec["size"],
            spec["color"],
            spec.get("display_text", canonical_text),
        )
        if getattr(widget, "_fx_inline_title_icon_signature", None) == signature:
            continue

        display_text = spec.get("display_text", canonical_text)
        try:
            widget.configure(
                text=display_text,
                image=_get_inline_title_icon_image(app, spec["icon"], spec["color"], spec["size"]),
                compound="left",
                anchor="w",
            )
            widget._fx_inline_title_icon_signature = signature
        except Exception:
            pass


def _normalize_input_path_value(value):
    text = _coerce_input_path_text(value).strip().strip('"')
    if not text:
        return ""
    return os.path.abspath(os.path.normpath(text))


def _sanitize_filename_component(value, fallback="task_result"):
    text = str(value or "").strip()
    if not text:
        return fallback
    text = text.replace("\r", " ").replace("\n", " ")
    text = re.sub(r'[<>:"/\\\\|?*\x00-\x1f]+', "_", text)
    text = re.sub(r"\s+", "_", text).strip(" ._")
    return text or fallback


def _decode_input_path_bytes(raw):
    candidates = []
    for encoding in ("utf-8", "mbcs", "gbk", sys.getfilesystemencoding()):
        if encoding and encoding not in candidates:
            candidates.append(encoding)
    for encoding in candidates:
        try:
            return raw.decode(encoding)
        except Exception:
            continue
    return raw.decode("utf-8", errors="replace")


def _coerce_input_path_text(value):
    if value is None:
        return ""
    if isinstance(value, (bytes, bytearray)):
        return _decode_input_path_bytes(bytes(value))

    text = str(value).strip()
    if not text:
        return ""

    for marker in ("b'", 'b"', "B'", 'B"'):
        marker_index = text.find(marker)
        if marker_index < 0:
            continue
        literal_text = text[marker_index:]
        try:
            literal_value = ast.literal_eval(literal_text)
        except Exception:
            continue
        if isinstance(literal_value, (bytes, bytearray)):
            return _decode_input_path_bytes(bytes(literal_value))

    return text


class _LazyImportSymbol:
    def __init__(self, module_name, attr_name):
        self._module_name = module_name
        self._attr_name = attr_name
        self._resolved = None

    def _resolve(self):
        if self._resolved is None:
            module = _resolve_lazy_runtime_module(self._module_name)
            self._resolved = getattr(module, self._attr_name)
        return self._resolved

    def __call__(self, *args, **kwargs):
        return self._resolve()(*args, **kwargs)

    def __getattr__(self, item):
        return getattr(self._resolve(), item)

    def __repr__(self):
        return f"<lazy import {self._module_name}.{self._attr_name}>"


class _LazyImportModule(types.ModuleType):
    def __init__(self, module_name):
        super().__init__(module_name)
        self.__dict__["_fx_lazy_module_name"] = module_name
        self.__dict__["__package__"] = module_name.rpartition(".")[0]
        if "." not in module_name:
            self.__dict__["__path__"] = []

    def _resolve(self):
        return _resolve_lazy_runtime_module(self._fx_lazy_module_name)

    def __getattr__(self, item):
        return getattr(self._resolve(), item)


def _has_lazy_runtime_module(module_name):
    try:
        return importlib.util.find_spec(module_name) is not None
    except Exception:
        return False


def _is_lazy_runtime_proxy(module_obj, module_name):
    return getattr(module_obj, "_fx_lazy_module_name", None) == module_name


def _resolve_lazy_runtime_module(module_name):
    cache = globals().setdefault("_fx_lazy_runtime_cache", {})
    if module_name in cache:
        return cache[module_name]

    removed = {}
    for name in LAZY_RUNTIME_IMPORT_SPECS:
        if name == module_name or name.startswith(module_name + ".") or module_name.startswith(name + "."):
            module_obj = sys.modules.get(name)
            if _is_lazy_runtime_proxy(module_obj, name):
                removed[name] = sys.modules.pop(name)

    spec = LAZY_RUNTIME_IMPORT_SPECS.get(module_name, {})
    candidate_names = [module_name]
    for fallback_name in spec.get("fallback_modules", ()):
        if fallback_name and fallback_name not in candidate_names:
            candidate_names.append(fallback_name)

    last_exc = None
    resolved_name = module_name
    module = None
    for candidate_name in candidate_names:
        try:
            module = importlib.import_module(candidate_name)
            resolved_name = candidate_name
            break
        except Exception as exc:
            last_exc = exc
    if module is None:
        for name, module_obj in removed.items():
            if name not in sys.modules:
                sys.modules[name] = module_obj
        raise last_exc

    if resolved_name != module_name:
        _debug(f"lazy_runtime_imports:fallback:{module_name}->{resolved_name}")

    cache[module_name] = module
    cache[resolved_name] = module
    return module


def _install_runtime_lazy_imports():
    previous = {}
    installed = {}

    for module_name, spec in LAZY_RUNTIME_IMPORT_SPECS.items():
        probe_name = spec.get("probe", module_name)
        if not _has_lazy_runtime_module(probe_name):
            continue
        proxy = _LazyImportModule(module_name)
        for attr_name, target in spec.get("symbols", {}).items():
            target_module, target_attr = target
            setattr(proxy, attr_name, _LazyImportSymbol(target_module, target_attr))
        installed[module_name] = proxy

    moviepy_proxy = installed.get("moviepy")
    moviepy_editor_proxy = installed.get("moviepy.editor")
    if moviepy_proxy is not None and moviepy_editor_proxy is not None:
        moviepy_proxy.editor = moviepy_editor_proxy

    for module_name, proxy in installed.items():
        if module_name in sys.modules:
            continue
        previous[module_name] = None
        sys.modules[module_name] = proxy

    if previous:
        _debug(f"lazy_runtime_imports:installed:{','.join(sorted(previous))}")
    return previous


def _restore_runtime_lazy_imports(previous):
    if not previous:
        return
    for module_name, prior in previous.items():
        current = sys.modules.get(module_name)
        if prior is None:
            if _is_lazy_runtime_proxy(current, module_name):
                sys.modules.pop(module_name, None)
        else:
            sys.modules[module_name] = prior
    _debug("lazy_runtime_imports:restored")


class _UnifiedInputPathPicker:
    def __init__(self, parent, initial_path=""):
        self.parent = parent
        self.result = ""
        self.current_dir = self._resolve_initial_dir(initial_path)
        self.path_var = tkinter.StringVar(value=self.current_dir)
        self.status_var = tkinter.StringVar(value="可直接双击文件确认，或选择文件夹后点确定。")
        self._item_paths = {}

        self.window = customtkinter.CTkToplevel(parent)
        _apply_window_icon(self.window)
        self.window.title("选择文件或文件夹")
        self.window.geometry("860x560")
        self.window.minsize(760, 480)
        self.window.transient(parent)
        self.window.protocol("WM_DELETE_WINDOW", self._cancel)
        self.window.grid_columnconfigure(0, weight=1)
        self.window.grid_rowconfigure(0, weight=1)

        card = customtkinter.CTkFrame(self.window, corner_radius=16)
        card.grid(row=0, column=0, sticky="nsew", padx=18, pady=18)
        card.grid_columnconfigure(0, weight=1)
        card.grid_rowconfigure(1, weight=1)

        toolbar = customtkinter.CTkFrame(card, fg_color="transparent")
        toolbar.grid(row=0, column=0, sticky="ew", padx=18, pady=(18, 10))
        toolbar.grid_columnconfigure(0, weight=1)

        self.path_entry = customtkinter.CTkEntry(toolbar, textvariable=self.path_var, height=38)
        self.path_entry.grid(row=0, column=0, sticky="ew", padx=(0, 10))

        customtkinter.CTkButton(toolbar, text="打开", width=72, height=38, command=self._open_typed_path).grid(
            row=0, column=1, padx=(0, 8)
        )
        customtkinter.CTkButton(toolbar, text="上一级", width=84, height=38, command=self._go_up).grid(
            row=0, column=2, padx=(0, 8)
        )
        customtkinter.CTkButton(toolbar, text="刷新", width=72, height=38, command=self._refresh_entries).grid(
            row=0, column=3
        )

        body = tkinter.Frame(card)
        body.grid(row=1, column=0, sticky="nsew", padx=18, pady=(0, 10))
        body.grid_columnconfigure(0, weight=1)
        body.grid_rowconfigure(0, weight=1)

        self.tree = tkinter.ttk.Treeview(body, columns=("kind", "size"), show="tree headings", selectmode="browse")
        self.tree.heading("#0", text="名称")
        self.tree.heading("kind", text="类型")
        self.tree.heading("size", text="大小")
        self.tree.column("#0", width=460, anchor="w")
        self.tree.column("kind", width=100, anchor="center")
        self.tree.column("size", width=120, anchor="e")
        self.tree.grid(row=0, column=0, sticky="nsew")

        scrollbar = tkinter.ttk.Scrollbar(body, orient="vertical", command=self.tree.yview)
        scrollbar.grid(row=0, column=1, sticky="ns")
        self.tree.configure(yscrollcommand=scrollbar.set)

        customtkinter.CTkLabel(card, textvariable=self.status_var, anchor="w", justify="left").grid(
            row=2, column=0, sticky="ew", padx=18, pady=(0, 10)
        )

        actions = customtkinter.CTkFrame(card, fg_color="transparent")
        actions.grid(row=3, column=0, sticky="ew", padx=18, pady=(0, 18))
        actions.grid_columnconfigure(0, weight=1)

        customtkinter.CTkButton(
            actions,
            text="选择当前文件夹",
            width=130,
            height=38,
            command=self._choose_current_directory,
        ).grid(row=0, column=0, sticky="w")
        customtkinter.CTkButton(actions, text="确定", width=88, height=38, command=self._confirm_selection).grid(
            row=0, column=1, padx=(8, 8)
        )
        customtkinter.CTkButton(actions, text="取消", width=88, height=38, command=self._cancel).grid(
            row=0, column=2
        )

        self.tree.bind("<<TreeviewSelect>>", self._on_tree_select)
        self.tree.bind("<Double-1>", self._on_tree_double_click)
        self.tree.bind("<Return>", self._on_tree_double_click)
        self.path_entry.bind("<Return>", lambda _event: self._open_typed_path())

        self._center_on_parent()
        self._refresh_entries()

    def _resolve_initial_dir(self, initial_path):
        normalized = _normalize_input_path_value(initial_path)
        if normalized and os.path.isfile(normalized):
            return os.path.dirname(normalized)
        if normalized and os.path.isdir(normalized):
            return normalized
        home_dir = _normalize_input_path_value(os.path.expanduser("~"))
        return home_dir if os.path.isdir(home_dir) else os.getcwd()

    def _center_on_parent(self):
        try:
            self.parent.update_idletasks()
            self.window.update_idletasks()
            width = 860
            height = 560
            x = max(self.parent.winfo_rootx() + (self.parent.winfo_width() - width) // 2, 80)
            y = max(self.parent.winfo_rooty() + (self.parent.winfo_height() - height) // 2, 60)
            self.window.geometry(f"{width}x{height}+{x}+{y}")
        except Exception:
            pass

    def _format_size(self, size):
        if size < 1024:
            return f"{size} B"
        if size < 1024 * 1024:
            return f"{size / 1024:.1f} KB"
        if size < 1024 * 1024 * 1024:
            return f"{size / (1024 * 1024):.1f} MB"
        return f"{size / (1024 * 1024 * 1024):.1f} GB"

    def _refresh_entries(self):
        target_dir = _normalize_input_path_value(self.path_var.get()) or self.current_dir
        if not os.path.isdir(target_dir):
            self.status_var.set("当前路径不是文件夹，请输入有效路径后再打开。")
            return

        try:
            entries = sorted(os.scandir(target_dir), key=lambda entry: (not entry.is_dir(), entry.name.lower()))
        except Exception as exc:
            self.status_var.set(f"无法读取目录: {exc}")
            return

        self.current_dir = target_dir
        self.path_var.set(target_dir)
        self.tree.delete(*self.tree.get_children())
        self._item_paths = {}

        for index, entry in enumerate(entries):
            item_id = f"item_{index}"
            item_path = os.path.abspath(entry.path)
            item_kind = "文件夹" if entry.is_dir() else "文件"
            item_size = ""
            if entry.is_file():
                try:
                    item_size = self._format_size(entry.stat().st_size)
                except Exception:
                    item_size = "-"
            self.tree.insert("", "end", iid=item_id, text=entry.name, values=(item_kind, item_size))
            self._item_paths[item_id] = item_path

        self.status_var.set(f"当前目录: {target_dir} | 共 {len(entries)} 项，可直接选择文件或文件夹。")

    def _get_selected_path(self):
        selection = self.tree.selection()
        if not selection:
            return ""
        return self._item_paths.get(selection[0], "")

    def _on_tree_select(self, _event=None):
        selected_path = self._get_selected_path()
        if not selected_path:
            self.status_var.set(f"当前目录: {self.current_dir} | 可双击文件确认，双击文件夹进入。")
            return
        selected_type = "文件夹" if os.path.isdir(selected_path) else "文件"
        self.status_var.set(f"已选中{selected_type}: {selected_path}")

    def _on_tree_double_click(self, _event=None):
        selected_path = self._get_selected_path()
        if not selected_path:
            return
        if os.path.isdir(selected_path):
            self.path_var.set(selected_path)
            self._refresh_entries()
            return
        self._finish(selected_path)

    def _open_typed_path(self):
        typed_path = _normalize_input_path_value(self.path_var.get())
        if not typed_path:
            tkinter.messagebox.showwarning("提示", "请输入文件或文件夹路径。", parent=self.window)
            return
        if os.path.isdir(typed_path):
            self.path_var.set(typed_path)
            self._refresh_entries()
            return
        if os.path.isfile(typed_path):
            self._finish(typed_path)
            return
        tkinter.messagebox.showwarning("提示", "路径不存在，请重新输入。", parent=self.window)

    def _go_up(self):
        current = _normalize_input_path_value(self.current_dir)
        if not current:
            return
        parent_dir = os.path.dirname(current.rstrip("\\/")) or current
        if not os.path.isdir(parent_dir):
            parent_dir = current
        self.path_var.set(parent_dir)
        self._refresh_entries()

    def _choose_current_directory(self):
        self._finish(self.current_dir)

    def _confirm_selection(self):
        selected_path = self._get_selected_path()
        if not selected_path:
            tkinter.messagebox.showwarning(
                "提示",
                "请先选择一个文件或文件夹，或使用“选择当前文件夹”。",
                parent=self.window,
            )
            return
        self._finish(selected_path)

    def _finish(self, path):
        self.result = _normalize_input_path_value(path)
        try:
            self.window.grab_release()
        except Exception:
            pass
        self.window.destroy()

    def _cancel(self):
        self.result = ""
        try:
            self.window.grab_release()
        except Exception:
            pass
        self.window.destroy()

    def show(self):
        try:
            self.window.lift()
            self.window.attributes("-topmost", True)
            self.window.after(150, lambda: self.window.attributes("-topmost", False))
        except Exception:
            pass
        try:
            self.window.grab_set()
        except Exception:
            pass
        self.path_entry.focus_set()
        self.parent.wait_window(self.window)
        return self.result


def _choose_input_path_via_shell_dialog(app, initial_dir=""):
    normalized_initial = _normalize_input_path_value(initial_dir)
    if normalized_initial and os.path.isfile(normalized_initial):
        normalized_initial = os.path.dirname(normalized_initial)
    if normalized_initial and not os.path.isdir(normalized_initial):
        normalized_initial = ""

    def _browse_callback(hwnd, msg, lparam, lpdata):
        if msg == BFFM_INITIALIZED and lpdata:
            try:
                win32gui.SendMessage(hwnd, BFFM_SETSELECTIONW, 1, lpdata)
            except Exception as exc:
                _debug(f"shell_picker:set_selection_error:{exc}")
        return 0

    try:
        app.update_idletasks()
    except Exception:
        pass

    try:
        owner_hwnd = int(app.winfo_id())
    except Exception:
        owner_hwnd = 0

    flags = BIF_USENEWUI | shellcon.BIF_BROWSEINCLUDEFILES | shellcon.BIF_VALIDATE
    try:
        result = shell.SHBrowseForFolder(
            owner_hwnd,
            None,
            "选择文件或文件夹",
            flags,
            _browse_callback,
            normalized_initial or None,
        )
    except Exception as exc:
        _debug(f"shell_picker:dialog_error:{exc}")
        return ""

    if not result:
        return ""

    pidl = result[0] if isinstance(result, (tuple, list)) else result
    if not pidl:
        return ""

    try:
        selected_path = shell.SHGetPathFromIDList(pidl)
    except Exception as exc:
        _debug(f"shell_picker:get_path_error:{exc}")
        return ""

    return _normalize_input_path_value(selected_path)


def _resolve_app_asset(name):
    try:
        base_dir = Path(sys.executable).resolve().parent if getattr(sys, "frozen", False) else Path(__file__).resolve().parent
    except Exception:
        base_dir = Path(__file__).resolve().parent
    return base_dir / "assets" / name


def _apply_native_window_icons(window, ico_path):
    """Set instance and class icons after Tk has created the native window."""
    try:
        try:
            top_window = window.winfo_toplevel()
        except Exception:
            top_window = window
        try:
            top_window.update_idletasks()
        except Exception:
            pass
        hwnd = int(top_window.winfo_id())
        user32 = ctypes.windll.user32
        user32.GetAncestor.argtypes = (ctypes.c_void_p, ctypes.c_uint)
        user32.GetAncestor.restype = ctypes.c_void_p
        user32.LoadImageW.argtypes = (
            ctypes.c_void_p,
            ctypes.c_wchar_p,
            ctypes.c_uint,
            ctypes.c_int,
            ctypes.c_int,
            ctypes.c_uint,
        )
        user32.LoadImageW.restype = ctypes.c_void_p
        user32.SendMessageW.argtypes = (
            ctypes.c_void_p,
            ctypes.c_uint,
            ctypes.c_void_p,
            ctypes.c_void_p,
        )
        user32.SendMessageW.restype = ctypes.c_void_p
        user32.SetClassLongPtrW.argtypes = (
            ctypes.c_void_p,
            ctypes.c_int,
            ctypes.c_void_p,
        )
        user32.SetClassLongPtrW.restype = ctypes.c_void_p
        load_from_file = 0x00000010
        image_icon = 1
        icon_big = user32.LoadImageW(None, str(ico_path), image_icon, 32, 32, load_from_file)
        icon_small = user32.LoadImageW(None, str(ico_path), image_icon, 16, 16, load_from_file)
        if not icon_big and not icon_small:
            return False
        root_value = user32.GetAncestor(ctypes.c_void_p(hwnd), 2)
        root_hwnd = int(root_value.value if root_value else hwnd)
        if icon_big:
            user32.SendMessageW(ctypes.c_void_p(root_hwnd), 0x0080, ctypes.c_void_p(1), icon_big)
        if icon_small:
            user32.SendMessageW(ctypes.c_void_p(root_hwnd), 0x0080, ctypes.c_void_p(0), icon_small)
        # Tk registers a shared TkTopLevel class. Keep the class icon aligned too,
        # otherwise Windows may keep showing the default icon in the title/task bar.
        if icon_big:
            user32.SetClassLongPtrW(ctypes.c_void_p(root_hwnd), -14, icon_big)
        if icon_small:
            user32.SetClassLongPtrW(ctypes.c_void_p(root_hwnd), -34, icon_small)
        window._fx_native_icon_handles = (icon_big, icon_small)
        window._fx_native_icon_hwnd = root_hwnd
        get_icon = user32.SendMessageW
        observed_big = get_icon(ctypes.c_void_p(root_hwnd), 0x007F, ctypes.c_void_p(1), 0)
        observed_small = get_icon(ctypes.c_void_p(root_hwnd), 0x007F, ctypes.c_void_p(0), 0)
        _debug(
            "native_window_icon:hwnd=%s,root=%s,big=%s,small=%s,observed_big=%s,observed_small=%s"
            % (hwnd, root_hwnd, icon_big, icon_small, observed_big, observed_small)
        )
        return True
    except Exception as exc:
        _debug(f"native_window_icon:error:{exc}")
        return False


def _apply_app_icon_bitmap(app):
    ico_path = _resolve_app_asset(APP_ICON_ICO)
    try:
        if ico_path.exists():
            app.iconbitmap(default=str(ico_path))
            _apply_native_window_icons(app, ico_path)
    except Exception as exc:
        _debug(f"app_icon:iconbitmap_error:{exc}")


def _apply_app_icon_png(app):
    png_path = _resolve_app_asset(APP_ICON_PNG)
    try:
        if png_path.exists():
            app._fx_window_icon = tkinter.PhotoImage(file=str(png_path))
            app.iconphoto(True, app._fx_window_icon)
    except Exception as exc:
        _debug(f"app_icon:iconphoto_error:{exc}")


def _apply_app_icon_native(app):
    ico_path = _resolve_app_asset(APP_ICON_ICO)
    try:
        if ico_path.exists():
            _apply_native_window_icons(app, ico_path)
    except Exception as exc:
        _debug(f"app_icon:native_icon_error:{exc}")


def _apply_app_icon(app):
    _apply_app_icon_bitmap(app)
    _apply_app_icon_png(app)


def _apply_app_icon_after_first_paint(app):
    """Apply the taskbar icon after the main window has had a chance to map."""
    started_at = time.perf_counter()
    try:
        _apply_app_icon_native(app)
        _apply_app_icon_png(app)
        _debug("startup:icon_png_applied_deferred")
    except Exception as exc:
        _debug(f"startup:icon_apply_deferred_error:{exc}")
    finally:
        _record_performance("main_icon_apply_deferred", started_at=started_at)


def _apply_window_icon(window):
    try:
        if hasattr(window, "_fx_window_icon") and getattr(window, "_fx_window_icon", None) is not None:
            window.iconphoto(True, window._fx_window_icon)
            return
    except Exception:
        pass
    ico_path = _resolve_app_asset(APP_ICON_ICO)
    png_path = _resolve_app_asset(APP_ICON_PNG)
    try:
        if ico_path.exists():
            window.iconbitmap(default=str(ico_path))
            _apply_native_window_icons(window, ico_path)
    except Exception as exc:
        _debug(f"window_icon:iconbitmap_error:{exc}")
    try:
        if png_path.exists():
            window._fx_window_icon = tkinter.PhotoImage(file=str(png_path))
            window.iconphoto(True, window._fx_window_icon)
    except Exception as exc:
        _debug(f"window_icon:iconphoto_error:{exc}")


def _apply_release_identity(app):
    try:
        app.title(f"风兮文件批量处理工具箱 {APP_DISPLAY_VERSION}")
    except Exception as exc:
        _debug(f"release_identity:title_error:{exc}")


def _style_sidebar_aux_button(button, text, image, font, variant):
    palette = SIDEBAR_AUX_STYLES.get(variant, SIDEBAR_AUX_STYLES["help"])
    button.configure(
        text=text,
        image=image,
        compound="left",
        border_spacing=12,
        width=208,
        height=40,
        anchor="w",
        corner_radius=14,
        border_width=1,
        font=font,
        fg_color=palette["fg_color"],
        hover_color=palette["hover_color"],
        border_color=palette["border_color"],
        text_color=palette["text_color"],
    )


def _run_with_fast_sidebar_button_construction(app, build_callable):
    try:
        original_init = customtkinter.CTkButton.__init__
    except Exception:
        return build_callable()

    if getattr(app, "_fx_sidebar_fast_build_active", False):
        return build_callable()

    def patched_init(button_self, *args, **kwargs):
        master = args[0] if args else kwargs.get("master")
        sidebar_frame = getattr(app, "sidebar_frame", None)
        if master is sidebar_frame and isinstance(kwargs.get("text"), str):
            # Sidebar buttons are restyled before the window becomes visible,
            # so we can use cheaper placeholder text/font during construction.
            kwargs["text"] = ""
            kwargs["font"] = FAST_SIDEBAR_BUILD_FONT
        return original_init(button_self, *args, **kwargs)

    app._fx_sidebar_fast_build_active = True
    customtkinter.CTkButton.__init__ = patched_init
    try:
        return build_callable()
    finally:
        customtkinter.CTkButton.__init__ = original_init
        app._fx_sidebar_fast_build_active = False


def _get_sidebar_button_font(app):
    font = getattr(app, "_fx_sidebar_button_font", None)
    if font is None:
        font = customtkinter.CTkFont(family="Microsoft YaHei UI", size=14, weight="bold")
        app._fx_sidebar_button_font = font
    return font


def _compact_main_tabview_header(app):
    tabview = getattr(app, "main_panel", None)
    if tabview is None:
        return

    try:
        tabview._outer_spacing = 0
        tabview._outer_button_overhang = 0
        tabview._button_height = 0
    except Exception:
        pass

    segmented_button = getattr(tabview, "_segmented_button", None)
    if segmented_button is not None:
        try:
            segmented_button.grid_forget()
            segmented_button.configure(height=1, width=1)
        except Exception:
            pass

    try:
        for row in (0, 1, 2):
            tabview.grid_rowconfigure(row, weight=0, minsize=0)
        tabview.grid_rowconfigure(3, weight=1, minsize=0)
    except Exception:
        pass

    canvas = getattr(tabview, "_canvas", None)
    if canvas is not None:
        try:
            canvas.grid(row=0, rowspan=4, column=0, sticky="nsew")
        except Exception:
            pass

    if not getattr(tabview, "_fx_compact_header_patch_ready", False):
        original_set_grid_current_tab = getattr(tabview, "_set_grid_current_tab", None)

        if callable(original_set_grid_current_tab):
            def compact_set_grid_current_tab(panel_self):
                original_set_grid_current_tab()
                try:
                    current_name = getattr(panel_self, "_current_name", "")
                    tab_dict = getattr(panel_self, "_tab_dict", {}) or {}
                    current_tab = tab_dict.get(current_name)
                    if current_tab is not None:
                        current_tab.grid(row=0, column=0, sticky="nsew", padx=0, pady=0)
                except Exception:
                    pass
                try:
                    segmented = getattr(panel_self, "_segmented_button", None)
                    if segmented is not None:
                        segmented.grid_forget()
                except Exception:
                    pass

            try:
                tabview._fx_original_set_grid_current_tab = original_set_grid_current_tab
                tabview._set_grid_current_tab = types.MethodType(compact_set_grid_current_tab, tabview)
                tabview._fx_compact_header_patch_ready = True
            except Exception:
                pass

    try:
        current_name = getattr(tabview, "_current_name", "")
        tab_dict = getattr(tabview, "_tab_dict", {}) or {}
        current_tab = tab_dict.get(current_name)
        if current_tab is not None:
            current_tab.grid(row=0, column=0, sticky="nsew", padx=0, pady=0)
    except Exception:
        pass


def _apply_shell_layout_tightening(app):
    if getattr(app, "_fx_shell_layout_tightened", False):
        return

    shell_fill_color = globals().get("COLOR_CARD_ALT", "#303030")
    try:
        app.configure(fg_color=shell_fill_color)
    except Exception:
        pass

    app.sidebar_frame.configure(width=300)
    app.grid_columnconfigure(0, minsize=300, weight=0)

    app.top_bar.grid_configure(pady=(8, 0))
    app.main_panel.grid_configure(pady=(0, 0))
    app.bottom_bar.grid_configure(pady=(0, 8))
    try:
        app.main_panel.configure(fg_color=shell_fill_color)
    except Exception:
        pass
    _compact_main_tabview_header(app)

    app.top_bar.configure(height=92)
    app.btn_browse.configure(height=40, text="浏览文件/文件夹")
    app.entry_path.configure(height=40)
    try:
        app.top_bar.grid_columnconfigure(0, weight=1)
        app.top_bar.grid_columnconfigure(1, weight=0)
        app.top_bar.grid_columnconfigure(2, weight=0)
    except Exception:
        pass
    for child in app.top_bar.winfo_children():
        if child is app.btn_browse:
            child.grid_configure(pady=(2, 8), padx=(0, 12))
        elif child is app.entry_path:
            child.grid_configure(pady=(2, 8), padx=(24, 16))
        elif child is getattr(app, "_fx_output_strategy_controls", None):
            child.grid_configure(pady=(8, 8), padx=(0, 24), sticky="e")
        else:
            child.grid_configure(pady=(6, 0), padx=24)

    app.bottom_bar.configure(height=228)
    try:
        app.bottom_bar.grid_propagate(False)
        app.bottom_bar.grid_columnconfigure(0, weight=1)
        app.bottom_bar.grid_columnconfigure(1, weight=1)
        app.bottom_bar.grid_rowconfigure(2, weight=0, minsize=128)
    except Exception:
        pass
    for child in app.bottom_bar.winfo_children():
        if child is app.progress_bar:
            child.grid_configure(row=0, column=0, columnspan=1, pady=(10, 8), padx=(24, 12), sticky="ew")
        elif child is getattr(app, "_fx_progress_status_label", None):
            try:
                child.configure(height=30, anchor="w", justify="left")
            except Exception:
                pass
            child.grid_configure(row=0, column=1, padx=(0, 24), pady=(8, 8), sticky="ew")
        elif child is app.log_box:
            try:
                child.configure(height=128)
            except Exception:
                pass
            child.grid_configure(pady=(6, 8), padx=30, sticky="ew")
        else:
            try:
                child.configure(height=42)
            except Exception:
                pass
            for action_child in child.winfo_children():
                try:
                    if isinstance(action_child, customtkinter.CTkButton):
                        action_child.configure(height=40)
                    elif isinstance(action_child, customtkinter.CTkSwitch):
                        action_child.configure(height=28)
                except Exception:
                    pass
            child.grid_configure(padx=30, pady=0, sticky="ew")

    sidebar_children = app.sidebar_frame.winfo_children()
    if sidebar_children:
        try:
            sidebar_children[0].grid_configure(padx=14, pady=(14, 8), sticky="ew")
        except Exception:
            pass
        try:
            header_children = sidebar_children[0].winfo_children()
            if len(header_children) >= 3:
                brand_image = _get_sidebar_brand_image(app)
                if brand_image is not None:
                    try:
                        header_children[0].configure(text="", image=brand_image, compound="center")
                    except Exception:
                        pass
                header_children[0].grid_configure(padx=(12, 10), pady=14, sticky="w")
                try:
                    header_children[1].configure(font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=17, weight="bold"))
                except Exception:
                    pass
                header_children[1].grid_configure(padx=(0, 14), pady=(12, 2), sticky="sw")
                header_children[2].grid_configure(padx=(0, 14), pady=(0, 8), sticky="nw")
        except Exception:
            pass

    sidebar_icon_images = _get_sidebar_icon_images(app)
    sidebar_button_font = _get_sidebar_button_font(app)
    for nav_name in (
        "btn_nav_wm",
        "btn_nav_rm_wm",
        "btn_nav_cv",
        "btn_nav_audio",
        "btn_nav_zip",
        "btn_nav_pdf",
        "btn_nav_img",
        "btn_nav_meta",
        "btn_nav_file",
    ):
        nav = getattr(app, nav_name, None)
        if nav is not None:
            try:
                nav_spec = SIDEBAR_BUTTON_SPECS.get(nav_name, {})
                nav.configure(
                    text=nav_spec.get("label", nav.cget("text")),
                    image=sidebar_icon_images.get(("nav", nav_spec.get("icon"))),
                    compound="left",
                    border_spacing=12,
                    font=sidebar_button_font,
                    height=40,
                    anchor="w",
                )
                nav.grid_configure(padx=12, pady=3, sticky="ew")
            except Exception:
                pass

    if getattr(app, "btn_help", None) is not None:
        try:
            app.btn_help.grid_remove()
        except Exception:
            pass
        if getattr(app, "btn_help_proxy", None) is None:
            app.btn_help_proxy = customtkinter.CTkButton(
                app.sidebar_frame,
                command=lambda target=app: _show_inline_help(target),
            )
        else:
            app.btn_help_proxy.configure(command=lambda target=app: _show_inline_help(target))
        _style_sidebar_aux_button(
            app.btn_help_proxy,
            SIDEBAR_AUX_BUTTON_SPECS["btn_help_proxy"]["label"],
            sidebar_icon_images.get(("help", SIDEBAR_AUX_BUTTON_SPECS["btn_help_proxy"]["icon"])),
            sidebar_button_font,
            "help",
        )
        app.sidebar_frame.grid_rowconfigure(10, minsize=0, weight=0)
        app.btn_help_proxy.grid(row=10, column=0, padx=12, pady=(7, 3), sticky="ew")
        app.sidebar_frame.grid_rowconfigure(11, minsize=0, weight=0)
    if getattr(app, "btn_donate", None) is not None:
        try:
            app.btn_donate.configure(command=lambda target=app: _show_inline_donate(target))
        except Exception:
            pass
        _style_sidebar_aux_button(
            app.btn_donate,
            SIDEBAR_AUX_BUTTON_SPECS["btn_donate"]["label"],
            sidebar_icon_images.get(("donate", SIDEBAR_AUX_BUTTON_SPECS["btn_donate"]["icon"])),
            sidebar_button_font,
            "donate",
        )
        app.btn_donate.grid_configure(row=11, padx=12, pady=(3, 7), sticky="ew")

    footer_candidates = app.sidebar_frame.winfo_children()
    if footer_candidates:
        app.sidebar_frame.grid_rowconfigure(12, minsize=0, weight=0)
        for footer in footer_candidates:
            if isinstance(footer, customtkinter.CTkLabel):
                try:
                    footer.grid_configure(row=12, padx=12, pady=(10, 8), sticky="ew")
                except Exception:
                    pass
                break

    app._fx_shell_layout_tightened = True


def _tighten_single_tab_layout(app, task_name):
    tab_attr = TAB_LAYOUT_ATTRS.get(task_name)
    if not tab_attr:
        return

    tab = getattr(app, tab_attr, None)
    if tab is None or getattr(tab, "_fx_layout_tightened", False):
        return

    children = tab.winfo_children()
    if children:
        try:
            children[0].pack_configure(padx=20, pady=18)
        except Exception:
            pass

    if task_name == "watermark" and children:
        _tighten_watermark_tab_layout(app, tab)

    if task_name == "pdf" and children:
        pdf_card = children[0]
        _tighten_pdf_tab_layout(pdf_card)

    if task_name == "meta" and children:
        meta_card = children[0]
        _tighten_meta_tab_layout(meta_card)

    if task_name == "audio" and children:
        audio_card = children[0]
        _tighten_audio_tab_layout(app, audio_card)

    if task_name == "zip":
        _patch_zip_tab_texts(app, tab)

    _apply_inline_title_icons(app, tab)

    tab._fx_layout_tightened = True


def _tighten_pdf_tab_layout(pdf_card):
    try:
        pdf_card.pack_configure(padx=18, pady=10)
    except Exception:
        pass

    pdf_sections = pdf_card.winfo_children()
    if len(pdf_sections) < 2:
        return

    try:
        pdf_sections[0].pack_configure(anchor="w", padx=24, pady=(16, 8))
    except Exception:
        pass
    try:
        pdf_sections[1].pack_configure(fill="both", expand=True, padx=24, pady=(0, 12))
    except Exception:
        pass

    content_sections = list(pdf_sections[1].winfo_children())
    if len(content_sections) == 1:
        try:
            nested_sections = list(content_sections[0].winfo_children())
        except Exception:
            nested_sections = []
        if len(nested_sections) >= 2:
            try:
                content_sections[0].pack_configure(fill="both", expand=True, padx=0, pady=0)
            except Exception:
                pass
            content_sections = nested_sections
    if len(content_sections) < 2:
        return

    base_panel, detail_shell = content_sections[0], content_sections[1]
    try:
        base_panel.configure(width=230)
        base_panel.pack_configure(side="left", fill="y", padx=(0, 10))
    except Exception:
        pass
    try:
        detail_shell.pack_configure(side="left", fill="both", expand=True, padx=0, pady=0)
    except Exception:
        pass

    base_children = list(base_panel.winfo_children())
    if base_children:
        try:
            base_children[0].pack_configure(anchor="w", pady=(0, 4))
        except Exception:
            pass

    shared_panel = None
    compact_button_font = customtkinter.CTkFont(size=11)
    for child in base_children[1:]:
        if isinstance(child, customtkinter.CTkButton):
            try:
                child.configure(height=28, font=compact_button_font)
                child.pack_configure(fill="x", pady=(0, 2))
            except Exception:
                pass
        else:
            shared_panel = child

    if shared_panel is not None:
        try:
            shared_panel.pack_configure(fill="x", pady=(4, 0))
        except Exception:
            pass
        shared_children = list(shared_panel.winfo_children())
        for index, child in enumerate(shared_children):
            try:
                if isinstance(child, customtkinter.CTkSwitch):
                    child.configure(height=24)
                    child.pack_configure(anchor="w", pady=(0, 6))
                elif isinstance(child, customtkinter.CTkLabel):
                    child.configure(font=customtkinter.CTkFont(size=10))
                    child.pack_configure(anchor="w", pady=(0, 2))
                elif isinstance(child, customtkinter.CTkEntry):
                    child.configure(height=30)
                    child.pack_configure(fill="x", pady=0)
                elif index == 0:
                    child.pack_configure(pady=(0, 6))
            except Exception:
                pass

    try:
        detail_children = list(detail_shell.winfo_children())
    except Exception:
        detail_children = []
    for panel in detail_children:
        try:
            for child in panel.winfo_children():
                if isinstance(child, customtkinter.CTkLabel):
                    try:
                        child.pack_configure(padx=6)
                    except Exception:
                        pass
        except Exception:
            pass


def _tighten_meta_tab_layout(meta_card):
    try:
        meta_card.pack_configure(padx=18, pady=10)
    except Exception:
        pass

    sections = meta_card.winfo_children()
    if len(sections) < 5:
        return

    try:
        sections[0].pack_configure(anchor="w", padx=24, pady=(8, 2))
    except Exception:
        pass

    spacer = sections[1]
    try:
        if not spacer.winfo_children():
            spacer.pack_forget()
        else:
            spacer.configure(height=8)
            spacer.pack_configure(fill="x", padx=24, pady=0)
    except Exception:
        pass

    try:
        sections[2].pack_configure(fill="x", padx=24, pady=(4, 4))
        for radio in sections[2].winfo_children():
            try:
                radio.configure(height=26)
            except Exception:
                pass
    except Exception:
        pass

    for section in sections[3:5]:
        try:
            section.configure(height=82)
            section.pack_propagate(False)
            section.pack_configure(fill="x", padx=24, pady=(0, 4))
            section_children = list(section.winfo_children())
            if len(section_children) >= 2:
                try:
                    section_children[0].pack_configure(anchor="w", pady=(0, 2))
                except Exception:
                    pass
                try:
                    section_children[1].configure(height=30)
                    section_children[1].pack_configure(fill="x", pady=0)
                except Exception:
                    pass
        except Exception:
            pass


def _tighten_audio_tab_layout(app, audio_card):
    try:
        audio_card.pack_configure(padx=18, pady=(0, 8))
    except Exception:
        pass

    sections = list(audio_card.winfo_children())
    if sections:
        try:
            sections[0].pack_configure(anchor="w", padx=45, pady=(0, 10))
        except Exception:
            pass
    if len(sections) > 1:
        try:
            sections[1].pack_configure(fill="x", padx=45, pady=(0, 12))
        except Exception:
            pass

    settings_frame = _find_audio_settings_frame(app)
    if settings_frame is None:
        return

    for child in settings_frame.winfo_children():
        try:
            if child is getattr(app, "_fx_audio_transcribe_preview_frame", None):
                child.pack_configure(fill="x", pady=(4, 6))
            elif child is getattr(app, "_fx_audio_transcribe_model_hint", None):
                child.pack_configure(fill="x", pady=(2, 0))
            else:
                child.pack_configure(pady=(0, 6))
        except Exception:
            pass

    preview_box = getattr(app, "_fx_audio_transcribe_preview_box", None)
    if preview_box is not None:
        try:
            preview_box.configure(height=150)
        except Exception:
            pass


def _layout_zip_depth_controls(app, tab):
    try:
        card = list(tab.winfo_children())[0]
    except Exception:
        return
    depth_frame = getattr(app, "_fx_zip_depth_frame", None)
    if depth_frame is None:
        return

    try:
        children = list(card.winfo_children())
    except Exception:
        return
    if len(children) < 3:
        return

    title = children[0]
    mode_frame = children[1]
    description = children[2]
    spacer = children[3] if len(children) > 3 and children[3] is not depth_frame else None

    try:
        for child in children:
            try:
                child.pack_forget()
            except Exception:
                pass
            try:
                child.grid_forget()
            except Exception:
                pass
        card.grid_columnconfigure(0, weight=3)
        card.grid_columnconfigure(1, weight=2)
        title.grid(row=0, column=0, columnspan=2, sticky="w", padx=45, pady=(22, 14))
        mode_frame.grid(row=1, column=0, sticky="nsew", padx=(45, 10), pady=(0, 12))
        depth_frame.grid(row=1, column=1, sticky="nsew", padx=(10, 45), pady=(0, 12))
        description.grid(row=2, column=0, columnspan=2, sticky="ew", padx=45, pady=(0, 10))
        if spacer is not None:
            spacer.grid(row=3, column=0, columnspan=2, sticky="ew", padx=45, pady=(0, 4))
        try:
            description.configure(wraplength=860)
        except Exception:
            pass
    except Exception as exc:
        _debug(f"zip_depth_layout_error:{exc}")


def _split_zip_depth_fields(value):
    text = str(value or "").strip()
    if not text:
        return "", ""
    try:
        min_depth, max_depth = normalize_zip_depth_range(text)
    except Exception:
        return "", text
    start_text = str(min_depth) if min_depth > 1 or max_depth is not None else ""
    end_text = "" if max_depth is None else str(max_depth)
    return start_text, end_text


def _set_zip_depth_fields(app, value):
    start_text, end_text = _split_zip_depth_fields(value)
    try:
        if getattr(app, "zip_min_depth_var", None) is not None:
            app.zip_min_depth_var.set(start_text)
        if getattr(app, "zip_max_depth_var", None) is not None:
            app.zip_max_depth_var.set(end_text)
    except Exception:
        pass


def _ensure_zip_depth_controls(app, tab):
    if getattr(app, "_fx_zip_depth_controls_ready", False):
        _layout_zip_depth_controls(app, tab)
        return
    if getattr(app, "zip_min_depth_var", None) is None:
        try:
            app.zip_min_depth_var = tkinter.StringVar(master=app, value="")
        except Exception:
            app.zip_min_depth_var = tkinter.StringVar(value="")
    if getattr(app, "zip_max_depth_var", None) is None:
        try:
            app.zip_max_depth_var = tkinter.StringVar(master=app, value="")
        except Exception:
            app.zip_max_depth_var = tkinter.StringVar(value="")
    else:
        try:
            legacy_value = str(app.zip_max_depth_var.get() or "").strip()
            if legacy_value and ("-" in legacy_value or "~" in legacy_value):
                _set_zip_depth_fields(app, legacy_value)
        except Exception:
            pass
    if getattr(app, "zip_archive_policy_var", None) is None:
        try:
            app.zip_archive_policy_var = tkinter.StringVar(master=app, value=ZIP_ARCHIVE_POLICY_REUSE_LABEL)
        except Exception:
            app.zip_archive_policy_var = tkinter.StringVar(value=ZIP_ARCHIVE_POLICY_REUSE_LABEL)
    else:
        _set_zip_archive_policy(app, _safe_var_get(app, "zip_archive_policy_var", ZIP_ARCHIVE_POLICY_REUSE_LABEL))

    try:
        children = list(tab.winfo_children())
    except Exception:
        children = []
    parent = children[0] if children else tab

    frame = customtkinter.CTkFrame(parent, fg_color="#151A20", corner_radius=12)
    frame.pack(fill="x", padx=45, pady=(0, 12))
    frame.grid_columnconfigure(1, weight=1)
    frame.grid_columnconfigure(3, weight=1)
    customtkinter.CTkLabel(
        frame,
        text="压缩层级范围",
        font=("Microsoft YaHei UI", 13, "bold"),
        text_color="#E8EDF5",
    ).grid(row=0, column=0, sticky="w", padx=(14, 10), pady=(12, 4))
    start_entry = customtkinter.CTkEntry(
        frame,
        textvariable=app.zip_min_depth_var,
        placeholder_text="起始",
        height=34,
        width=72,
    )
    start_entry.grid(row=0, column=1, sticky="ew", padx=(0, 8), pady=(12, 4))
    customtkinter.CTkLabel(
        frame,
        text="-",
        font=("Microsoft YaHei UI", 16, "bold"),
        text_color="#C9D3E1",
    ).grid(row=0, column=2, sticky="ew", padx=(0, 8), pady=(12, 4))
    entry = customtkinter.CTkEntry(
        frame,
        textvariable=app.zip_max_depth_var,
        placeholder_text="留空=不限，4=1-4，2-4=只压第2到4层",
        height=34,
        width=72,
    )
    try:
        entry.configure(placeholder_text="结束")
    except Exception:
        pass
    entry.grid(row=0, column=3, sticky="ew", padx=(0, 14), pady=(12, 4))
    customtkinter.CTkLabel(
        frame,
        text="适用于“全层级递归”和“智能混合”；根目录为第 1 层，填 4 表示 1-4，填 2-4 表示只压第 2 到第 4 层。",
        font=("Microsoft YaHei UI", 11),
        text_color="#AEB7C4",
        wraplength=360,
        justify="left",
    ).grid(row=1, column=0, columnspan=4, sticky="w", padx=14, pady=(0, 12))
    customtkinter.CTkLabel(
        frame,
        text="已有压缩包",
        font=("Microsoft YaHei UI", 12, "bold"),
        text_color="#E8EDF5",
    ).grid(row=2, column=0, sticky="w", padx=(14, 10), pady=(0, 6))
    policy_combo = customtkinter.CTkComboBox(
        frame,
        values=list(ZIP_ARCHIVE_POLICY_LABELS),
        variable=app.zip_archive_policy_var,
        height=34,
        state="readonly",
        font=("Microsoft YaHei UI", 12),
        dropdown_font=("Microsoft YaHei UI", 12),
    )
    policy_combo.grid(row=2, column=1, columnspan=3, sticky="ew", padx=(0, 14), pady=(0, 6))
    customtkinter.CTkLabel(
        frame,
        text="递归/智能模式遇到本次计划输出位置已有同名 zip 时：可复用旧包继续跑，也可删除旧包后重新生成。",
        font=("Microsoft YaHei UI", 11),
        text_color="#AEB7C4",
        wraplength=360,
        justify="left",
    ).grid(row=3, column=0, columnspan=4, sticky="w", padx=14, pady=(0, 12))
    try:
        for child in frame.winfo_children():
            if isinstance(child, customtkinter.CTkLabel):
                label_text = str(child.cget("text") or "")
                if "1-4" in label_text or "2-4" in label_text:
                    child.configure(
                        text="适用于“全层级递归”和“智能混合”；根目录为第 1 层。留空表示不限；起始填 2、结束填 4 表示只压第 2 到第 4 层。"
                    )
    except Exception:
        pass
    app.zip_min_depth_entry = start_entry
    app.zip_max_depth_entry = entry
    app.zip_archive_policy_combo = policy_combo
    app._fx_zip_depth_frame = frame
    app._fx_zip_depth_controls_ready = True
    _install_zip_last_settings_memory(app)
    _layout_zip_depth_controls(app, tab)


def _patch_zip_tab_texts(app, tab):
    try:
        stack = list(tab.winfo_children())
    except Exception:
        stack = []
    while stack:
        widget = stack.pop()
        try:
            text_value = widget.cget("text")
        except Exception:
            text_value = None
        if isinstance(text_value, str):
            new_label = ZIP_MODE_LABEL_TEXTS.get(text_value)
            if new_label is not None and text_value != new_label:
                try:
                    widget.configure(text=new_label)
                except Exception:
                    pass
            elif text_value.startswith("功能说明：") and text_value != ZIP_MODE_DESCRIPTION_TEXT:
                try:
                    widget.configure(text=ZIP_MODE_DESCRIPTION_TEXT, justify="left")
                except Exception:
                    pass
        try:
            stack.extend(widget.winfo_children())
        except Exception:
            pass
    _ensure_zip_depth_controls(app, tab)


def _refresh_visible_tab_layout(app, task_name):
    tab_attr = TAB_LAYOUT_ATTRS.get(task_name)
    if not tab_attr:
        return

    tab = getattr(app, tab_attr, None)
    if tab is None:
        return
    children = list(tab.winfo_children())
    if not children:
        return

    card = children[0]
    if task_name == "pdf":
        _tighten_pdf_tab_layout(card)
    elif task_name == "meta":
        _tighten_meta_tab_layout(card)
    elif task_name == "zip":
        _patch_zip_tab_texts(app, tab)
    _apply_inline_title_icons(app, tab)


def _widget_direct_child_under(root_widget, target_widget):
    if root_widget is None or target_widget is None:
        return None
    current = target_widget
    previous = None
    while current is not None and current is not root_widget:
        previous = current
        current = getattr(current, "master", None)
    return previous if current is root_widget else None


def _find_watermark_text_panel(app, tab):
    text_widget = getattr(app, "wm_text", None)
    direct_child = _widget_direct_child_under(tab, text_widget)
    if direct_child is not None:
        return direct_child
    for child in list(tab.winfo_children()):
        stack = list(child.winfo_children())
        while stack:
            widget = stack.pop()
            if widget is text_widget:
                return child
            try:
                stack.extend(widget.winfo_children())
            except Exception:
                pass
    children = list(tab.winfo_children())
    return children[0] if children else tab


def _find_watermark_settings_panel(app, tab, left_panel):
    for attr_name in ("slider_size", "slider_opacity", "slider_angle", "wm_font_combo", "wm_font_dropdown"):
        widget = getattr(app, attr_name, None)
        direct_child = _widget_direct_child_under(tab, widget)
        if direct_child is not None and direct_child is not left_panel:
            return direct_child
    for child in list(tab.winfo_children()):
        try:
            grid_info = child.grid_info()
        except Exception:
            grid_info = {}
        if child is not left_panel and str(grid_info.get("column", "")) == "1":
            return child
    for child in list(tab.winfo_children()):
        if child is not left_panel:
            return child
    return None


def _tighten_watermark_tab_layout(app, tab):
    try:
        tab.grid_rowconfigure(0, weight=1, minsize=0)
        tab.grid_columnconfigure(0, weight=3, minsize=0)
        tab.grid_columnconfigure(1, weight=2, minsize=0)
    except Exception:
        pass

    left_panel = _find_watermark_text_panel(app, tab)
    right_panel = _find_watermark_settings_panel(app, tab, left_panel)
    if left_panel is None or right_panel is None:
        return
    _ensure_active_watermark_filename_rule_controls(app, right_panel)

    try:
        left_panel.grid_configure(row=0, column=0, padx=(0, 14), pady=0, sticky="nsew")
        right_panel.grid_configure(row=0, column=1, padx=0, pady=0, sticky="nsew")
        right_panel.grid_configure(padx=0, pady=0, sticky="nsew")
        left_panel.configure(height=520)
        right_panel.configure(height=520)
        left_panel.grid_rowconfigure(1, weight=1)
        right_panel.grid_rowconfigure(10, weight=1)
    except Exception:
        pass

    left_children = list(left_panel.winfo_children())
    preview_frame = next((child for child in left_children if getattr(child, "_fx_wm_color_preview_controls", False)), None)
    text_area = next(
        (
            child
            for child in left_children[1:]
            if child is not preview_frame and not getattr(child, "_fx_wm_color_preview_controls", False)
        ),
        None,
    )
    if left_children:
        try:
            left_children[0].pack_configure(anchor="w", padx=24, pady=(18, 8))
        except Exception:
            pass
    if preview_frame is not None:
        try:
            preview_frame.configure(height=132)
            preview_frame.pack_propagate(False)
            preview_frame.pack_configure(fill="x", padx=24, pady=(0, 8))
        except Exception:
            pass
    if text_area is not None:
        try:
            text_area.configure(height=292 if preview_frame is not None else 390)
            text_area.pack_configure(fill="both", expand=True, padx=24, pady=(0, 18))
        except Exception:
            pass

    right_children = list(right_panel.winfo_children())
    filename_rule_controls = next(
        (child for child in right_children if getattr(child, "_fx_wm_filename_rule_controls", False)),
        None,
    )
    for child in right_children:
        try:
            current = child.pack_info()
        except Exception:
            continue
        side = current.get("side", None)
        fill = current.get("fill", None)
        expand = bool(int(current.get("expand", 0))) if str(current.get("expand", "0")).isdigit() else current.get("expand", False)
        padx = current.get("padx", 0)
        try:
            child.pack_configure(side=side, fill=fill, expand=expand, padx=padx, pady=(0, 3))
        except Exception:
            pass

    if right_children:
        try:
            right_children[0].configure(height=28)
            right_children[0].pack_configure(anchor="w", padx=24, pady=(6, 1))
        except Exception:
            pass
    for index in (1, 2):
        if index < len(right_children):
            try:
                right_children[index].configure(height=28)
                right_children[index].pack_propagate(False)
                right_children[index].pack_configure(fill="x", padx=24, pady=(0, 1))
                for radio in right_children[index].winfo_children():
                    try:
                        radio.configure(height=28)
                    except Exception:
                        pass
            except Exception:
                pass
    for index in (3, 5, 6, 7):
        if index < len(right_children):
            try:
                if right_children[index] is filename_rule_controls:
                    continue
                right_children[index].configure(height=30)
                right_children[index].pack_configure(anchor="w", padx=24, pady=(0, 1))
            except Exception:
                pass
    if len(right_children) > 4:
        try:
            if right_children[4] is not filename_rule_controls:
                right_children[4].configure(height=38)
                right_children[4].pack_configure(fill="x", padx=24, pady=(0, 2))
        except Exception:
            pass
    for index in (8, 9, 10):
        if index < len(right_children):
            try:
                if right_children[index] is filename_rule_controls:
                    continue
                right_children[index].configure(height=36)
                right_children[index].pack_propagate(False)
                right_children[index].pack_configure(fill="x", padx=24, pady=(0, 1))
                slider_parts = list(right_children[index].winfo_children())
                if slider_parts:
                    try:
                        slider_parts[0].configure(height=18)
                        slider_parts[0].pack_propagate(False)
                        slider_parts[0].pack_configure(fill="x", pady=(0, 0))
                        for item in slider_parts[0].winfo_children():
                            try:
                                item.configure(height=18, font=customtkinter.CTkFont(size=11))
                            except Exception:
                                pass
                    except Exception:
                        pass
                if len(slider_parts) > 1:
                    try:
                        slider_parts[1].configure(height=16)
                        slider_parts[1].pack_configure(fill="x", pady=(0, 0))
                    except Exception:
                        pass
            except Exception:
                pass
    if filename_rule_controls is not None:
        try:
            filename_rule_controls.configure(height=82)
            filename_rule_controls.pack_propagate(False)
            filename_rule_controls.pack_configure(fill="x", padx=24, pady=(0, 3))
        except Exception:
            pass
    elif len(right_children) > 11:
        try:
            controls_height = 82 if getattr(right_children[11], "_fx_wm_filename_rule_controls", False) else 30
            right_children[11].configure(height=controls_height)
            right_children[11].pack_propagate(False)
            right_children[11].pack_configure(fill="x", padx=24, pady=(0, 2))
            for child in right_children[11].winfo_children():
                try:
                    if getattr(right_children[11], "_fx_wm_filename_rule_controls", False):
                        continue
                    child.configure(height=30)
                except Exception:
                    pass
        except Exception:
            pass


def _tighten_layout(app, task_name=None):
    try:
        _apply_shell_layout_tightening(app)
        if task_name is None:
            for name in TAB_LAYOUT_ATTRS:
                _tighten_single_tab_layout(app, name)
        else:
            _tighten_single_tab_layout(app, task_name)
    except Exception as exc:
        _debug(f"layout:tighten_error:{exc}")


def _load_runtime_namespace():
    started_at = time.perf_counter()
    _debug("load_runtime:start")
    if not RUNTIME_BIN.exists():
        _debug(f"load_runtime:missing:{RUNTIME_BIN}")
        raise FileNotFoundError(f"Missing runtime payload: {RUNTIME_BIN}")
    try:
        code = marshal.loads(RUNTIME_BIN.read_bytes())
        namespace = {
            "__name__": "fengxi_embedded_runtime",
            "__file__": __file__,
            "__package__": None,
        }
        lazy_import_state = _install_runtime_lazy_imports()
        try:
            exec(code, namespace)
        finally:
            _restore_runtime_lazy_imports(lazy_import_state)
    except Exception:
        _record_performance("runtime_load", started_at=started_at, details={"status": "error"})
        raise
    else:
        _record_performance(
            "runtime_load",
            started_at=started_at,
            details={"status": "success", "namespace_size": len(namespace)},
        )
        _debug("load_runtime:done")
        return namespace


def _locate_ffmpeg():
    env_bin = os.environ.get("FFMPEG_BINARY")
    if env_bin and os.path.exists(env_bin):
        return env_bin
    try:
        import imageio_ffmpeg
        ffmpeg_bin = imageio_ffmpeg.get_ffmpeg_exe()
        if ffmpeg_bin and os.path.exists(ffmpeg_bin):
            os.environ["FFMPEG_BINARY"] = ffmpeg_bin
            return ffmpeg_bin
    except Exception:
        pass
    return None


def _ffmpeg_convert(src, dst, target_fmt, bitrate="192k"):
    ffmpeg_bin = _locate_ffmpeg()
    if not ffmpeg_bin:
        return "MISSING_LIB:ffmpeg backend not available"

    codec_map = {
        "mp3": ["-acodec", "libmp3lame"],
        "wav": ["-acodec", "pcm_s16le"],
        "flac": ["-acodec", "flac"],
        "m4a": ["-acodec", "aac"],
        "aac": ["-acodec", "aac"],
        "ogg": ["-acodec", "libvorbis"],
    }

    target_fmt = target_fmt.lower()
    cmd = [ffmpeg_bin, "-y", "-i", src, "-vn"]
    if bitrate and target_fmt not in ("wav", "flac"):
        cmd.extend(["-b:a", bitrate])
    cmd.extend(codec_map.get(target_fmt, []))
    cmd.append(dst)

    try:
        result = subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8", errors="ignore")
    except Exception as exc:
        return f"ERROR:{exc}"

    if result.returncode != 0:
        detail = (result.stderr or result.stdout or "").strip()
        return f"ERROR:{detail or 'ffmpeg convert failed'}"
    if not os.path.exists(dst) or os.path.getsize(dst) <= 0:
        return "ERROR:output file not created"
    return "SUCCESS"


_debug("bootstrap:start")
_ns = _load_runtime_namespace()
_ns["convert_audio_format"] = _ffmpeg_convert
_ns["HAS_MOVIEPY"] = bool(_locate_ffmpeg()) or bool(_ns.get("HAS_MOVIEPY"))
_debug(f"bootstrap:patched:has_moviepy={_ns['HAS_MOVIEPY']}")

for _key, _value in _ns.items():
    if _key in {"__name__", "__file__", "__package__", "__builtins__"}:
        continue
    globals()[_key] = _value
_debug("bootstrap:globals_loaded")


def create_watermark_packet(content, font_name, font_size, opacity, angle, color=None):
    return _watermark_core_create_watermark_packet(
        content,
        font_name,
        font_size,
        opacity,
        angle,
        font_path_resolver=get_font_path_by_name,
        color=color,
    )


def add_watermark_to_pdf(
    src,
    dst,
    watermark_packet,
    page_range="all",
    check_text=None,
    force_mode=False,
    copy_guard=False,
    copy_guard_strength="standard",
):
    return _watermark_core_add_watermark_to_pdf(
        src,
        dst,
        watermark_packet,
        page_range=_watermark_core_normalize_page_range(page_range),
        check_text=check_text,
        force_mode=force_mode,
        copy_guard=copy_guard,
        copy_guard_strength=copy_guard_strength,
    )


def add_watermark_to_word(
    word_app,
    src,
    dst,
    text,
    raw_font_name,
    font_size,
    opacity,
    angle,
    page_range="all",
    force_mode=False,
    color=None,
    copy_guard=False,
    copy_guard_strength="standard",
):
    return _watermark_core_add_watermark_to_word(
        word_app,
        src,
        dst,
        text,
        raw_font_name,
        font_size,
        opacity,
        angle,
        page_range=_watermark_core_normalize_page_range(page_range),
        force_mode=force_mode,
        word_font_resolver=get_word_compatible_font_name,
        com_context_factory=_DisableWin32ComGenCache,
        color=color,
        copy_guard=copy_guard,
        copy_guard_strength=copy_guard_strength,
    )


globals()["create_watermark_packet"] = create_watermark_packet
globals()["add_watermark_to_pdf"] = add_watermark_to_pdf
globals()["add_watermark_to_word"] = add_watermark_to_word
_ns["create_watermark_packet"] = create_watermark_packet
_ns["add_watermark_to_pdf"] = add_watermark_to_pdf
_ns["add_watermark_to_word"] = add_watermark_to_word
_debug("bootstrap:watermark_core_wrapped")


def modify_file_timestamp(src, dst, timestamp):
    return _meta_core_modify_file_timestamp(src, dst, timestamp)


def modify_office_meta(app, src, dst, author_name, app_type="word"):
    return _meta_core_modify_office_meta(app, src, dst, author_name, app_type=app_type)


def modify_pdf_author(src, dst, author_name):
    return _meta_core_modify_pdf_author(src, dst, author_name)


globals()["modify_file_timestamp"] = modify_file_timestamp
globals()["modify_office_meta"] = modify_office_meta
globals()["modify_pdf_author"] = modify_pdf_author
_ns["modify_file_timestamp"] = modify_file_timestamp
_ns["modify_office_meta"] = modify_office_meta
_ns["modify_pdf_author"] = modify_pdf_author
_debug("bootstrap:meta_core_wrapped")


def _patch_task_routing():
    try:
        original = FengxiToolboxApp.resolve_task_config
    except Exception as exc:
        _debug(f"patch_task_routing:missing:{exc}")
        return
    if getattr(original, "__fx_task_routing_patch__", False):
        return

    def patched(self, task_type, all_files):
        args, force_single_thread, extra = original(self, task_type, all_files)
        try:
            if task_type == "remove_wm":
                force_single_thread = True
                _debug("patch_task_routing:force_single_thread:remove_wm")
            elif task_type == "pdf" and self.pdf_mode_var.get() == "merge":
                force_single_thread = True
                _debug("patch_task_routing:force_single_thread:pdf_merge")
            elif task_type == "file" and self.file_mode_var.get() == "dedup":
                force_single_thread = True
                _debug("patch_task_routing:force_single_thread:file_dedup")
        except Exception as exc:
            _debug(f"patch_task_routing:error:{exc}")
        return args, force_single_thread, extra

    patched.__fx_task_routing_patch__ = True
    FengxiToolboxApp.resolve_task_config = patched
    _debug("patch_task_routing:installed")


_patch_task_routing()


def _choose_input_path_interactive(app):
    current_value = ""
    try:
        current_value = app.input_path.get()
    except Exception:
        current_value = ""

    current_path = _normalize_input_path_value(current_value)
    last_input_dir = getattr(app, "_fx_last_input_dir", "")

    initial_dir = current_path
    if initial_dir and os.path.isfile(initial_dir):
        initial_dir = os.path.dirname(initial_dir)
    if not initial_dir:
        initial_dir = _normalize_input_path_value(last_input_dir)
    if initial_dir and not os.path.isdir(initial_dir):
        initial_dir = os.path.dirname(initial_dir)
    if not initial_dir:
        home_dir = _normalize_input_path_value(os.path.expanduser("~"))
        initial_dir = home_dir if os.path.isdir(home_dir) else os.getcwd()

    selected_path = _choose_input_path_via_shell_dialog(app, initial_dir)
    if not selected_path:
        _debug("shell_picker:cancelled_or_failed")
        return ""

    normalized_path = _normalize_input_path_value(selected_path)
    if not normalized_path:
        return ""

    app._fx_input_pick_mode = "file" if os.path.isfile(normalized_path) else "folder"
    app._fx_last_input_dir = normalized_path if os.path.isdir(normalized_path) else os.path.dirname(normalized_path)
    try:
        app.input_path.set(normalized_path)
    except Exception:
        pass
    try:
        app.log(
            f"🎯 [目标锁定] 已载入{'文件' if os.path.isfile(normalized_path) else '文件夹'}: {normalized_path}"
        )
    except Exception:
        pass
    return normalized_path


def _patch_drag_drop_input_support():
    try:
        original = FengxiToolboxApp.accept_drag_drop
    except Exception as exc:
        _debug(f"patch_drag_drop:missing:{exc}")
        return

    if getattr(original, "__fx_drag_drop_patch__", False):
        return

    def patched(self, filenames):
        if not filenames:
            return

        raw_items = list(filenames)
        normalized_items = []
        for item in raw_items:
            normalized = _normalize_input_path_value(item)
            if normalized:
                normalized_items.append(normalized)

        if not normalized_items:
            try:
                self.log("❌ [错误] 拖拽目标路径无效")
            except Exception:
                pass
            return

        selected_path = normalized_items[0]
        is_file = os.path.isfile(selected_path)
        is_dir = os.path.isdir(selected_path)
        if not (is_file or is_dir):
            try:
                self.log(f"❌ [错误] 拖拽目标不存在: {selected_path}")
            except Exception:
                pass
            return

        self._fx_input_pick_mode = "file" if is_file else "folder"
        self._fx_last_input_dir = os.path.dirname(selected_path) if is_file else selected_path

        try:
            self.input_path.set(selected_path)
        except Exception:
            pass

        if len(normalized_items) > 1:
            try:
                self.log(f"⚠️ [拖拽] 检测到 {len(normalized_items)} 个项目，当前仅使用第 1 个: {selected_path}")
            except Exception:
                pass

        try:
            self.log(f"🎯 [目标锁定] 已拖入{'文件' if is_file else '文件夹'}: {selected_path}")
        except Exception:
            pass

    patched.__fx_drag_drop_patch__ = True
    FengxiToolboxApp.accept_drag_drop = patched
    _debug("patch_drag_drop:installed")


_patch_drag_drop_input_support()


def _run_single_file_zip_via_staging(app, input_file, original_run_process):
    return _run_zip_task_with_core(app, input_file)


def _patch_single_input_support():
    try:
        original_select_folder = FengxiToolboxApp.select_folder
        original_collect_input_files = FengxiToolboxApp.collect_input_files
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_single_input:missing:{exc}")
        return

    if getattr(original_run_process, "__fx_single_input_patch__", False):
        return

    def patched_select_folder(self):
        _choose_input_path_interactive(self)

    def patched_collect_input_files(self, input_folder, task_type):
        normalized_input = _normalize_input_path_value(input_folder)
        single_target = _normalize_input_path_value(getattr(self, "_fx_single_input_target", ""))
        target = ""
        if normalized_input and os.path.isfile(normalized_input):
            target = normalized_input
        elif single_target and os.path.isfile(single_target):
            parent_dir = os.path.dirname(single_target)
            if normalized_input and normalized_input == parent_dir:
                target = single_target
        if target:
            log_key = (task_type, target)
            if getattr(self, "_fx_single_input_logged", None) != log_key:
                try:
                    self.log(f"📄 [单文件模式] 仅处理: {os.path.basename(target)}")
                except Exception:
                    pass
                self._fx_single_input_logged = log_key
            return [target]
        return original_collect_input_files(self, input_folder, task_type)

    def patched_run_process(self, input_folder, task_type):
        normalized_input = _normalize_input_path_value(input_folder)
        prev_single_target = getattr(self, "_fx_single_input_target", None)
        prev_single_log = getattr(self, "_fx_single_input_logged", None)
        enable_multithread_var = getattr(self, "enable_multithread", None)
        prev_multithread = None
        try:
            if normalized_input and os.path.isfile(normalized_input):
                self._fx_last_input_dir = os.path.dirname(normalized_input)
                self._fx_input_pick_mode = "file"
                self._fx_single_input_logged = None
                if enable_multithread_var is not None:
                    try:
                        prev_multithread = bool(enable_multithread_var.get())
                        enable_multithread_var.set(False)
                        _debug(f"patch_single_input:force_single_thread:{task_type}")
                    except Exception as exc:
                        _debug(f"patch_single_input:toggle_multithread_error:{exc}")
                if task_type == "zip":
                    return _run_single_file_zip_via_staging(self, normalized_input, original_run_process)
                self._fx_single_input_target = normalized_input
                return original_run_process(self, os.path.dirname(normalized_input), task_type)
            if normalized_input and os.path.isdir(normalized_input):
                self._fx_last_input_dir = normalized_input
                self._fx_input_pick_mode = "folder"
            self._fx_single_input_target = None
            self._fx_single_input_logged = None
            return original_run_process(self, input_folder, task_type)
        finally:
            if prev_multithread is not None and enable_multithread_var is not None:
                try:
                    enable_multithread_var.set(prev_multithread)
                except Exception:
                    pass
            self._fx_single_input_target = prev_single_target
            self._fx_single_input_logged = prev_single_log

    def patched_on_start_click(self):
        if getattr(self, "is_running", False):
            return
        if getattr(self, "current_task", None) in {"help", "donate"}:
            self.log("ℹ️ [说明页面] 当前页面仅用于查看内容，请先切换到具体功能页。")
            return
        selected_input = _normalize_input_path_value(self.input_path.get())
        if not selected_input:
            self.log("❌ [错误] 目标路径未定义")
            tkinter.messagebox.showwarning("提示", "请先选择文件或文件夹！")
            return
        if self.current_task == "watermark" and not self.font_list:
            tkinter.messagebox.showerror("资源错误", "fonts 文件夹为空！\n无法添加水印。")
            return
        if self.current_task == "audio" and not HAS_MOVIEPY:
            tkinter.messagebox.showerror("依赖缺失", "未安装 moviepy 库，无法进行音频处理。\n请运行: pip install moviepy")
            return

        self.is_running = True
        self.stop_event = False
        self.progress_bar.set(0)
        _set_progress_status(self, stage="准备中", fraction=0.0)
        self.btn_run.configure(state="disabled", text="⚡ 执行中...", fg_color="#455A64")
        self.btn_stop.configure(state="normal")
        _set_watermark_checkpoint_action_state(self, "disabled")
        threading.Thread(target=self.run_process, args=(selected_input, self.current_task), daemon=True).start()

    patched_select_folder.__fx_single_input_patch__ = True
    patched_collect_input_files.__fx_single_input_patch__ = True
    patched_run_process.__fx_single_input_patch__ = True
    patched_on_start_click.__fx_single_input_patch__ = True
    FengxiToolboxApp.select_folder = patched_select_folder
    FengxiToolboxApp.collect_input_files = patched_collect_input_files
    FengxiToolboxApp.run_process = patched_run_process
    FengxiToolboxApp.on_start_click = patched_on_start_click
    _debug("patch_single_input:installed")


_patch_single_input_support()


def _patch_file_manager_core():
    try:
        original_process_single = FengxiToolboxApp.process_single_file
    except Exception as exc:
        _debug(f"patch_file_manager_core:missing:{exc}")
        return

    if getattr(original_process_single, "__fx_file_manager_core_patch__", False):
        return

    def patched_process_single_file(self, src, input_folder, output_folder, task_type, args, failed_list):
        if task_type == "file":
            values = list(args or [])
            mode = str(values[0] if values else "").strip().lower()
            if mode == "rename":
                try:
                    _file_core_apply_rename_to_file(
                        src,
                        input_folder,
                        output_folder,
                        args,
                        copy_file_safe=copy_file_safe,
                        log=getattr(self, "log", None),
                    )
                    return None
                except Exception as exc:
                    try:
                        failed_list.append(f"{src}: {exc}")
                    except Exception:
                        pass
                    return None
        return original_process_single(self, src, input_folder, output_folder, task_type, args, failed_list)

    patched_process_single_file.__fx_file_manager_core_patch__ = True
    FengxiToolboxApp.process_single_file = patched_process_single_file
    _debug("patch_file_manager_core:installed")


_patch_file_manager_core()


def _patch_meta_core():
    try:
        original_process_single = FengxiToolboxApp.process_single_file
    except Exception as exc:
        _debug(f"patch_meta_core:missing:{exc}")
        return

    if getattr(original_process_single, "__fx_meta_core_patch__", False):
        return

    def patched_process_single_file(self, src, input_folder, output_folder, task_type, args, failed_list):
        if task_type == "meta":
            return _meta_core_process_meta_file(
                src,
                input_folder,
                output_folder,
                args,
                failed_list,
                copy_file_safe=copy_file_safe,
                log=getattr(self, "log", None),
            )
        return original_process_single(self, src, input_folder, output_folder, task_type, args, failed_list)

    patched_process_single_file.__fx_meta_core_patch__ = True
    patched_process_single_file.__fx_file_manager_core_patch__ = bool(
        getattr(original_process_single, "__fx_file_manager_core_patch__", False)
    )
    patched_process_single_file.__wrapped__ = original_process_single
    FengxiToolboxApp.process_single_file = patched_process_single_file
    _debug("patch_meta_core:installed")


_patch_meta_core()


def _patch_convert_file_adapter():
    try:
        original_process_single = FengxiToolboxApp.process_single_file
    except Exception as exc:
        _debug(f"patch_convert_file_adapter:missing:{exc}")
        return

    if getattr(original_process_single, "__fx_convert_file_adapter_patch__", False):
        return

    def patched_process_single_file(self, src, input_folder, output_folder, task_type, args, failed_list):
        if task_type == "convert":
            values = list(args or [])
            mode = _convert_core_normalize_mode(values[0] if values else _get_convert_mode(self))
            if mode != "imgs2pdf":
                skip_complex = bool(values[1]) if len(values) > 1 else False
                result = process_convert_file(
                    src,
                    input_folder,
                    output_folder,
                    mode,
                    ConvertFileContext(
                        word_app=values[2] if len(values) > 2 else None,
                        ppt_app=values[3] if len(values) > 3 else None,
                        skip_complex=skip_complex,
                        convert_doc_to_pdf=lambda word_app, src_path, dst_path: _convert_doc_to_pdf_safely(word_app, src_path, dst_path),
                        convert_pdf_to_word=convert_pdf_to_word,
                        convert_ppt_to_pdf=lambda ppt_app, src_path, dst_path: _convert_ppt_to_pdf_safely(ppt_app, src_path, dst_path),
                        convert_pdf_to_ppt=lambda ppt_app, src_path, dst_path: _convert_pdf_to_ppt_safely(ppt_app, src_path, dst_path),
                        convert_txt_to_word=convert_txt_to_word_file,
                        convert_md_to_pdf=convert_md_to_pdf_file,
                        convert_pdf_to_md=convert_pdf_to_md_file,
                        check_pdf_complexity=check_pdf_complexity,
                        copy_file_safe=copy_file_safe,
                        log=getattr(self, "log", None),
                    ),
                )
                if not result.get("ok"):
                    try:
                        failed_list.append(str(src))
                    except Exception:
                        pass
                return None
        return original_process_single(self, src, input_folder, output_folder, task_type, args, failed_list)

    patched_process_single_file.__fx_convert_file_adapter_patch__ = True
    patched_process_single_file.__fx_meta_core_patch__ = bool(
        getattr(original_process_single, "__fx_meta_core_patch__", False)
    )
    patched_process_single_file.__fx_file_manager_core_patch__ = bool(
        getattr(original_process_single, "__fx_file_manager_core_patch__", False)
    )
    patched_process_single_file.__wrapped__ = original_process_single
    FengxiToolboxApp.process_single_file = patched_process_single_file
    _debug("patch_convert_file_adapter:installed")


_patch_convert_file_adapter()


def _build_generic_resume_output_path(src, input_folder, output_folder, task_type, args):
    values = list(args or [])
    task = str(task_type or "").strip().lower()
    try:
        if task == "file":
            mode = str(values[0] if values else "").strip().lower()
            if mode == "rename":
                spec = _file_core_normalize_rename_spec(values)
                return _file_core_plan_renamed_output_path(src, output_folder, spec)

        if task == "meta":
            return _meta_core_build_output_path(src, input_folder, output_folder)[1]

        if task == "convert":
            mode = _convert_core_normalize_mode(values[0] if values else "word2pdf")
            if mode != "imgs2pdf":
                return _convert_core_plan_output_path(src, input_folder, output_folder, mode)

        if task == "image":
            mode = str(values[0] if values else "").strip().lower()
            if mode in {"convert", "compress"}:
                target_format = str(values[2] if len(values) > 2 else Path(src).suffix.lstrip(".") or "jpg").strip().lower()
                target_format = target_format.lstrip(".") or Path(src).suffix.lstrip(".") or "jpg"
                try:
                    rel = os.path.relpath(str(src), str(input_folder))
                except Exception:
                    rel = os.path.basename(str(src))
                if rel.startswith(".."):
                    rel = os.path.basename(str(src))
                return str((Path(output_folder) / rel).with_suffix(f".{target_format}"))

        if task == "pdf":
            mode = str(values[0] if values else "").strip().lower()
            if mode == "encrypt":
                try:
                    rel = os.path.relpath(str(src), str(input_folder))
                except Exception:
                    rel = os.path.basename(str(src))
                if rel.startswith(".."):
                    rel = os.path.basename(str(src))
                return str(Path(output_folder) / rel)
    except Exception:
        return ""
    return ""


def _generic_resume_outputs_complete(src, input_folder, output_folder, task_type, args):
    values = list(args or [])
    task = str(task_type or "").strip().lower()
    try:
        if task == "pdf":
            mode = str(values[0] if values else "").strip().lower()
            if mode == "split":
                split_dir = Path(output_folder) / Path(src).stem
                if not split_dir.is_dir():
                    return False, ""
                try:
                    page_count = len(pypdf.PdfReader(str(src)).pages)
                except Exception:
                    page_count = 0
                outputs = [item for item in split_dir.glob("*.pdf") if is_nonempty_file(item)]
                return page_count > 0 and len(outputs) >= page_count, str(split_dir)
    except Exception:
        return False, ""

    output_path = _build_generic_resume_output_path(src, input_folder, output_folder, task_type, args)
    if not output_path:
        return False, ""
    try:
        same_file_output = Path(output_path).resolve() == Path(src).resolve()
    except Exception:
        same_file_output = False
    if same_file_output:
        return False, output_path
    return is_nonempty_file(output_path), output_path


def _patch_generic_process_resume():
    try:
        original_process_single = FengxiToolboxApp.process_single_file
    except Exception as exc:
        _debug(f"patch_process_resume:missing:{exc}")
        return

    if getattr(original_process_single, "__fx_process_resume_patch__", False):
        return

    def patched_process_single_file(self, src, input_folder, output_folder, task_type, args, failed_list):
        complete, output_path = _generic_resume_outputs_complete(src, input_folder, output_folder, task_type, args)
        if complete:
            try:
                log = getattr(self, "log", None)
                if callable(log):
                    log(format_resume_skip_message("断点续跑", src, output_path))
            except Exception:
                pass
            return None
        return original_process_single(self, src, input_folder, output_folder, task_type, args, failed_list)

    patched_process_single_file.__fx_process_resume_patch__ = True
    patched_process_single_file.__fx_convert_file_adapter_patch__ = bool(
        getattr(original_process_single, "__fx_convert_file_adapter_patch__", False)
    )
    patched_process_single_file.__fx_meta_core_patch__ = bool(
        getattr(original_process_single, "__fx_meta_core_patch__", False)
    )
    patched_process_single_file.__fx_file_manager_core_patch__ = bool(
        getattr(original_process_single, "__fx_file_manager_core_patch__", False)
    )
    patched_process_single_file.__wrapped__ = original_process_single
    FengxiToolboxApp.process_single_file = patched_process_single_file
    _debug("patch_process_resume:installed")


_patch_generic_process_resume()


def _run_file_dedup_task(app, input_folder):
    normalized_input = _normalize_input_path_value(input_folder)
    result = _run_file_dedup_task_core(
        app,
        normalized_input,
        collect_input_files=getattr(app, "collect_input_files", None),
        progress_tracker=_get_active_progress_tracker(app),
        progress_bar=getattr(app, "progress_bar", None),
        get_last_task_result=_get_last_task_result,
        start_task_result=_start_task_result,
        set_task_result_output_strategy=_set_task_result_output_strategy,
        set_task_result_output_root=_set_task_result_output_root,
        add_task_result_output=_add_task_result_output,
        set_task_result_counts=_set_task_result_counts,
        set_task_result_finished=_set_task_result_finished,
        normalize_input_path=_normalize_input_path_value,
        get_task_output_strategy=_get_task_output_strategy,
        clamp_progress_value=_clamp_progress_value,
        set_progress_status=_set_progress_status,
        log=getattr(app, "log", None),
    )
    return result


def _patch_file_dedup_core_task():
    try:
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_file_dedup_core:missing:{exc}")
        return

    if getattr(original_run_process, "__fx_file_dedup_core_patch__", False):
        return

    def patched_run_process(self, input_folder, task_type):
        if task_type == "file":
            try:
                mode = str(self.file_mode_var.get() if getattr(self, "file_mode_var", None) is not None else "")
            except Exception:
                mode = ""
            if mode.strip().lower() == "dedup":
                try:
                    _run_file_dedup_task(self, input_folder)
                except Exception as exc:
                    self.log(f"🔥 [严重错误] {exc}")
                finally:
                    self.reset_ui()
                return
        return original_run_process(self, input_folder, task_type)

    patched_run_process.__fx_file_dedup_core_patch__ = True
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_file_dedup_core:installed")


_patch_file_dedup_core_task()


def _safe_float(value, default=0.0):
    try:
        return float(value)
    except Exception:
        return default


def _safe_getattr(obj, attr_name, default=None):
    try:
        return getattr(obj, attr_name)
    except Exception:
        return default


def _shape_rotation_distance(rotation):
    angle = _safe_float(rotation) % 360.0
    candidates = (0.0, 90.0, 180.0, 270.0, 360.0)
    return min(abs(angle - candidate) for candidate in candidates)


def _coerce_remove_wm_mode(value):
    normalized = str(value or "").strip()
    if normalized in REMOVE_WM_MODE_VALUES:
        return normalized
    mapped = REMOVE_WM_MODE_LABEL_TO_VALUE.get(normalized)
    if mapped in REMOVE_WM_MODE_VALUES:
        return mapped
    return REMOVE_WM_MODE_DEFAULT


def _get_remove_wm_mode_label(value):
    return REMOVE_WM_MODE_VALUE_TO_LABEL.get(
        _coerce_remove_wm_mode(value),
        REMOVE_WM_MODE_VALUE_TO_LABEL[REMOVE_WM_MODE_DEFAULT],
    )


def _get_remove_wm_mode_hint(value):
    return REMOVE_WM_MODE_HINTS.get(_coerce_remove_wm_mode(value), REMOVE_WM_MODE_HINTS[REMOVE_WM_MODE_DEFAULT])


def _get_effective_remove_wm_mode(mode=None):
    if mode is not None:
        return _coerce_remove_wm_mode(mode)
    runtime_mode = getattr(_REMOVE_WM_MODE_CONTEXT, "mode", None)
    return _coerce_remove_wm_mode(runtime_mode)


def _get_remove_wm_profile(mode=None):
    normalized = _get_effective_remove_wm_mode(mode)
    return REMOVE_WM_MODE_PROFILES.get(normalized, REMOVE_WM_MODE_PROFILES[REMOVE_WM_MODE_DEFAULT])


def _push_remove_wm_runtime_mode(mode):
    previous = getattr(_REMOVE_WM_MODE_CONTEXT, "mode", None)
    _REMOVE_WM_MODE_CONTEXT.mode = _coerce_remove_wm_mode(mode)
    return previous


def _pop_remove_wm_runtime_mode(previous):
    if previous is None:
        try:
            delattr(_REMOVE_WM_MODE_CONTEXT, "mode")
        except Exception:
            pass
    else:
        _REMOVE_WM_MODE_CONTEXT.mode = previous


def _contains_watermark_keyword(text):
    normalized = (text or "").strip().lower()
    if not normalized:
        return False
    keywords = (
        "xmu_done",
        "watermark",
        "confidential",
        "draft",
        "sample",
        "do not copy",
        "internal use",
        "for review",
        "water mark",
        "water-mark",
        "水印",
        "机密",
        "保密",
        "绝密",
        "内部",
        "样稿",
        "草稿",
        "作废",
        "仅供",
        "测试章",
    )
    return any(keyword in normalized for keyword in keywords)


def _collect_shape_marker_text(shape):
    parts = []
    for attr_name in ("Name", "AlternativeText", "Title"):
        try:
            value = getattr(shape, attr_name, "")
        except Exception:
            value = ""
        if value:
            parts.append(str(value))
    try:
        text_effect = getattr(shape, "TextEffect", None)
        text_value = getattr(text_effect, "Text", "")
        if text_value:
            parts.append(str(text_value))
    except Exception:
        pass
    try:
        text_frame = getattr(shape, "TextFrame", None)
        if text_frame and getattr(text_frame, "HasText", False):
            text_value = text_frame.TextRange.Text
            if text_value:
                parts.append(str(text_value))
    except Exception:
        pass
    return " ".join(part for part in parts if part).strip()


def _collect_inline_shape_marker_text(ishape):
    parts = []
    for attr_name in ("AlternativeText", "Title"):
        try:
            value = getattr(ishape, attr_name, "")
        except Exception:
            value = ""
        if value:
            parts.append(str(value))
    try:
        range_obj = getattr(ishape, "Range", None)
        text_value = getattr(range_obj, "Text", "")
        if text_value:
            parts.append(str(text_value))
    except Exception:
        pass
    return " ".join(part for part in parts if part).strip()


def _shape_looks_like_watermark(shape, page_width, page_height, preserve_mine=False, mode=None):
    marker_text = _collect_shape_marker_text(shape)
    normalized_marker = marker_text.lower()
    if "xmu_done" in normalized_marker:
        return not preserve_mine
    if _contains_watermark_keyword(marker_text):
        return True

    width = _safe_float(_safe_getattr(shape, "Width", 0.0))
    height = _safe_float(_safe_getattr(shape, "Height", 0.0))
    left = _safe_float(_safe_getattr(shape, "Left", 0.0))
    top = _safe_float(_safe_getattr(shape, "Top", 0.0))
    transparency = _safe_float(_safe_getattr(_safe_getattr(shape, "Fill", None), "Transparency", 0.0))
    rotation_distance = _shape_rotation_distance(_safe_getattr(shape, "Rotation", 0.0))

    if page_width <= 0 or page_height <= 0:
        return False

    profile = _get_remove_wm_profile(mode)
    large_enough = (
        width >= page_width * profile["shape_width_ratio"]
        or height >= page_height * profile["shape_height_ratio"]
    )
    center_x = left + width / 2.0
    center_y = top + height / 2.0
    near_center = (
        abs(center_x - page_width / 2.0) <= page_width * profile["shape_center_ratio"]
        and abs(center_y - page_height / 2.0) <= page_height * profile["shape_center_ratio"]
    )
    off_canvas = left < 0 or top < 0
    diagonal = rotation_distance > profile["shape_rotation_distance"]
    translucent = transparency >= profile["shape_transparency"]
    near_watermark_zone = near_center or off_canvas

    if large_enough and near_watermark_zone and (diagonal or translucent):
        return True

    if profile.get("allow_plain_very_large"):
        very_large = (
            width >= page_width * profile["very_large_width_ratio"]
            or height >= page_height * profile["very_large_height_ratio"]
        )
        return very_large and near_watermark_zone

    return False


def _inline_shape_looks_like_watermark(ishape, page_width, page_height, preserve_mine=False, mode=None):
    marker_text = _collect_inline_shape_marker_text(ishape)
    normalized_marker = marker_text.lower()
    if "xmu_done" in normalized_marker:
        return not preserve_mine
    if _contains_watermark_keyword(marker_text):
        return True

    width = _safe_float(_safe_getattr(ishape, "Width", 0.0))
    height = _safe_float(_safe_getattr(ishape, "Height", 0.0))
    if page_width <= 0 or page_height <= 0:
        return False

    profile = _get_remove_wm_profile(mode)
    return width >= page_width * profile["inline_width_ratio"] and height >= page_height * profile["inline_height_ratio"]


def _remove_watermark_from_word_safely(word_app, src, dst, preserve_mine=False, is_pdf_source=False, mode=None):
    doc = None
    removed_header_shapes = 0
    removed_doc_shapes = 0
    removed_header_inline = 0
    effective_mode = _get_effective_remove_wm_mode(mode)
    try:
        src_abs = os.path.abspath(src)
        dst_abs = os.path.abspath(dst)
        with _DisableWin32ComGenCache():
            doc = _watermark_core_open_word_document_safely(word_app, src_abs)
            page_width = _safe_float(_safe_getattr(_safe_getattr(doc, "PageSetup", None), "PageWidth", 0.0))
            page_height = _safe_float(_safe_getattr(_safe_getattr(doc, "PageSetup", None), "PageHeight", 0.0))

            for section in doc.Sections:
                for header in section.Headers:
                    try:
                        header_shapes = header.Shapes
                        for index in range(header_shapes.Count, 0, -1):
                            shape = header.Shapes(index)
                            if not _shape_looks_like_watermark(
                                shape,
                                page_width,
                                page_height,
                                preserve_mine=preserve_mine,
                                mode=effective_mode,
                            ):
                                continue
                            shape.Delete()
                            removed_header_shapes += 1
                    except Exception as exc:
                        _debug(f"patch_remove_wm:header_shape_iter_error:{exc}")

                    try:
                        inline_shapes = header.Range.InlineShapes
                        for index in range(inline_shapes.Count, 0, -1):
                            ishape = inline_shapes(index)
                            if not _inline_shape_looks_like_watermark(
                                ishape,
                                page_width,
                                page_height,
                                preserve_mine=preserve_mine,
                                mode=effective_mode,
                            ):
                                continue
                            ishape.Delete()
                            removed_header_inline += 1
                    except Exception as exc:
                        _debug(f"patch_remove_wm:header_inline_iter_error:{exc}")

            if is_pdf_source and not preserve_mine:
                try:
                    for index in range(doc.Shapes.Count, 0, -1):
                        shape = doc.Shapes(index)
                        try:
                            if not _shape_looks_like_watermark(
                                shape,
                                page_width,
                                page_height,
                                preserve_mine=False,
                                mode=effective_mode,
                            ):
                                continue
                        except Exception:
                            continue
                        shape.Delete()
                        removed_doc_shapes += 1
                except Exception as exc:
                    _debug(f"patch_remove_wm:doc_shape_iter_error:{exc}")

            doc.SaveAs(dst_abs)
            doc.Close()
            doc = None
        _debug(
            "patch_remove_wm:safe_cleanup:"
            f"header_shapes={removed_header_shapes}:header_inline={removed_header_inline}:"
            f"doc_shapes={removed_doc_shapes}:mode={effective_mode}:pdf_source={is_pdf_source}:dst={dst_abs}"
        )
        return "SUCCESS"
    except Exception as exc:
        _debug(f"patch_remove_wm:safe_cleanup_error:{exc}")
        try:
            if doc is not None:
                doc.Close(False)
        except Exception:
            pass
        return f"ERROR:{exc}"


def _patch_remove_watermark_robustness():
    try:
        original = remove_watermark_from_word
    except Exception as exc:
        _debug(f"patch_remove_wm:missing:{exc}")
        return

    if getattr(original, "__fx_remove_wm_patch__", False):
        return

    def patched(word_app, src, dst, preserve_mine=False, is_pdf_source=False, mode=None):
        return _remove_watermark_from_word_safely(
            word_app,
            src,
            dst,
            preserve_mine=preserve_mine,
            is_pdf_source=is_pdf_source,
            mode=mode,
        )

    patched.__fx_remove_wm_patch__ = True
    patched.__fx_remove_wm_original__ = original
    globals()["remove_watermark_from_word"] = patched
    _debug("patch_remove_wm:installed")


_patch_remove_watermark_robustness()


def _write_failed_report(output_folder, failed_list):
    if not failed_list:
        return None
    report_path = os.path.join(output_folder, "!失败文件清单.txt")
    with open(report_path, "w", encoding="utf-8") as fh:
        fh.write("以下文件处理失败（已原样保留）：\n")
        for item in failed_list:
            fh.write(f"{item}\n")
    return report_path


def _copy_tree_contents(src_dir, dst_dir, skip_names=None):
    skip_names = set(skip_names or [])
    src_root = Path(src_dir)
    dst_root = Path(dst_dir)
    if not src_root.exists():
        return
    for item in src_root.rglob("*"):
        rel = item.relative_to(src_root)
        if rel.name in skip_names:
            continue
        target = dst_root / rel
        if item.is_dir():
            target.mkdir(parents=True, exist_ok=True)
        else:
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(item, target)


def _read_failed_report_items(report_path):
    if not report_path or not os.path.exists(report_path):
        return []
    lines = Path(report_path).read_text(encoding="utf-8", errors="ignore").splitlines()
    return [line.strip() for line in lines[1:] if line.strip()]


def _dispatch_com_app_dynamic(progid):
    try:
        if str(progid).lower() == "word.application":
            import pywintypes
            import win32com.client.dynamic

            word_clsid = pywintypes.IID("{000209FF-0000-0000-C000-000000000046}")
            dispatch = pythoncom.CoCreateInstance(
                word_clsid,
                None,
                pythoncom.CLSCTX_LOCAL_SERVER,
                pythoncom.IID_IDispatch,
            )
            return win32com.client.dynamic.Dispatch(dispatch)
        import win32com.client.dynamic

        return win32com.client.dynamic.Dispatch(progid)
    except Exception as exc:
        _debug(f"com_dynamic_dispatch:fallback:{progid}:{exc}")
        return _FX_ORIGINAL_WIN32COM_DISPATCH_EX(progid)


class _DisableWin32ComGenCache:
    def __init__(self):
        self._original_get_class = None
        self._original_get_module = None

    def __enter__(self):
        try:
            import win32com.client.gencache as gencache

            self._original_get_class = gencache.GetClassForCLSID
            self._original_get_module = gencache.GetModuleForCLSID
            gencache.GetClassForCLSID = lambda *_args, **_kwargs: None
            gencache.GetModuleForCLSID = lambda *_args, **_kwargs: None
        except Exception as exc:
            _debug(f"com_gencache:disable_error:{exc}")
        return self

    def __exit__(self, exc_type, exc, tb):
        if self._original_get_class is None:
            return False
        try:
            import win32com.client.gencache as gencache

            if self._original_get_class is not None:
                gencache.GetClassForCLSID = self._original_get_class
            if self._original_get_module is not None:
                gencache.GetModuleForCLSID = self._original_get_module
        except Exception as restore_exc:
            _debug(f"com_gencache:restore_error:{restore_exc}")
        return False


def _is_win32com_gen_cache_error(exc):
    text = str(exc or "").lower()
    return (
        "win32com.gen_py" in text
        and (
            "clsidtoclassmap" in text
            or "clsidtopackagemap" in text
            or "vtablestopackagemap" in text
            or "has no attribute" in text
        )
    )


def _safe_office_dispatch_ex(progid, *args, **kwargs):
    progid_text = str(progid or "")
    if progid_text.lower() == "word.application":
        return _dispatch_com_app_dynamic(progid_text)
    try:
        return _FX_ORIGINAL_WIN32COM_DISPATCH_EX(progid, *args, **kwargs)
    except Exception as exc:
        if progid_text.lower() in {"word.application", "powerpoint.application"} and _is_win32com_gen_cache_error(exc):
            _debug(f"com_dispatch_ex:gen_cache_fallback:{progid_text}:{exc}")
            return _dispatch_com_app_dynamic(progid_text)
        raise


def _safe_office_dispatch(progid, *args, **kwargs):
    progid_text = str(progid or "")
    if progid_text.lower() == "word.application":
        return _dispatch_com_app_dynamic(progid_text)
    try:
        return _FX_ORIGINAL_WIN32COM_DISPATCH(progid, *args, **kwargs)
    except Exception as exc:
        if progid_text.lower() in {"word.application", "powerpoint.application"} and _is_win32com_gen_cache_error(exc):
            _debug(f"com_dispatch:gen_cache_fallback:{progid_text}:{exc}")
            return _dispatch_com_app_dynamic(progid_text)
        raise


def _install_safe_office_dispatch_patch():
    try:
        current_dispatch = getattr(win32com.client, "Dispatch", None)
        current_dispatch_ex = getattr(win32com.client, "DispatchEx", None)
        dispatch_installed = getattr(current_dispatch, "__fx_safe_office_dispatch_patch__", False)
        dispatch_ex_installed = getattr(current_dispatch_ex, "__fx_safe_office_dispatch_patch__", False)
        if dispatch_installed and dispatch_ex_installed:
            return
        _safe_office_dispatch.__fx_safe_office_dispatch_patch__ = True
        _safe_office_dispatch.__wrapped__ = _FX_ORIGINAL_WIN32COM_DISPATCH
        _safe_office_dispatch_ex.__fx_safe_office_dispatch_patch__ = True
        _safe_office_dispatch_ex.__wrapped__ = _FX_ORIGINAL_WIN32COM_DISPATCH_EX
        win32com.client.Dispatch = _safe_office_dispatch
        win32com.client.DispatchEx = _safe_office_dispatch_ex
        _debug("com_dispatch_ex:safe_office_patch_installed")
    except Exception as exc:
        _debug(f"com_dispatch_ex:safe_office_patch_error:{exc}")


def _create_hidden_word_app():
    with _DisableWin32ComGenCache():
        word_app = _dispatch_com_app_dynamic("Word.Application")
    try:
        word_app.Visible = False
    except Exception:
        pass
    try:
        word_app.DisplayAlerts = 0
    except Exception:
        pass
    return word_app


_FX_RUNTIME_CONVERT_DOC_TO_PDF = convert_doc_to_pdf
_FX_RUNTIME_CONVERT_PPT_TO_PDF = convert_ppt_to_pdf


def _convert_doc_to_pdf_safely(word_app, src, dst):
    status = ""
    try:
        with _DisableWin32ComGenCache():
            status = _FX_RUNTIME_CONVERT_DOC_TO_PDF(word_app, src, dst)
    except Exception as exc:
        status = f"ERROR:{exc}"
    if str(status).strip() == "SUCCESS" and os.path.exists(dst) and os.path.getsize(dst) > 0:
        return "SUCCESS"

    fallback_status = _export_word_docx_to_pdf_safely(src, dst)
    if str(fallback_status).strip() == "SUCCESS":
        return "SUCCESS"
    return status or fallback_status


def _convert_ppt_to_pdf_safely(ppt_app, src, dst):
    with _DisableWin32ComGenCache():
        return _FX_RUNTIME_CONVERT_PPT_TO_PDF(ppt_app, src, dst)


def _convert_pdf_to_ppt_safely(ppt_app, src, dst):
    temp_dir = None
    try:
        import fitz
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        src_abs = os.path.abspath(src)
        dst_abs = os.path.abspath(dst)
        os.makedirs(os.path.dirname(dst_abs), exist_ok=True)
        temp_dir = tempfile.mkdtemp(prefix="fx_pdf2ppt_")
        with fitz.open(src_abs) as doc:
            if doc.is_encrypted and not doc.authenticate(""):
                return "ERROR:encrypted_pdf_requires_password"
            if len(doc) <= 0:
                return "ERROR:empty_pdf"

            first_rect = doc[0].rect
            slide_width_pt = max(120.0, float(first_rect.width or 720.0))
            slide_height_pt = max(120.0, float(first_rect.height or 540.0))
            pt_to_emu = 12700

            presentation = Presentation()
            presentation.slide_width = int(round(slide_width_pt * pt_to_emu))
            presentation.slide_height = int(round(slide_height_pt * pt_to_emu))
            blank_layout = presentation.slide_layouts[6]

            def to_emu(value):
                return int(round(float(value) * pt_to_emu))

            def page_mapper(page):
                rect = page.rect
                scale = min(slide_width_pt / max(1.0, float(rect.width)), slide_height_pt / max(1.0, float(rect.height)))
                offset_x = (slide_width_pt - float(rect.width) * scale) / 2.0
                offset_y = (slide_height_pt - float(rect.height) * scale) / 2.0

                def map_bbox(bbox):
                    x0, y0, x1, y1 = [float(v) for v in bbox]
                    left = offset_x + x0 * scale
                    top = offset_y + y0 * scale
                    width = max(1.0, (x1 - x0) * scale)
                    height = max(1.0, (y1 - y0) * scale)
                    return to_emu(left), to_emu(top), to_emu(width), to_emu(height), scale

                return map_bbox

            def set_run_font(run, span, scale):
                font_name = str(span.get("font") or "").strip()
                if font_name and len(font_name) < 64:
                    run.font.name = font_name
                size = span.get("size")
                try:
                    run.font.size = Pt(max(1.0, min(96.0, float(size) * max(0.2, scale))))
                except Exception:
                    run.font.size = Pt(10)
                lowered_font = font_name.lower()
                run.font.bold = bool("bold" in lowered_font or "black" in lowered_font)
                run.font.italic = bool("italic" in lowered_font or "oblique" in lowered_font)
                try:
                    color = int(span.get("color") or 0)
                    run.font.color.rgb = RGBColor((color >> 16) & 255, (color >> 8) & 255, color & 255)
                except Exception:
                    pass

            for page_index, page in enumerate(doc, start=1):
                slide = presentation.slides.add_slide(blank_layout)
                mapper = page_mapper(page)
                page_dict = page.get_text("dict")

                for block_index, block in enumerate(page_dict.get("blocks") or [], start=1):
                    block_type = int(block.get("type", 0) or 0)
                    bbox = block.get("bbox")
                    if not bbox:
                        continue
                    left, top, width, height, scale = mapper(bbox)
                    if block_type == 1:
                        image_bytes = block.get("image")
                        if not image_bytes:
                            continue
                        ext = str(block.get("ext") or "png").lower().lstrip(".") or "png"
                        image_path = os.path.join(temp_dir, f"page_{page_index:04d}_image_{block_index:03d}.{ext}")
                        with open(image_path, "wb") as image_file:
                            image_file.write(image_bytes)
                        if os.path.exists(image_path) and os.path.getsize(image_path) > 0:
                            slide.shapes.add_picture(image_path, left, top, width, height)
                        continue

                    if block_type != 0:
                        continue
                    lines = block.get("lines") or []
                    if not lines:
                        continue
                    shape = slide.shapes.add_textbox(left, top, width, max(height, to_emu(8)))
                    frame = shape.text_frame
                    frame.clear()
                    frame.margin_left = 0
                    frame.margin_right = 0
                    frame.margin_top = 0
                    frame.margin_bottom = 0
                    frame.word_wrap = False
                    paragraph = frame.paragraphs[0]
                    used_paragraph = False
                    for line in lines:
                        spans = [span for span in (line.get("spans") or []) if str(span.get("text") or "")]
                        if not spans:
                            continue
                        if used_paragraph:
                            paragraph = frame.add_paragraph()
                        used_paragraph = True
                        paragraph.space_after = Pt(0)
                        paragraph.space_before = Pt(0)
                        for span in spans:
                            run = paragraph.add_run()
                            run.text = str(span.get("text") or "")
                            set_run_font(run, span, scale)

            presentation.save(dst_abs)
        if os.path.exists(dst_abs) and os.path.getsize(dst_abs) > 0:
            return "SUCCESS"
        return "ERROR:no_output"
    except Exception as exc:
        return f"ERROR:{exc}"
    finally:
        try:
            if temp_dir:
                shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass


_install_safe_office_dispatch_patch()


def _export_word_docx_to_pdf_safely(docx_path, pdf_path):
    doc = None
    word_app = None
    docx_abs = os.path.abspath(docx_path)
    pdf_abs = os.path.abspath(pdf_path)
    try:
        os.makedirs(os.path.dirname(pdf_abs), exist_ok=True)
        word_app = _create_hidden_word_app()
        with _DisableWin32ComGenCache():
            doc = _watermark_core_open_word_document_safely(word_app, docx_abs)
            doc.ExportAsFixedFormat(pdf_abs, 17)
            doc.Close(False)
            doc = None
        if os.path.exists(pdf_abs) and os.path.getsize(pdf_abs) > 0:
            return "SUCCESS"
        return "ERROR:no_output"
    except Exception as exc:
        _debug(f"remove_wm_pdf_export_direct:error:{exc}")
        return f"ERROR:{exc}"
    finally:
        try:
            if doc is not None:
                doc.Close(False)
        except Exception:
            pass
        try:
            if word_app is not None:
                word_app.Quit()
        except Exception:
            pass


globals()["convert_doc_to_pdf"] = _convert_doc_to_pdf_safely
globals()["convert_ppt_to_pdf"] = _convert_ppt_to_pdf_safely
_ns["convert_doc_to_pdf"] = _convert_doc_to_pdf_safely
_ns["convert_ppt_to_pdf"] = _convert_ppt_to_pdf_safely


def _get_remove_wm_overwrite_original(app):
    if getattr(app, "rm_wm_overwrite_original", None) is None:
        return False
    try:
        return bool(app.rm_wm_overwrite_original.get())
    except Exception:
        return False


def _build_single_remove_wm_output_path(source_path):
    source = Path(source_path)
    candidate = source.with_name(f"{source.stem}_去水印{source.suffix}")
    counter = 2
    while candidate.exists() and os.path.normcase(str(candidate)) != os.path.normcase(str(source)):
        candidate = source.with_name(f"{source.stem}_去水印_{counter}{source.suffix}")
        counter += 1
    return str(candidate)


def _build_single_remove_wm_base_output_path(source_path):
    source = Path(source_path)
    suffix = Path(_build_single_remove_wm_output_path(source)).suffix
    return str(source.with_name(f"{source.stem}_去水印{suffix}"))


def _replace_file_safely(src, dst):
    dst_path = Path(dst)
    temp_path = dst_path.with_name(f".{dst_path.name}.fx_replace_tmp")
    try:
        if temp_path.exists():
            temp_path.unlink()
    except Exception:
        pass
    try:
        shutil.copy2(src, temp_path)
        os.replace(str(temp_path), str(dst_path))
    finally:
        try:
            if temp_path.exists():
                temp_path.unlink()
        except Exception:
            pass


def _finalize_single_remove_wm_output(app, source_file, staged_output_file, overwrite_original=False):
    source_file = _normalize_input_path_value(source_file)
    staged_output_file = _normalize_input_path_value(staged_output_file)
    if not source_file or not staged_output_file or not os.path.exists(staged_output_file):
        return None

    final_path = source_file if overwrite_original else _build_single_remove_wm_output_path(source_file)
    if overwrite_original:
        _replace_file_safely(staged_output_file, final_path)
        try:
            app.log(f"✅ [单文件去水印] 已直接覆盖原文件：{os.path.basename(final_path)}")
        except Exception:
            pass
    else:
        shutil.copy2(staged_output_file, final_path)
        try:
            app.log(f"✅ [单文件去水印] 已输出新文件：{final_path}")
        except Exception:
            pass
    return final_path


def _resolve_result_output_folder(input_value):
    normalized_input = _normalize_input_path_value(input_value)
    if normalized_input and os.path.isfile(normalized_input):
        input_root = os.path.dirname(normalized_input)
    else:
        input_root = normalized_input
    return normalized_input, input_root, os.path.join(input_root, RESULT_FOLDER_NAME)


def _resolve_output_root_for_task(input_value, task_type, strategy_value=None):
    normalized_input = _normalize_input_path_value(input_value)
    if normalized_input and os.path.isfile(normalized_input):
        input_root = os.path.dirname(normalized_input)
    else:
        input_root = normalized_input

    resolved_strategy = _resolve_output_strategy(task_type, strategy_value)
    if resolved_strategy == "overwrite":
        output_root = input_root
    elif resolved_strategy == "same_dir":
        output_root = input_root
    else:
        output_root = os.path.join(input_root, RESULT_FOLDER_NAME)

    return normalized_input, input_root, output_root, resolved_strategy


def _task_result_now():
    return time.time()


def _new_task_result(input_value="", task_type=""):
    normalized_input = _normalize_input_path_value(input_value)
    started_at = _task_result_now()
    return {
        "task_type": task_type or "",
        "input": normalized_input,
        "output_strategy": "",
        "output_strategy_label": "",
        "status": "unknown",
        "success": False,
        "stopped": False,
        "skipped": False,
        "message": "",
        "detail": "",
        "error": "",
        "outputs": [],
        "output_root": "",
        "failed_items": [],
        "processed_count": 0,
        "success_count": 0,
        "failed_count": 0,
        "skipped_count": 0,
        "started_at": started_at,
        "finished_at": None,
        "duration_seconds": 0.0,
    }


def _normalize_task_output_path(path_value):
    normalized = _normalize_input_path_value(path_value)
    return normalized if normalized else ""


def _add_task_result_output(result, path_value):
    if not isinstance(result, dict):
        return
    normalized = _normalize_task_output_path(path_value)
    if not normalized:
        return
    outputs = result.setdefault("outputs", [])
    if normalized not in outputs:
        outputs.append(normalized)


def _set_task_result_output_root(result, path_value):
    if not isinstance(result, dict):
        return
    normalized = _normalize_task_output_path(path_value)
    if normalized:
        result["output_root"] = normalized


def _set_task_result_output_strategy(result, task_type, strategy_value):
    if not isinstance(result, dict):
        return
    _apply_output_strategy_to_result(result, task_type, strategy_value)


def _set_task_result_counts(result, *, processed=None, success=None, failed=None, skipped=None):
    if not isinstance(result, dict):
        return
    if processed is not None:
        result["processed_count"] = max(0, int(processed))
    if success is not None:
        result["success_count"] = max(0, int(success))
    if failed is not None:
        result["failed_count"] = max(0, int(failed))
    if skipped is not None:
        result["skipped_count"] = max(0, int(skipped))


def _set_task_result_finished(result, status, message="", detail="", error="", stopped=False, skipped=False):
    if not isinstance(result, dict):
        return result
    finished_at = _task_result_now()
    result["finished_at"] = finished_at
    result["duration_seconds"] = max(0.0, finished_at - float(result.get("started_at") or finished_at))
    normalized_status = str(status or "unknown").strip().lower() or "unknown"
    result["status"] = normalized_status
    result["success"] = normalized_status == "success"
    result["stopped"] = bool(stopped or normalized_status == "stopped")
    result["skipped"] = bool(skipped or normalized_status == "skipped")
    result["message"] = str(message or result.get("message") or "")
    result["detail"] = str(detail or result.get("detail") or "")
    result["error"] = str(error or result.get("error") or "")
    result["failed_count"] = max(0, int(result.get("failed_count") or 0))
    result["success_count"] = max(0, int(result.get("success_count") or 0))
    result["processed_count"] = max(0, int(result.get("processed_count") or 0))
    result["skipped_count"] = max(0, int(result.get("skipped_count") or 0))
    result["failed_items"] = list(result.get("failed_items") or [])
    result["outputs"] = [item for item in result.get("outputs", []) if item]
    return result


def _schedule_task_completion_notice(app, result):
    """Show one completion notice after a task is finalized successfully."""

    if app is None or not isinstance(result, dict):
        return
    if str(result.get("status") or "").strip().lower() != "success":
        return
    if result.get("_fx_completion_notice_scheduled"):
        return
    result["_fx_completion_notice_scheduled"] = True

    task_label = _get_feature_label(result.get("task_type"), fallback="任务")
    success_count = int(result.get("success_count") or 0)
    skipped_count = int(result.get("skipped_count") or 0)
    failed_count = int(result.get("failed_count") or 0)
    detail_lines = [f"{task_label}已完成。", f"成功处理：{success_count} 个文件"]
    if skipped_count:
        detail_lines.append(f"跳过：{skipped_count} 个文件")
    if failed_count:
        detail_lines.append(f"失败：{failed_count} 个文件")
    output_root = str(result.get("output_root") or "").strip()
    if output_root:
        detail_lines.extend(["", f"输出位置：{output_root}"])
    message = "\n".join(detail_lines)

    def show_notice():
        try:
            tkinter.messagebox.showinfo("任务完成", message, parent=app)
        except Exception as exc:
            _debug(f"task_completion_notice:error:{exc}")

    try:
        app.after(0, show_notice)
    except Exception:
        show_notice()


def _task_result_snapshot(result):
    if not isinstance(result, dict):
        return {}
    return {
        "task_type": result.get("task_type", ""),
        "input": result.get("input", ""),
        "output_strategy_requested": result.get("output_strategy_requested", ""),
        "output_strategy": result.get("output_strategy", ""),
        "output_strategy_label": result.get("output_strategy_label", ""),
        "status": result.get("status", ""),
        "success": bool(result.get("success", False)),
        "stopped": bool(result.get("stopped", False)),
        "skipped": bool(result.get("skipped", False)),
        "message": result.get("message", ""),
        "detail": result.get("detail", ""),
        "error": result.get("error", ""),
        "outputs": list(result.get("outputs") or []),
        "output_root": result.get("output_root", ""),
        "failed_items": list(result.get("failed_items") or []),
        "processed_count": int(result.get("processed_count") or 0),
        "success_count": int(result.get("success_count") or 0),
        "failed_count": int(result.get("failed_count") or 0),
        "skipped_count": int(result.get("skipped_count") or 0),
        "started_at": result.get("started_at"),
        "finished_at": result.get("finished_at"),
        "duration_seconds": float(result.get("duration_seconds") or 0.0),
    }


def _attach_task_result(app, result):
    if app is None or not isinstance(result, dict):
        return result
    try:
        app._fx_last_task_result = result
    except Exception:
        pass
    return result


def _start_task_result(app, input_value, task_type):
    result = _new_task_result(input_value=input_value, task_type=task_type)
    _set_task_result_output_strategy(result, task_type, _get_task_output_strategy(app, task_type))
    return _attach_task_result(app, result)


def _get_last_task_result(app):
    try:
        result = getattr(app, "_fx_last_task_result", None)
    except Exception:
        result = None
    return result if isinstance(result, dict) else None


def _clear_last_task_result(app):
    try:
        app._fx_last_task_result = None
    except Exception:
        pass


_FX_ES_CONTINUOUS = 0x80000000
_FX_ES_SYSTEM_REQUIRED = 0x00000001
_FX_ES_AWAYMODE_REQUIRED = 0x00000040


def _fx_background_guard_begin(app, reason="task"):
    """Keep Windows from idling/sleeping while long batch work is running."""

    if app is None:
        return False
    try:
        count = int(getattr(app, "_fx_background_guard_count", 0) or 0)
    except Exception:
        count = 0
    try:
        app._fx_background_guard_count = count + 1
    except Exception:
        pass
    if count > 0:
        return True
    try:
        flags = _FX_ES_CONTINUOUS | _FX_ES_SYSTEM_REQUIRED | _FX_ES_AWAYMODE_REQUIRED
        result = ctypes.windll.kernel32.SetThreadExecutionState(flags)
        try:
            app._fx_background_guard_active = bool(result)
        except Exception:
            pass
        _debug(f"background_guard:begin:{reason}:{bool(result)}")
        return bool(result)
    except Exception as exc:
        _debug(f"background_guard:begin_error:{reason}:{exc}")
        return False


def _fx_background_guard_end(app, reason="task"):
    if app is None:
        return False
    try:
        count = int(getattr(app, "_fx_background_guard_count", 0) or 0)
    except Exception:
        count = 0
    next_count = max(0, count - 1)
    try:
        app._fx_background_guard_count = next_count
    except Exception:
        pass
    if next_count > 0:
        return True
    try:
        result = ctypes.windll.kernel32.SetThreadExecutionState(_FX_ES_CONTINUOUS)
        try:
            app._fx_background_guard_active = False
        except Exception:
            pass
        _debug(f"background_guard:end:{reason}:{bool(result)}")
        return bool(result)
    except Exception as exc:
        _debug(f"background_guard:end_error:{reason}:{exc}")
        return False


def _export_task_result(result, output_path):
    if not isinstance(result, dict) or not output_path:
        return False
    try:
        path = Path(_normalize_input_path_value(output_path))
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
        return True
    except Exception as exc:
        _debug(f"task_result:export_error:{exc}")
        return False


def _finalize_current_task_result(app, status, message="", detail="", error="", stopped=False, skipped=False):
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, "", "")
    return _set_task_result_finished(
        result,
        status=status,
        message=message,
        detail=detail,
        error=error,
        stopped=stopped,
        skipped=skipped,
    )


def _summarize_logs_for_task_result(logs):
    entries = [str(item).strip() for item in (logs or []) if str(item).strip()]
    if not entries:
        return "", "", False, False

    error_lines = [item for item in entries if "❌" in item or "🔥" in item or "错误" in item or "失败" in item]
    stop_lines = [item for item in entries if "⏹️" in item or "中止" in item or "停止" in item]
    success_lines = [item for item in entries if "🎉 [完成]" in item or "✅ [完成]" in item or "已全部完成" in item]
    if error_lines:
        return "failed", error_lines[-1], False, False
    if stop_lines:
        return "stopped", stop_lines[-1], True, False
    if success_lines:
        return "success", success_lines[-1], False, False
    return "", entries[-1], False, False


def _infer_task_result_from_context(app, input_value, task_type, return_value=None, logs=None, exception=None):
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, input_value, task_type)
    if task_type and not result.get("task_type"):
        result["task_type"] = task_type
    if input_value and not result.get("input"):
        result["input"] = _normalize_input_path_value(input_value)

    if exception is not None:
        return _set_task_result_finished(
            result,
            status="failed",
            message=str(exception),
            detail=str(exception),
            error=str(exception),
        )

    if result.get("finished_at"):
        return result

    inferred_status = ""
    inferred_message = ""
    stopped = False
    skipped = False

    if getattr(app, "stop_event", False):
        inferred_status = "stopped"
        inferred_message = "用户停止"
        stopped = True

    if not inferred_status:
        log_status, log_message, log_stopped, log_skipped = _summarize_logs_for_task_result(logs)
        if log_status:
            inferred_status = log_status
            inferred_message = log_message
            stopped = log_stopped
            skipped = log_skipped

    if not inferred_status and isinstance(return_value, str):
        normalized_return = return_value.strip()
        if normalized_return == "SUCCESS":
            inferred_status = "success"
            inferred_message = normalized_return
        elif normalized_return.startswith("ERROR"):
            inferred_status = "failed"
            inferred_message = normalized_return
        elif normalized_return.startswith("SKIP"):
            inferred_status = "skipped"
            inferred_message = normalized_return
            skipped = True

    if not inferred_status:
        failed_items = list(result.get("failed_items") or [])
        if failed_items or int(result.get("failed_count") or 0) > 0:
            inferred_status = "failed"
            inferred_message = result.get("detail") or "存在失败项目"
        elif result.get("outputs") or result.get("output_root"):
            inferred_status = "success"
            inferred_message = result.get("detail") or "执行完成"
        else:
            inferred_status = "success"
            inferred_message = result.get("detail") or "执行完成"

    return _set_task_result_finished(
        result,
        status=inferred_status,
        message=inferred_message,
        detail=result.get("detail") or inferred_message,
        error=result.get("error") or (inferred_message if inferred_status == "failed" else ""),
        stopped=stopped,
        skipped=skipped,
    )


def _collect_result_file_paths(output_root):
    normalized_root = _normalize_input_path_value(output_root)
    if not normalized_root or not os.path.exists(normalized_root):
        return []
    root_path = Path(normalized_root)
    results = []
    if root_path.is_file():
        return [str(root_path)]
    skip_names = {"!失败文件清单.txt"}
    for item in root_path.rglob("*"):
        if not item.is_file():
            continue
        if item.name in skip_names:
            continue
        results.append(str(item))
    return results


def _coerce_progress_value(value, default=0.0):
    try:
        return float(value)
    except Exception:
        return default


def _clamp_progress_value(value):
    return max(0.0, min(1.0, _coerce_progress_value(value)))


def _get_active_progress_tracker(app):
    return getattr(app, "_fx_progress_tracker", None)


def _format_progress_eta(seconds):
    if seconds is None:
        return "--"
    try:
        seconds = float(seconds)
    except Exception:
        return "--"
    if seconds < 0 or seconds == float("inf"):
        return "--"
    if seconds < 1:
        return "<1秒"
    seconds = int(round(seconds))
    minutes, sec = divmod(seconds, 60)
    hours, minutes = divmod(minutes, 60)
    if hours:
        return f"{hours}:{minutes:02d}:{sec:02d}"
    return f"{minutes:02d}:{sec:02d}"


def _shorten_progress_name(value, max_chars=34):
    text = str(value or "").strip()
    if not text:
        return ""
    try:
        name = os.path.basename(text.rstrip("\\/")) or text
    except Exception:
        name = text
    if len(name) <= max_chars:
        return name
    return name[: max(1, max_chars - 1)] + "…"


def _set_progress_status(
    app,
    *,
    current_file="",
    stage="",
    fraction=None,
    completed=None,
    total=None,
    eta_seconds=None,
):
    if app is None:
        return ""
    parts = []
    if current_file:
        parts.append(f"当前：{_shorten_progress_name(current_file)}")
    if stage:
        parts.append(f"阶段：{stage}")
    if total:
        try:
            completed_value = int(max(0, completed or 0))
            total_value = int(max(1, total or 1))
            parts.append(f"文件：{min(completed_value, total_value)}/{total_value}")
        except Exception:
            pass
    if fraction is not None:
        percent = int(round(_clamp_progress_value(fraction) * 100))
        parts.append(f"总进度：{percent}%")
    if eta_seconds is not None:
        parts.append(f"预计剩余：{_format_progress_eta(eta_seconds)}")
    text = " | ".join(parts) if parts else PROGRESS_STATUS_IDLE_TEXT
    try:
        app._fx_last_progress_status = text
    except Exception:
        pass
    status_var = getattr(app, "_fx_progress_status_var", None)
    if status_var is not None:
        try:
            status_var.set(text)
        except Exception:
            pass
    return text


def _extract_progress_path_from_call(args, kwargs):
    candidates = []
    if isinstance(kwargs, dict):
        candidates.extend(kwargs.values())
    candidates.extend(args or ())
    for candidate in candidates:
        text = str(candidate or "").strip()
        if not text:
            continue
        try:
            if os.path.exists(text) or os.path.splitext(os.path.basename(text))[1]:
                return text
        except Exception:
            continue
    return ""


def _unwrap_runtime_run_process(callable_obj):
    current = callable_obj
    seen = set()
    while isinstance(current, types.FunctionType) and id(current) not in seen:
        seen.add(id(current))
        if getattr(current, "__module__", "") == "fengxi_embedded_runtime" and current.__name__ == "run_process":
            return current
        next_func = None
        for cell in current.__closure__ or ():
            try:
                value = cell.cell_contents
            except ValueError:
                continue
            if isinstance(value, types.FunctionType):
                next_func = value
                break
        if next_func is None:
            break
        current = next_func
    return None


def _build_runtime_progress_site_map(runtime_run_process):
    if runtime_run_process is None:
        return {
            "before_process_single": set(),
            "pass_through": set(),
            "direct_pre": set(),
        }

    cached = getattr(runtime_run_process, "__fx_progress_site_map__", None)
    if cached is not None:
        return cached

    instructions = list(dis.get_instructions(runtime_run_process))
    site_map = {
        "before_process_single": set(),
        "pass_through": set(),
        "direct_pre": set(),
    }

    for index, instr in enumerate(instructions):
        if instr.opname != "LOAD_ATTR" or instr.argval != "progress_bar":
            continue

        method_index = None
        last_call_index = None
        for follow_index in range(index + 1, min(index + 40, len(instructions))):
            follow = instructions[follow_index]
            if follow.opname == "LOAD_METHOD" and follow.argval == "set":
                method_index = follow_index
            elif method_index is not None and follow.opname == "CALL":
                last_call_index = follow_index
            elif method_index is not None and follow.opname == "POP_TOP":
                break
        call_index = last_call_index
        if method_index is None or call_index is None:
            continue
        if call_index + 1 >= len(instructions) or instructions[call_index + 1].opname != "POP_TOP":
            continue

        call_offset = instructions[call_index].offset
        look_back = instructions[max(0, index - 12):index]
        look_ahead = instructions[call_index + 1:min(len(instructions), call_index + 18)]

        if any(item.opname == "LOAD_METHOD" and item.argval == "process_single_file" for item in look_ahead):
            site_map["before_process_single"].add(call_offset)
            continue

        if any(item.opname == "STORE_FAST" and item.argval == "processed" for item in look_back):
            site_map["pass_through"].add(call_offset)
            continue

        meaningful_next = None
        for item in look_ahead:
            if item.opname in {"POP_TOP", "EXTENDED_ARG", "NOP"}:
                continue
            meaningful_next = item
            break
        if meaningful_next is not None and meaningful_next.opname.startswith("JUMP"):
            site_map["pass_through"].add(call_offset)
            continue

        site_map["direct_pre"].add(call_offset)

    runtime_run_process.__fx_progress_site_map__ = site_map
    return site_map


def _get_zip_max_depth(app):
    start_value = ""
    end_value = ""
    start_var = getattr(app, "zip_min_depth_var", None)
    end_var = getattr(app, "zip_max_depth_var", None)
    if start_var is not None:
        try:
            start_value = start_var.get()
        except Exception:
            start_value = ""
    if end_var is not None:
        try:
            end_value = end_var.get()
        except Exception:
            end_value = ""
    start_value = str(start_value or "").strip()
    end_value = str(end_value or "").strip()
    if start_value and end_value:
        return f"{start_value}-{end_value}"
    if start_value and not end_value:
        return start_value
    return end_value


def _zip_archive_policy_label(value):
    normalized = normalize_zip_archive_policy(ZIP_ARCHIVE_POLICY_LABEL_TO_VALUE.get(str(value or ""), value))
    return ZIP_ARCHIVE_POLICY_VALUE_TO_LABEL.get(normalized, ZIP_ARCHIVE_POLICY_REUSE_LABEL)


def _set_zip_archive_policy(app, value):
    label = _zip_archive_policy_label(value)
    try:
        if getattr(app, "zip_archive_policy_var", None) is not None:
            app.zip_archive_policy_var.set(label)
    except Exception:
        pass
    return label


def _get_zip_archive_policy(app):
    value = ZIP_ARCHIVE_POLICY_REUSE_LABEL
    var = getattr(app, "zip_archive_policy_var", None)
    if var is not None:
        try:
            value = var.get()
        except Exception:
            value = ZIP_ARCHIVE_POLICY_REUSE_LABEL
    return normalize_zip_archive_policy(ZIP_ARCHIVE_POLICY_LABEL_TO_VALUE.get(str(value or ""), value))


def _count_zip_progress_units(input_value, mode, max_depth=None):
    normalized = _normalize_input_path_value(input_value)
    if not normalized:
        return 0
    try:
        return estimate_zip_progress_units(normalized, normalize_zip_mode(mode), max_depth=max_depth)
    except Exception:
        return 1 if os.path.exists(normalized) else 0


def _estimate_progress_total_units(app, input_value, task_type):
    normalized = _normalize_input_path_value(input_value)
    try:
        all_files = list(app.collect_input_files(normalized, task_type))
    except Exception:
        all_files = []

    if task_type == "pdf":
        mode = ""
        if getattr(app, "pdf_mode_var", None) is not None:
            try:
                mode = app.pdf_mode_var.get()
            except Exception:
                mode = ""
        pdf_count = sum(1 for item in all_files if str(item).lower().endswith(".pdf"))
        if mode in {"merge", "split", "encrypt", "ocr", "compress"}:
            return pdf_count

    if task_type == "zip":
        mode = ""
        if getattr(app, "zip_mode_var", None) is not None:
            try:
                mode = app.zip_mode_var.get()
            except Exception:
                mode = ""
        return _count_zip_progress_units(normalized, mode, _get_zip_max_depth(app))

    if task_type == "image":
        mode = _get_image_pdf_mode(app)
        if mode in {"to_pdf", "merge_pdf"}:
            return max(1, len(_collect_image_to_pdf_files(app, normalized)))

    return len(all_files)


class _FxRunProgressTracker:
    def __init__(self, app, total_units, original_progress_set, runtime_run_process, site_map):
        self.app = app
        self.total_units = max(1, int(total_units or 0))
        self.original_progress_set = original_progress_set
        self.runtime_code = getattr(runtime_run_process, "__code__", None)
        self.site_map = site_map or {}
        self.lock = threading.Lock()
        self.completed_units = 0
        self.pending_direct_offset = None
        self.seen_activity = False
        self.keep_final_on_reset = False
        self.started_at = time.time()
        self.current_file = ""
        self.current_stage = "准备中"

    def _apply_progress_locked(self, value):
        fraction = _clamp_progress_value(value)
        self.original_progress_set(fraction)
        self._emit_status_locked(fraction)

    def _estimate_eta_locked(self, fraction):
        if fraction <= 0.01 or fraction >= 0.999:
            return None
        elapsed = max(0.0, time.time() - self.started_at)
        if elapsed <= 0:
            return None
        return elapsed * (1.0 - fraction) / fraction

    def _emit_status_locked(self, fraction):
        _set_progress_status(
            self.app,
            current_file=self.current_file,
            stage=self.current_stage,
            fraction=fraction,
            completed=self.completed_units,
            total=self.total_units,
            eta_seconds=self._estimate_eta_locked(fraction),
        )

    def current_fraction(self):
        with self.lock:
            return self.completed_units / self.total_units

    def set_current_item(self, current_file="", stage="处理中"):
        with self.lock:
            if current_file:
                self.current_file = current_file
            if stage:
                self.current_stage = stage
            self.seen_activity = True
            self._apply_progress_locked(self.completed_units / self.total_units)

    def set_current_item_fraction(self, item_fraction, stage=None, current_file=None):
        with self.lock:
            if current_file:
                self.current_file = current_file
            if stage:
                self.current_stage = stage
            self.seen_activity = True
            overall = (self.completed_units + max(0.0, min(1.0, _coerce_progress_value(item_fraction)))) / self.total_units
            self._apply_progress_locked(overall)

    def complete_units(self, count=1, stage="完成"):
        with self.lock:
            if stage:
                self.current_stage = stage
            self.seen_activity = True
            self.completed_units = min(self.total_units, self.completed_units + max(0, int(count or 0)))
            self.pending_direct_offset = None
            self._apply_progress_locked(self.completed_units / self.total_units)

    def finalize_pending(self):
        with self.lock:
            if self.pending_direct_offset is not None:
                self.completed_units = min(self.total_units, self.completed_units + 1)
                self.pending_direct_offset = None
                self.seen_activity = True
                self.current_stage = "完成"
                self._apply_progress_locked(self.completed_units / self.total_units)

    def note_process_single_complete(self):
        self.complete_units(1)

    def handle_runtime_progress_call(self, call_offset, raw_value):
        raw_fraction = _clamp_progress_value(raw_value)
        with self.lock:
            if call_offset in self.site_map.get("before_process_single", set()):
                self.seen_activity = True
                if not self.current_stage or self.current_stage == "准备中":
                    self.current_stage = "处理中"
                self._apply_progress_locked(self.completed_units / self.total_units)
                return True

            if call_offset in self.site_map.get("direct_pre", set()):
                if self.pending_direct_offset is not None:
                    self.completed_units = min(self.total_units, self.completed_units + 1)
                self.pending_direct_offset = call_offset
                self.seen_activity = True
                self.current_stage = "处理中"
                self._apply_progress_locked(self.completed_units / self.total_units)
                return True

            if call_offset in self.site_map.get("pass_through", set()):
                self.seen_activity = True
                self.current_stage = "收尾"
                estimated_completed = int(round(raw_fraction * self.total_units))
                self.completed_units = min(self.total_units, max(self.completed_units, estimated_completed))
                self._apply_progress_locked(raw_fraction)
                return True

        return False

    def handle_reset(self):
        self.finalize_pending()
        if getattr(self.app, "stop_event", False) or not self.seen_activity:
            return False
        self.keep_final_on_reset = True
        self.current_stage = "已完成"
        self.original_progress_set(1.0)
        self._emit_status_locked(1.0)
        return True

    def should_ignore_reset_zero(self, frame_code, value):
        if not self.keep_final_on_reset:
            return False
        if frame_code is not getattr(type(self.app), "reset_ui").__code__:
            return False
        return _coerce_progress_value(value, default=1.0) <= 0.0


def _install_run_progress_tracker(app, input_value, task_type, runtime_run_process, site_map):
    progress_bar = getattr(app, "progress_bar", None)
    original_progress_set = getattr(progress_bar, "set", None)
    if progress_bar is None or not callable(original_progress_set):
        return None

    tracker = _FxRunProgressTracker(
        app=app,
        total_units=_estimate_progress_total_units(app, input_value, task_type),
        original_progress_set=original_progress_set,
        runtime_run_process=runtime_run_process,
        site_map=site_map,
    )

    original_process_single = getattr(app, "process_single_file", None)
    original_reset_ui = getattr(app, "reset_ui", None)
    tracker.set_current_item(stage="准备中")

    def patched_progress_set(value, *args, **kwargs):
        frame = inspect.currentframe().f_back
        try:
            if frame is not None and tracker.runtime_code is not None and frame.f_code is tracker.runtime_code:
                if tracker.handle_runtime_progress_call(frame.f_lasti, value):
                    return None
            if frame is not None and tracker.should_ignore_reset_zero(frame.f_code, value):
                return None
            return original_progress_set(value, *args, **kwargs)
        finally:
            del frame

    def patched_process_single_file(*args, **kwargs):
        current_path = _extract_progress_path_from_call(args, kwargs)
        tracker.set_current_item(current_path, "处理中")
        try:
            return original_process_single(*args, **kwargs)
        finally:
            tracker.note_process_single_complete()

    def patched_reset_ui(*args, **kwargs):
        tracker.handle_reset()
        return original_reset_ui(*args, **kwargs)

    progress_bar.set = patched_progress_set
    if callable(original_process_single):
        app.process_single_file = patched_process_single_file
    if callable(original_reset_ui):
        app.reset_ui = patched_reset_ui
    app._fx_progress_tracker = tracker

    def restore():
        try:
            progress_bar.set = original_progress_set
        except Exception:
            pass
        try:
            if callable(original_process_single):
                app.process_single_file = original_process_single
        except Exception:
            pass
        try:
            if callable(original_reset_ui):
                app.reset_ui = original_reset_ui
        except Exception:
            pass
        try:
            if getattr(app, "_fx_progress_tracker", None) is tracker:
                app._fx_progress_tracker = None
        except Exception:
            pass

    tracker.restore = restore
    return tracker


def _get_zip_mode(app):
    mode = ""
    mode_var = getattr(app, "zip_mode_var", None)
    if mode_var is not None:
        try:
            mode = mode_var.get()
        except Exception:
            mode = ""
    return normalize_zip_mode(mode)


def _build_zip_plan_messages(input_path, mode, max_depth=None, archive_policy="reuse_existing"):
    source = Path(str(input_path or "")).resolve()
    mode = normalize_zip_mode(mode)
    archive_policy = normalize_zip_archive_policy(archive_policy)
    min_depth, normalized_max_depth = normalize_zip_depth_range(max_depth)
    try:
        jobs = plan_zip_archives(source, mode, max_depth=max_depth)
    except Exception:
        jobs = []

    messages = []
    if mode in {"recursive", "smart_recursive"}:
        if min_depth <= 1 and normalized_max_depth is None:
            messages.append("ℹ️ [ZIP] 压缩层级范围：不限。")
        elif normalized_max_depth is None:
            messages.append(f"ℹ️ [ZIP] 压缩层级范围：第 {min_depth} 层及更深（根目录为第 1 层）。")
        elif min_depth <= 1:
            messages.append(f"ℹ️ [ZIP] 压缩层级范围：1-{normalized_max_depth} 层（根目录为第 1 层）。")
        else:
            messages.append(f"ℹ️ [ZIP] 压缩层级范围：{min_depth}-{normalized_max_depth} 层（根目录为第 1 层）。")
    if mode in {"recursive", "smart_recursive"}:
        if archive_policy == "rebuild_existing":
            messages.append("ℹ️ [ZIP] 已有同名压缩包策略：删除旧包并重新压缩。")
        else:
            messages.append("ℹ️ [ZIP] 已有同名压缩包策略：复用已有压缩包（断点续跑）。")
    if mode == "smart_recursive":
        messages.append(
            "ℹ️ [ZIP] 智能混合：每层先生成压缩包；遇到“普通文件 + 子文件夹”的层会停止继续下钻。"
        )

    if jobs:
        messages.append(f"📦 [ZIP] 本次预计生成 {len(jobs)} 个压缩包。")
        preview_limit = 5
        for job in jobs[:preview_limit]:
            output = Path(job.get("output", ""))
            try:
                output_text = str(output.relative_to(source))
            except Exception:
                output_text = str(output)
            messages.append(f"   → {output_text}")
        if len(jobs) > preview_limit:
            messages.append(f"   … 还有 {len(jobs) - preview_limit} 个压缩包未展开显示")
    return messages


def _run_zip_task_with_core(app, input_folder):
    normalized_input = _normalize_input_path_value(input_folder)
    mode = _get_zip_mode(app)
    max_depth = _get_zip_max_depth(app)
    archive_policy = _get_zip_archive_policy(app)
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "zip")
    _set_task_result_output_strategy(result, "zip", _get_task_output_strategy(app, "zip"))

    input_path = Path(normalized_input) if normalized_input else Path(".")
    output_root = input_path.parent if input_path.is_file() else input_path
    _set_task_result_output_root(result, output_root)

    tracker = _get_active_progress_tracker(app)
    progress_bar = getattr(app, "progress_bar", None)
    last_completed = {"value": 0}
    stage_labels = {
        "preparing": "准备压缩",
        "compressing": "正在压缩",
        "done": "压缩完成",
    }

    def report_progress(completed=0, total=1, current_path="", stage="compressing", fraction=0.0):
        stage_label = stage_labels.get(str(stage or ""), str(stage or "正在压缩"))
        current_name = os.path.basename(str(current_path or "").rstrip("\\/")) or str(current_path or "")
        if tracker is not None:
            try:
                tracker.total_units = max(1, int(total or tracker.total_units or 1))
            except Exception:
                pass
            if stage == "done":
                delta = max(0, int(completed or 0) - int(last_completed.get("value") or 0))
                if current_name:
                    tracker.set_current_item(current_name, stage_label)
                if delta > 0:
                    tracker.complete_units(delta, stage_label)
                    last_completed["value"] = int(completed or 0)
            else:
                tracker.set_current_item(current_name, stage_label)
            return

        try:
            if progress_bar is not None:
                progress_bar.set(_clamp_progress_value(fraction))
        except Exception:
            pass
        try:
            _set_progress_status(
                app,
                current_file=current_name,
                stage=stage_label,
                fraction=_clamp_progress_value(fraction),
                completed=int(completed or 0),
                total=max(1, int(total or 1)),
            )
        except Exception:
            pass

    def log_message(message):
        try:
            app.log(str(message))
        except Exception:
            pass

    try:
        log_message(f"[ZIP] mode={mode} max_depth={max_depth or 'unlimited'} archive_policy={archive_policy} input={normalized_input}")
        for message in _build_zip_plan_messages(normalized_input, mode, max_depth=max_depth, archive_policy=archive_policy):
            log_message(message)
        core_result = _zip_core_run_zip_task(
            normalized_input,
            mode,
            max_depth=max_depth,
            archive_policy=archive_policy,
            progress=report_progress,
            stop_requested=lambda: bool(getattr(app, "stop_event", False)),
            log=log_message,
        )
    except Exception as exc:
        _set_task_result_counts(result, processed=0, success=0, failed=1, skipped=0)
        result["failed_items"] = [normalized_input]
        _set_task_result_finished(result, "failed", message=str(exc), detail=str(exc), error=str(exc))
        raise

    outputs = list(core_result.get("outputs") or [])
    for output in outputs:
        _add_task_result_output(result, output)
    failed_items = list(core_result.get("failed_items") or [])
    result["failed_items"] = failed_items
    _set_task_result_counts(
        result,
        processed=core_result.get("processed_count", len(outputs) + len(failed_items)),
        success=core_result.get("success_count", len(outputs)),
        failed=core_result.get("failed_count", len(failed_items)),
        skipped=core_result.get("skipped_count", 0),
    )

    status = str(core_result.get("status") or "unknown")
    message = str(core_result.get("message") or "")
    if status == "success":
        log_message(f"[ZIP] completed: {len(outputs)} archive(s)")
        _set_task_result_finished(result, "success", message=message or "ZIP completed", detail=message or "ZIP completed")
    elif status == "stopped":
        log_message("[ZIP] stopped")
        _set_task_result_finished(result, "stopped", message=message or "ZIP stopped", detail=message or "ZIP stopped", stopped=True)
    elif status == "skipped":
        log_message("[ZIP] skipped")
        _set_task_result_finished(result, "skipped", message=message or "ZIP skipped", detail=message or "ZIP skipped", skipped=True)
    else:
        log_message(f"[ZIP] failed: {message}")
        _set_task_result_finished(
            result,
            "failed",
            message=message or "ZIP failed",
            detail=message or "ZIP failed",
            error=message or "ZIP failed",
        )
    return outputs


def _run_remove_wm_pdf_roundtrip(app, pdf_files, input_folder, output_folder, mode=None):
    failed_list = []
    skipped_count = 0
    preserve_mine = False
    remove_mode = _coerce_remove_wm_mode(mode or _get_remove_wm_mode(app))
    tracker = _get_active_progress_tracker(app)
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, input_folder, "remove_wm")
    _set_task_result_output_root(result, output_folder)
    if getattr(app, "rm_wm_preserve_mine", None) is not None:
        try:
            preserve_mine = bool(app.rm_wm_preserve_mine.get())
        except Exception:
            preserve_mine = False

    try:
        app.log(f"[去水印强度] 当前模式：{_get_remove_wm_mode_label(remove_mode)}")
    except Exception:
        pass

    temp_root = Path(tempfile.mkdtemp(prefix="fx_rm_pdf_"))
    pythoncom.CoInitialize()
    try:
        total = max(1, len(pdf_files))
        for index, src in enumerate(pdf_files):
            if getattr(app, "stop_event", False):
                try:
                    app.log("⏹️ [停止] 去水印任务已被用户中止")
                except Exception:
                    pass
                break

            rel = os.path.relpath(src, input_folder)
            dst_pdf = os.path.join(output_folder, rel)
            os.makedirs(os.path.dirname(dst_pdf), exist_ok=True)
            if is_nonempty_file(dst_pdf):
                try:
                    app.log(format_resume_skip_message("去水印断点续跑", src, dst_pdf))
                except Exception:
                    pass
                _add_task_result_output(result, dst_pdf)
                skipped_count += 1
                if tracker is not None:
                    tracker.complete_units(1)
                else:
                    try:
                        app.progress_bar.set((index + 1) / total)
                    except Exception:
                        pass
                continue
            stage_dir = temp_root / f"job_{index:03d}"
            stage_dir.mkdir(parents=True, exist_ok=True)
            stage_docx = stage_dir / "from_pdf.docx"
            cleaned_docx = stage_dir / "cleaned.docx"
            if tracker is not None:
                tracker.set_current_item(src, "PDF 转 Word")

            try:
                app.log(f"🧼 [PDF去水印] 正在处理：{os.path.basename(src)}")
            except Exception:
                pass

            try:
                convert_status = convert_pdf_to_word(src, str(stage_docx))
                if convert_status != "SUCCESS" or not stage_docx.exists():
                    raise RuntimeError(f"convert_pdf_to_word:{convert_status}")
                if tracker is not None:
                    tracker.set_current_item_fraction(0.35, stage="去除水印", current_file=src)

                chosen_docx = stage_docx
                word_for_remove = _create_hidden_word_app()
                try:
                    remove_status = remove_watermark_from_word(
                        word_for_remove,
                        str(stage_docx),
                        str(cleaned_docx),
                        preserve_mine=preserve_mine,
                        is_pdf_source=True,
                        mode=remove_mode,
                    )
                finally:
                    try:
                        word_for_remove.Quit()
                    except Exception:
                        pass

                if remove_status != "SUCCESS":
                    raise RuntimeError(f"remove_watermark_from_word:{remove_status}")
                if not cleaned_docx.exists():
                    raise RuntimeError("remove_watermark_from_word:no_output")
                chosen_docx = cleaned_docx
                if tracker is not None:
                    tracker.set_current_item_fraction(0.70, stage="导出 PDF", current_file=src)

                word_for_pdf = _create_hidden_word_app()
                try:
                    with _DisableWin32ComGenCache():
                        pdf_status = convert_doc_to_pdf(word_for_pdf, str(chosen_docx), dst_pdf)
                finally:
                    try:
                        word_for_pdf.Quit()
                    except Exception:
                        pass
                if pdf_status != "SUCCESS" or not os.path.exists(dst_pdf):
                    try:
                        app.log("[提示] [PDF去水印] Word 导出 PDF 兼容模式重试中...")
                    except Exception:
                        pass
                    pdf_status = _export_word_docx_to_pdf_safely(str(chosen_docx), dst_pdf)
                if pdf_status != "SUCCESS" or not os.path.exists(dst_pdf):
                    raise RuntimeError(f"export_docx_to_pdf:{pdf_status}")
            except Exception as exc:
                failed_list.append(rel)
                try:
                    app.log(f"❌ [失败] PDF 去水印错误: {os.path.basename(src)}: {exc}")
                except Exception:
                    pass
            else:
                try:
                    app.log(f"✅ [PDF去水印] 已输出：{os.path.basename(dst_pdf)}")
                except Exception:
                    pass
                _add_task_result_output(result, dst_pdf)
            finally:
                if tracker is not None:
                    tracker.complete_units(1)
                else:
                    try:
                        app.progress_bar.set((index + 1) / total)
                    except Exception:
                        pass
    finally:
        pythoncom.CoUninitialize()
        shutil.rmtree(temp_root, ignore_errors=True)
    _set_task_result_counts(result, processed=len(pdf_files), success=len(pdf_files) - len(failed_list), failed=len(failed_list), skipped=skipped_count)
    if failed_list:
        result["failed_items"] = list(failed_list)
    return failed_list


def _run_remove_wm_task(app, input_folder, original_run_process):
    normalized_input = _normalize_input_path_value(input_folder)
    input_root = os.path.dirname(normalized_input) if normalized_input and os.path.isfile(normalized_input) else normalized_input
    output_folder = os.path.join(input_root, RESULT_FOLDER_NAME)
    os.makedirs(output_folder, exist_ok=True)
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "remove_wm")
    _set_task_result_output_root(result, output_folder)

    all_files = app.collect_input_files(normalized_input, "remove_wm")
    pdf_files = [path for path in all_files if path.lower().endswith(".pdf")]
    other_files = [path for path in all_files if not path.lower().endswith(".pdf")]

    if not pdf_files:
        original_run_process(app, input_folder, "remove_wm")
        return _infer_task_result_from_context(app, normalized_input, "remove_wm")

    failed_list = []

    if other_files:
        staging_root = Path(tempfile.mkdtemp(prefix="fx_rm_mix_"))
        staging_output = staging_root / RESULT_FOLDER_NAME
        original_reset_ui = getattr(app, "reset_ui", None)
        try:
            for src in other_files:
                rel = os.path.relpath(src, input_root)
                staged = staging_root / rel
                staged.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(src, staged)

            if callable(original_reset_ui):
                app.reset_ui = lambda: None
            original_run_process(app, str(staging_root), "remove_wm")
            _copy_tree_contents(staging_output, output_folder, skip_names={"!澶辫触鏂囦欢娓呭崟.txt"})
            failed_list.extend(_read_failed_report_items(staging_output / "!澶辫触鏂囦欢娓呭崟.txt"))
        finally:
            if callable(original_reset_ui):
                app.reset_ui = original_reset_ui
            shutil.rmtree(staging_root, ignore_errors=True)

    failed_list.extend(_run_remove_wm_pdf_roundtrip(app, pdf_files, input_root, output_folder))

    if failed_list:
        deduped_failed = []
        seen = set()
        for item in failed_list:
            if item in seen:
                continue
            seen.add(item)
            deduped_failed.append(item)
        app.log("\n========= ❌ 失败清单 =========")
        for item in deduped_failed:
            app.log(f"• {item}")
        report_path = _write_failed_report(output_folder, deduped_failed)
        if report_path:
            app.log(f"\n📄 [报告] 已生成报告: {report_path}")
        app.log("完成 (含错误)")
        app.log(f"去水印任务结束，但有 {len(deduped_failed)} 个文件处理失败。")
    elif not getattr(app, "stop_event", False):
        app.log("\n🎉 [完成] 去水印已全部处理完成！")


def _run_remove_wm_task(app, input_folder, original_run_process):
    normalized_input, input_root, output_folder, resolved_strategy = _resolve_output_root_for_task(
        input_folder,
        "remove_wm",
        _get_task_output_strategy(app, "remove_wm"),
    )
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "remove_wm")
    is_single_input = bool(normalized_input and os.path.isfile(normalized_input))
    overwrite_requested = _get_remove_wm_overwrite_original(app)
    overwrite_original = is_single_input and overwrite_requested
    remove_mode = _get_remove_wm_mode(app)
    previous_remove_mode = _push_remove_wm_runtime_mode(remove_mode)
    tracker = _get_active_progress_tracker(app)
    total_items = 0

    single_output_root = None
    actual_strategy = resolved_strategy
    if is_single_input:
        actual_strategy = "overwrite" if overwrite_original else "same_dir"
        single_output_root = Path(tempfile.mkdtemp(prefix="fx_rm_single_out_"))
        output_folder = str(single_output_root / RESULT_FOLDER_NAME)
    _set_task_result_output_strategy(result, "remove_wm", actual_strategy)
    _set_task_result_output_root(result, input_root if is_single_input else output_folder)
    os.makedirs(output_folder, exist_ok=True)

    try:
        all_files = app.collect_input_files(normalized_input, "remove_wm")
        pdf_files = [path for path in all_files if path.lower().endswith(".pdf")]
        other_files = [path for path in all_files if not path.lower().endswith(".pdf")]
        total_items = len(pdf_files) + len(other_files)

        if not pdf_files and not is_single_input:
            return original_run_process(app, input_folder, "remove_wm")

        failed_list = []

        if is_single_input and not overwrite_original:
            resume_output = _build_single_remove_wm_base_output_path(normalized_input)
            if is_nonempty_file(resume_output):
                try:
                    app.log(format_resume_skip_message("去水印断点续跑", normalized_input, resume_output))
                except Exception:
                    pass
                result["outputs"] = []
                result["failed_items"] = []
                _add_task_result_output(result, resume_output)
                _set_task_result_output_root(result, os.path.dirname(resume_output))
                _set_task_result_counts(result, processed=max(1, total_items), success=1, failed=0, skipped=1)
                _set_task_result_finished(
                    result,
                    "success",
                    message="单文件去水印已复用已有结果",
                    detail=f"输出: {resume_output}",
                )
                return

        if other_files:
            staging_root = Path(tempfile.mkdtemp(prefix="fx_rm_mix_"))
            staging_output = staging_root / RESULT_FOLDER_NAME
            failed_report_path = staging_output / "!失败文件清单.txt"
            original_reset_ui = getattr(app, "reset_ui", None)
            try:
                for src in other_files:
                    rel = os.path.relpath(src, input_root)
                    staged = staging_root / rel
                    staged.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(src, staged)

                if callable(original_reset_ui):
                    app.reset_ui = lambda: None
                original_run_process(app, str(staging_root), "remove_wm")
                if tracker is not None:
                    tracker.finalize_pending()
                _copy_tree_contents(staging_output, output_folder, skip_names={failed_report_path.name})
                failed_list.extend(_read_failed_report_items(failed_report_path))
            finally:
                if callable(original_reset_ui):
                    app.reset_ui = original_reset_ui
                shutil.rmtree(staging_root, ignore_errors=True)

        if pdf_files:
            failed_list.extend(_run_remove_wm_pdf_roundtrip(app, pdf_files, input_root, output_folder, mode=remove_mode))
        skipped_count = int(result.get("skipped_count") or 0)

        if is_single_input:
            rel = os.path.relpath(normalized_input, input_root)
            staged_output_file = os.path.join(output_folder, rel)
            deduped_failed = list(dict.fromkeys(item for item in failed_list if item))
            if not deduped_failed and not os.path.exists(staged_output_file):
                deduped_failed.append(rel)

            if deduped_failed:
                app.log("\n========= ❌ 失败清单 =========")
                for item in deduped_failed:
                    app.log(f"• {item}")
                app.log("去水印未生成有效结果，原文件保持不变。")
                result["outputs"] = []
                result["failed_items"] = list(deduped_failed)
                _set_task_result_output_root(result, input_root)
                _set_task_result_counts(result, processed=max(1, total_items), success=0, failed=len(deduped_failed), skipped=0)
                _set_task_result_finished(
                    result,
                    "failed",
                    message="去水印未生成有效结果，原文件保持不变",
                    detail=f"失败 {len(deduped_failed)} 个文件",
                    error=f"失败 {len(deduped_failed)} 个文件",
                )
                return

            final_path = _finalize_single_remove_wm_output(
                app,
                normalized_input,
                staged_output_file,
                overwrite_original=overwrite_original,
            )
            if final_path:
                result["outputs"] = []
                result["failed_items"] = []
                _add_task_result_output(result, final_path)
                _set_task_result_output_root(result, os.path.dirname(final_path))
                _set_task_result_counts(result, processed=max(1, total_items), success=1, failed=0, skipped=0)
                _set_task_result_finished(
                    result,
                    "success",
                    message="单文件去水印已处理完成",
                    detail=f"输出: {final_path}",
                )
            else:
                result["outputs"] = []
                result["failed_items"] = [rel]
                _set_task_result_output_root(result, input_root)
                _set_task_result_counts(result, processed=max(1, total_items), success=0, failed=1, skipped=0)
                _set_task_result_finished(
                    result,
                    "failed",
                    message="去水印未生成有效结果",
                    detail="去水印未生成有效结果",
                    error="去水印未生成有效结果",
                )
            if final_path and not getattr(app, "stop_event", False):
                app.log("\n🎉 [完成] 单文件去水印已处理完成！")
            return

        if failed_list:
            deduped_failed = []
            seen = set()
            for item in failed_list:
                if item in seen:
                    continue
                seen.add(item)
                deduped_failed.append(item)
            app.log("\n========= ❌ 失败清单 =========")
            for item in deduped_failed:
                app.log(f"• {item}")
            report_path = _write_failed_report(output_folder, deduped_failed)
            if report_path:
                app.log(f"\n📄 [报告] 已生成报告: {report_path}")
                _add_task_result_output(result, report_path)
            app.log("完成 (含错误)")
            app.log(f"去水印任务结束，但有 {len(deduped_failed)} 个文件处理失败。")
            result["failed_items"] = list(deduped_failed)
            _set_task_result_counts(result, processed=total_items, success=max(0, total_items - len(deduped_failed)), failed=len(deduped_failed), skipped=skipped_count)
            _set_task_result_finished(
                result,
                "failed",
                message=f"去水印任务结束，但有 {len(deduped_failed)} 个文件处理失败。",
                detail=f"失败 {len(deduped_failed)} 个文件",
                error=f"失败 {len(deduped_failed)} 个文件",
            )
        elif not getattr(app, "stop_event", False):
            app.log("\n🎉 [完成] 去水印已全部处理完成！")
            _set_task_result_counts(result, processed=total_items, success=total_items, failed=0, skipped=skipped_count)
            _set_task_result_finished(
                result,
                "success",
                message="去水印已全部处理完成",
                detail=f"成功处理 {len(pdf_files) + len(other_files)} 个文件",
            )
        else:
            _set_task_result_counts(result, processed=total_items, success=max(0, total_items - len(failed_list)), failed=len(failed_list), skipped=skipped_count)
            _set_task_result_finished(result, "stopped", message="用户停止去水印任务", detail="用户停止去水印任务", stopped=True)
    finally:
        _pop_remove_wm_runtime_mode(previous_remove_mode)
        if single_output_root is not None:
            shutil.rmtree(single_output_root, ignore_errors=True)


def _run_pdf_ocr_task(app, input_folder):
    from tools.fx_pdf_ocr import (
        default_model_root,
        get_default_backend_key,
        get_default_profile_key,
        get_default_preprocess_key,
        resolve_model_root,
    )
    from tools.fx_pdf_ocr_task import (
        PdfOcrTaskCallbacks,
        PdfOcrTaskOptions,
        run_pdf_ocr_task_core,
    )

    normalized_input, input_root, output_folder, resolved_strategy = _resolve_output_root_for_task(
        input_folder,
        "pdf",
        _get_task_output_strategy(app, "pdf"),
    )
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "pdf")
    _set_task_result_output_strategy(result, "pdf", resolved_strategy)
    _set_task_result_output_root(result, output_folder)
    os.makedirs(output_folder, exist_ok=True)
    all_files = app.collect_input_files(normalized_input, "pdf")
    pdf_files = [f for f in all_files if f.lower().endswith(".pdf")]

    if not pdf_files:
        app.log("⚠️ [提示] 未找到可处理的 PDF 文件")
        _set_task_result_counts(result, processed=0, success=0, failed=0, skipped=1)
        _set_task_result_finished(result, "skipped", message="未找到可处理的 PDF 文件", detail="未找到可处理的 PDF 文件", skipped=True)
        return

    model_root = default_model_root()
    if getattr(app, "pdf_ocr_model_root", None) is not None:
        model_root = app.pdf_ocr_model_root.get().strip() or model_root
    resolved_model_root = resolve_model_root(model_root)

    language_display = ""
    if getattr(app, "pdf_ocr_language", None) is not None:
        language_display = app.pdf_ocr_language.get().strip()
    config_map = getattr(app, "_fx_pdf_ocr_lang_map", {})
    language_config = config_map.get(language_display, get_default_profile_key())

    backend_display = ""
    if getattr(app, "pdf_ocr_backend", None) is not None:
        backend_display = app.pdf_ocr_backend.get().strip()
    backend_map = getattr(app, "_fx_pdf_ocr_backend_map", {})
    backend_key = backend_map.get(backend_display, get_default_backend_key())

    mode_display = "mixed"
    if getattr(app, "pdf_ocr_mode", None) is not None:
        mode_display = app.pdf_ocr_mode.get().strip()
    mode_map = getattr(app, "_fx_pdf_ocr_mode_map", {})
    extraction_mode = mode_map.get(mode_display, mode_display or "mixed")
    if "|" in extraction_mode:
        extraction_mode = extraction_mode.split("|", 1)[0].strip()

    use_cls = False
    if getattr(app, "pdf_ocr_cls", None) is not None:
        use_cls = bool(app.pdf_ocr_cls.get())

    compare_report = False
    if getattr(app, "pdf_ocr_compare_report", None) is not None:
        compare_report = bool(app.pdf_ocr_compare_report.get())

    normalize_page_rotation = True
    if getattr(app, "pdf_ocr_normalize_rotation", None) is not None:
        normalize_page_rotation = bool(app.pdf_ocr_normalize_rotation.get())

    preprocess_display = ""
    if getattr(app, "pdf_ocr_preprocess", None) is not None:
        preprocess_display = app.pdf_ocr_preprocess.get().strip()
    preprocess_map = getattr(app, "_fx_pdf_ocr_preprocess_map", {})
    preprocess_mode = preprocess_map.get(preprocess_display, get_default_preprocess_key())

    password = ""
    if getattr(app, "pdf_pwd_entry", None) is not None:
        password = app.pdf_pwd_entry.get().strip()

    cpu_threads = max(1, min(os.cpu_count() or 4, 8))
    options = PdfOcrTaskOptions(
        model_root=resolved_model_root,
        profile_key=language_config,
        backend_key=backend_key,
        extraction_mode=extraction_mode,
        cls=use_cls,
        compare_report=compare_report,
        password=password,
        limit_side_len=2880,
        cpu_threads=cpu_threads,
        preprocess_mode=preprocess_mode,
        layered=True,
        normalize_page_rotation=normalize_page_rotation,
    )
    tracker = _get_active_progress_tracker(app)
    total = len(pdf_files)
    _clear_pdf_ocr_preview(app, "实时预览：正在准备 OCR 搜索版 PDF 任务...\n")

    def _on_engine_ready(engine_backend_key):
        app.log(f"🛡️ [安全模式] OCR 搜索版 PDF 使用单线程稳定处理，共 {total} 个文件...")
        app.log(f"🤖 [OCR] 风兮模型目录：{resolved_model_root}")
        app.log(f"🧩 [OCR] 后端：{engine_backend_key}{' (自动选择)' if backend_key == 'auto' else ''}")
        app.log(
            f"🧠 [OCR] 模型：{language_config} | 模式：{extraction_mode} | 图像增强：{preprocess_mode} | 修正文本方向（OCR）：{'开' if use_cls else '关'}"
            f" | 仅修正文本方向（不进行OCR）：{'开' if normalize_page_rotation else '关'} | 对比报告：{'开' if compare_report else '关'}"
        )

    def _on_file_started(src, _dst, _index, _total):
        app.log(f"🔎 [OCR] 正在处理：{os.path.basename(src)}")
        _set_pdf_ocr_preview_text(app, f"\n[{os.path.basename(src)}] 开始 OCR，识别内容会逐页显示在下方。\n")
        if tracker is not None:
            tracker.set_current_item(src, "OCR 准备")
            tracker.set_current_item_fraction(0.02, stage="OCR 准备", current_file=src)

    def _mark_compare_stage(src):
        if tracker is not None:
            tracker.set_current_item_fraction(0.08, stage="后端对比报告", current_file=src)

    def _on_compare_report(src, _report_path, report_result):
        app.log(f"🧪 [OCR] 已生成后端对比报告：{report_result['report_path']}")
        _mark_compare_stage(src)

    def _on_compare_report_failed(src, report_exc):
        app.log(f"⚠️ [OCR] 后端对比报告生成失败，已跳过：{report_exc}")
        _mark_compare_stage(src)

    def _on_page_progress(src, index, total_count, page_done, total_pages):
        if total_pages <= 0:
            return
        overall_fraction = page_done / total_pages
        if tracker is not None:
            tracker.set_current_item_fraction(
                overall_fraction,
                stage=f"OCR 第 {page_done}/{total_pages} 页",
                current_file=src,
            )
        else:
            app.progress_bar.set((index + overall_fraction) / total_count)

    def _on_page_preview(src, page_done, total_pages, page_payload):
        _append_pdf_ocr_preview_page(app, src, page_done, total_pages, page_payload)

    def _on_file_finished(src, dst, ocr_result):
        _add_task_result_output(result, dst)
        usage_text = ", ".join(f"{key}:{value}" for key, value in sorted((ocr_result.get("backend_usage") or {}).items()))
        if usage_text:
            app.log(f"🧭 [OCR] 实际后端使用：{usage_text}")
        rotation_result = dict(ocr_result.get("rotation_normalization") or {})
        changed_pages = int(rotation_result.get("changed_pages") or 0)
        if changed_pages:
            app.log(f"🧭 [OCR] 仅修正文本方向（不进行OCR）已修正 {changed_pages} 页，兼容 Noteful 等阅读器")
        app.log(f"✅ [OCR] 已生成可搜索 PDF：{os.path.basename(src)}")

    def _on_file_failed(src, _dst, _rel, exc):
        app.log(f"❌ [失败] OCR 错误: {os.path.basename(src)}: {exc}")

    def _on_file_completed(_src, _dst, index, total_count):
        if tracker is not None:
            tracker.complete_units(1)
        else:
            app.progress_bar.set((index + 1) / total_count)

    core_result = run_pdf_ocr_task_core(
        pdf_files,
        input_root,
        output_folder,
        options,
        PdfOcrTaskCallbacks(
            log=app.log,
            stop_requested=lambda: app.stop_event,
            on_engine_ready=_on_engine_ready,
            on_file_started=_on_file_started,
            on_page_progress=_on_page_progress,
            on_page_preview=_on_page_preview,
            on_file_finished=_on_file_finished,
            on_file_failed=_on_file_failed,
            on_file_completed=_on_file_completed,
            on_compare_report=_on_compare_report,
            on_compare_report_failed=_on_compare_report_failed,
        ),
    )
    if core_result.get("stopped"):
        app.log("⏹️ [停止] OCR 任务已被用户中止")
    failed_list = list(core_result.get("failed_items") or [])
    success_count = int(core_result.get("success_count") or 0)
    processed_count = int(core_result.get("processed_count") or 0)
    skipped_count = int(core_result.get("skipped_count") or 0)

    if failed_list:
        result["failed_items"] = list(failed_list)
        _set_task_result_counts(result, processed=len(pdf_files), success=success_count, failed=len(failed_list), skipped=skipped_count)
        app.log("\n========= ❌ 失败清单 =========")
        for item in failed_list:
            app.log(f"• {item}")
        report_path = _write_failed_report(output_folder, failed_list)
        if report_path:
            app.log(f"\n📄 [报告] 已生成报告: {report_path}")
            _add_task_result_output(result, report_path)
        app.log(f"完成 (含错误)")
        app.log(f"OCR 搜索版 PDF 任务结束，但有 {len(failed_list)} 个文件处理失败。")
        _set_task_result_finished(
            result,
            "failed",
            message=f"OCR 搜索版 PDF 任务结束，但有 {len(failed_list)} 个文件处理失败。",
            detail=f"失败 {len(failed_list)} 个文件",
            error=f"失败 {len(failed_list)} 个文件",
        )
    elif not app.stop_event:
        _set_task_result_counts(result, processed=len(pdf_files), success=success_count, failed=0, skipped=skipped_count)
        app.log("\n🎉 [完成] OCR 搜索版 PDF 已全部生成！")
        _set_task_result_finished(
            result,
            "success",
            message="OCR 搜索版 PDF 已全部生成",
            detail=f"成功处理 {success_count} 个文件",
        )
    else:
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=len(failed_list), skipped=skipped_count)
        _set_task_result_finished(result, "stopped", message="用户停止 OCR 任务", detail="用户停止 OCR 任务", stopped=True)


def _run_pdf_rotation_task(app, input_folder):
    from tools.fx_pdf_rotation_task import (
        PdfRotationTaskCallbacks,
        PdfRotationTaskOptions,
        run_pdf_rotation_task_core,
    )

    normalized_input, input_root, output_folder, resolved_strategy = _resolve_output_root_for_task(
        input_folder,
        "pdf",
        _get_task_output_strategy(app, "pdf"),
    )
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "pdf")
    _set_task_result_output_strategy(result, "pdf", resolved_strategy)
    _set_task_result_output_root(result, output_folder)
    os.makedirs(output_folder, exist_ok=True)
    all_files = app.collect_input_files(normalized_input, "pdf")
    pdf_files = [f for f in all_files if f.lower().endswith(".pdf")]
    if not pdf_files:
        app.log("⚠️ [提示] 未找到可修正角度的 PDF 文件")
        _set_task_result_counts(result, processed=0, success=0, failed=0, skipped=1)
        _set_task_result_finished(result, "skipped", message="未找到可修正角度的 PDF 文件", detail="未找到可修正角度的 PDF 文件", skipped=True)
        return

    password = ""
    if getattr(app, "pdf_pwd_entry", None) is not None:
        try:
            password = app.pdf_pwd_entry.get().strip()
        except Exception:
            password = ""
    delete_source = False
    if getattr(app, "pdf_delete_var", None) is not None:
        try:
            delete_source = bool(app.pdf_delete_var.get())
        except Exception:
            delete_source = False

    tracker = _get_active_progress_tracker(app)
    total = len(pdf_files)
    app.log(f"🧭 [仅修正文本方向（不进行OCR）] 共 {total} 个 PDF，输出到结果文件夹；原文件保持不变。")

    def _on_file_started(src, _dst, _index, _total):
        app.log(f"📄 [仅修正文本方向（不进行OCR）] 正在处理：{os.path.basename(src)}")
        if tracker is not None:
            tracker.set_current_item(src, "仅修正文本方向（不进行OCR）")
            tracker.set_current_item_fraction(0.02, stage="仅修正文本方向（不进行OCR）", current_file=src)

    def _on_file_finished(src, dst, payload):
        _add_task_result_output(result, dst)
        rotation_result = dict(payload.get("rotation_normalization") or {})
        changed_pages = int(rotation_result.get("changed_pages") or 0)
        if payload.get("resumed"):
            app.log(f"⏭️ [仅修正文本方向（不进行OCR）] 已复用结果：{os.path.basename(dst)}")
        elif changed_pages:
            app.log(f"✅ [仅修正文本方向（不进行OCR）] {os.path.basename(src)}：修正 {changed_pages} 页")
        else:
            app.log(f"✅ [仅修正文本方向（不进行OCR）] {os.path.basename(src)}：未发现页面旋转，已复制结果")

    def _on_file_failed(src, _dst, _rel, exc):
        app.log(f"❌ [失败] 仅修正文本方向（不进行OCR）错误：{os.path.basename(src)}：{exc}")

    def _on_file_completed(_src, _dst, index, total_count):
        if tracker is not None:
            tracker.complete_units(1)
        else:
            app.progress_bar.set((index + 1) / total_count)

    core_result = run_pdf_rotation_task_core(
        pdf_files,
        input_root,
        output_folder,
        PdfRotationTaskOptions(password=password, delete_source=delete_source),
        PdfRotationTaskCallbacks(
            stop_requested=lambda: app.stop_event,
            on_file_started=_on_file_started,
            on_file_finished=_on_file_finished,
            on_file_failed=_on_file_failed,
            on_file_completed=_on_file_completed,
        ),
    )
    failed_list = list(core_result.get("failed_items") or [])
    success_count = int(core_result.get("success_count") or 0)
    processed_count = int(core_result.get("processed_count") or 0)
    skipped_count = int(core_result.get("skipped_count") or 0)

    if failed_list:
        result["failed_items"] = list(failed_list)
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=len(failed_list), skipped=skipped_count)
        report_path = _write_failed_report(output_folder, failed_list)
        if report_path:
            _add_task_result_output(result, report_path)
        _set_task_result_finished(
            result,
            "failed",
            message=f"仅修正文本方向（不进行OCR）结束，但有 {len(failed_list)} 个文件处理失败。",
            detail=f"失败 {len(failed_list)} 个文件",
            error=f"失败 {len(failed_list)} 个文件",
        )
    elif core_result.get("stopped") or app.stop_event:
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=0, skipped=skipped_count)
        _set_task_result_finished(result, "stopped", message="用户停止仅修正文本方向（不进行OCR）任务", detail="用户停止仅修正文本方向（不进行OCR）任务", stopped=True)
    else:
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=0, skipped=skipped_count)
        app.log(f"🎉 [完成] 仅修正文本方向（不进行OCR）完成，共修正 {core_result.get('changed_pages', 0)} 页。")
        _set_task_result_finished(result, "success", message="仅修正文本方向（不进行OCR）已完成", detail=f"成功处理 {success_count} 个文件")


def _is_parallel_enabled(app):
    var = getattr(app, "enable_multithread", None)
    if var is None:
        return False
    try:
        return bool(var.get())
    except Exception:
        return False


def _get_parallel_worker_count(item_count):
    if item_count <= 1:
        return 1
    cpu_count = os.cpu_count() or 4
    return max(1, min(int(item_count), PARALLEL_MAX_WORKERS, max(2, cpu_count)))


def _get_watermark_pdf_worker_count(item_count):
    """Keep PDF rewrites parallel without exhausting RAM on large documents."""

    if item_count <= 1:
        return 1
    cpu_count = os.cpu_count() or 4
    return max(1, min(int(item_count), WATERMARK_PDF_MAX_WORKERS, max(2, cpu_count)))


def _reserve_unique_output_path(src, output_folder, builder, reserved):
    target = Path(builder(src, output_folder))
    target_dir = target.parent
    stem = target.stem
    suffix = target.suffix
    counter = 2
    normalized = os.path.normcase(str(target))
    while normalized in reserved:
        target = target_dir / f"{stem}_{counter}{suffix}"
        normalized = os.path.normcase(str(target))
        counter += 1
    reserved.add(normalized)
    return str(target)


def _get_pdf_compress_profile(app):
    level_var = getattr(app, "pdf_compress_level_var", None)
    image_var = getattr(app, "pdf_image_compress_level_var", None)
    try:
        level = str(level_var.get() or "标准").strip() if level_var is not None else "标准"
    except Exception:
        level = "标准"
    try:
        image_level = str(image_var.get() or "标准").strip() if image_var is not None else "标准"
    except Exception:
        image_level = "标准"
    return (
        level if level in PDF_COMPRESS_LEVELS else "标准",
        image_level if image_level in PDF_IMAGE_COMPRESS_LEVELS else "标准",
    )


def _build_pdf_compress_output_path(src, output_folder):
    return _pdf_compress_core_build_output_path(src, output_folder)


def _build_pdf_compress_resume_output_path(src, output_folder):
    source = Path(src)
    return str(Path(output_folder) / f"{source.stem}_压缩{source.suffix}")


def compress_pdf_file(src, dst, compress_level="标准", image_level="标准", password=""):
    return _pdf_compress_core_compress_pdf_file(src, dst, compress_level, image_level, password=password)


def _run_pdf_compress_task(app, input_folder):
    normalized_input, input_root, output_folder, resolved_strategy = _resolve_output_root_for_task(
        input_folder,
        "pdf",
        _get_task_output_strategy(app, "pdf"),
    )
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "pdf")
    _set_task_result_output_strategy(result, "pdf", resolved_strategy)
    _set_task_result_output_root(result, output_folder)
    all_files = app.collect_input_files(normalized_input, "pdf")
    pdf_files = [f for f in all_files if f.lower().endswith(".pdf")]
    if not pdf_files:
        app.log("⚠️ 未找到可压缩的 PDF 文件。")
        _set_task_result_counts(result, processed=0, success=0, failed=0, skipped=1)
        _set_task_result_finished(result, "skipped", message="未找到可压缩的 PDF 文件", detail="未找到可压缩的 PDF 文件", skipped=True)
        return

    os.makedirs(output_folder, exist_ok=True)
    compress_level, image_level = _get_pdf_compress_profile(app)
    password = ""
    if getattr(app, "pdf_pwd_entry", None) is not None:
        try:
            password = app.pdf_pwd_entry.get().strip()
        except Exception:
            password = ""

    tracker = _get_active_progress_tracker(app)
    failed_list = []
    success_count = 0
    resumed_count = 0
    total = len(pdf_files)
    app.log(f"📉 [PDF 压缩] 共 {total} 个 PDF，压缩程度：{compress_level}，图片压缩：{image_level}")
    profile_stamp_by_src = {
        src: _pdf_compress_core_build_profile_stamp(src, compress_level, image_level, password=password)
        for src in pdf_files
    }

    reserved_outputs = set()
    stem_counts = {}
    for src in pdf_files:
        stem = Path(src).stem.lower()
        stem_counts[stem] = stem_counts.get(stem, 0) + 1
    jobs = []
    for src in pdf_files:
        resume_output = _build_pdf_compress_resume_output_path(src, output_folder)
        if (
            stem_counts.get(Path(src).stem.lower(), 0) == 1
            and is_nonempty_file(resume_output)
            and _pdf_compress_core_meta_matches(resume_output, profile_stamp_by_src.get(src))
        ):
            jobs.append((src, resume_output))
        else:
            jobs.append((src, _reserve_unique_output_path(src, output_folder, _build_pdf_compress_output_path, reserved_outputs)))
    should_delete_source = False
    if getattr(app, "pdf_delete_var", None) is not None:
        try:
            should_delete_source = bool(app.pdf_delete_var.get())
        except Exception:
            should_delete_source = False

    def process_one_pdf_compress(job):
        src, dst = job
        if is_nonempty_file(dst):
            return {"src": src, "dst": dst, "ok": True, "status": "RESUME", "ratio": 0, "image_changes": "0", "resumed": True}
        before_size = os.path.getsize(src)
        status = compress_pdf_file(src, dst, compress_level, image_level, password=password)
        if not status.startswith("SUCCESS"):
            return {"src": src, "dst": dst, "ok": False, "status": status}
        after_size = os.path.getsize(dst)
        ratio = 0 if before_size <= 0 else max(0, round((1 - after_size / before_size) * 100, 1))
        status_parts = status.split(":")
        image_changes = status_parts[1] if len(status_parts) > 1 else "0"
        engine = status_parts[2] if len(status_parts) > 2 else "pymupdf"
        _pdf_compress_core_write_meta(
            dst,
            profile_stamp_by_src.get(src),
            {
                "status": status,
                "engine": engine,
                "source_size": before_size,
                "output_size": after_size,
                "ratio": ratio,
                "image_changes": image_changes,
            },
        )
        return {
            "src": src,
            "dst": dst,
            "ok": True,
            "status": status,
            "ratio": ratio,
            "image_changes": image_changes,
            "engine": engine,
        }

    def _log_pdf_compress_success(item):
        dst = item["dst"]
        engine = item.get("engine") or "pymupdf"
        if engine in {"original", "kept_original"}:
            suffix = " | 保留原大小，避免越压越大"
        elif engine == "ghostscript":
            suffix = " | 引擎 ghostscript"
        elif engine == "optimized":
            suffix = " | 引擎 内置优化"
        elif engine == "pikepdf":
            suffix = " | 引擎 对象流优化"
        elif engine == "pymupdf":
            suffix = " | 引擎 内置"
        elif engine == "rasterized":
            suffix = " | 引擎 图片化"
        else:
            suffix = f" | 引擎 {engine}"
        app.log(
            f"✅ [PDF 压缩] {os.path.basename(dst)} | 减少 {item.get('ratio', 0)}% | 图片 {item.get('image_changes', '0')} 项{suffix}"
        )

    parallel_workers = _get_parallel_worker_count(total) if _is_parallel_enabled(app) else 1
    if parallel_workers > 1:
        app.log(f"🚀 [批量并行] PDF 压缩启用 {parallel_workers} 个线程。")
        futures = {}
        with concurrent.futures.ThreadPoolExecutor(max_workers=parallel_workers) as executor:
            for job in jobs:
                if getattr(app, "stop_event", False):
                    app.log("⏹️ [停止] PDF 压缩任务已被用户中止")
                    break
                src, _dst = job
                if tracker is not None:
                    tracker.set_current_item(src, "PDF 压缩")
                app.log(f"📄 [PDF 压缩] 已加入并行处理：{os.path.basename(src)}")
                futures[executor.submit(process_one_pdf_compress, job)] = job
            for future in concurrent.futures.as_completed(futures):
                src, _dst = futures[future]
                try:
                    item = future.result()
                except Exception as exc:
                    item = {"src": src, "dst": _dst, "ok": False, "status": str(exc)}
                if not item.get("ok"):
                    failed_list.append(f"{item.get('src', src)}: {item.get('status')}")
                    app.log(f"❌ [失败] {os.path.basename(str(item.get('src', src)))}: {item.get('status')}")
                else:
                    _log_pdf_compress_success(item)
                    success_count += 1
                    if item.get("resumed"):
                        resumed_count += 1
                    _add_task_result_output(result, dst)
                    if should_delete_source and not item.get("resumed"):
                        try:
                            os.remove(item["src"])
                            app.log(f"🗑️ 已删除源文件：{os.path.basename(item['src'])}")
                        except Exception as exc:
                            failed_list.append(f"{item['src']}: 删除源文件失败: {exc}")
                if tracker is not None:
                    tracker.complete_units(1)
                else:
                    app.progress_bar.set(min(1.0, (success_count + len(failed_list)) / total))
    else:
        for index, (src, dst) in enumerate(jobs):
            if getattr(app, "stop_event", False):
                app.log("⏹️ [停止] PDF 压缩任务已被用户中止")
                break
            try:
                if tracker is not None:
                    tracker.set_current_item(src, "PDF 压缩")
                app.log(f"📄 [PDF 压缩] 正在处理：{os.path.basename(src)}")
                item = process_one_pdf_compress((src, dst))
                if not item.get("ok"):
                    failed_list.append(f"{src}: {item.get('status')}")
                    app.log(f"❌ [失败] {os.path.basename(src)}: {item.get('status')}")
                    continue
                _log_pdf_compress_success(item)
                success_count += 1
                if item.get("resumed"):
                    resumed_count += 1
                _add_task_result_output(result, dst)
                if should_delete_source and not item.get("resumed"):
                    try:
                        os.remove(src)
                        app.log(f"🗑️ 已删除源文件：{os.path.basename(src)}")
                    except Exception as exc:
                        failed_list.append(f"{src}: 删除源文件失败: {exc}")
            except Exception as exc:
                failed_list.append(f"{src}: {exc}")
                app.log(f"❌ [失败] {os.path.basename(src)}: {exc}")
            finally:
                if tracker is not None:
                    tracker.complete_units(1)
                else:
                    app.progress_bar.set((index + 1) / total)

    if failed_list:
        result["failed_items"] = list(failed_list)
        _set_task_result_counts(result, processed=total, success=success_count, failed=len(failed_list), skipped=resumed_count)
        app.log("\n========= ❌ 失败清单 =========")
        for item in failed_list:
            app.log(f"• {item}")
        report_path = _write_failed_report(output_folder, failed_list)
        if report_path:
            app.log(f"\n📄 [报告] 已生成报告: {report_path}")
            _add_task_result_output(result, report_path)
        app.log(f"PDF 压缩任务结束，但有 {len(failed_list)} 个文件处理失败。")
        _set_task_result_finished(
            result,
            "failed",
            message=f"PDF 压缩任务结束，但有 {len(failed_list)} 个文件处理失败。",
            detail=f"失败 {len(failed_list)} 个文件",
            error=f"失败 {len(failed_list)} 个文件",
        )
    elif not getattr(app, "stop_event", False):
        _set_task_result_counts(result, processed=total, success=success_count, failed=0, skipped=resumed_count)
        app.log("\n🎉 [完成] PDF 压缩已全部完成！")
        _set_task_result_finished(
            result,
            "success",
            message="PDF 压缩已全部完成",
            detail=f"成功处理 {success_count} 个文件",
        )
    else:
        _set_task_result_counts(result, processed=success_count + len(failed_list), success=success_count, failed=len(failed_list), skipped=resumed_count)
        _set_task_result_finished(result, "stopped", message="用户停止 PDF 压缩任务", detail="用户停止 PDF 压缩任务", stopped=True)


def _collect_image_to_pdf_files(app, input_value):
    normalized_input = _normalize_input_path_value(input_value)
    return _image_pdf_task_collect_files(
        normalized_input,
        collect_input_files=getattr(app, "collect_input_files", None),
        valid_exts=IMAGE_TO_PDF_EXTS,
    )


def _build_image_pdf_output_path(src, output_folder):
    return _image_pdf_task_build_output_path(src, output_folder)


def _image_file_to_pdf(src, dst):
    return _image_pdf_task_image_file_to_pdf(src, dst)


def _get_image_pdf_mode(app):
    mode_var = getattr(app, "img_mode_var", None)
    try:
        return str(mode_var.get() or "").strip() if mode_var is not None else ""
    except Exception:
        return ""


def _get_convert_mode(app):
    mode_var = getattr(app, "cv_mode", None)
    try:
        raw_mode = mode_var.get() if mode_var is not None else ""
    except Exception:
        raw_mode = ""
    return _convert_core_normalize_mode(raw_mode)


def _get_convert_preview_detail(app):
    try:
        return _convert_core_describe_mode(_get_convert_mode(app), fallback="格式转换")
    except Exception:
        return "格式转换"


def _collect_convert_files(app, input_value):
    normalized_input = _normalize_input_path_value(input_value)
    return _convert_core_collect_files(
        normalized_input,
        _get_convert_mode(app),
        collect_input_files=getattr(app, "collect_input_files", None),
    )


def _run_convert_imgs_to_pdf_task(app, input_folder):
    normalized_input, input_root, output_folder, resolved_strategy = _resolve_output_root_for_task(
        input_folder,
        "image",
        "result_folder",
    )
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "convert")
    _set_task_result_output_strategy(result, "image", resolved_strategy)
    _set_task_result_output_root(result, output_folder)

    tracker = _get_active_progress_tracker(app)
    completed_units = {"value": 0}
    expected_total = len(_collect_convert_files(app, normalized_input))

    def _on_merge_started(dst, total):
        if tracker is not None:
            tracker.total_units = max(1, int(total or expected_total or 1))
            tracker.set_current_item(dst, "格式转换 · 多图合并 PDF")

    def _on_item_finished(src, dst, item):
        _add_task_result_output(result, dst)
        app.log(f"✅ [格式转换] 多图合并 PDF 已输出：{os.path.basename(dst)}")

    def _on_item_failed(src, _dst, status):
        app.log(f"❌ [失败] 多图合并 PDF: {status}")

    def _on_item_completed(units):
        completed_units["value"] += int(units or 0)
        if tracker is not None:
            tracker.complete_units(units)
        else:
            app.progress_bar.set(min(1.0, completed_units["value"] / max(1, expected_total)))

    core_result = run_convert_imgs_to_pdf_task_core(
        normalized_input,
        input_root=input_root,
        output_folder=output_folder,
        collect_input_files=getattr(app, "collect_input_files", None),
        merge_images_to_pdf=merge_images_to_pdf,
        callbacks=ConvertImgsToPdfCallbacks(
            log=app.log,
            stop_requested=lambda: getattr(app, "stop_event", False),
            on_merge_started=_on_merge_started,
            on_item_finished=_on_item_finished,
            on_item_failed=_on_item_failed,
            on_item_completed=_on_item_completed,
        ),
    )

    failed_list = list(core_result.get("failed_items") or [])
    success_count = int(core_result.get("success_count") or 0)
    processed_count = int(core_result.get("processed_count") or expected_total or 0)
    skipped_count = int(core_result.get("skipped_count") or 0)

    if core_result.get("status") == "skipped":
        message = core_result.get("message") or "未找到可合并为 PDF 的图片文件"
        app.log(f"⚠️ {message}")
        _set_task_result_counts(result, processed=0, success=0, failed=0, skipped=skipped_count or 1)
        _set_task_result_finished(result, "skipped", message=message, detail=message, skipped=True)
    elif failed_list:
        result["failed_items"] = list(failed_list)
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=len(failed_list), skipped=skipped_count)
        app.log("\n========= ❌ 失败清单 =========")
        for item in failed_list:
            app.log(f"• {item}")
        report_path = _write_failed_report(output_folder, failed_list)
        if report_path:
            app.log(f"\n📄 [报告] 已生成报告: {report_path}")
            _add_task_result_output(result, report_path)
        _set_task_result_finished(
            result,
            "failed",
            message=f"多图合并 PDF 任务结束，但有 {len(failed_list)} 个文件处理失败。",
            detail=f"失败 {len(failed_list)} 个文件",
            error=f"失败 {len(failed_list)} 个文件",
        )
    elif getattr(app, "stop_event", False):
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=0, skipped=skipped_count)
        _set_task_result_finished(result, "stopped", message="用户停止多图合并 PDF 任务", detail="用户停止多图合并 PDF 任务", stopped=True)
    else:
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=0, skipped=skipped_count)
        app.log("\n🎉 [完成] 格式转换多图合并 PDF 已全部完成！")
        _set_task_result_finished(
            result,
            "success",
            message="格式转换多图合并 PDF 已全部完成",
            detail=f"成功合并 {success_count} 张图片",
        )
    return result


CONVERT_EXTENDED_MODES = {"pdf2ppt", "txt2word", "md2pdf", "pdf2md"}


def _run_convert_extended_task(app, input_folder):
    mode = _get_convert_mode(app)
    normalized_input, input_root, output_folder, resolved_strategy = _resolve_output_root_for_task(
        input_folder,
        "convert",
        _get_task_output_strategy(app, "convert"),
    )
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "convert")
    _set_task_result_output_strategy(result, "convert", resolved_strategy)
    _set_task_result_output_root(result, output_folder)

    files = _convert_core_collect_files(
        normalized_input,
        mode,
        collect_input_files=getattr(app, "collect_input_files", None),
    )
    total = len(files)
    if total <= 0:
        message = f"未找到可用于{_convert_core_describe_mode(mode, fallback='格式转换')}的文件"
        app.log(f"⚠️ [格式转换] {message}")
        _set_task_result_counts(result, processed=0, success=0, failed=0, skipped=1)
        _set_task_result_finished(result, "skipped", message=message, detail=message, skipped=True)
        return result

    os.makedirs(output_folder, exist_ok=True)
    tracker = _get_active_progress_tracker(app)
    if tracker is not None:
        tracker.total_units = max(1, total)

    app.log(f"🔄 [格式转换] {_convert_core_describe_mode(mode, fallback='格式转换')} | 共 {total} 个文件")
    failed_items = []
    outputs = []
    skipped_count = 0

    try:
        context = ConvertFileContext(
            ppt_app=None,
            convert_pdf_to_ppt=lambda ppt_app_obj, src_path, dst_path: _convert_pdf_to_ppt_safely(ppt_app_obj, src_path, dst_path),
            convert_txt_to_word=convert_txt_to_word_file,
            convert_md_to_pdf=convert_md_to_pdf_file,
            convert_pdf_to_md=convert_pdf_to_md_file,
            copy_file_safe=copy_file_safe,
            log=getattr(app, "log", None),
        )

        for index, src in enumerate(files, start=1):
            if getattr(app, "stop_event", False):
                break
            if tracker is not None:
                tracker.set_current_item(src, _convert_core_describe_mode(mode, fallback="格式转换"))
            status_text = "准备"
            try:
                expected_output = _convert_core_plan_output_path(src, input_root, output_folder, mode)
                if is_nonempty_file(expected_output):
                    outputs.append(expected_output)
                    _add_task_result_output(result, expected_output)
                    skipped_count += 1
                    app.log(f"⏭️ [断点续跑] 已存在结果，跳过: {os.path.basename(src)}")
                    status_text = "断点续跑"
                else:
                    item_result = process_convert_file(src, input_root, output_folder, mode, context)
                    status = str(item_result.get("status") or "")
                    output_path = str(item_result.get("output") or expected_output)
                    if item_result.get("ok") and is_nonempty_file(output_path):
                        outputs.append(output_path)
                        _add_task_result_output(result, output_path)
                        if item_result.get("skipped"):
                            skipped_count += 1
                        app.log(f"✅ [格式转换] {os.path.basename(src)} -> {output_path}")
                        status_text = "完成"
                    else:
                        failed_items.append(str(src))
                        app.log(f"❌ [格式转换失败] {os.path.basename(src)} | {status or item_result.get('message') or 'unknown'}")
                        status_text = "失败"
            except Exception as exc:
                failed_items.append(str(src))
                app.log(f"❌ [格式转换失败] {os.path.basename(src)} | {exc}")
                status_text = "失败"
            finally:
                if tracker is not None:
                    tracker.complete_units(1, stage=status_text, current_file=src)
                else:
                    app.progress_bar.set(min(1.0, index / max(1, total)))
    finally:
        pass

    processed_count = len(outputs) + len(failed_items)
    success_count = len(outputs)
    if failed_items:
        result["failed_items"] = list(failed_items)
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=len(failed_items), skipped=skipped_count)
        app.log("\n========= ❌ 失败清单 =========")
        for item in failed_items:
            app.log(f"• {item}")
        report_path = _write_failed_report(output_folder, failed_items)
        if report_path:
            _add_task_result_output(result, report_path)
            app.log(f"\n📄 [报告] 已生成报告: {report_path}")
        _set_task_result_finished(
            result,
            "failed",
            message=f"格式转换完成，但有 {len(failed_items)} 个文件失败。",
            detail=f"成功 {success_count} 个，失败 {len(failed_items)} 个",
            error=f"失败 {len(failed_items)} 个文件",
        )
    elif getattr(app, "stop_event", False):
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=0, skipped=skipped_count)
        _set_task_result_finished(result, "stopped", message="用户停止格式转换任务", detail="用户停止格式转换任务", stopped=True)
    else:
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=0, skipped=skipped_count)
        app.log("\n🎉 [完成] 格式转换已全部完成！")
        _set_task_result_finished(
            result,
            "success",
            message="格式转换已全部完成",
            detail=f"成功处理 {success_count} 个文件",
        )
    return result


def _run_image_to_pdf_task(app, input_folder, merge=False):
    normalized_input, input_root, output_folder, resolved_strategy = _resolve_output_root_for_task(
        input_folder,
        "image",
        _get_task_output_strategy(app, "image"),
    )
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "image")
    _set_task_result_output_strategy(result, "image", resolved_strategy)
    _set_task_result_output_root(result, output_folder)
    image_files = _collect_image_to_pdf_files(app, normalized_input)
    if not image_files:
        app.log("⚠️ 未找到可转 PDF 的图片文件。")
        _set_task_result_counts(result, processed=0, success=0, failed=0, skipped=1)
        _set_task_result_finished(result, "skipped", message="未找到可转 PDF 的图片文件", detail="未找到可转 PDF 的图片文件", skipped=True)
        return

    tracker = _get_active_progress_tracker(app)
    should_delete = False
    if getattr(app, "img_delete_var", None) is not None:
        try:
            should_delete = bool(app.img_delete_var.get())
        except Exception:
            should_delete = False

    total = len(image_files)
    if merge:
        app.log(f"🧩 [多图合并PDF] 共 {total} 张图片，正在合并...")
    else:
        app.log(f"📄 [图片转PDF] 共 {total} 张图片，逐张生成 PDF...")

    def _on_merge_started(dst, _total):
        if tracker is not None:
            tracker.set_current_item(dst, "多图合并 PDF")

    def _on_item_started(src, _dst, _index, _total):
        if tracker is not None:
            tracker.set_current_item(src, "图片转 PDF")
        if _is_parallel_enabled(app) and not merge:
            app.log(f"🖼️ [图片转PDF] 已加入并行处理：{os.path.basename(src)}")

    def _on_item_finished(src, dst, item):
        _add_task_result_output(result, dst)
        if item.get("merge"):
            app.log(f"✅ [多图合并PDF] 已输出：{os.path.basename(dst)}")
        else:
            app.log(f"✅ [图片转PDF] {os.path.basename(src)} -> {os.path.basename(dst)}")

    def _on_item_failed(src, _dst, status):
        if merge:
            app.log(f"❌ [失败] 多图合并 PDF: {status}")
        else:
            app.log(f"❌ [失败] {os.path.basename(src)}: {status}")

    completed_units = {"value": 0}

    def _on_item_completed(units):
        completed_units["value"] += int(units or 0)
        if tracker is not None:
            tracker.complete_units(units)
        else:
            app.progress_bar.set(min(1.0, completed_units["value"] / max(1, total)))

    parallel_workers = _get_parallel_worker_count(total) if (_is_parallel_enabled(app) and not merge) else 1
    core_result = run_image_pdf_task_core(
        image_files,
        input_root,
        normalized_input,
        output_folder,
        ImagePdfTaskOptions(
            merge=merge,
            delete_source=should_delete,
            parallel_workers=parallel_workers,
            executor_factory=concurrent.futures.ThreadPoolExecutor,
            image_to_pdf=_image_file_to_pdf,
            merge_images_to_pdf=merge_images_to_pdf,
        ),
        ImagePdfTaskCallbacks(
            log=app.log,
            stop_requested=lambda: getattr(app, "stop_event", False),
            on_item_started=_on_item_started,
            on_item_finished=_on_item_finished,
            on_item_failed=_on_item_failed,
            on_item_completed=_on_item_completed,
            on_merge_started=_on_merge_started,
        ),
    )
    failed_list = list(core_result.get("failed_items") or [])
    success_count = int(core_result.get("success_count") or 0)
    processed_count = int(core_result.get("processed_count") or 0)
    skipped_count = int(core_result.get("skipped_count") or 0)

    if failed_list:
        result["failed_items"] = list(failed_list)
        _set_task_result_counts(result, processed=len(image_files), success=success_count, failed=len(failed_list), skipped=skipped_count)
        app.log("\n========= ❌ 失败清单 =========")
        for item in failed_list:
            app.log(f"• {item}")
        report_path = _write_failed_report(output_folder, failed_list)
        if report_path:
            app.log(f"\n📄 [报告] 已生成报告: {report_path}")
            _add_task_result_output(result, report_path)
        app.log(f"图片转 PDF 任务结束，但有 {len(failed_list)} 个文件处理失败。")
        _set_task_result_finished(
            result,
            "failed",
            message=f"图片转 PDF 任务结束，但有 {len(failed_list)} 个文件处理失败。",
            detail=f"失败 {len(failed_list)} 个文件",
            error=f"失败 {len(failed_list)} 个文件",
        )
    elif not getattr(app, "stop_event", False):
        _set_task_result_counts(result, processed=len(image_files), success=success_count, failed=0, skipped=skipped_count)
        app.log("\n🎉 [完成] 图片 PDF 任务已全部完成！")
        _set_task_result_finished(
            result,
            "success",
            message="图片 PDF 任务已全部完成",
            detail=f"成功处理 {success_count} 个文件",
        )
    else:
        _set_task_result_counts(result, processed=processed_count, success=success_count, failed=len(failed_list), skipped=skipped_count)
        _set_task_result_finished(result, "stopped", message="用户停止图片 PDF 任务", detail="用户停止图片 PDF 任务", stopped=True)


def _patch_pdf_ocr_mode():
    try:
        original_init_pdf_ui = FengxiToolboxApp.init_pdf_ui
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_pdf_ocr_mode:missing:{exc}")
        return

    if getattr(original_init_pdf_ui, "__fx_pdf_ocr_patch__", False):
        return

    def patched_init_pdf_ui(self):
        original_init_pdf_ui(self)
        if getattr(self, "_fx_pdf_ocr_ui_ready", False):
            return
        try:
            from tools.fx_pdf_ocr import (
                build_backend_display_map,
                build_backend_status_text,
                build_preprocess_display_map,
                build_profile_display_map,
                default_model_root,
                get_default_backend_display,
                get_default_preprocess_display,
                get_default_profile_display,
            )

            card = self.tab_pdf.winfo_children()[0]
            card_children = card.winfo_children()
            header = card_children[0]
            body = card_children[1]

            try:
                header.pack_configure(anchor="w", padx=28, pady=(24, 14))
            except Exception:
                pass
            try:
                body.pack_configure(fill="both", expand=True, padx=28, pady=(0, 20))
            except Exception:
                pass

            if len(body.winfo_children()) >= 5:
                try:
                    body.winfo_children()[4].configure(text="设置密码 (加密模式 / OCR 解密文档可用):")
                except Exception:
                    pass

            self.pdf_ocr_model_root = tkinter.StringVar(value=default_model_root())
            backend_map = build_backend_display_map()
            if not backend_map:
                fallback_backend = get_default_backend_display()
                backend_map = {fallback_backend: "auto"}
            self._fx_pdf_ocr_backend_map = backend_map
            self.pdf_ocr_backend = tkinter.StringVar(
                value=get_default_backend_display()
            )
            lang_map = build_profile_display_map()
            if not lang_map:
                fallback = get_default_profile_display()
                lang_map = {fallback: "general"}
            self._fx_pdf_ocr_lang_map = lang_map
            self.pdf_ocr_language = tkinter.StringVar(
                value=get_default_profile_display()
            )
            self._fx_pdf_ocr_mode_map = {
                "mixed | 混合 OCR / 原文本 (推荐)": "mixed",
                "fullPage | 整页强制 OCR": "fullPage",
                "imageOnly | 仅 OCR 图片": "imageOnly",
            }
            self.pdf_ocr_mode = tkinter.StringVar(value="mixed | 混合 OCR / 原文本 (推荐)")
            preprocess_map = build_preprocess_display_map()
            if not preprocess_map:
                preprocess_map = {get_default_preprocess_display(): "auto"}
            self._fx_pdf_ocr_preprocess_map = preprocess_map
            self.pdf_ocr_preprocess = tkinter.StringVar(value=get_default_preprocess_display())
            self.pdf_ocr_cls = tkinter.BooleanVar(value=False)
            self.pdf_ocr_compare_report = tkinter.BooleanVar(value=False)
            self.pdf_ocr_normalize_rotation = tkinter.BooleanVar(value=True)

            base_controls = list(body.winfo_children())
            merge_text = base_controls[0].cget("text")
            split_text = base_controls[1].cget("text")
            encrypt_text = base_controls[2].cget("text")
            delete_text = base_controls[3].cget("text")
            pwd_label_text = base_controls[4].cget("text")
            pwd_placeholder = base_controls[5].cget("placeholder_text")
            pwd_value = ""
            try:
                pwd_value = base_controls[5].get()
            except Exception:
                pass
            for widget in base_controls:
                try:
                    widget.destroy()
                except Exception:
                    pass
            self.pdf_pwd_var = tkinter.StringVar(value=pwd_value)

            content_row = customtkinter.CTkFrame(body, fg_color="transparent")
            content_row.pack(fill="both", expand=True, padx=0, pady=(4, 0))

            base_panel = customtkinter.CTkFrame(content_row, fg_color="transparent", width=250)
            base_panel.pack(side="left", fill="y", padx=(0, 14))
            base_panel.pack_propagate(False)

            detail_shell = customtkinter.CTkFrame(content_row, **self._get_panel_style())
            detail_shell.pack(side="left", fill="both", expand=True)
            detail_shell.grid_columnconfigure(0, weight=1)
            detail_shell.grid_rowconfigure(0, weight=1)

            nav_label = customtkinter.CTkLabel(
                base_panel,
                text="PDF 功能",
                text_color="#E6EEF2",
                font=customtkinter.CTkFont(size=14, weight="bold"),
            )
            nav_label.pack(anchor="w", pady=(4, 10))

            mode_buttons = {}

            def select_pdf_mode(mode):
                try:
                    self.pdf_mode_var.set(mode)
                except Exception:
                    pass
                for key, button in mode_buttons.items():
                    try:
                        if key == mode:
                            button.configure(fg_color="#7A695B", border_color="#D0B38A", text_color="#FFFFFF")
                        else:
                            button.configure(fg_color="transparent", border_color="#44515A", text_color="#DDE6EA")
                    except Exception:
                        pass
                for key, panel in self._fx_pdf_detail_panels.items():
                    try:
                        if key == mode:
                            panel.grid(row=0, column=0, sticky="nsew", padx=10, pady=10)
                            panel.tkraise()
                        else:
                            panel.grid_remove()
                    except Exception:
                        pass

            def make_mode_button(mode, title, hint):
                frame = customtkinter.CTkButton(
                    base_panel,
                    text=title,
                    command=lambda selected=mode: select_pdf_mode(selected),
                    height=28,
                    anchor="w",
                    corner_radius=8,
                    border_width=1,
                    fg_color="transparent",
                    hover_color="#303030",
                    border_color="#44515A",
                    text_color="#DDE6EA",
                    font=customtkinter.CTkFont(size=11),
                )
                frame._fx_pdf_mode_hint = hint
                frame.pack(fill="x", pady=(0, 2))
                mode_buttons[mode] = frame

            make_mode_button("merge", merge_text, "多份合成一份")
            make_mode_button("split", split_text, "逐页拆分输出")
            make_mode_button("encrypt", encrypt_text, "设置打开密码")
            make_mode_button("compress", "PDF 压缩", "压缩体积和图片")
            make_mode_button("rotation", "仅修正文本方向（不进行OCR）", "不进行 OCR，仅修正页面方向")
            make_mode_button("ocr", "OCR 搜索版 PDF", "生成可搜索文字层")

            shared_panel = customtkinter.CTkFrame(base_panel, fg_color="transparent")
            shared_panel.pack(fill="x", pady=(6, 0))

            self.chk_delete = customtkinter.CTkSwitch(
                shared_panel,
                text=delete_text,
                variable=self.pdf_delete_var,
                **self._get_switch_style(),
            )
            self.chk_delete.pack(anchor="w", pady=(0, 10))

            customtkinter.CTkLabel(
                shared_panel,
                text=pwd_label_text,
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
            ).pack(anchor="w", pady=(0, 4))

            self.pdf_pwd_entry = customtkinter.CTkEntry(
                shared_panel,
                textvariable=self.pdf_pwd_var,
                placeholder_text=pwd_placeholder,
                **self._get_entry_style(),
            )
            self.pdf_pwd_entry.pack(fill="x")
            self._fx_pdf_shared_pwd_entry = self.pdf_pwd_entry

            self.pdf_compress_level_var = tkinter.StringVar(value="标准")
            self.pdf_image_compress_level_var = tkinter.StringVar(value="标准")
            self._fx_pdf_detail_panels = {}

            def create_detail_panel(key, title):
                panel = customtkinter.CTkFrame(detail_shell, fg_color="transparent")
                panel.grid_columnconfigure(0, weight=1)
                customtkinter.CTkLabel(
                    panel,
                    text=title,
                    text_color="#E6EEF2",
                    font=customtkinter.CTkFont(size=15, weight="bold"),
                    height=22,
                ).pack(anchor="w", padx=8, pady=(0, 8))
                self._fx_pdf_detail_panels[key] = panel
                return panel

            def add_panel_note(parent, text):
                customtkinter.CTkLabel(
                    parent,
                    text=text,
                    text_color=COLOR_TEXT_SOFT,
                    font=customtkinter.CTkFont(size=12),
                    justify="left",
                    wraplength=620,
                ).pack(anchor="w", fill="x", padx=8, pady=(0, 8))

            merge_panel = create_detail_panel("merge", "PDF 合并")
            add_panel_note(merge_panel, "把输入中的 PDF 按文件名顺序合并为一个 PDF。适合先把文件放进同一个文件夹后统一处理。")

            split_panel = create_detail_panel("split", "PDF 拆分")
            add_panel_note(split_panel, "把每份 PDF 按页面拆成单页文件，并在结果目录中按原文件名建立子文件夹。")

            encrypt_panel = create_detail_panel("encrypt", "PDF 加密")
            add_panel_note(encrypt_panel, "填写打开密码后开始处理。这个密码也兼容加密 PDF 的 OCR 和压缩读取。")
            encrypt_pwd_field = customtkinter.CTkFrame(encrypt_panel, fg_color="transparent")
            encrypt_pwd_field.pack(fill="x", padx=8, pady=(6, 10))
            customtkinter.CTkLabel(
                encrypt_pwd_field,
                text="打开密码：",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
            ).pack(anchor="w", pady=(0, 4))
            self._fx_pdf_encrypt_pwd_entry = customtkinter.CTkEntry(
                encrypt_pwd_field,
                textvariable=self.pdf_pwd_var,
                placeholder_text=pwd_placeholder,
                **self._get_entry_style(),
            )
            self._fx_pdf_encrypt_pwd_entry.pack(fill="x")
            self.pdf_pwd_entry = self._fx_pdf_encrypt_pwd_entry

            compress_panel = create_detail_panel("compress", "PDF 压缩")
            add_panel_note(compress_panel, "PDF 压缩程度控制对象清理、字体和数据流压缩；图片压缩程度控制内嵌图片的重压缩和降采样。图片化压缩会把整页转成图片以换取更小体积。")
            compress_grid = customtkinter.CTkFrame(compress_panel, fg_color="transparent")
            compress_grid.pack(fill="x", padx=8, pady=(2, 10))
            compress_grid.grid_columnconfigure(0, weight=1)
            compress_grid.grid_columnconfigure(1, weight=1)

            pdf_level_field = customtkinter.CTkFrame(compress_grid, fg_color="transparent")
            pdf_level_field.grid(row=0, column=0, sticky="ew", padx=(0, 8))
            customtkinter.CTkLabel(
                pdf_level_field,
                text="PDF 压缩程度：",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
            ).pack(anchor="w", pady=(0, 4))
            customtkinter.CTkComboBox(
                pdf_level_field,
                values=list(PDF_COMPRESS_LEVELS.keys()),
                variable=self.pdf_compress_level_var,
                height=32,
                **self._get_combo_style(),
            ).pack(fill="x")

            image_level_field = customtkinter.CTkFrame(compress_grid, fg_color="transparent")
            image_level_field.grid(row=0, column=1, sticky="ew")
            customtkinter.CTkLabel(
                image_level_field,
                text="图片压缩程度：",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
            ).pack(anchor="w", pady=(0, 4))
            customtkinter.CTkComboBox(
                image_level_field,
                values=list(PDF_IMAGE_COMPRESS_LEVELS.keys()),
                variable=self.pdf_image_compress_level_var,
                height=32,
                **self._get_combo_style(),
            ).pack(fill="x")

            add_panel_note(
                compress_panel,
                "提示：如果 PDF 主要由扫描图片组成，调高图片压缩更有效；如果 PDF 主要是文字，PDF 压缩程度通常更关键。若选择“图片化压缩”，页面会转成图片，适合上传分享，不适合复制、搜索或编辑文本。",
            )

            rotation_panel = create_detail_panel("rotation", "仅修正文本方向（不进行OCR）")
            add_panel_note(
                rotation_panel,
                "把 PDF 页面依赖的 90/180/270 度旋转写入页面内容，并将页面旋转值归零。适合修复 OCR 文件在 Noteful 等笔记软件中画面与文字层方向不一致的问题；本功能不需要 OCR，不重新识别文字。",
            )
            customtkinter.CTkLabel(
                rotation_panel,
                text="处理结果：复制到结果文件夹后再修正，原 PDF 默认保持不变。若输入 PDF 已加密，请在左侧填写密码。",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                justify="left",
                wraplength=620,
            ).pack(anchor="w", fill="x", padx=8, pady=(2, 8))

            ocr_panel = create_detail_panel("ocr", "OCR 配置")

            customtkinter.CTkLabel(
                ocr_panel,
                text="OCR 模型目录：",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                height=18,
            ).pack(anchor="w", padx=8, pady=(0, 2))

            customtkinter.CTkEntry(
                ocr_panel,
                textvariable=self.pdf_ocr_model_root,
                **self._get_entry_style(),
            ).pack(fill="x", padx=8, pady=(0, 6))

            ocr_top_fields = customtkinter.CTkFrame(ocr_panel, fg_color="transparent")
            ocr_top_fields.pack(fill="x", padx=8, pady=(0, 6))

            backend_field = customtkinter.CTkFrame(ocr_top_fields, fg_color="transparent")
            backend_field.pack(side="left", fill="x", expand=True, padx=(0, 8))

            customtkinter.CTkLabel(
                backend_field,
                text="OCR 后端：",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                height=18,
            ).pack(anchor="w", pady=(0, 2))

            customtkinter.CTkComboBox(
                backend_field,
                values=list(backend_map.keys()),
                variable=self.pdf_ocr_backend,
                height=32,
                **self._get_combo_style(),
            ).pack(fill="x")

            self.pdf_ocr_backend_status_var = tkinter.StringVar(value="")

            def refresh_pdf_ocr_backend_status(detailed=True):
                try:
                    raw_status = build_backend_status_text(detailed=detailed)
                    summary_parts = []
                    for line in raw_status.splitlines():
                        line = line.strip()
                        if not line.startswith("- "):
                            continue
                        name, _, desc = line[2:].partition(":")
                        short_desc = desc.split("|", 1)[0].strip() if desc else ""
                        summary_parts.append(f"{name.strip()}: {short_desc}")
                    if summary_parts:
                        self.pdf_ocr_backend_status_var.set("后端状态： " + " | ".join(summary_parts))
                    else:
                        self.pdf_ocr_backend_status_var.set(raw_status)
                except Exception as inner_exc:
                    self.pdf_ocr_backend_status_var.set(f"后端状态读取失败：{inner_exc}")

            lang_field = customtkinter.CTkFrame(ocr_top_fields, fg_color="transparent")
            lang_field.pack(side="left", fill="x", expand=True)

            customtkinter.CTkLabel(
                lang_field,
                text="OCR 识别配置：",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                height=18,
            ).pack(anchor="w", pady=(0, 2))

            self.pdf_ocr_lang_combo = customtkinter.CTkComboBox(
                lang_field,
                values=list(lang_map.keys()),
                variable=self.pdf_ocr_language,
                height=32,
                **self._get_combo_style(),
            )
            self.pdf_ocr_lang_combo.pack(fill="x")

            ocr_mid_fields = customtkinter.CTkFrame(ocr_panel, fg_color="transparent")
            ocr_mid_fields.pack(fill="x", padx=8, pady=(0, 6))

            mode_field = customtkinter.CTkFrame(ocr_mid_fields, fg_color="transparent")
            mode_field.pack(side="left", fill="x", expand=True, padx=(0, 8))

            customtkinter.CTkLabel(
                mode_field,
                text="提取模式：",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                height=18,
            ).pack(anchor="w", pady=(0, 2))

            customtkinter.CTkComboBox(
                mode_field,
                values=list(self._fx_pdf_ocr_mode_map.keys()),
                variable=self.pdf_ocr_mode,
                height=32,
                **self._get_combo_style(),
            ).pack(fill="x")

            preprocess_field = customtkinter.CTkFrame(ocr_mid_fields, fg_color="transparent")
            preprocess_field.pack(side="left", fill="x", expand=True, padx=(0, 8))

            customtkinter.CTkLabel(
                preprocess_field,
                text="图像增强：",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                height=18,
            ).pack(anchor="w", pady=(0, 2))

            customtkinter.CTkComboBox(
                preprocess_field,
                values=list(self._fx_pdf_ocr_preprocess_map.keys()),
                variable=self.pdf_ocr_preprocess,
                height=32,
                **self._get_combo_style(),
            ).pack(fill="x")

            refresh_field = customtkinter.CTkFrame(ocr_mid_fields, fg_color="transparent")
            refresh_field.pack(side="left", fill="y")

            customtkinter.CTkLabel(
                refresh_field,
                text="状态检测：",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                height=18,
            ).pack(anchor="w", pady=(0, 2))

            customtkinter.CTkButton(
                refresh_field,
                text="刷新状态",
                command=lambda: refresh_pdf_ocr_backend_status(True),
                height=32,
                width=120,
                corner_radius=8,
                border_width=1,
                fg_color="transparent",
                hover_color="#303030",
                border_color="#7A695B",
                text_color="#E7E9EE",
            ).pack(anchor="w")

            customtkinter.CTkLabel(
                ocr_panel,
                textvariable=self.pdf_ocr_backend_status_var,
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=10),
                justify="left",
                wraplength=560,
            ).pack(anchor="w", fill="x", padx=8, pady=(0, 6))

            direction_switch_row = customtkinter.CTkFrame(ocr_panel, fg_color="transparent")
            direction_switch_row.pack(fill="x", padx=8, pady=(0, 4))

            customtkinter.CTkSwitch(
                direction_switch_row,
                text="修正文本方向（OCR）",
                variable=self.pdf_ocr_cls,
                **self._get_switch_style(),
            ).pack(side="left", anchor="w")

            rotation_switch_style = dict(self._get_switch_style())
            rotation_switch_style.update(
                {
                    "fg_color": "#5B4936",
                    "progress_color": "#B98245",
                    "button_color": "#E6BF7C",
                    "button_hover_color": "#F2D59B",
                    "text_color": "#F0D09A",
                }
            )
            customtkinter.CTkSwitch(
                direction_switch_row,
                text="仅修正文本方向（不进行OCR）",
                variable=self.pdf_ocr_normalize_rotation,
                **rotation_switch_style,
            ).pack(side="right", anchor="e")

            customtkinter.CTkSwitch(
                ocr_panel,
                text="生成后端对比报告（首页）",
                variable=self.pdf_ocr_compare_report,
                **self._get_switch_style(),
            ).pack(anchor="w", padx=8, pady=(0, 4))

            ocr_preview_frame = customtkinter.CTkFrame(ocr_panel, fg_color="#202426", corner_radius=10)
            ocr_preview_frame.pack(fill="x", padx=8, pady=(2, 6))
            self._fx_pdf_ocr_preview_frame = ocr_preview_frame

            ocr_preview_header = customtkinter.CTkFrame(ocr_preview_frame, fg_color="transparent")
            ocr_preview_header.pack(fill="x", padx=10, pady=(8, 2))
            customtkinter.CTkLabel(
                ocr_preview_header,
                text="实时 OCR 预览",
                text_color="#E6EEF2",
                font=customtkinter.CTkFont(size=11, weight="bold"),
            ).pack(side="left", anchor="w")
            customtkinter.CTkButton(
                ocr_preview_header,
                text="清空",
                command=lambda target=self: _clear_pdf_ocr_preview(target),
                height=24,
                width=58,
                corner_radius=7,
                fg_color="#2A6DA8",
                hover_color="#1F5C91",
                text_color="#FFFFFF",
                font=customtkinter.CTkFont(size=10),
            ).pack(side="right")

            ocr_preview_box = customtkinter.CTkTextbox(
                ocr_preview_frame,
                height=118,
                wrap="word",
                activate_scrollbars=True,
                fg_color="#262A2C",
                text_color="#DDE6EA",
                border_width=0,
                corner_radius=8,
                font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=11),
            )
            ocr_preview_box.pack(fill="x", padx=10, pady=(0, 10))
            ocr_preview_box.configure(state="disabled")
            self._fx_pdf_ocr_preview_box = ocr_preview_box
            _clear_pdf_ocr_preview(self)

            customtkinter.CTkLabel(
                ocr_panel,
                text="说明：生成双层可搜索 PDF，保留原页面画面并叠加透明文字层；仅修正文本方向（不进行OCR）会把 90/180/270 度显示旋转写入页面内容，避免导入 Noteful 后文字层方向错位。",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=10),
                justify="left",
                wraplength=560,
            ).pack(anchor="w", fill="x", padx=8, pady=(0, 2))

            self.pdf_ocr_backend_status_var.set("后端状态：按需检测，可直接运行 OCR；如需查看详细可用性再点刷新。")
            self._fx_select_pdf_mode = select_pdf_mode
            try:
                select_pdf_mode(self.pdf_mode_var.get())
            except Exception:
                select_pdf_mode("merge")

            self._fx_pdf_ocr_ui_ready = True
        except Exception as exc:
            _debug(f"patch_pdf_ocr_mode:init_ui_error:{exc}")

    def patched_run_process(self, input_folder, task_type):
        if task_type == "pdf":
            try:
                pdf_mode = self.pdf_mode_var.get() if getattr(self, "pdf_mode_var", None) is not None else ""
                if pdf_mode == "compress":
                    try:
                        _run_pdf_compress_task(self, input_folder)
                    except Exception as exc:
                        self.log(f"🔥 [严重错误] {exc}")
                    finally:
                        self.reset_ui()
                    return
                if pdf_mode == "rotation":
                    try:
                        _run_pdf_rotation_task(self, input_folder)
                    except Exception as exc:
                        self.log(f"🔥 [严重错误] {exc}")
                    finally:
                        self.reset_ui()
                    return
                if pdf_mode == "ocr":
                    try:
                        _run_pdf_ocr_task(self, input_folder)
                    except Exception as exc:
                        self.log(f"🔥 [严重错误] {exc}")
                    finally:
                        self.reset_ui()
                    return
            except Exception as exc:
                _debug(f"patch_pdf_ocr_mode:run_process_error:{exc}")
        return original_run_process(self, input_folder, task_type)

    patched_init_pdf_ui.__fx_pdf_ocr_patch__ = True
    patched_run_process.__fx_pdf_ocr_patch__ = True
    FengxiToolboxApp.init_pdf_ui = patched_init_pdf_ui
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_pdf_ocr_mode:installed")


_patch_pdf_ocr_mode()


def _patch_pdf_merge_resume():
    try:
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_pdf_merge_resume:missing:{exc}")
        return

    if getattr(original_run_process, "__fx_pdf_merge_resume_patch__", False):
        return

    def patched_run_process(self, input_folder, task_type):
        if task_type == "pdf":
            try:
                pdf_mode = self.pdf_mode_var.get() if getattr(self, "pdf_mode_var", None) is not None else ""
            except Exception:
                pdf_mode = ""
            if pdf_mode == "merge":
                normalized_input, _input_root, output_folder, resolved_strategy = _resolve_output_root_for_task(
                    input_folder,
                    "pdf",
                    _get_task_output_strategy(self, "pdf"),
                )
                merged_output = os.path.join(output_folder, "Merged_All.pdf")
                if is_nonempty_file(merged_output):
                    result = _get_last_task_result(self)
                    if result is None:
                        result = _start_task_result(self, normalized_input, "pdf")
                    _set_task_result_output_strategy(result, "pdf", resolved_strategy)
                    _set_task_result_output_root(result, output_folder)
                    _add_task_result_output(result, merged_output)
                    try:
                        all_files = self.collect_input_files(normalized_input, "pdf")
                    except Exception:
                        all_files = []
                    pdf_count = sum(1 for item in all_files if str(item).lower().endswith(".pdf"))
                    count = max(1, pdf_count)
                    _set_task_result_counts(result, processed=count, success=count, failed=0, skipped=count)
                    _set_task_result_finished(
                        result,
                        "success",
                        message="PDF 合并已复用已有结果",
                        detail=f"输出: {merged_output}",
                    )
                    try:
                        self.log(format_resume_skip_message("PDF 合并断点续跑", normalized_input, merged_output))
                    except Exception:
                        pass
                    try:
                        self.reset_ui()
                    except Exception:
                        pass
                    return
        return original_run_process(self, input_folder, task_type)

    patched_run_process.__fx_pdf_merge_resume_patch__ = True
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_pdf_merge_resume:installed")


_patch_pdf_merge_resume()


def _patch_image_pdf_modes():
    try:
        original_init_img_ui = FengxiToolboxApp.init_img_ui
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_image_pdf_modes:missing:{exc}")
        return

    if getattr(original_init_img_ui, "__fx_image_pdf_patch__", False):
        return

    def patched_init_img_ui(self):
        original_init_img_ui(self)
        if getattr(self, "_fx_image_pdf_ui_ready", False):
            return
        try:
            card = self.tab_img.winfo_children()[0]
            body = card.winfo_children()[1]
            children = list(body.winfo_children())
            insert_after = children[2] if len(children) > 2 else None
            pdf_modes_frame = customtkinter.CTkFrame(body, fg_color="transparent")
            if insert_after is not None:
                pdf_modes_frame.pack(after=insert_after, fill="x", pady=(0, 8))
            else:
                pdf_modes_frame.pack(fill="x", pady=(0, 8))

            customtkinter.CTkRadioButton(
                pdf_modes_frame,
                text="图片转 PDF (Single Image PDF)",
                variable=self.img_mode_var,
                value="to_pdf",
                **self._get_radio_style(),
            ).pack(anchor="w", pady=(0, 6))

            customtkinter.CTkRadioButton(
                pdf_modes_frame,
                text="多图合并 PDF (Merge Images PDF)",
                variable=self.img_mode_var,
                value="merge_pdf",
                **self._get_radio_style(),
            ).pack(anchor="w", pady=(0, 2))

            customtkinter.CTkLabel(
                pdf_modes_frame,
                text="图片转 PDF 会为每张图片生成一个 PDF；多图合并 PDF 会按文件名顺序合成一份 PDF。",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                justify="left",
                wraplength=560,
            ).pack(anchor="w", fill="x", pady=(4, 0))
            self._fx_image_pdf_ui_ready = True
        except Exception as exc:
            _debug(f"patch_image_pdf_modes:init_ui_error:{exc}")

    def patched_run_process(self, input_folder, task_type):
        if task_type == "image":
            try:
                image_mode = _get_image_pdf_mode(self)
                if image_mode == "to_pdf":
                    try:
                        _run_image_to_pdf_task(self, input_folder, merge=False)
                    except Exception as exc:
                        self.log(f"🔥 [严重错误] {exc}")
                    finally:
                        self.reset_ui()
                    return
                if image_mode == "merge_pdf":
                    try:
                        _run_image_to_pdf_task(self, input_folder, merge=True)
                    except Exception as exc:
                        self.log(f"🔥 [严重错误] {exc}")
                    finally:
                        self.reset_ui()
                    return
            except Exception as exc:
                _debug(f"patch_image_pdf_modes:run_process_error:{exc}")
        return original_run_process(self, input_folder, task_type)

    patched_init_img_ui.__fx_image_pdf_patch__ = True
    patched_run_process.__fx_image_pdf_patch__ = True
    FengxiToolboxApp.init_img_ui = patched_init_img_ui
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_image_pdf_modes:installed")


_patch_image_pdf_modes()


def _patch_convert_imgs_to_pdf_task():
    try:
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_convert_imgs2pdf:missing:{exc}")
        return

    if getattr(original_run_process, "__fx_convert_imgs2pdf_patch__", False):
        return

    def patched_run_process(self, input_folder, task_type):
        if task_type == "convert" and _get_convert_mode(self) == "imgs2pdf":
            try:
                _run_convert_imgs_to_pdf_task(self, input_folder)
            except Exception as exc:
                self.log(f"🔥 [严重错误] {exc}")
            finally:
                self.reset_ui()
            return
        return original_run_process(self, input_folder, task_type)

    patched_run_process.__fx_convert_imgs2pdf_patch__ = True
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_convert_imgs2pdf:installed")


_patch_convert_imgs_to_pdf_task()


def _iter_child_widgets(widget):
    try:
        children = list(widget.winfo_children())
    except Exception:
        return
    for child in children:
        yield child
        yield from _iter_child_widgets(child)


def _install_convert_mode_grid(app):
    if getattr(app, "cv_mode", None) is None:
        app.cv_mode = tkinter.StringVar(value="word2pdf")
    card = app.tab_cv.winfo_children()[0]
    body = card.winfo_children()[1] if len(card.winfo_children()) > 1 else card
    mode_values = tuple(CONVERT_MODE_SPECS.keys())
    mode_widgets = []
    old_grids = []
    for widget in _iter_child_widgets(body):
        try:
            if getattr(widget, "_fx_convert_mode_grid", False):
                old_grids.append(widget)
                continue
            if isinstance(widget, customtkinter.CTkRadioButton) and str(widget.cget("value")) in mode_values:
                mode_widgets.append(widget)
        except Exception:
            continue

    insert_before = mode_widgets[0].master if mode_widgets else None
    parents_to_destroy = []
    for widget in mode_widgets:
        try:
            parent = widget.master
            if parent is not body and parent not in parents_to_destroy:
                parents_to_destroy.append(parent)
        except Exception:
            pass

    for widget in old_grids:
        try:
            widget.destroy()
        except Exception:
            pass

    section = customtkinter.CTkFrame(
        body,
        fg_color=globals().get("COLOR_CARD_ALT", "#303030"),
        corner_radius=8,
        border_width=1,
        border_color=globals().get("COLOR_BORDER", "#3A3A3A"),
    )
    section._fx_convert_mode_grid = True
    pack_options = {"fill": "x", "pady": (8, 10)}
    if insert_before is not None:
        pack_options["before"] = insert_before
    try:
        section.pack(**pack_options)
    except Exception:
        section.pack(fill="x", pady=(8, 10))
    section.grid_columnconfigure(0, weight=1)
    section.grid_columnconfigure(1, weight=1)

    customtkinter.CTkLabel(
        section,
        text="转换类型",
        text_color=COLOR_TEXT_SOFT,
        font=customtkinter.CTkFont(size=12, weight="bold"),
    ).grid(row=0, column=0, columnspan=2, sticky="w", padx=12, pady=(10, 6))

    mode_items = (
        ("word2pdf", "Word -> PDF"),
        ("pdf2word", "PDF -> Word"),
        ("ppt2pdf", "PPT -> PDF"),
        ("pdf2ppt", "PDF -> PPT"),
        ("txt2word", "TXT -> Word"),
        ("md2pdf", "Markdown -> PDF"),
        ("pdf2md", "PDF -> Markdown"),
        ("imgs2pdf", "多图合并 PDF"),
    )
    for index, (value, label) in enumerate(mode_items):
        radio = customtkinter.CTkRadioButton(
            section,
            text=label,
            variable=app.cv_mode,
            value=value,
            **app._get_radio_style(),
        )
        radio.grid(row=index // 2 + 1, column=index % 2, sticky="w", padx=12, pady=(0, 8))

    for widget in mode_widgets:
        try:
            widget.destroy()
        except Exception:
            pass
    for parent in parents_to_destroy:
        try:
            if parent is not section:
                parent.destroy()
        except Exception:
            pass


def _patch_convert_extended_modes():
    try:
        original_init_convert_ui = FengxiToolboxApp.init_convert_ui
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_convert_extended_modes:missing:{exc}")
        return

    if getattr(original_init_convert_ui, "__fx_convert_extended_modes_patch__", False):
        return

    def patched_init_convert_ui(self):
        original_init_convert_ui(self)
        if getattr(self, "_fx_convert_extended_ui_ready", False):
            return
        try:
            if getattr(self, "cv_mode", None) is None:
                self.cv_mode = tkinter.StringVar(value="word2pdf")
            card = self.tab_cv.winfo_children()[0]
            body = card.winfo_children()[1] if len(card.winfo_children()) > 1 else card
            section = customtkinter.CTkFrame(body, fg_color="transparent")
            section.pack(fill="x", pady=(8, 4))
            customtkinter.CTkLabel(
                section,
                text="新增转换",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=12, weight="bold"),
            ).pack(anchor="w", pady=(0, 6))
            for value, label in (
                ("pdf2ppt", "PDF 转 PPT (PDF to PPT)"),
                ("txt2word", "TXT 转 Word (TXT to Word)"),
                ("md2pdf", "Markdown 转 PDF (MD to PDF)"),
                ("pdf2md", "PDF 转 Markdown (PDF to MD)"),
            ):
                customtkinter.CTkRadioButton(
                    section,
                    text=label,
                    variable=self.cv_mode,
                    value=value,
                    **self._get_radio_style(),
                ).pack(anchor="w", pady=(0, 6))
            customtkinter.CTkLabel(
                section,
                text="PDF 转 PPT 会将每页 PDF 作为图片放入幻灯片；PDF 转 Markdown 只提取可复制文本，扫描件需先做 OCR。",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                justify="left",
                wraplength=560,
            ).pack(anchor="w", fill="x", pady=(2, 0))
            self._fx_convert_extended_ui_ready = True
        except Exception as exc:
            _debug(f"patch_convert_extended_modes:init_ui_error:{exc}")

    def patched_run_process(self, input_folder, task_type):
        if task_type == "convert" and _get_convert_mode(self) in CONVERT_EXTENDED_MODES:
            try:
                _run_convert_extended_task(self, input_folder)
            except Exception as exc:
                self.log(f"🔥 [严重错误] {exc}")
            finally:
                self.reset_ui()
            return
        return original_run_process(self, input_folder, task_type)

    patched_init_convert_ui.__fx_convert_extended_modes_patch__ = True
    patched_run_process.__fx_convert_extended_modes_patch__ = True
    FengxiToolboxApp.init_convert_ui = patched_init_convert_ui
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_convert_extended_modes:installed")


_patch_convert_extended_modes()


def _patch_convert_mode_grid_cleanup():
    try:
        original_init_convert_ui = FengxiToolboxApp.init_convert_ui
    except Exception as exc:
        _debug(f"patch_convert_mode_grid_cleanup:missing:{exc}")
        return

    if getattr(original_init_convert_ui, "__fx_convert_mode_grid_cleanup_patch__", False):
        return

    def patched_init_convert_ui(self):
        original_init_convert_ui(self)
        try:
            _install_convert_mode_grid(self)
        except Exception as exc:
            _debug(f"patch_convert_mode_grid_cleanup:init_ui_error:{exc}")

    patched_init_convert_ui.__fx_convert_mode_grid_cleanup_patch__ = True
    FengxiToolboxApp.init_convert_ui = patched_init_convert_ui
    _debug("patch_convert_mode_grid_cleanup:installed")


_patch_convert_mode_grid_cleanup()


def _collect_audio_files(app, input_value):
    normalized_input = _normalize_input_path_value(input_value)
    return _audio_task_collect_files(normalized_input, collect_input_files=getattr(app, "collect_input_files", None))


def _get_audio_task_args(app):
    return _audio_task_get_args(app)


def _get_audio_transcribe_cache_dir():
    return str(_get_user_pref_root() / "models" / "faster-whisper")


def _get_audio_transcribe_args(app):
    args = _audio_task_get_transcribe_args(app)
    if not str(args.get("cache_dir") or "").strip():
        args["cache_dir"] = _get_audio_transcribe_cache_dir()
    return args


def _build_audio_output_path(src, input_root, output_folder, target_fmt):
    return _audio_task_build_output_path(src, input_root, output_folder, target_fmt)


def _process_one_audio_file(job):
    return _audio_task_process_one_file(job, convert_audio_format, copy_file_safe, _speech_transcribe_media_file)


def _set_pdf_ocr_preview_text(app, text, *, clear=False, force_bottom=False):
    box = getattr(app, "_fx_pdf_ocr_preview_box", None)
    if box is None:
        return False

    def update():
        try:
            should_follow = True
            if not clear and not force_bottom:
                try:
                    should_follow = float(box.yview()[1]) >= 0.98
                except Exception:
                    should_follow = True
            box.configure(state="normal")
            if clear:
                box.delete("1.0", "end")
            box.insert("end", str(text or ""))
            box.configure(state="disabled")
            if clear or force_bottom or should_follow:
                box.see("end")
            return True
        except Exception as exc:
            _debug(f"pdf_ocr:preview_update_error:{exc}")
            return False

    try:
        app.after(0, update)
        return True
    except Exception:
        return bool(update())


def _clear_pdf_ocr_preview(app, message=None):
    text = message or "实时预览：开始 OCR 后，这里会按页显示已识别出的文字。\n"
    return _set_pdf_ocr_preview_text(app, text, clear=True, force_bottom=True)


def _append_pdf_ocr_preview_page(app, src, page_done, total_pages, payload):
    payload = payload if isinstance(payload, dict) else {}
    basename = os.path.basename(str(src or ""))
    try:
        page_number = int(payload.get("page_number") or page_done or 0)
    except Exception:
        page_number = int(page_done or 0)
    try:
        page_total = int(payload.get("total_pages") or total_pages or 0)
    except Exception:
        page_total = int(total_pages or 0)
    lines = payload.get("lines") if isinstance(payload.get("lines"), list) else []
    output = [f"\n[{basename}] 第 {page_number}/{page_total} 页\n"]
    clean_lines = [str(line).strip() for line in lines if str(line or "").strip()]
    if clean_lines:
        for line in clean_lines[:12]:
            output.append(f"  {line}\n")
        if len(clean_lines) > 12:
            output.append(f"  ... 还有 {len(clean_lines) - 12} 行未显示\n")
    else:
        text_count = int(payload.get("text_count") or 0)
        ocr_count = int(payload.get("ocr_count") or 0)
        output.append(f"  本页暂未识别到可预览文本。（原文 {text_count} 行，OCR {ocr_count} 行）\n")
    return _set_pdf_ocr_preview_text(app, "".join(output))


def _set_audio_transcribe_preview_text(app, text, *, clear=False, force_bottom=False):
    box = getattr(app, "_fx_audio_transcribe_preview_box", None)
    if box is None:
        return False

    def update():
        try:
            should_follow = True
            if not clear and not force_bottom:
                try:
                    should_follow = float(box.yview()[1]) >= 0.98
                except Exception:
                    should_follow = True
            box.configure(state="normal")
            if clear:
                box.delete("1.0", "end")
            box.insert("end", str(text or ""))
            box.configure(state="disabled")
            if clear or force_bottom or should_follow:
                box.see("end")
            return True
        except Exception as exc:
            _debug(f"audio_transcribe:preview_update_error:{exc}")
            return False

    try:
        app.after(0, update)
        return True
    except Exception:
        return bool(update())


def _clear_audio_transcribe_preview(app, message=None):
    text = message or "实时预览：开始识别后，这里会滚动显示已经识别出来的内容。\n"
    return _set_audio_transcribe_preview_text(app, text, clear=True, force_bottom=True)


def _append_audio_transcribe_progress(app, src, payload):
    payload = payload if isinstance(payload, dict) else {}
    event_type = str(payload.get("type") or "")
    basename = os.path.basename(str(src or payload.get("src") or ""))
    lines = []
    if event_type == "stage":
        stage = str(payload.get("stage") or "")
        if stage == "load_model":
            lines.append(f"\n[{basename}] 加载识别模型...\n")
        elif stage == "transcribe":
            lines.append(f"[{basename}] 开始识别，片段会实时显示在下方。\n")
        elif stage == "write_outputs":
            lines.append(f"[{basename}] 正在写出转写文件...\n")
    elif event_type == "segment":
        segment = payload.get("segment") if isinstance(payload.get("segment"), dict) else {}
        text = str(segment.get("text", "")).strip()
        if text:
            start = _speech_format_timestamp(segment.get("start", 0)).replace(",", ".")
            end = _speech_format_timestamp(segment.get("end", 0)).replace(",", ".")
            lines.append(f"[{start} -> {end}] {text}\n")
    elif event_type == "done":
        count = int(payload.get("segments") or 0)
        lines.append(f"[{basename}] 识别完成，共 {count} 段。\n")
    if not lines:
        return False
    return _set_audio_transcribe_preview_text(app, "".join(lines))


def _run_audio_task(app, input_folder):
    normalized_input, input_root, output_folder = _resolve_result_output_folder(input_folder)
    audio_files = _collect_audio_files(app, normalized_input)
    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, input_folder, "audio")
    try:
        mode, _target_fmt, _bitrate, _delete_source = _get_audio_task_args(app)
    except Exception:
        mode = ""
    if mode == "transcribe":
        _clear_audio_transcribe_preview(app, "实时预览：正在准备语音转文字任务...\n")
    callbacks = AudioTaskCallbacks(
        log=getattr(app, "log", None),
        stop_requested=lambda: bool(getattr(app, "stop_event", False)),
        on_transcript_progress=lambda src, payload: _append_audio_transcribe_progress(app, src, payload),
    )
    return run_audio_task_core(
        app,
        input_folder,
        normalized_input=normalized_input,
        input_root=input_root,
        output_folder=output_folder,
        audio_files=audio_files,
        result=result,
        tracker=_get_active_progress_tracker(app),
        is_parallel_enabled=_is_parallel_enabled,
        get_parallel_worker_count=_get_parallel_worker_count,
        convert_audio_format=convert_audio_format,
        copy_file_safe=copy_file_safe,
        set_task_result_counts=_set_task_result_counts,
        set_task_result_finished=_set_task_result_finished,
        set_task_result_output_root=_set_task_result_output_root,
        add_task_result_output=_add_task_result_output,
        write_failed_report=_write_failed_report,
        log=getattr(app, "log", None),
        progress_bar=getattr(app, "progress_bar", None),
        stop_requested=lambda: bool(getattr(app, "stop_event", False)),
        executor_factory=concurrent.futures.ThreadPoolExecutor,
        get_audio_task_args=_get_audio_task_args,
        get_audio_transcribe_args=_get_audio_transcribe_args,
        transcribe_media_file=_speech_transcribe_media_file,
        callbacks=callbacks,
    )


def _patch_audio_parallel_task():
    try:
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_audio_parallel:missing:{exc}")
        return

    if getattr(original_run_process, "__fx_audio_parallel_patch__", False):
        return

    def patched_run_process(self, input_folder, task_type):
        if task_type == "audio":
            try:
                _run_audio_task(self, input_folder)
            except Exception as exc:
                self.log(f"🔥 [严重错误] {exc}")
            finally:
                self.reset_ui()
            return
        return original_run_process(self, input_folder, task_type)

    patched_run_process.__fx_audio_parallel_patch__ = True
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_audio_parallel:installed")


_patch_audio_parallel_task()


def _find_audio_settings_frame(app):
    tab = getattr(app, "tab_audio", None)
    if tab is None:
        return None
    stack = list(tab.winfo_children())
    while stack:
        widget = stack.pop(0)
        try:
            children = list(widget.winfo_children())
        except Exception:
            children = []
        texts = []
        for child in children:
            try:
                texts.append(str(child.cget("text")))
            except Exception:
                pass
        if any("视频提取音频" in text for text in texts) and any("音频格式互转" in text for text in texts):
            return widget
        stack.extend(children)
    return None


def _install_audio_transcribe_ui(app):
    if getattr(app, "_fx_audio_transcribe_ui_ready", False):
        return
    frame = _find_audio_settings_frame(app)
    if frame is None:
        return

    try:
        app.audio_transcribe_model = tkinter.StringVar(master=app, value="base")
        app.audio_transcribe_language = tkinter.StringVar(master=app, value="自动识别")
        app.audio_transcribe_format = tkinter.StringVar(master=app, value="txt")
        app.audio_transcribe_cache_dir = tkinter.StringVar(master=app, value=_get_audio_transcribe_cache_dir())
    except Exception as exc:
        _debug(f"audio_transcribe:variables_error:{exc}")
        return

    try:
        customtkinter.CTkRadioButton(
            frame,
            text="语音转文字 (音频/视频生成文本)",
            variable=app.audio_mode_var,
            value="transcribe",
        ).pack(anchor="w", pady=8)
    except Exception as exc:
        _debug(f"audio_transcribe:radio_error:{exc}")

    controls = customtkinter.CTkFrame(frame, fg_color="transparent")
    controls.pack(fill="x", pady=(4, 8))
    try:
        customtkinter.CTkLabel(controls, text="识别模型:").pack(side="left", padx=(0, 8))
        customtkinter.CTkComboBox(
            controls,
            variable=app.audio_transcribe_model,
            values=list(SPEECH_MODEL_OPTIONS),
            width=86,
        ).pack(side="left", padx=(0, 12))
        customtkinter.CTkLabel(controls, text="语言:").pack(side="left", padx=(0, 8))
        customtkinter.CTkComboBox(
            controls,
            variable=app.audio_transcribe_language,
            values=list(SPEECH_LANGUAGE_OPTIONS.keys()),
            width=92,
        ).pack(side="left", padx=(0, 12))
        customtkinter.CTkLabel(controls, text="输出:").pack(side="left", padx=(0, 8))
        customtkinter.CTkComboBox(
            controls,
            variable=app.audio_transcribe_format,
            values=list(SPEECH_OUTPUT_FORMATS),
            width=86,
        ).pack(side="left")
    except Exception as exc:
        _debug(f"audio_transcribe:controls_error:{exc}")

    preview_frame = customtkinter.CTkFrame(frame, fg_color="transparent")
    preview_frame.pack(fill="x", pady=(2, 4))
    app._fx_audio_transcribe_preview_frame = preview_frame
    preview_header = customtkinter.CTkFrame(preview_frame, fg_color="transparent")
    preview_header.pack(fill="x", pady=(0, 2))
    customtkinter.CTkLabel(
        preview_header,
        text="实时识别预览",
        text_color=globals().get("COLOR_TEXT_MAIN", "#E8EEF2"),
        font=customtkinter.CTkFont(size=12, weight="bold"),
        anchor="w",
    ).pack(side="left")
    customtkinter.CTkButton(
        preview_header,
        text="清空",
        width=54,
        height=24,
        command=lambda target=app: _clear_audio_transcribe_preview(target),
    ).pack(side="right")
    preview_box = customtkinter.CTkTextbox(
        preview_frame,
        height=150,
        wrap="word",
        font=customtkinter.CTkFont(size=12),
    )
    preview_box.pack(fill="x")
    app._fx_audio_transcribe_preview_box = preview_box
    _clear_audio_transcribe_preview(app)

    hint = customtkinter.CTkLabel(
        frame,
        text=(
            "模型说明：base 为默认推荐；tiny 最快但错字更多；small 更稳；"
            "medium 准确率最高但更慢。首次使用会缓存模型到风兮本地目录。"
        ),
        text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
        font=customtkinter.CTkFont(size=11),
        justify="left",
        anchor="w",
        wraplength=620,
    )
    hint.pack(fill="x", pady=(2, 4))
    app._fx_audio_transcribe_model_hint = hint

    app._fx_audio_transcribe_controls = controls
    app._fx_audio_transcribe_ui_ready = True
    try:
        _install_audio_last_settings_memory(app)
    except Exception as exc:
        _debug(f"audio_transcribe:last_settings_install_error:{exc}")


def _schedule_audio_last_settings_persistence(app, delay_ms=350):
    if getattr(app, "_fx_last_settings_loading", False):
        return
    after_id = getattr(app, "_fx_audio_last_settings_after_id", None)
    if after_id:
        try:
            app.after_cancel(after_id)
        except Exception:
            pass

    def persist(target=app):
        try:
            target._fx_audio_last_settings_after_id = None
        except Exception:
            pass
        try:
            if _last_settings_category_ready(target, "audio"):
                _save_last_settings_category(target, "audio", update_active=(getattr(target, "current_task", "") == "audio"))
        except Exception as exc:
            _debug(f"last_settings:audio_auto_save_error:{exc}")

    try:
        app._fx_audio_last_settings_after_id = app.after(delay_ms, persist)
    except Exception:
        persist()


def _install_audio_last_settings_memory(app):
    if getattr(app, "_fx_audio_last_settings_memory_ready", False):
        return

    def changed(*_args, target=app):
        _schedule_audio_last_settings_persistence(target)

    trace_ids = []
    for name in (
        "audio_mode_var",
        "audio_target_fmt",
        "audio_bitrate",
        "audio_delete_var",
        "audio_transcribe_model",
        "audio_transcribe_language",
        "audio_transcribe_format",
    ):
        var = getattr(app, name, None)
        if not isinstance(var, tkinter.Variable):
            continue
        try:
            trace_ids.append((name, var.trace_add("write", changed)))
        except Exception:
            pass
    app._fx_audio_last_settings_trace_ids = trace_ids
    app._fx_audio_last_settings_memory_ready = True


def _schedule_zip_last_settings_persistence(app, delay_ms=350):
    if getattr(app, "_fx_last_settings_loading", False):
        return
    after_id = getattr(app, "_fx_zip_last_settings_after_id", None)
    if after_id:
        try:
            app.after_cancel(after_id)
        except Exception:
            pass

    def persist(target=app):
        try:
            target._fx_zip_last_settings_after_id = None
        except Exception:
            pass
        try:
            if _last_settings_category_ready(target, "zip"):
                _save_last_settings_category(target, "zip", update_active=(getattr(target, "current_task", "") == "zip"))
        except Exception as exc:
            _debug(f"last_settings:zip_auto_save_error:{exc}")

    try:
        app._fx_zip_last_settings_after_id = app.after(delay_ms, persist)
    except Exception:
        persist()


def _install_zip_last_settings_memory(app):
    if getattr(app, "_fx_zip_last_settings_memory_ready", False):
        return

    def changed(*_args, target=app):
        _schedule_zip_last_settings_persistence(target)

    trace_ids = []
    for name in ("zip_mode_var", "zip_min_depth_var", "zip_max_depth_var", "zip_archive_policy_var"):
        var = getattr(app, name, None)
        if not isinstance(var, tkinter.Variable):
            continue
        try:
            trace_ids.append((name, var.trace_add("write", changed)))
        except Exception:
            pass
    app._fx_zip_last_settings_trace_ids = trace_ids
    app._fx_zip_last_settings_memory_ready = True


def _patch_audio_transcribe_ui():
    try:
        original_init_audio_ui = FengxiToolboxApp.init_audio_ui
    except Exception as exc:
        _debug(f"audio_transcribe:patch_missing:{exc}")
        return
    if getattr(original_init_audio_ui, "__fx_audio_transcribe_ui_patch__", False):
        return

    def patched_init_audio_ui(self):
        original_init_audio_ui(self)
        try:
            _install_audio_transcribe_ui(self)
        except Exception as exc:
            _debug(f"audio_transcribe:init_error:{exc}")

    patched_init_audio_ui.__fx_audio_transcribe_ui_patch__ = True
    FengxiToolboxApp.init_audio_ui = patched_init_audio_ui
    _debug("audio_transcribe:ui_patch_installed")


_patch_audio_transcribe_ui()


def _patch_zip_core_task():
    try:
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_zip_core:missing:{exc}")
        return

    if getattr(original_run_process, "__fx_zip_core_patch__", False):
        return

    def patched_run_process(self, input_folder, task_type):
        if task_type == "zip":
            try:
                _run_zip_task_with_core(self, input_folder)
            except Exception as exc:
                self.log(f"[ZIP] ERROR: {exc}")
            finally:
                self.reset_ui()
            return
        return original_run_process(self, input_folder, task_type)

    patched_run_process.__fx_zip_core_patch__ = True
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_zip_core:installed")


_patch_zip_core_task()


def _iter_widget_tree(widget):
    if widget is None:
        return
    yield widget
    try:
        children = widget.winfo_children()
    except Exception:
        children = []
    for child in children:
        yield from _iter_widget_tree(child)


def _find_watermark_skip_switch(root_widget):
    candidate = None
    for widget in _iter_widget_tree(root_widget):
        if not isinstance(widget, customtkinter.CTkSwitch):
            continue
        try:
            label = str(widget.cget("text") or "")
        except Exception:
            continue
        if ("文件名" in label and "跳过" in label) or ("-" in label and "结尾" in label):
            candidate = widget
    return candidate


def _get_option_menu_style(combo_style):
    allowed_keys = {
        "corner_radius",
        "fg_color",
        "button_color",
        "button_hover_color",
        "text_color",
        "text_color_disabled",
        "dropdown_fg_color",
        "dropdown_hover_color",
        "dropdown_text_color",
        "font",
        "dropdown_font",
        "state",
        "hover",
        "dynamic_resizing",
        "anchor",
    }
    return {key: value for key, value in dict(combo_style or {}).items() if key in allowed_keys}


def _create_watermark_filename_rule_controls(app, controls_parent, option_menu_style=None):
    if controls_parent is None:
        return None
    if getattr(app, "wm_skip_name_position_var", None) is None:
        app.wm_skip_name_position_var = tkinter.StringVar(value="结尾")
    if getattr(app, "wm_skip_name_text_var", None) is None:
        app.wm_skip_name_text_var = tkinter.StringVar(value="-")
    if getattr(app, "wm_copy_skipped_var", None) is None:
        app.wm_copy_skipped_var = tkinter.BooleanVar(value=False)
    _ensure_watermark_type_skip_vars(app)

    option_menu_style = dict(option_menu_style or {})
    controls_row = customtkinter.CTkFrame(
        controls_parent,
        height=82,
        fg_color="transparent",
    )
    controls_row._fx_wm_filename_rule_controls = True
    try:
        controls_row.pack_propagate(False)
    except Exception:
        pass

    fields_row = customtkinter.CTkFrame(controls_row, fg_color="transparent", height=30)
    fields_row.pack(fill="x", pady=(0, 2))

    customtkinter.CTkLabel(
        fields_row,
        text="匹配位置",
        text_color=globals().get("COLOR_TEXT_SOFT"),
        font=customtkinter.CTkFont(size=11),
        height=30,
    ).pack(side="left", padx=(0, 8))

    customtkinter.CTkOptionMenu(
        fields_row,
        variable=app.wm_skip_name_position_var,
        values=["结尾", "开头"],
        width=82,
        height=30,
        **option_menu_style,
    ).pack(side="left", padx=(0, 8))

    app.wm_skip_name_entry = customtkinter.CTkEntry(
        fields_row,
        width=126,
        height=30,
        textvariable=app.wm_skip_name_text_var,
        placeholder_text="-",
    )
    app.wm_skip_name_entry.pack(side="left", fill="x", expand=True, padx=(0, 0))
    controls_row._fx_wm_filename_rule_entry = app.wm_skip_name_entry

    customtkinter.CTkLabel(
        controls_row,
        text='留空默认 “-”，可填写任意开头或结尾字符',
        text_color=globals().get("COLOR_TEXT_SOFT"),
        font=customtkinter.CTkFont(size=11),
        height=18,
        anchor="w",
    ).pack(anchor="w", fill="x")
    customtkinter.CTkCheckBox(
        controls_row,
        text="跳过文件复制到输出文件夹",
        variable=app.wm_copy_skipped_var,
        height=22,
        font=customtkinter.CTkFont(size=11),
    ).pack(anchor="w", fill="x", pady=(1, 0))
    return controls_row


def _ensure_active_watermark_filename_rule_controls(app, right_panel=None):
    try:
        skip_switch = _find_watermark_skip_switch(getattr(app, "tab_wm", None))
    except Exception:
        skip_switch = None
    controls_parent = right_panel or (getattr(skip_switch, "master", None) if skip_switch is not None else None)
    if controls_parent is None:
        return None

    active_controls = None
    duplicate_controls = []
    try:
        for child in controls_parent.winfo_children():
            if getattr(child, "_fx_wm_filename_rule_controls", False):
                if active_controls is None:
                    active_controls = child
                else:
                    duplicate_controls.append(child)
    except Exception:
        active_controls = None
        duplicate_controls = []

    for stale in duplicate_controls:
        try:
            stale.pack_forget()
        except Exception:
            pass
        try:
            stale.grid_forget()
        except Exception:
            pass

    if skip_switch is not None:
        try:
            skip_switch.configure(text="按文件名规则跳过")
        except Exception:
            pass

    if active_controls is None:
        combo_style = {}
        try:
            combo_style = app._get_combo_style()
        except Exception:
            combo_style = {}
        active_controls = _create_watermark_filename_rule_controls(
            app,
            controls_parent,
            _get_option_menu_style(combo_style),
        )
        if active_controls is not None:
            try:
                active_controls.pack(fill="x", padx=24, pady=(0, 2))
            except Exception:
                pass
    if active_controls is not None and skip_switch is not None:
        try:
            pack_info = active_controls.pack_info()
        except Exception:
            pack_info = {}
        try:
            if str(pack_info.get("after", "")) != str(skip_switch):
                active_controls.pack_forget()
                active_controls.pack(after=skip_switch, fill="x", padx=24, pady=(0, 3))
            else:
                active_controls.pack_configure(fill="x", padx=24, pady=(0, 3))
        except Exception:
            pass
    try:
        active_entry = getattr(active_controls, "_fx_wm_filename_rule_entry", None)
        if active_entry is not None:
            app.wm_skip_name_entry = active_entry
    except Exception:
        pass
    try:
        _install_watermark_filename_rule_memory(app)
    except Exception:
        pass
    try:
        _ensure_watermark_type_skip_controls(app, controls_parent, after_widget=active_controls)
    except Exception as exc:
        _debug(f"wm_type_skip_controls_error:{exc}")
    return active_controls


def _normalize_watermark_filename_rule_mode(value):
    text = str(value or "").strip().lower()
    if text in ("prefix", "start", "starts_with", "startswith"):
        return "prefix"
    if text in ("suffix", "end", "ends_with", "endswith"):
        return "suffix"
    if "开头" in text or "起始" in text or "前缀" in text:
        return "prefix"
    if "结尾" in text or "末尾" in text or "后缀" in text:
        return "suffix"
    return "suffix"


def _watermark_filename_rule_mode_to_label(mode):
    return "开头" if mode == "prefix" else "结尾"


def _get_watermark_filename_rule(app):
    skip_var = getattr(app, "wm_skip_hyphen_var", None)
    if skip_var is None:
        return None
    try:
        enabled = bool(skip_var.get())
    except Exception:
        enabled = False
    if not enabled:
        return None

    mode_var = getattr(app, "wm_skip_name_position_var", None)
    marker_var = getattr(app, "wm_skip_name_text_var", None)

    try:
        mode_label = str(mode_var.get() or "结尾").strip() if mode_var is not None else "结尾"
    except Exception:
        mode_label = "结尾"
    try:
        marker_text = str(marker_var.get() or "").strip() if marker_var is not None else ""
    except Exception:
        marker_text = ""

    if not marker_text:
        marker_text = "-"
    mode = _normalize_watermark_filename_rule_mode(mode_label)
    return mode, marker_text


def _get_watermark_copy_skipped_to_output(app):
    var = getattr(app, "wm_copy_skipped_var", None)
    if var is None:
        return False
    try:
        return bool(var.get())
    except Exception:
        return False


def _ensure_watermark_type_skip_vars(app):
    for name in ("wm_skip_pdf_type_var", "wm_skip_word_type_var", "wm_skip_ppt_type_var"):
        if getattr(app, name, None) is None:
            try:
                setattr(app, name, tkinter.BooleanVar(value=False))
            except Exception:
                setattr(app, name, tkinter.BooleanVar())


def _get_watermark_type_skip_groups(app):
    groups = []
    for var_name, group_name in (
        ("wm_skip_pdf_type_var", "pdf"),
        ("wm_skip_word_type_var", "word"),
        ("wm_skip_ppt_type_var", "ppt"),
    ):
        var = getattr(app, var_name, None)
        try:
            if var is not None and bool(var.get()):
                groups.append(group_name)
        except Exception:
            pass
    return set(groups)


def _watermark_type_group_for_suffix(suffix):
    suffix_text = str(suffix or "").lower()
    if suffix_text in WATERMARK_PDF_EXTS:
        return "pdf"
    if suffix_text in WATERMARK_WORD_EXTS:
        return "word"
    if suffix_text in WATERMARK_PPT_EXTS:
        return "ppt"
    return ""


def _filter_watermark_files_by_type_skip(app, files):
    skipped_groups = _get_watermark_type_skip_groups(app)
    if not skipped_groups:
        return list(files or []), []
    kept = []
    skipped = []
    for path in files or []:
        group = _watermark_type_group_for_suffix(Path(str(path)).suffix)
        if group and group in skipped_groups:
            skipped.append(path)
        else:
            kept.append(path)
    return kept, skipped


def _ensure_watermark_type_skip_controls(app, controls_parent, after_widget=None):
    if controls_parent is None:
        return None
    _ensure_watermark_type_skip_vars(app)
    active = None
    try:
        for child in controls_parent.winfo_children():
            if getattr(child, "_fx_wm_type_skip_controls", False):
                active = child
                break
    except Exception:
        active = None
    if active is None:
        active = customtkinter.CTkFrame(controls_parent, height=58, fg_color="transparent")
        active._fx_wm_type_skip_controls = True
        try:
            active.pack_propagate(False)
        except Exception:
            pass
        customtkinter.CTkLabel(
            active,
            text="\u4e0d\u6dfb\u52a0\u6c34\u5370\u7684\u6587\u4ef6\u7c7b\u578b",
            text_color=globals().get("COLOR_TEXT_SOFT"),
            font=customtkinter.CTkFont(size=11),
            height=18,
            anchor="w",
        ).pack(anchor="w", fill="x")
        row = customtkinter.CTkFrame(active, fg_color="transparent", height=28)
        row.pack(fill="x", pady=(1, 0))
        for label, var_name in (
            ("PDF", "wm_skip_pdf_type_var"),
            ("Word", "wm_skip_word_type_var"),
            ("PPT", "wm_skip_ppt_type_var"),
        ):
            customtkinter.CTkCheckBox(
                row,
                text=label,
                variable=getattr(app, var_name),
                width=64,
                height=24,
                font=customtkinter.CTkFont(size=11),
            ).pack(side="left", padx=(0, 8))
    try:
        active.pack_forget()
        if after_widget is not None:
            active.pack(after=after_widget, fill="x", padx=24, pady=(0, 3))
        else:
            active.pack(fill="x", padx=24, pady=(0, 3))
    except Exception:
        pass
    return active


def _watermark_filename_matches_rule(name_no_ext, mode, marker):
    normalized_name = str(name_no_ext or "")
    normalized_marker = str(marker or "")
    if not normalized_name or not normalized_marker:
        return False
    if mode == "prefix":
        return normalized_name.startswith(normalized_marker)
    return normalized_name.endswith(normalized_marker)


def _get_preview_mode_detail(app, task_type):
    try:
        if task_type == "pdf" and getattr(app, "pdf_mode_var", None) is not None:
            mode = str(app.pdf_mode_var.get() or "")
            return _get_feature_preview_mode_label(task_type, mode, mode)
        if task_type == "image":
            mode = _get_image_pdf_mode(app)
            return _get_feature_preview_mode_label(task_type, mode, mode)
        if task_type == "convert":
            return _get_convert_preview_detail(app)
        if task_type == "zip" and getattr(app, "zip_mode_var", None) is not None:
            mode = str(app.zip_mode_var.get() or "")
            return _get_feature_preview_mode_label(task_type, mode, mode)
        if task_type == "file" and getattr(app, "file_mode_var", None) is not None:
            mode = str(app.file_mode_var.get() or "")
            return _get_feature_preview_mode_label(task_type, mode, mode)
        if task_type == "remove_wm":
            mode = _get_remove_wm_mode(app)
            return _get_feature_preview_mode_label(task_type, mode, _get_remove_wm_mode_label(mode))
        if task_type == "audio":
            mode, target_fmt, bitrate, _delete_source = _get_audio_task_args(app)
            if mode == "transcribe":
                args = _get_audio_transcribe_args(app)
                return (
                    "语音转文字 "
                    f"{args.get('model_name', 'base')} / "
                    f"{args.get('language', '自动识别')} / "
                    f"{args.get('output_format', 'txt')}"
                )
            return f"{mode} -> {target_fmt} ({bitrate})"
        if task_type == "watermark":
            detail = _get_feature_preview_mode_label(task_type, "default", "添加水印")
            if bool(_safe_var_get(app, "wm_copy_guard_enabled_var", False)):
                strength = WATERMARK_COPY_GUARD_STRENGTH_LABELS.get(
                    _get_watermark_copy_guard_strength(app),
                    "标准",
                )
                detail += f" + 复制干扰层（{strength}）"
            return detail
    except Exception:
        return ""
    return ""


def _collect_preview_files(app, input_value, task_type):
    normalized_input = _normalize_input_path_value(input_value)
    if not normalized_input:
        return []
    try:
        if task_type == "zip":
            mode = _get_zip_mode(app)
            max_depth = _get_zip_max_depth(app)
            return [
                str(job.get("source") or "")
                for job in plan_zip_archives(normalized_input, mode, max_depth=max_depth)
                if job.get("source")
            ]
        if task_type == "pdf":
            all_files = app.collect_input_files(normalized_input, "pdf")
            return [item for item in all_files if str(item).lower().endswith(".pdf")]
        if task_type == "image" and _get_image_pdf_mode(app) in {"to_pdf", "merge_pdf"}:
            return _collect_image_to_pdf_files(app, normalized_input)
        if task_type == "convert":
            return _collect_convert_files(app, normalized_input)
        if task_type == "audio":
            return _collect_audio_files(app, normalized_input)
        return list(app.collect_input_files(normalized_input, task_type))
    except Exception:
        if normalized_input and os.path.isfile(normalized_input):
            return [normalized_input]
    return []


def _split_watermark_preview_files(app, files):
    rule = _get_watermark_filename_rule(app)
    kept_by_type, type_skipped = _filter_watermark_files_by_type_skip(app, files)
    skipped_keys = {
        _normalize_watermark_file_key(path)
        for path in type_skipped
        if _normalize_watermark_file_key(path)
    }
    if not rule:
        return list(kept_by_type), len(skipped_keys)

    mode, marker = rule
    eligible = []
    for path in kept_by_type:
        path_key = _normalize_watermark_file_key(path)
        if path_key and path_key in skipped_keys:
            continue
        try:
            name_no_ext = os.path.splitext(os.path.basename(str(path)))[0]
        except Exception:
            eligible.append(path)
            continue
        if _watermark_filename_matches_rule(name_no_ext, mode, marker):
            if path_key:
                skipped_keys.add(path_key)
            continue
        eligible.append(path)
    return eligible, len(skipped_keys)


def _count_watermark_preview_skips(app, files):
    _eligible, skipped = _split_watermark_preview_files(app, files)
    return skipped


def _count_watermark_checkpoint_preview_skips(app, input_value, files):
    """Count reusable checkpoint entries without changing checkpoint state."""

    normalized_input = _normalize_input_path_value(input_value)
    if not normalized_input or not files:
        return 0
    try:
        is_single_input = os.path.isfile(normalized_input)
        input_root = os.path.dirname(normalized_input) if is_single_input else normalized_input
        requested_strategy = _get_task_output_strategy(app, "watermark")
        actual_strategy = requested_strategy if is_single_input else "result_folder"
        if actual_strategy == "overwrite":
            return 0
        output_root = (
            input_root
            if actual_strategy == "same_dir"
            else os.path.join(input_root, RESULT_FOLDER_NAME)
        )
        settings = _get_watermark_settings(app)
        checkpoint_path = _watermark_checkpoint_path(normalized_input, output_root)
        checkpoint = _watermark_checkpoint_load(checkpoint_path)
        header = checkpoint.get("header") if isinstance(checkpoint, dict) else None
        expected_identity = _watermark_checkpoint_identity(
            normalized_input,
            actual_strategy,
            settings,
        )
        if not isinstance(header, dict) or str(header.get("identity") or "") != expected_identity:
            return 0

        eligible_files, _rule_skipped = _split_watermark_preview_files(app, files)
        reusable_count = 0
        for source in eligible_files:
            output = _watermark_planned_output_path(
                source,
                input_root,
                output_root,
                actual_strategy,
                settings,
                is_single_input=is_single_input,
            )
            if output and _watermark_checkpoint_item_reusable(checkpoint, source, output):
                reusable_count += 1
        return reusable_count
    except Exception as exc:
        _debug(f"start_preview:watermark_checkpoint_error:{exc}")
        return 0


def _get_start_preview_risks(app, task_type, output_strategy):
    risks = []
    try:
        if output_strategy == "overwrite":
            risks.append("将覆盖原文件")
        if task_type == "remove_wm" and _get_remove_wm_overwrite_original(app):
            risks.append("去水印将直接覆盖原文件")
        if task_type == "watermark" and bool(_safe_var_get(app, "wm_delete_var", False)):
            risks.append("水印完成后会删除源文件")
        if task_type == "pdf" and bool(_safe_var_get(app, "pdf_delete_var", False)):
            risks.append("PDF 处理完成后会删除源文件")
        if task_type == "image" and bool(_safe_var_get(app, "img_delete_var", False)):
            risks.append("图片处理完成后会删除源文件")
        if task_type == "audio":
            _mode, _fmt, _bitrate, delete_source = _get_audio_task_args(app)
            if delete_source:
                risks.append("音视频转换完成后会删除源文件")
        if task_type == "file" and str(_safe_var_get(app, "file_mode_var", "")) == "dedup":
            risks.append("文件去重可能移动或删除重复文件")
    except Exception as exc:
        _debug(f"start_preview:risk_error:{exc}")
    return risks


def _build_start_preview(app, input_value=None, task_type=None):
    task_type = str(task_type or getattr(app, "current_task", "") or "")
    normalized_input = _normalize_input_path_value(input_value if input_value is not None else _safe_var_get(app, "input_path", ""))
    files = _collect_preview_files(app, normalized_input, task_type)
    rule_skipped = _count_watermark_preview_skips(app, files) if task_type == "watermark" else 0
    checkpoint_skipped = (
        _count_watermark_checkpoint_preview_skips(app, normalized_input, files)
        if task_type == "watermark"
        else 0
    )
    skipped = rule_skipped + checkpoint_skipped
    effective_count = max(0, len(files) - skipped)
    output_strategy = _get_task_output_strategy(app, task_type)
    mode_detail = _get_preview_mode_detail(app, task_type)
    risks = _get_start_preview_risks(app, task_type, output_strategy)
    return {
        "task_type": task_type,
        "task_label": _get_feature_label(task_type),
        "mode_detail": mode_detail,
        "input": normalized_input,
        "input_kind": "单文件" if normalized_input and os.path.isfile(normalized_input) else "文件夹",
        "total_count": len(files),
        "skipped_count": skipped,
        "rule_skipped_count": rule_skipped,
        "checkpoint_skipped_count": checkpoint_skipped,
        "effective_count": effective_count,
        "output_strategy": output_strategy,
        "output_strategy_label": _get_output_strategy_label(output_strategy),
        "risks": risks,
    }


def _format_start_preview_message(preview):
    lines = [
        "请确认本次任务：",
        "",
        f"功能：{preview.get('task_label')}",
    ]
    if preview.get("mode_detail"):
        lines.append(f"模式：{preview.get('mode_detail')}")
    lines.extend(
        [
            f"输入类型：{preview.get('input_kind')}",
            f"将处理：{preview.get('effective_count')} 个文件",
        ]
    )
    rule_skipped = int(preview.get("rule_skipped_count") or 0)
    checkpoint_skipped = int(preview.get("checkpoint_skipped_count") or 0)
    if rule_skipped > 0:
        lines.append(f"预计按规则跳过：{rule_skipped} 个文件")
    if checkpoint_skipped > 0:
        lines.append(f"预计断点续传跳过：{checkpoint_skipped} 个文件")
    lines.append(f"输出策略：{preview.get('output_strategy_label')}")
    risks = list(preview.get("risks") or [])
    if risks:
        lines.extend(["", "风险提示："])
        lines.extend(f"- {item}" for item in risks)
    lines.extend(["", "是否开始处理？"])
    return "\n".join(lines)


def _confirm_start_preview(app, input_value=None, task_type=None):
    if getattr(app, "_fx_start_via_queue", False):
        return True
    preview = _build_start_preview(app, input_value=input_value, task_type=task_type)
    if int(preview.get("effective_count") or 0) <= 0:
        app.log("⚠️ [预览] 未找到可处理的文件")
        tkinter.messagebox.showwarning("任务预览", "未找到可处理的文件，请检查输入路径或当前功能模式。")
        return False
    app._fx_last_start_preview = preview
    try:
        risk_text = "；".join(preview.get("risks") or []) or "无高风险选项"
        app.log(
            f"🔎 [任务预览] {preview.get('task_label')} {preview.get('mode_detail') or ''} | "
            f"将处理 {preview.get('effective_count')} 个文件 | 输出策略：{preview.get('output_strategy_label')} | {risk_text}"
        )
    except Exception:
        pass
    return bool(
        tkinter.messagebox.askokcancel(
            "任务预览",
            _format_start_preview_message(preview),
            parent=app,
        )
    )


def _patch_start_preview_confirmation():
    try:
        original_on_start_click = FengxiToolboxApp.on_start_click
    except Exception as exc:
        _debug(f"start_preview:missing:{exc}")
        return
    if getattr(original_on_start_click, "__fx_start_preview_patch__", False):
        return

    def patched_on_start_click(self):
        if getattr(self, "is_running", False):
            return original_on_start_click(self)
        if getattr(self, "current_task", None) in {"help", "donate"}:
            return original_on_start_click(self)
        selected_input = _normalize_input_path_value(self.input_path.get())
        if selected_input and os.path.exists(selected_input):
            if not _confirm_start_preview(self, selected_input, getattr(self, "current_task", "")):
                return
        return original_on_start_click(self)

    patched_on_start_click.__fx_start_preview_patch__ = True
    FengxiToolboxApp.on_start_click = patched_on_start_click
    _debug("start_preview:patch_installed")


_patch_start_preview_confirmation()


def _get_user_pref_file():
    return _get_user_pref_root() / "user_prefs.json"


def _get_queue_history_file():
    return _get_user_pref_root() / "queue_history.json"


def _user_prefs_context():
    return UserPrefsContext(
        pref_file=_get_user_pref_file,
        output_strategy_values=OUTPUT_STRATEGY_VALUES,
        output_strategy_default=OUTPUT_STRATEGY_DEFAULT,
        remove_wm_values=REMOVE_WM_MODE_VALUES,
        remove_wm_default=REMOVE_WM_MODE_DEFAULT,
        remove_wm_label_to_value=REMOVE_WM_MODE_LABEL_TO_VALUE,
        preset_categories=tuple(PRESET_CATEGORY_LABELS.keys()),
        preset_category_labels=PRESET_CATEGORY_LABELS,
        preset_category_to_task=PRESET_CATEGORY_TO_TASK,
        preset_label_to_category=PRESET_LABEL_TO_CATEGORY,
        debug=_debug,
    )


def _load_user_prefs():
    return _prefs_load_user_prefs(_user_prefs_context())


def _save_user_prefs(data):
    _prefs_save_user_prefs(data, _user_prefs_context())


def _get_saved_output_strategy():
    return _prefs_get_saved_output_strategy(_user_prefs_context())


def _save_output_strategy(value):
    _prefs_save_output_strategy(value, _user_prefs_context())


def _get_saved_remove_wm_mode():
    return _prefs_get_saved_remove_wm_mode(_user_prefs_context())


def _save_remove_wm_mode(value):
    _prefs_save_remove_wm_mode(value, _user_prefs_context())


def _get_remove_wm_mode(app=None):
    var = getattr(app, "rm_wm_mode_var", None) if app is not None else None
    if var is not None:
        try:
            return _coerce_remove_wm_mode(var.get())
        except Exception:
            pass
    return _get_saved_remove_wm_mode()


def _refresh_remove_wm_mode_hint(app):
    hint_var = getattr(app, "rm_wm_mode_hint_var", None)
    if hint_var is None:
        return
    try:
        mode = _get_remove_wm_mode(app)
        hint_var.set(_get_remove_wm_mode_hint(mode))
    except Exception as exc:
        _debug(f"remove_wm_mode:hint_error:{exc}")


def _install_remove_wm_mode_memory(app):
    if getattr(app, "_fx_remove_wm_mode_memory_ready", False):
        return
    var = getattr(app, "rm_wm_mode_var", None)
    if not isinstance(var, tkinter.Variable):
        return

    saved = _get_saved_remove_wm_mode()
    try:
        app._fx_remove_wm_mode_loading = True
        var.set(_get_remove_wm_mode_label(saved))
    except Exception as exc:
        _debug(f"remove_wm_mode:load_error:{exc}")
    finally:
        app._fx_remove_wm_mode_loading = False

    def on_change(*_args, target=app):
        if getattr(target, "_fx_remove_wm_mode_loading", False):
            return
        try:
            _save_remove_wm_mode(_coerce_remove_wm_mode(var.get()))
        except Exception as exc:
            _debug(f"remove_wm_mode:save_error:{exc}")
        _refresh_remove_wm_mode_hint(target)

    try:
        var.trace_add("write", on_change)
    except Exception:
        pass

    app._fx_remove_wm_mode_memory_ready = True
    _refresh_remove_wm_mode_hint(app)


def _get_output_strategy_label(value):
    return OUTPUT_STRATEGY_VALUE_TO_LABEL.get(str(value or ""), OUTPUT_STRATEGY_VALUE_TO_LABEL[OUTPUT_STRATEGY_DEFAULT])


def _coerce_output_strategy_value(value):
    normalized = str(value or "").strip()
    if normalized in OUTPUT_STRATEGY_VALUES:
        return normalized
    mapped = OUTPUT_STRATEGY_LABEL_TO_VALUE.get(normalized)
    if mapped in OUTPUT_STRATEGY_VALUES:
        return mapped
    return OUTPUT_STRATEGY_DEFAULT


def _resolve_output_strategy(task_type, requested_value=None):
    normalized_task = str(task_type or "")
    requested = _coerce_output_strategy_value(requested_value or _get_saved_output_strategy() or OUTPUT_STRATEGY_DEFAULT)
    if _feature_forces_result_folder(normalized_task):
        return "result_folder"
    if not _feature_supports_output_strategy(normalized_task):
        return "same_dir"
    if normalized_task == "zip" and requested == "overwrite":
        return "same_dir"
    return requested


def _apply_output_strategy_to_result(result, task_type, strategy_value):
    if not isinstance(result, dict):
        return result
    requested = _coerce_output_strategy_value(strategy_value or OUTPUT_STRATEGY_DEFAULT)
    normalized = _resolve_output_strategy(task_type, requested)
    result["output_strategy_requested"] = requested
    result["output_strategy"] = normalized
    result["output_strategy_label"] = _get_output_strategy_label(normalized)
    return result


def _get_task_output_strategy(app, task_type):
    var = getattr(app, "output_strategy_var", None)
    requested = ""
    try:
        if var is not None:
            requested = str(var.get() or "").strip()
    except Exception:
        requested = ""
    return _resolve_output_strategy(task_type, _coerce_output_strategy_value(requested or None))


def _refresh_output_strategy_hint(app):
    hint_var = getattr(app, "output_strategy_hint_var", None)
    if hint_var is None:
        return
    task_type = str(getattr(app, "current_task", "") or "")
    requested = ""
    try:
        requested = str(getattr(app, "output_strategy_var").get() or "").strip()
    except Exception:
        requested = ""
    requested_value = _coerce_output_strategy_value(requested or None)
    resolved = _resolve_output_strategy(task_type, requested_value)
    label = _get_output_strategy_label(resolved)
    if _feature_forces_result_folder(task_type) and resolved != requested_value:
        text = f"当前功能将自动使用：{label}"
    elif task_type == "zip" and requested_value == "overwrite":
        text = "批量压缩不支持覆盖原文件，已自动改为同目录新文件。"
    else:
        text = f"当前输出策略：{label}"
    try:
        hint_var.set(text)
    except Exception:
        pass


def _install_output_strategy_memory(app):
    if getattr(app, "_fx_output_strategy_memory_ready", False):
        return
    var = getattr(app, "output_strategy_var", None)
    if not isinstance(var, tkinter.Variable):
        return

    saved = _get_saved_output_strategy()
    try:
        app._fx_output_strategy_loading = True
        var.set(_get_output_strategy_label(saved))
    except Exception as exc:
        _debug(f"output_strategy:load_error:{exc}")
    finally:
        app._fx_output_strategy_loading = False

    def on_change(*_args):
        if getattr(app, "_fx_output_strategy_loading", False):
            return
        try:
            _save_output_strategy(_coerce_output_strategy_value(var.get()))
        except Exception as exc:
            _debug(f"output_strategy:save_error:{exc}")
        _refresh_output_strategy_hint(app)

    try:
        var.trace_add("write", on_change)
    except Exception:
        pass

    app._fx_output_strategy_memory_ready = True
    _refresh_output_strategy_hint(app)


def _get_saved_watermark_text():
    return _prefs_get_saved_watermark_text(_user_prefs_context())


def _save_watermark_text(value):
    _prefs_save_watermark_text(value, _user_prefs_context())


def _get_saved_watermark_filename_rule_settings():
    return _prefs_get_saved_watermark_filename_rule_settings(_user_prefs_context())


def _save_watermark_filename_rule_settings(app):
    skip_var = getattr(app, "wm_skip_hyphen_var", None)
    mode_var = getattr(app, "wm_skip_name_position_var", None)
    marker_var = getattr(app, "wm_skip_name_text_var", None)
    copy_var = getattr(app, "wm_copy_skipped_var", None)
    if skip_var is None and mode_var is None and marker_var is None:
        return

    enabled = False
    position = "结尾"
    marker = "-"

    try:
        if skip_var is not None:
            enabled = bool(skip_var.get())
    except Exception:
        enabled = False

    try:
        if mode_var is not None:
            raw_position = str(mode_var.get() or "").strip()
            position = _watermark_filename_rule_mode_to_label(
                _normalize_watermark_filename_rule_mode(raw_position)
            )
    except Exception:
        position = "结尾"

    try:
        if marker_var is not None:
            marker = str(marker_var.get() or "")
    except Exception:
        marker = "-"

    copy_skipped = False
    try:
        if copy_var is not None:
            copy_skipped = bool(copy_var.get())
    except Exception:
        copy_skipped = False

    _prefs_save_watermark_filename_rule_settings(
        _user_prefs_context(),
        enabled=enabled,
        position=position,
        marker=marker,
        copy_skipped=copy_skipped,
    )


def _flush_watermark_filename_rule_persistence(app):
    if getattr(app, "_fx_wm_filename_rule_loading", False):
        return
    after_id = getattr(app, "_fx_wm_filename_rule_save_after_id", None)
    if after_id is not None:
        try:
            app.after_cancel(after_id)
        except Exception:
            pass
        app._fx_wm_filename_rule_save_after_id = None
    _save_watermark_filename_rule_settings(app)


def _schedule_watermark_filename_rule_persistence(app, delay_ms=300):
    if getattr(app, "_fx_wm_filename_rule_loading", False):
        return
    after_id = getattr(app, "_fx_wm_filename_rule_save_after_id", None)
    if after_id is not None:
        try:
            app.after_cancel(after_id)
        except Exception:
            pass

    def persist_later(target=app):
        target._fx_wm_filename_rule_save_after_id = None
        _save_watermark_filename_rule_settings(target)

    try:
        app._fx_wm_filename_rule_save_after_id = app.after(delay_ms, persist_later)
    except Exception:
        _flush_watermark_filename_rule_persistence(app)


def _install_watermark_filename_rule_memory(app):
    if getattr(app, "_fx_wm_filename_rule_memory_ready", False):
        return

    skip_var = getattr(app, "wm_skip_hyphen_var", None)
    mode_var = getattr(app, "wm_skip_name_position_var", None)
    marker_var = getattr(app, "wm_skip_name_text_var", None)
    copy_var = getattr(app, "wm_copy_skipped_var", None)

    saved = _get_saved_watermark_filename_rule_settings()
    try:
        app._fx_wm_filename_rule_loading = True
        if "enabled" in saved and skip_var is not None:
            skip_var.set(bool(saved["enabled"]))
        if "position" in saved and mode_var is not None:
            mode_var.set(
                _watermark_filename_rule_mode_to_label(
                    _normalize_watermark_filename_rule_mode(saved["position"])
                )
            )
        if "marker" in saved and marker_var is not None:
            marker_var.set(saved["marker"])
        if "copy_skipped" in saved and copy_var is not None:
            copy_var.set(bool(saved["copy_skipped"]))
    except Exception as exc:
        _debug(f"wm_filename_rule_memory:load_error:{exc}")
    finally:
        app._fx_wm_filename_rule_loading = False

    def on_var_change(*_args, target=app):
        _schedule_watermark_filename_rule_persistence(target)

    def on_focus_out(_event=None, target=app):
        _flush_watermark_filename_rule_persistence(target)

    for var in (skip_var, mode_var, marker_var, copy_var):
        try:
            if var is not None:
                var.trace_add("write", on_var_change)
        except Exception:
            pass

    entry = getattr(app, "wm_skip_name_entry", None)
    if entry is not None:
        try:
            entry.bind("<FocusOut>", on_focus_out, add="+")
            entry.bind("<Return>", on_focus_out, add="+")
        except Exception:
            pass

    app._fx_wm_filename_rule_memory_ready = True


def _read_watermark_text_widget(app):
    widget = getattr(app, "wm_text", None)
    if widget is None:
        return ""
    try:
        return widget.get("1.0", "end-1c")
    except Exception:
        return ""


def _flush_watermark_text_persistence(app):
    if getattr(app, "_fx_wm_text_loading", False):
        return
    after_id = getattr(app, "_fx_wm_text_save_after_id", None)
    if after_id is not None:
        try:
            app.after_cancel(after_id)
        except Exception:
            pass
        app._fx_wm_text_save_after_id = None
    _save_watermark_text(_read_watermark_text_widget(app))


def _schedule_watermark_text_persistence(app, delay_ms=400):
    if getattr(app, "_fx_wm_text_loading", False):
        return
    after_id = getattr(app, "_fx_wm_text_save_after_id", None)
    if after_id is not None:
        try:
            app.after_cancel(after_id)
        except Exception:
            pass

    def persist_later(target=app):
        target._fx_wm_text_save_after_id = None
        _save_watermark_text(_read_watermark_text_widget(target))

    try:
        app._fx_wm_text_save_after_id = app.after(delay_ms, persist_later)
    except Exception:
        _flush_watermark_text_persistence(app)


def _install_watermark_text_memory(app):
    widget = getattr(app, "wm_text", None)
    if widget is None or getattr(app, "_fx_wm_text_memory_ready", False):
        return

    saved_text = _get_saved_watermark_text()
    if saved_text:
        current_text = _read_watermark_text_widget(app)
        if current_text != saved_text:
            try:
                app._fx_wm_text_loading = True
                widget.delete("1.0", "end")
                widget.insert("1.0", saved_text)
            except Exception as exc:
                _debug(f"wm_text_memory:load_error:{exc}")
            finally:
                app._fx_wm_text_loading = False

    def on_change(_event=None, target=app):
        _schedule_watermark_text_persistence(target)

    def on_focus_out(_event=None, target=app):
        _flush_watermark_text_persistence(target)

    try:
        widget.bind("<KeyRelease>", on_change, add="+")
        widget.bind("<FocusOut>", on_focus_out, add="+")
    except Exception:
        pass

    inner_text = getattr(widget, "_textbox", None)
    if inner_text is not None:
        try:
            inner_text.bind("<KeyRelease>", on_change, add="+")
            inner_text.bind("<FocusOut>", on_focus_out, add="+")
            inner_text.bind("<<Paste>>", on_change, add="+")
        except Exception:
            pass

    app._fx_wm_text_memory_ready = True


def _filter_watermark_result_folder_inputs(input_folder, files):
    """Exclude generated result folders before the runtime watermark scan."""

    values = list(files or [])
    normalized_input = _normalize_input_path_value(input_folder)
    if not normalized_input or os.path.isfile(normalized_input) or not os.path.isdir(normalized_input):
        return values, []

    result_name = str(RESULT_FOLDER_NAME).strip().rstrip(" .").casefold()
    input_root = Path(normalized_input)
    try:
        input_root_resolved = input_root.resolve()
    except Exception:
        input_root_resolved = input_root.absolute()

    if str(input_root.name).strip().rstrip(" .").casefold() == result_name:
        return [], values

    filtered = []
    excluded = []
    for item in values:
        try:
            source_path = Path(_normalize_input_path_value(item)).resolve()
            relative_parts = source_path.relative_to(input_root_resolved).parts
            if any(str(part).strip().rstrip(" .").casefold() == result_name for part in relative_parts[:-1]):
                excluded.append(item)
                continue
        except Exception:
            pass
        filtered.append(item)
    return filtered, excluded


def _patch_watermark_filename_rule_ui():
    try:
        original_init_watermark_ui = FengxiToolboxApp.init_watermark_ui
        original_collect_input_files = FengxiToolboxApp.collect_input_files
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_watermark_filename_rule:missing:{exc}")
        return

    if getattr(original_init_watermark_ui, "__fx_watermark_filename_rule_patch__", False):
        return

    def patched_init_watermark_ui(self):
        original_init_watermark_ui(self)
        if getattr(self, "_fx_wm_filename_rule_ui_ready", False):
            return
        try:
            self.wm_skip_name_position_var = tkinter.StringVar(value="结尾")
            self.wm_skip_name_text_var = tkinter.StringVar(value="-")
            self.wm_copy_skipped_var = tkinter.BooleanVar(value=False)
            _ensure_watermark_type_skip_vars(self)

            skip_switch = _find_watermark_skip_switch(getattr(self, "tab_wm", None))
            controls_parent = getattr(skip_switch, "master", None) if skip_switch is not None else None
            controls_row = customtkinter.CTkFrame(
                controls_parent if controls_parent is not None else self.tab_wm,
                height=82,
                fg_color="transparent",
            )
            controls_row._fx_wm_filename_rule_controls = True
            try:
                controls_row.pack_propagate(False)
            except Exception:
                pass

            if skip_switch is not None:
                try:
                    skip_switch.configure(text="按文件名规则跳过")
                except Exception:
                    pass
                controls_row.pack(after=skip_switch, fill="x", padx=0, pady=(0, 4))
            else:
                controls_row.grid(row=99, column=0, columnspan=2, sticky="ew", padx=18, pady=(4, 8))

            combo_style = {}
            try:
                combo_style = self._get_combo_style()
            except Exception:
                combo_style = {}
            option_menu_style = _get_option_menu_style(combo_style)

            fields_row = customtkinter.CTkFrame(controls_row, fg_color="transparent", height=30)
            fields_row.pack(fill="x", pady=(0, 2))

            customtkinter.CTkLabel(
                fields_row,
                text="匹配位置",
                text_color=globals().get("COLOR_TEXT_SOFT"),
                font=customtkinter.CTkFont(size=11),
                height=30,
            ).pack(side="left", padx=(0, 8))

            customtkinter.CTkOptionMenu(
                fields_row,
                variable=self.wm_skip_name_position_var,
                values=["结尾", "开头"],
                width=82,
                height=30,
                **option_menu_style,
            ).pack(side="left", padx=(0, 8))

            self.wm_skip_name_entry = customtkinter.CTkEntry(
                fields_row,
                width=126,
                height=30,
                textvariable=self.wm_skip_name_text_var,
                placeholder_text="-",
            )
            self.wm_skip_name_entry.pack(side="left", fill="x", expand=True, padx=(0, 0))
            controls_row._fx_wm_filename_rule_entry = self.wm_skip_name_entry

            customtkinter.CTkLabel(
                controls_row,
                text="留空默认 “-”，可填写任意开头或结尾字符",
                text_color=globals().get("COLOR_TEXT_SOFT"),
                font=customtkinter.CTkFont(size=11),
                height=18,
                anchor="w",
            ).pack(anchor="w", fill="x")
            customtkinter.CTkCheckBox(
                controls_row,
                text="跳过文件复制到输出文件夹",
                variable=self.wm_copy_skipped_var,
                height=22,
                font=customtkinter.CTkFont(size=11),
            ).pack(anchor="w", fill="x", pady=(1, 0))
            _ensure_watermark_type_skip_controls(self, controls_parent if controls_parent is not None else self.tab_wm, after_widget=controls_row)
            _install_watermark_filename_rule_memory(self)
        except Exception as exc:
            _debug(f"patch_watermark_filename_rule:init_ui_error:{exc}")
        self._fx_wm_filename_rule_ui_ready = True

    def patched_collect_input_files(self, input_folder, task_type):
        files = original_collect_input_files(self, input_folder, task_type)
        if task_type != "watermark":
            return files

        files, result_folder_files = _filter_watermark_result_folder_inputs(input_folder, files)
        if result_folder_files:
            try:
                log_key = (_normalize_input_path_value(input_folder), len(result_folder_files))
                if getattr(self, "_fx_wm_result_folder_filter_log_key", None) != log_key:
                    self._fx_wm_result_folder_filter_log_key = log_key
                    if os.path.basename(str(_normalize_input_path_value(input_folder)).rstrip("\\/")) == RESULT_FOLDER_NAME:
                        self.log(
                            f"[批量水印] 已忽略结果目录自身的 {len(result_folder_files)} 个文件，避免生成嵌套结果目录。"
                        )
                    else:
                        self.log(
                            f"[批量水印] 已排除已有结果目录中的 {len(result_folder_files)} 个文件，避免重复加水印。"
                        )
            except Exception:
                pass

        rule = getattr(self, "_fx_wm_filename_rule_runtime", None)
        if not rule:
            try:
                self._fx_wm_filename_rule_skipped_files = []
            except Exception:
                pass
            return files

        mode, marker = rule
        filtered_files = []
        skipped_files = []
        for path in files:
            try:
                name_no_ext = os.path.splitext(os.path.basename(str(path)))[0]
            except Exception:
                filtered_files.append(path)
                continue
            if _watermark_filename_matches_rule(name_no_ext, mode, marker):
                skipped_files.append(path)
            else:
                filtered_files.append(path)

        if skipped_files:
            try:
                self._fx_wm_filename_rule_skipped_files = list(skipped_files)
            except Exception:
                pass
            try:
                direction = "开头" if mode == "prefix" else "结尾"
                preview = ", ".join(os.path.basename(item) for item in skipped_files[:5])
                if len(skipped_files) > 5:
                    preview = f"{preview} 等 {len(skipped_files)} 个文件"
                self.log(f"⏭️ [智能水印] 已跳过文件名{direction}为“{marker}”的文件：{preview}")
            except Exception:
                pass
        else:
            try:
                self._fx_wm_filename_rule_skipped_files = []
            except Exception:
                pass
        return filtered_files

    def patched_run_process(self, input_folder, task_type):
        if task_type != "watermark":
            return original_run_process(self, input_folder, task_type)

        try:
            _flush_watermark_filename_rule_persistence(self)
        except Exception as exc:
            _debug(f"patch_watermark_filename_rule:flush_before_run_error:{exc}")

        rule = _get_watermark_filename_rule(self)
        if not rule:
            return original_run_process(self, input_folder, task_type)

        skip_var = getattr(self, "wm_skip_hyphen_var", None)
        previous_runtime_rule = getattr(self, "_fx_wm_filename_rule_runtime", None)
        previous_skip_value = None
        restore_skip_value = False
        previous_rule_loading = getattr(self, "_fx_wm_filename_rule_loading", False)
        try:
            self._fx_wm_filename_rule_runtime = rule
            self._fx_wm_filename_rule_skipped_files = []
            if skip_var is not None:
                try:
                    self._fx_wm_filename_rule_loading = True
                    previous_skip_value = skip_var.get()
                    skip_var.set(False)
                    restore_skip_value = True
                except Exception as exc:
                    _debug(f"patch_watermark_filename_rule:disable_builtin_skip_error:{exc}")
            return original_run_process(self, input_folder, task_type)
        finally:
            self._fx_wm_filename_rule_runtime = previous_runtime_rule
            if restore_skip_value and skip_var is not None:
                try:
                    skip_var.set(previous_skip_value)
                except Exception:
                    pass
            self._fx_wm_filename_rule_loading = previous_rule_loading

    patched_init_watermark_ui.__fx_watermark_filename_rule_patch__ = True
    patched_collect_input_files.__fx_watermark_filename_rule_patch__ = True
    patched_run_process.__fx_watermark_filename_rule_patch__ = True
    FengxiToolboxApp.init_watermark_ui = patched_init_watermark_ui
    FengxiToolboxApp.collect_input_files = patched_collect_input_files
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_watermark_filename_rule:installed")


_patch_watermark_filename_rule_ui()


def _patch_watermark_text_memory_ui():
    try:
        original_init_watermark_ui = FengxiToolboxApp.init_watermark_ui
        original_on_start_click = FengxiToolboxApp.on_start_click
    except Exception as exc:
        _debug(f"patch_watermark_text_memory:missing:{exc}")
        return

    if getattr(original_init_watermark_ui, "__fx_wm_text_memory_patch__", False):
        return

    def patched_init_watermark_ui(self):
        original_init_watermark_ui(self)
        try:
            _install_watermark_text_memory(self)
        except Exception as exc:
            _debug(f"patch_watermark_text_memory:init_error:{exc}")

    def patched_on_start_click(self):
        if getattr(self, "current_task", None) == "watermark":
            try:
                _flush_watermark_text_persistence(self)
            except Exception as exc:
                _debug(f"patch_watermark_text_memory:flush_before_run_error:{exc}")
        return original_on_start_click(self)

    patched_init_watermark_ui.__fx_wm_text_memory_patch__ = True
    patched_on_start_click.__fx_wm_text_memory_patch__ = True
    FengxiToolboxApp.init_watermark_ui = patched_init_watermark_ui
    FengxiToolboxApp.on_start_click = patched_on_start_click
    _debug("patch_watermark_text_memory:installed")


_patch_watermark_text_memory_ui()


def _patch_remove_wm_output_ui():
    try:
        original_init_remove_wm_ui = FengxiToolboxApp.init_remove_wm_ui
    except Exception as exc:
        _debug(f"patch_remove_wm_output_ui:missing:{exc}")
        return

    if getattr(original_init_remove_wm_ui, "__fx_remove_wm_output_ui_patch__", False):
        return

    def patched_init_remove_wm_ui(self):
        original_init_remove_wm_ui(self)
        if getattr(self, "_fx_remove_wm_output_ui_ready", False):
            return
        try:
            self.rm_wm_overwrite_original = tkinter.BooleanVar(value=False)
            self.rm_wm_mode_var = tkinter.StringVar(value=_get_remove_wm_mode_label(REMOVE_WM_MODE_DEFAULT))
            self.rm_wm_mode_hint_var = tkinter.StringVar(value="")
            card = self.tab_rm_wm.winfo_children()[0]
            body = card.winfo_children()[1]
            widgets = body.winfo_children()
            separator = widgets[4] if len(widgets) > 4 else None

            combo_style = {}
            try:
                combo_style = self._get_combo_style()
            except Exception:
                combo_style = {}
            option_menu_style = _get_option_menu_style(combo_style)

            mode_row = customtkinter.CTkFrame(body, fg_color="transparent")
            mode_row.pack(anchor="w", fill="x", padx=0, pady=(2, 4))
            if separator is not None:
                mode_row.pack_configure(before=separator)
            mode_row.grid_columnconfigure(1, weight=1)

            customtkinter.CTkLabel(
                mode_row,
                text="去水印强度",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=12),
                anchor="w",
            ).grid(row=0, column=0, sticky="w", padx=(0, 10))

            customtkinter.CTkOptionMenu(
                mode_row,
                variable=self.rm_wm_mode_var,
                values=[REMOVE_WM_MODE_VALUE_TO_LABEL[key] for key in REMOVE_WM_MODE_VALUES],
                width=150,
                height=30,
                command=lambda _value=None, target=self: _refresh_remove_wm_mode_hint(target),
                **option_menu_style,
            ).grid(row=0, column=1, sticky="w")

            mode_hint = customtkinter.CTkLabel(
                body,
                textvariable=self.rm_wm_mode_hint_var,
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                justify="left",
                wraplength=620,
            )
            mode_hint.pack(anchor="w", padx=0, pady=(0, 8))
            if separator is not None:
                mode_hint.pack_configure(before=separator)
            _install_remove_wm_mode_memory(self)

            switch = customtkinter.CTkSwitch(
                body,
                text="单文件时直接覆盖原文件（谨慎）",
                variable=self.rm_wm_overwrite_original,
                **self._get_switch_style(),
            )
            switch.pack(anchor="w", padx=0, pady=(4, 10))
            if separator is not None:
                switch.pack_configure(before=separator)

            note = customtkinter.CTkLabel(
                body,
                text="说明：单文件默认在同目录生成新的“_去水印”文件；开启后仅在处理成功时替换原文件。文件夹模式仍输出到【处理完成】结果文件夹。",
                text_color=COLOR_TEXT_SOFT,
                font=customtkinter.CTkFont(size=11),
                justify="left",
                wraplength=620,
            )
            note.pack(anchor="w", padx=0, pady=(0, 8))
            if separator is not None:
                note.pack_configure(before=separator)
        except Exception as exc:
            _debug(f"patch_remove_wm_output_ui:error:{exc}")
        self._fx_remove_wm_output_ui_ready = True

    patched_init_remove_wm_ui.__fx_remove_wm_output_ui_patch__ = True
    FengxiToolboxApp.init_remove_wm_ui = patched_init_remove_wm_ui
    _debug("patch_remove_wm_output_ui:installed")


_patch_remove_wm_output_ui()


def _patch_remove_wm_pdf_fallback():
    try:
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_remove_wm_pdf_fallback:missing:{exc}")
        return

    if getattr(original_run_process, "__fx_remove_wm_pdf_patch__", False):
        return

    def patched_run_process(self, input_folder, task_type):
        if task_type == "remove_wm":
            previous_remove_mode = _push_remove_wm_runtime_mode(_get_remove_wm_mode(self))
            try:
                try:
                    all_files = self.collect_input_files(input_folder, task_type)
                    if any(path.lower().endswith(".pdf") for path in all_files):
                        try:
                            _run_remove_wm_task(self, input_folder, original_run_process)
                        except Exception as exc:
                            self.log(f"🔥 [严重错误] {exc}")
                        finally:
                            self.reset_ui()
                        return
                except Exception as exc:
                    _debug(f"patch_remove_wm_pdf_fallback:run_process_error:{exc}")
                return original_run_process(self, input_folder, task_type)
            finally:
                _pop_remove_wm_runtime_mode(previous_remove_mode)
        return original_run_process(self, input_folder, task_type)

    patched_run_process.__fx_remove_wm_pdf_patch__ = True
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_remove_wm_pdf_fallback:installed")


_patch_remove_wm_pdf_fallback()


def _make_queue_task_id():
    return f"task_{int(time.time() * 1000)}_{os.getpid()}"


def _format_queue_time(timestamp=None):
    try:
        return time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(float(timestamp or time.time())))
    except Exception:
        return time.strftime("%Y-%m-%d %H:%M:%S")


def _safe_widget_get(widget):
    try:
        if isinstance(widget, customtkinter.CTkTextbox):
            return widget.get("1.0", "end-1c")
        if hasattr(widget, "get"):
            return widget.get()
    except Exception:
        return None
    return None


def _safe_widget_set(widget, value):
    try:
        if isinstance(widget, customtkinter.CTkTextbox):
            widget.delete("1.0", "end")
            widget.insert("1.0", str(value or ""))
            return True
        if isinstance(widget, customtkinter.CTkEntry):
            widget.delete(0, "end")
            widget.insert(0, str(value or ""))
            return True
        if hasattr(widget, "set"):
            widget.set(value)
            return True
    except Exception:
        return False
    return False


def _safe_var_get(app, name, default=None):
    try:
        var = getattr(app, name, None)
        if isinstance(var, tkinter.Variable):
            return var.get()
    except Exception:
        pass
    return default


def _safe_var_set(app, name, value):
    try:
        var = getattr(app, name, None)
        if isinstance(var, tkinter.Variable):
            var.set(value)
            return True
    except Exception:
        pass
    return False


def _safe_named_widget_get(app, name, default=""):
    try:
        widget = getattr(app, name, None)
        value = _safe_widget_get(widget)
        if value is not None:
            return value
    except Exception:
        pass
    return default


def _safe_named_widget_set(app, name, value):
    try:
        return _safe_widget_set(getattr(app, name, None), value)
    except Exception:
        return False


WATERMARK_DEFAULT_COLOR = "#C0C0C0"
WATERMARK_COPY_GUARD_STRENGTH_LABELS = {
    "light": "轻度",
    "standard": "标准",
    "strong": "强力",
}
WATERMARK_COPY_GUARD_STRENGTH_HINTS = {
    "light": "轻度：PDF 每页 4 个；Word 每个文件 4 段。",
    "standard": "标准：PDF 每页 8 个；Word 每个文件 8 段。",
    "strong": "强力：PDF 每页 14 个；Word 每个文件 14 段。",
}
WATERMARK_COPY_GUARD_STRENGTH_VALUES = tuple(WATERMARK_COPY_GUARD_STRENGTH_LABELS)
WATERMARK_COPY_GUARD_LABEL_TO_VALUE = {
    label: value for value, label in WATERMARK_COPY_GUARD_STRENGTH_LABELS.items()
}


def _get_watermark_copy_guard_strength(app):
    raw_value = str(_safe_var_get(app, "wm_copy_guard_strength_var", "标准") or "标准").strip()
    if raw_value in WATERMARK_COPY_GUARD_STRENGTH_VALUES:
        return raw_value
    return WATERMARK_COPY_GUARD_LABEL_TO_VALUE.get(raw_value, "standard")
WATERMARK_WORD_EXTS = {".doc", ".docx"}
WATERMARK_PPT_EXTS = {".ppt", ".pptx"}
WATERMARK_PDF_EXTS = {".pdf"}
WATERMARK_SUPPORTED_EXTS = WATERMARK_PDF_EXTS | WATERMARK_WORD_EXTS | WATERMARK_PPT_EXTS


def _split_watermark_supported_files(files):
    """Partition inputs before they enter watermark processing or checkpoint work."""

    supported = []
    unsupported = []
    for item in files or []:
        try:
            suffix = Path(str(item)).suffix.lower()
        except Exception:
            suffix = ""
        (supported if suffix in WATERMARK_SUPPORTED_EXTS else unsupported).append(item)
    return supported, unsupported


def _normalize_watermark_color(value, default=WATERMARK_DEFAULT_COLOR):
    try:
        return _watermark_core_color_to_hex(value, default=default)
    except Exception:
        return str(default or WATERMARK_DEFAULT_COLOR).upper()


def _get_watermark_color(app):
    return _normalize_watermark_color(_safe_var_get(app, "wm_color_var", WATERMARK_DEFAULT_COLOR))


def _set_watermark_color(app, color):
    normalized = _normalize_watermark_color(color)
    if str(_safe_var_get(app, "wm_color_var", "") or "").strip().upper() == normalized:
        return normalized
    if not _safe_var_set(app, "wm_color_var", normalized):
        try:
            app.wm_color_var = tkinter.StringVar(value=normalized)
        except Exception:
            pass
    return normalized


def _watermark_color_rgb_tuple(color):
    hex_color = _normalize_watermark_color(color)
    return tuple(int(hex_color[index : index + 2], 16) for index in (1, 3, 5))


def _watermark_log(app, logs, message):
    text = str(message)
    logs.append(text)
    pending = getattr(app, "_fx_watermark_ui_log_pending", None)
    if pending is None:
        pending = []
        app._fx_watermark_ui_log_pending = pending
    pending.append(text)
    _flush_watermark_ui_logs(app)


def _flush_watermark_ui_logs(app, force=False):
    pending = list(getattr(app, "_fx_watermark_ui_log_pending", []) or [])
    if not pending:
        return
    now = time.monotonic()
    last_flush = float(getattr(app, "_fx_watermark_ui_log_last_flush", 0.0) or 0.0)
    if not force and len(pending) < WATERMARK_UI_LOG_BATCH_SIZE and now - last_flush < WATERMARK_UI_LOG_FLUSH_INTERVAL_SECONDS:
        return
    try:
        app.log("\n".join(pending))
    except Exception:
        pass
    finally:
        try:
            app._fx_watermark_ui_log_pending = []
            app._fx_watermark_ui_log_last_flush = now
        except Exception:
            pass


def _watermark_status_kind(status):
    text = str(status or "").strip()
    if text == "SUCCESS":
        return "success"
    lowered = text.lower()
    if lowered.startswith("skip"):
        return "skipped"
    if lowered.startswith("error") or not text:
        return "failed"
    return "failed"


def _watermark_should_preserve_original(status, suffix=""):
    text = str(status or "").strip().lower()
    if text.startswith("skip:protected") or "protected pdf requires password" in text:
        return True
    if "protected word document requires password" in text:
        return True
    if text.startswith("skip:damaged word source") or "damaged word source" in text:
        return True
    suffix_text = str(suffix or "").strip().lower()
    if suffix_text in WATERMARK_WORD_EXTS:
        unreadable_markers = [
            "word 无法读取",
            "文档可能已损坏",
            "文件似乎已经损坏",
            "cannot read",
            "may be corrupted",
            "appears to be corrupted",
        ]
        return any(marker in text for marker in unreadable_markers)
    return False


def _watermark_normalize_path_part(value):
    return str(value or "").strip().rstrip(" .").casefold()


def _resolve_watermark_source_path(src, input_root):
    src_path = Path(src)
    if src_path.exists():
        return src_path
    root = Path(input_root)
    if not root.exists():
        return src_path
    try:
        rel = Path(os.path.relpath(str(src_path), str(root)))
        if str(rel) not in {"", "."} and not str(rel).startswith(".."):
            stripped_candidate = root / Path(*[str(part).strip().rstrip(" .") for part in rel.parts])
            if stripped_candidate.exists():
                return stripped_candidate
    except Exception:
        rel = None

    try:
        wanted_parts = []
        if rel is not None and str(rel) not in {"", "."} and not str(rel).startswith(".."):
            wanted_parts = [_watermark_normalize_path_part(part) for part in rel.parts]
        wanted_name = _watermark_normalize_path_part(src_path.name)
        for candidate in root.rglob("*"):
            try:
                if not candidate.is_file():
                    continue
                if _watermark_normalize_path_part(candidate.name) != wanted_name:
                    continue
                if wanted_parts:
                    try:
                        candidate_rel_parts = [
                            _watermark_normalize_path_part(part)
                            for part in candidate.relative_to(root).parts
                        ]
                        if candidate_rel_parts[-len(wanted_parts) :] != wanted_parts:
                            continue
                    except Exception:
                        pass
                return candidate
            except Exception:
                continue
    except Exception:
        pass
    return src_path


def _unique_watermark_path(path_value):
    path = Path(path_value)
    if not path.exists():
        return str(path)
    counter = 2
    while True:
        candidate = path.with_name(f"{path.stem}_{counter}{path.suffix}")
        if not candidate.exists():
            return str(candidate)
        counter += 1


def _watermark_relative_parent(src_path, input_root):
    try:
        rel_parent = Path(os.path.relpath(str(src_path.parent), str(input_root)))
        if str(rel_parent) in {"", "."} or str(rel_parent).startswith(".."):
            return Path()
        return rel_parent
    except Exception:
        return Path()


def _watermark_safe_relative_parent(src_path, input_root):
    rel_parent = _watermark_relative_parent(src_path, input_root)
    safe_parts = []
    for part in rel_parent.parts:
        safe = str(part).strip().rstrip(" .")
        if not safe:
            safe = _sanitize_filename_component(part, fallback="folder")
        safe_parts.append(safe)
    return Path(*safe_parts) if safe_parts else Path()


def _build_watermark_output_path(src, input_root, output_root, strategy, *, convert_to_pdf=False, single_input=False):
    src_path = Path(src)
    input_root_path = Path(input_root)
    output_root_path = Path(output_root)
    source_suffix = src_path.suffix.lower()
    output_suffix = ".pdf" if convert_to_pdf and source_suffix in WATERMARK_WORD_EXTS | WATERMARK_PPT_EXTS else src_path.suffix

    if strategy == "overwrite" and output_suffix.lower() == src_path.suffix.lower():
        return str(src_path)

    if strategy == "result_folder":
        target_dir = output_root_path / _watermark_safe_relative_parent(src_path, input_root_path)
        target_name = src_path.name if output_suffix.lower() == src_path.suffix.lower() else f"{src_path.stem}{output_suffix}"
        return str(target_dir / target_name)

    target_dir = src_path.parent if single_input or strategy == "same_dir" else output_root_path
    return _unique_watermark_path(target_dir / f"{src_path.stem}_加水印{output_suffix}")


def _watermark_replace_original(staged_path, source_path):
    staged = Path(staged_path)
    source = Path(source_path)
    if not staged.exists():
        return False
    os.replace(str(staged), str(source))
    return True


def _get_slider_value_for_preview(app, name, fallback):
    value = _safe_named_widget_get(app, name, fallback)
    return _safe_float(value, fallback)


def _resolve_preview_font(app, size):
    font_name = str(_safe_var_get(app, "selected_font", "") or "").strip()
    font_path = ""
    if font_name:
        try:
            font_path = str(get_font_path_by_name(font_name) or "")
        except Exception:
            font_path = ""
    for candidate in (font_path, "C:/Windows/Fonts/msyh.ttc", "C:/Windows/Fonts/simhei.ttf"):
        try:
            if candidate and os.path.exists(candidate):
                return ImageFont.truetype(candidate, max(10, int(size)))
        except Exception:
            pass
    try:
        return ImageFont.truetype("arial.ttf", max(10, int(size)))
    except Exception:
        return ImageFont.load_default()


def _make_watermark_preview_image(app):
    width, height = 360, 92
    image = PILImage.new("RGBA", (width, height), (248, 250, 252, 255))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((10, 10, width - 10, height - 10), radius=14, fill=(255, 255, 255, 255), outline=(218, 225, 232, 255), width=2)
    for x in range(28, width - 20, 36):
        draw.line((x, 16, x - 72, height - 16), fill=(237, 241, 245, 255), width=1)

    text = (_read_watermark_text_widget(app) or "Fengxi Watermark").strip() or "Fengxi Watermark"
    preview_lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not preview_lines:
        preview_lines = ["Fengxi Watermark"]
    if len(preview_lines) > 4:
        preview_lines = preview_lines[:4]
        preview_lines[-1] = f"{preview_lines[-1][:24]}..."
    preview_text = "\n".join(preview_lines)
    requested_font_size = max(16.0, min(44.0, _get_slider_value_for_preview(app, "slider_size", 60) * 0.45))
    opacity = max(0.08, min(1.0, _get_slider_value_for_preview(app, "slider_opacity", 0.18)))
    angle = _get_slider_value_for_preview(app, "slider_angle", 45)
    red, green, blue = _watermark_color_rgb_tuple(_get_watermark_color(app))
    alpha = max(24, min(230, int(opacity * 255)))

    text_layer = PILImage.new("RGBA", (width, height), (0, 0, 0, 0))
    text_draw = ImageDraw.Draw(text_layer)
    font_size = requested_font_size
    spacing = max(2, int(round(font_size * 0.16)))
    max_text_width = width - 40
    max_text_height = height - 42
    for _ in range(12):
        font = _resolve_preview_font(app, font_size)
        try:
            bbox = text_draw.multiline_textbbox(
                (0, 0),
                preview_text,
                font=font,
                spacing=spacing,
                align="center",
            )
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
        except Exception:
            line_lengths = [len(line) for line in preview_lines] or [1]
            text_width = max(line_lengths) * int(font_size * 0.6)
            text_height = len(preview_lines) * int(font_size) + max(0, len(preview_lines) - 1) * spacing
            bbox = (0, 0, text_width, text_height)
        if text_width <= max_text_width and text_height <= max_text_height:
            break
        font_size = max(10.0, font_size - 2.0)
        spacing = max(2, int(round(font_size * 0.16)))
    text_x = (width - text_width) / 2 - bbox[0]
    text_y = (height - 42 - text_height) / 2 - bbox[1] + 6
    text_draw.multiline_text(
        (text_x, text_y),
        preview_text,
        font=font,
        fill=(red, green, blue, alpha),
        spacing=spacing,
        align="center",
    )
    rotated = text_layer.rotate(angle, resample=PILImage.Resampling.BICUBIC, center=(width / 2, height / 2))
    image.alpha_composite(rotated)

    label = f"{_get_watermark_color(app)} | {int(round(opacity * 100))}% | {int(round(angle))}°"
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((18, height - 31, 178, height - 15), radius=7, fill=(246, 248, 250, 235))
    draw.text((26, height - 30), label, fill=(92, 103, 115, 255), font=ImageFont.load_default())
    return image


def _refresh_watermark_preview(app):
    label = getattr(app, "wm_preview_label", None)
    if label is None:
        return False
    try:
        if not bool(label.winfo_exists()):
            return False
    except Exception:
        return False
    try:
        image = _make_watermark_preview_image(app)
        ctk_image = customtkinter.CTkImage(light_image=image, dark_image=image, size=(360, 92))
        app._fx_wm_preview_image = ctk_image
        label.configure(image=ctk_image, text="")
        swatch = getattr(app, "wm_color_swatch", None)
        if swatch is not None:
            swatch.configure(fg_color=_get_watermark_color(app))
        return True
    except Exception as exc:
        _debug(f"watermark_preview:refresh_error:{exc}")
        try:
            label.configure(text="预览暂不可用")
        except Exception:
            pass
        return False


def _schedule_watermark_preview_refresh(app, delay_ms=120):
    try:
        after_id = getattr(app, "_fx_wm_preview_after_id", None)
        if after_id:
            app.after_cancel(after_id)
    except Exception:
        pass

    def refresh_later(target=app):
        try:
            target._fx_wm_preview_after_id = None
        except Exception:
            pass
        try:
            if hasattr(target, "winfo_exists") and not bool(target.winfo_exists()):
                return
        except Exception:
            return
        _refresh_watermark_preview(target)

    try:
        app._fx_wm_preview_after_id = app.after(delay_ms, refresh_later)
    except Exception:
        refresh_later()


def _choose_watermark_color(app):
    current = _get_watermark_color(app)
    try:
        _rgb, selected = tkinter.colorchooser.askcolor(color=current, title="选择水印颜色")
    except Exception as exc:
        _debug(f"watermark_color:chooser_error:{exc}")
        selected = None
    if selected:
        _set_watermark_color(app, selected)
        _refresh_watermark_preview(app)


def _iter_watermark_preview_frames(tab):
    stack = list(tab.winfo_children()) if tab is not None else []
    while stack:
        widget = stack.pop()
        if getattr(widget, "_fx_wm_color_preview_controls", False):
            yield widget
        try:
            stack.extend(widget.winfo_children())
        except Exception:
            pass


def _destroy_stale_watermark_preview_frames(app, tab, keep_frame=None):
    for frame in list(_iter_watermark_preview_frames(tab)):
        if frame is keep_frame:
            continue
        try:
            frame.destroy()
        except Exception:
            pass
    if keep_frame is None and getattr(app, "_fx_wm_color_preview_frame", None) is not None:
        try:
            if not app._fx_wm_color_preview_frame.winfo_exists():
                app._fx_wm_color_preview_frame = None
        except Exception:
            app._fx_wm_color_preview_frame = None


def _pack_watermark_color_preview_frame(app, frame, left_panel):
    text_widget = getattr(app, "wm_text", None)
    before_widget = text_widget if getattr(text_widget, "master", None) is left_panel else None
    try:
        if frame.winfo_manager():
            frame.pack_forget()
    except Exception:
        pass
    try:
        if before_widget is not None:
            frame.pack(fill="x", padx=24, pady=(0, 8), before=before_widget)
        else:
            frame.pack(fill="x", padx=24, pady=(0, 8))
    except Exception:
        frame.pack(fill="x", padx=24, pady=(0, 8))


def _install_watermark_color_preview_ui(app):
    tab = getattr(app, "tab_wm", None)
    if tab is None:
        return

    left_panel = _find_watermark_text_panel(app, tab)
    existing_frame = getattr(app, "_fx_wm_color_preview_frame", None)
    if getattr(app, "_fx_wm_color_preview_ready", False):
        try:
            existing_ok = existing_frame is not None and existing_frame.winfo_exists() and existing_frame.master is left_panel
        except Exception:
            existing_ok = False
        if existing_ok:
            _destroy_stale_watermark_preview_frames(app, tab, keep_frame=existing_frame)
            _pack_watermark_color_preview_frame(app, existing_frame, left_panel)
            _tighten_watermark_tab_layout(app, tab)
            _schedule_watermark_preview_refresh(app, delay_ms=120)
            return
        app._fx_wm_color_preview_ready = False

    _destroy_stale_watermark_preview_frames(app, tab)
    _set_watermark_color(app, _safe_var_get(app, "wm_color_var", WATERMARK_DEFAULT_COLOR))

    frame = customtkinter.CTkFrame(left_panel, fg_color="transparent", height=132)
    frame._fx_wm_color_preview_controls = True
    app._fx_wm_color_preview_frame = frame
    try:
        frame.pack_propagate(False)
    except Exception:
        pass
    _pack_watermark_color_preview_frame(app, frame, left_panel)

    row = customtkinter.CTkFrame(frame, fg_color="transparent", height=30)
    row.pack(fill="x", pady=(0, 4))
    customtkinter.CTkLabel(
        row,
        text="水印颜色",
        text_color=globals().get("COLOR_TEXT_SOFT"),
        font=customtkinter.CTkFont(size=12),
        width=62,
        anchor="w",
    ).pack(side="left", padx=(0, 8))

    app.wm_color_swatch = customtkinter.CTkFrame(row, width=28, height=22, corner_radius=8, fg_color=_get_watermark_color(app))
    app.wm_color_swatch.pack(side="left", padx=(0, 8))
    try:
        app.wm_color_swatch.pack_propagate(False)
    except Exception:
        pass

    app.wm_color_entry = customtkinter.CTkEntry(row, width=88, height=28, textvariable=app.wm_color_var)
    app.wm_color_entry.pack(side="left", padx=(0, 8))
    customtkinter.CTkButton(row, text="选择", width=56, height=28, command=lambda target=app: _choose_watermark_color(target)).pack(side="left", padx=(0, 8))
    customtkinter.CTkButton(row, text="刷新预览", width=78, height=28, command=lambda target=app: _refresh_watermark_preview(target)).pack(side="left")

    app.wm_preview_label = customtkinter.CTkLabel(frame, text="", width=360, height=92, corner_radius=14)
    app.wm_preview_label.pack(fill="x", expand=True)

    def color_changed(*_args):
        _schedule_watermark_preview_refresh(app)

    try:
        app.wm_color_var.trace_add("write", color_changed)
    except Exception:
        pass

    def preview_changed(_event=None, target=app):
        _schedule_watermark_preview_refresh(target)

    for widget_name in ("wm_text", "slider_size", "slider_opacity", "slider_angle"):
        widget = getattr(app, widget_name, None)
        if widget is None:
            continue
        try:
            widget.bind("<KeyRelease>", preview_changed, add="+")
            widget.bind("<B1-Motion>", preview_changed, add="+")
            widget.bind("<ButtonRelease-1>", preview_changed, add="+")
        except Exception:
            pass
        inner_text = getattr(widget, "_textbox", None)
        if inner_text is not None:
            try:
                inner_text.bind("<KeyRelease>", preview_changed, add="+")
                inner_text.bind("<<Paste>>", preview_changed, add="+")
            except Exception:
                pass

    app._fx_wm_color_preview_ready = True
    _tighten_watermark_tab_layout(app, tab)
    _schedule_watermark_preview_refresh(app, delay_ms=450)


def _patch_watermark_color_preview_ui():
    try:
        original_init_watermark_ui = FengxiToolboxApp.init_watermark_ui
    except Exception as exc:
        _debug(f"patch_watermark_color_preview:missing:{exc}")
        return

    if getattr(original_init_watermark_ui, "__fx_wm_color_preview_patch__", False):
        return

    def patched_init_watermark_ui(self):
        original_init_watermark_ui(self)
        try:
            _install_watermark_color_preview_ui(self)
        except Exception as exc:
            _debug(f"patch_watermark_color_preview:init_error:{exc}")

    patched_init_watermark_ui.__fx_wm_color_preview_patch__ = True
    FengxiToolboxApp.init_watermark_ui = patched_init_watermark_ui
    _debug("patch_watermark_color_preview:installed")


_patch_watermark_color_preview_ui()


WATERMARK_RANGE_UI_LABELS = {
    "all": "\u6bcf\u4e00\u9875",
    "first": "\u4ec5\u7b2c\u4e00\u9875",
    "first_random": "\u7b2c\u4e00\u9875 + \u968f\u673a\u4e00\u9875",
}
WATERMARK_RANGE_UI_VALUE_BY_LABEL = {label: value for value, label in WATERMARK_RANGE_UI_LABELS.items()}


def _watermark_range_ui_label(value):
    return WATERMARK_RANGE_UI_LABELS.get(_watermark_core_normalize_page_range(value), WATERMARK_RANGE_UI_LABELS["all"])


def _set_watermark_range_var(app, value):
    normalized = WATERMARK_RANGE_UI_VALUE_BY_LABEL.get(str(value or "").strip())
    if normalized is None:
        normalized = _watermark_core_normalize_page_range(value)
    var = getattr(app, "wm_range_var", None)
    if isinstance(var, tkinter.Variable):
        try:
            var.set(normalized)
            return normalized
        except Exception:
            pass
    return normalized


def _schedule_watermark_range_last_settings_save(app):
    scheduler = globals().get("_schedule_watermark_last_settings_persistence")
    if callable(scheduler):
        try:
            scheduler(app)
        except Exception as exc:
            _debug(f"watermark_range:last_settings_save_error:{exc}")


def _select_first_random_watermark_range(app):
    _set_watermark_range_var(app, "first_random")
    _schedule_watermark_range_last_settings_save(app)


def _iter_watermark_range_radios(app):
    tab = getattr(app, "tab_wm", None)
    target_var = getattr(app, "wm_range_var", None)
    if tab is None or target_var is None:
        return []
    try:
        target_name = str(target_var)
    except Exception:
        target_name = ""
    radios = []
    stack = list(tab.winfo_children())
    while stack:
        widget = stack.pop(0)
        if isinstance(widget, customtkinter.CTkRadioButton):
            try:
                variable_name = str(widget.cget("variable"))
            except Exception:
                variable_name = ""
            if variable_name and target_name and variable_name == target_name:
                radios.append(widget)
        try:
            stack.extend(widget.winfo_children())
        except Exception:
            pass
    return radios


def _find_active_watermark_range_parent(app):
    radios = _iter_watermark_range_radios(app)
    if not radios:
        return None
    parent_groups = []
    for radio in radios:
        parent = getattr(radio, "master", None)
        if parent is not None and parent not in parent_groups:
            parent_groups.append(parent)
    if not parent_groups:
        return None
    mapped_groups = []
    for parent in parent_groups:
        try:
            if parent.winfo_ismapped():
                mapped_groups.append(parent)
        except Exception:
            pass
    return (mapped_groups or parent_groups)[-1]


def _install_watermark_range_options(app):
    if getattr(app, "_fx_wm_range_options_ready", False):
        return
    if getattr(app, "wm_range_var", None) is None:
        return
    try:
        _set_watermark_range_var(app, _safe_var_get(app, "wm_range_var", "all"))
        parent = _find_active_watermark_range_parent(app)
        if parent is not None:
            existing = [
                child
                for child in parent.winfo_children()
                if isinstance(child, customtkinter.CTkRadioButton)
                and str(child.cget("value")) == "first_random"
            ]
            if not existing:
                style = {}
                try:
                    style = app._get_radio_style()
                except Exception:
                    style = {}
                radio = customtkinter.CTkRadioButton(
                    parent,
                    text=WATERMARK_RANGE_UI_LABELS["first_random"],
                    variable=app.wm_range_var,
                    value="first_random",
                    command=lambda target=app: _select_first_random_watermark_range(target),
                    **style,
                )
                radio.pack(anchor="w", pady=(0, 2))
                app._fx_wm_range_random_radio = radio
            else:
                app._fx_wm_range_random_radio = existing[-1]
                try:
                    app._fx_wm_range_random_radio.configure(command=lambda target=app: _select_first_random_watermark_range(target))
                except Exception:
                    pass
        app._fx_wm_range_options_ready = True
    except Exception as exc:
        _debug(f"watermark_range:install_error:{exc}")


def _patch_watermark_range_options_ui():
    try:
        original_init_watermark_ui = FengxiToolboxApp.init_watermark_ui
    except Exception as exc:
        _debug(f"patch_watermark_range_options:missing:{exc}")
        return
    if getattr(original_init_watermark_ui, "__fx_wm_range_options_patch__", False):
        return

    def patched_init_watermark_ui(self):
        original_init_watermark_ui(self)
        try:
            _install_watermark_range_options(self)
        except Exception as exc:
            _debug(f"watermark_range:init_error:{exc}")

    patched_init_watermark_ui.__fx_wm_range_options_patch__ = True
    FengxiToolboxApp.init_watermark_ui = patched_init_watermark_ui
    _debug("patch_watermark_range_options:installed")


_patch_watermark_range_options_ui()


def _install_watermark_checkpoint_action_ui(app):
    if getattr(app, "_fx_wm_checkpoint_action_ui_ready", False):
        return getattr(app, "wm_clear_checkpoint_button", None)
    stop_button = getattr(app, "btn_stop", None)
    parent = getattr(stop_button, "master", None)
    if stop_button is None or parent is None:
        return None

    button = customtkinter.CTkButton(
        parent,
        text="清除断点",
        command=lambda target=app: _clear_watermark_checkpoint_for_app(target),
        width=128,
        height=40,
        corner_radius=8,
        fg_color="transparent",
        border_width=1,
        border_color=globals().get("COLOR_BORDER", "#3A3A3A"),
        text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
        hover_color=globals().get("COLOR_CARD_ALT", "#303030"),
        font=customtkinter.CTkFont(size=12, weight="bold"),
    )

    manager = ""
    try:
        manager = str(stop_button.winfo_manager() or "")
    except Exception:
        pass
    if manager == "grid":
        try:
            stop_info = stop_button.grid_info()
            stop_column = int(stop_info.get("column", 0))
            stop_row = int(stop_info.get("row", 0))
            for sibling in parent.winfo_children():
                if sibling is stop_button:
                    continue
                try:
                    sibling_info = sibling.grid_info()
                    sibling_column = int(sibling_info.get("column", -1))
                    if sibling_column > stop_column:
                        sibling.grid_configure(column=sibling_column + 1)
                except Exception:
                    pass
            button.grid(
                row=stop_row,
                column=stop_column + 1,
                padx=(8, 8),
                pady=0,
                sticky="ew",
            )
        except Exception:
            button.pack(side="left", padx=(8, 0))
    elif manager == "pack":
        try:
            stop_pack_info = stop_button.pack_info()
            button.pack(
                side=stop_pack_info.get("side", "left"),
                padx=(8, 0),
                pady=stop_pack_info.get("pady", 0),
                after=stop_button,
            )
        except Exception:
            button.pack(side="left", padx=(8, 0))
    else:
        button.pack(side="left", padx=(8, 0))

    try:
        button.configure(state="disabled" if getattr(app, "is_running", False) else "normal")
    except Exception:
        pass
    app.wm_clear_checkpoint_button = button
    app._fx_wm_checkpoint_action_ui_ready = True
    return button


def _set_watermark_checkpoint_action_state(app, state=None):
    button = getattr(app, "wm_clear_checkpoint_button", None)
    if button is None:
        return
    if state is None:
        state = "disabled" if getattr(app, "is_running", False) else "normal"
    try:
        button.configure(state=state)
    except Exception:
        pass


def _patch_watermark_checkpoint_action_ui():
    try:
        original_init_watermark_ui = FengxiToolboxApp.init_watermark_ui
    except Exception as exc:
        _debug(f"watermark_checkpoint_ui:missing:{exc}")
        return
    if getattr(original_init_watermark_ui, "__fx_wm_checkpoint_action_patch__", False):
        return

    def patched_init_watermark_ui(self):
        original_init_watermark_ui(self)
        try:
            _install_watermark_checkpoint_action_ui(self)
        except Exception as exc:
            _debug(f"watermark_checkpoint_ui:init_error:{exc}")

    patched_init_watermark_ui.__fx_wm_checkpoint_action_patch__ = True
    FengxiToolboxApp.init_watermark_ui = patched_init_watermark_ui
    _debug("patch_watermark_checkpoint_ui:installed")

    try:
        original_setup_main_area = FengxiToolboxApp.setup_main_area
    except Exception as exc:
        _debug(f"watermark_checkpoint_ui:main_area_missing:{exc}")
        return
    if getattr(original_setup_main_area, "__fx_wm_checkpoint_action_patch__", False):
        return

    def patched_setup_main_area(self, *args, **kwargs):
        result = original_setup_main_area(self, *args, **kwargs)
        try:
            _install_watermark_checkpoint_action_ui(self)
        except Exception as exc:
            _debug(f"watermark_checkpoint_ui:main_area_error:{exc}")
        return result

    patched_setup_main_area.__fx_wm_checkpoint_action_patch__ = True
    FengxiToolboxApp.setup_main_area = patched_setup_main_area
    _debug("patch_watermark_checkpoint_ui:main_area_installed")


_patch_watermark_checkpoint_action_ui()


def _install_watermark_copy_guard_ui(app):
    if getattr(app, "_fx_wm_copy_guard_ui_ready", False):
        return getattr(app, "_fx_wm_copy_guard_frame", None)
    text_widget = getattr(app, "wm_text", None)
    parent = getattr(text_widget, "master", None)
    if text_widget is None or parent is None:
        return None

    if getattr(app, "wm_copy_guard_enabled_var", None) is None:
        app.wm_copy_guard_enabled_var = tkinter.BooleanVar(value=False)
    if getattr(app, "wm_copy_guard_strength_var", None) is None:
        app.wm_copy_guard_strength_var = tkinter.StringVar(value="标准")

    frame = customtkinter.CTkFrame(parent, fg_color="transparent", height=104)
    frame._fx_wm_copy_guard_module = True
    try:
        frame.pack_propagate(False)
    except Exception:
        pass

    top_row = customtkinter.CTkFrame(frame, fg_color="transparent", height=34)
    top_row.pack(fill="x")
    top_row.pack_propagate(False)

    strength_control = customtkinter.CTkSegmentedButton(
        top_row,
        values=[WATERMARK_COPY_GUARD_STRENGTH_LABELS[key] for key in WATERMARK_COPY_GUARD_STRENGTH_VALUES],
        variable=app.wm_copy_guard_strength_var,
        height=30,
        font=customtkinter.CTkFont(size=11),
    )

    def refresh_state(*_args):
        enabled = bool(_safe_var_get(app, "wm_copy_guard_enabled_var", False))
        try:
            strength_control.configure(state="normal" if enabled else "disabled")
        except Exception:
            pass
        try:
            refresh_hint()
        except Exception:
            pass
        try:
            _schedule_watermark_last_settings_persistence(app)
        except Exception:
            pass

    switch = customtkinter.CTkSwitch(
        top_row,
        text="复制干扰层（PDF / Word）",
        variable=app.wm_copy_guard_enabled_var,
        command=refresh_state,
        height=30,
        font=customtkinter.CTkFont(size=12, weight="bold"),
    )
    switch.pack(side="left", fill="x", expand=True)
    strength_control.pack(side="right", padx=(12, 0))

    hint = customtkinter.CTkLabel(
        frame,
        text="",
        anchor="w",
        justify="left",
        wraplength=430,
        font=customtkinter.CTkFont(size=11),
        text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
    )
    hint.pack(fill="x", pady=(4, 0))

    def refresh_hint(*_args):
        strength = _get_watermark_copy_guard_strength(app)
        strength_hint = WATERMARK_COPY_GUARD_STRENGTH_HINTS.get(
            strength,
            WATERMARK_COPY_GUARD_STRENGTH_HINTS["standard"],
        )
        hint.configure(
            text=(
                "PDF 插入不可见文本块，Word 在段间插入 1pt 白色乱码段落；"
                "干扰块含字母、数字、符号和乱码。单段复制正常，跨段或整份复制会混入干扰。"
                f"{strength_hint}"
            )
        )

    try:
        frame.pack(before=text_widget, fill="x", padx=24, pady=(0, 8))
    except Exception:
        frame.pack(fill="x", padx=24, pady=(0, 8))
    app.wm_copy_guard_switch = switch
    app.wm_copy_guard_strength_control = strength_control
    app.wm_copy_guard_hint_label = hint
    app._fx_wm_copy_guard_frame = frame
    app._fx_wm_copy_guard_ui_ready = True
    try:
        app._fx_wm_copy_guard_enabled_trace_id = app.wm_copy_guard_enabled_var.trace_add(
            "write",
            refresh_state,
        )
    except Exception:
        app._fx_wm_copy_guard_enabled_trace_id = None
    try:
        app._fx_wm_copy_guard_strength_trace_id = app.wm_copy_guard_strength_var.trace_add(
            "write",
            refresh_state,
        )
    except Exception:
        app._fx_wm_copy_guard_strength_trace_id = None
    refresh_state()
    return frame


def _patch_watermark_copy_guard_ui():
    try:
        original_init_watermark_ui = FengxiToolboxApp.init_watermark_ui
    except Exception as exc:
        _debug(f"watermark_copy_guard:patch_missing:{exc}")
        return
    if getattr(original_init_watermark_ui, "__fx_wm_copy_guard_patch__", False):
        return

    def patched_init_watermark_ui(self):
        original_init_watermark_ui(self)
        try:
            _install_watermark_copy_guard_ui(self)
        except Exception as exc:
            _debug(f"watermark_copy_guard:init_error:{exc}")

    patched_init_watermark_ui.__fx_wm_copy_guard_patch__ = True
    FengxiToolboxApp.init_watermark_ui = patched_init_watermark_ui
    _debug("patch_watermark_copy_guard:installed")


_patch_watermark_copy_guard_ui()


def _get_watermark_settings(app):
    text = _read_watermark_text_widget(app) or str(_safe_named_widget_get(app, "wm_text", "") or "")
    font_name = str(_safe_var_get(app, "selected_font", "") or "").strip()
    if not font_name:
        try:
            font_name = str((getattr(app, "font_list", []) or [""])[0] or "")
        except Exception:
            font_name = ""
    font_size = _safe_float(_safe_named_widget_get(app, "slider_size", _safe_var_get(app, "wm_size", 60)), 60.0)
    opacity = _safe_float(_safe_named_widget_get(app, "slider_opacity", _safe_var_get(app, "wm_opacity", 0.08)), 0.08)
    angle = _safe_float(_safe_named_widget_get(app, "slider_angle", _safe_var_get(app, "wm_angle", 45)), 45.0)
    page_range = _watermark_core_normalize_page_range(_safe_var_get(app, "wm_range_var", "all"))
    overwrite_mode = str(_safe_var_get(app, "wm_overwrite_var", "smart") or "smart").strip().lower()
    color = _get_watermark_color(app)
    return {
        "text": text,
        "font_name": font_name,
        "font_size": font_size,
        "opacity": opacity,
        "angle": angle,
        "color": color,
        "page_range": page_range,
        "force_mode": overwrite_mode == "force",
        "convert_pdf": bool(_safe_var_get(app, "wm_convert_pdf", False)),
        "delete_source": bool(_safe_var_get(app, "wm_delete_var", False)),
        "copy_guard_enabled": bool(_safe_var_get(app, "wm_copy_guard_enabled_var", False)),
        "copy_guard_strength": _get_watermark_copy_guard_strength(app),
    }


def _watermark_planned_output_path(src, input_root, output_root, actual_strategy, settings, is_single_input=False):
    src_path = Path(src)
    suffix = src_path.suffix.lower()
    convert_to_pdf = bool(settings.get("convert_pdf")) and suffix in WATERMARK_WORD_EXTS | WATERMARK_PPT_EXTS
    return Path(
        _build_watermark_output_path(
            src_path,
            input_root,
            output_root,
            actual_strategy,
            convert_to_pdf=convert_to_pdf,
            single_input=is_single_input,
        )
    )


def _watermark_checkpoint_context(input_value, input_root, output_root, actual_strategy, settings, total):
    # In-place overwrite changes the source signature, so output-based resume is
    # intentionally disabled for that mode.
    if actual_strategy == "overwrite":
        return None
    try:
        checkpoint_path = _watermark_checkpoint_path(input_value, output_root)
        identity = _watermark_checkpoint_identity(input_value, actual_strategy, settings)
        checkpoint = _watermark_checkpoint_prepare(
            checkpoint_path,
            input_value,
            actual_strategy,
            identity,
            total,
        )
        return {"path": checkpoint_path, "data": checkpoint}
    except Exception as exc:
        _debug(f"watermark_checkpoint:prepare_error:{exc}")
        return None


def _watermark_checkpoint_split(checkpoint, all_files, input_root, output_root, actual_strategy, settings, is_single_input):
    if not checkpoint:
        return list(all_files), []

    def output_getter(src):
        return _watermark_planned_output_path(
            src,
            input_root,
            output_root,
            actual_strategy,
            settings,
            is_single_input=is_single_input,
        )

    return _watermark_checkpoint_split_files(all_files, checkpoint.get("data"), output_getter)


def _watermark_checkpoint_mark(checkpoint, source, output, status):
    if not checkpoint:
        return
    try:
        _watermark_checkpoint_mark_item(checkpoint["path"], source, output, status)
    except Exception as exc:
        _debug(f"watermark_checkpoint:item_error:{exc}")


def _watermark_checkpoint_finish(checkpoint, state, processed, total, message=""):
    if not checkpoint:
        return
    try:
        _watermark_checkpoint_mark_state(checkpoint["path"], state, processed, total, message)
        if state == "completed":
            _watermark_checkpoint_clear(checkpoint["path"])
    except Exception as exc:
        _debug(f"watermark_checkpoint:state_error:{exc}")


def _clear_watermark_checkpoint_for_app(app):
    if getattr(app, "is_running", False):
        try:
            app.log("[断点续传] 当前任务仍在运行，请先按停止，再清除断点。")
        except Exception:
            pass
        return False
    try:
        input_value = _normalize_input_path_value(app.input_path.get())
    except Exception:
        input_value = ""
    if not input_value:
        try:
            app.log("[断点续传] 尚未选择输入文件或文件夹。")
        except Exception:
            pass
        return False
    is_single_input = os.path.isfile(input_value)
    input_root = os.path.dirname(input_value) if is_single_input else input_value
    requested_strategy = _get_task_output_strategy(app, "watermark")
    actual_strategy = requested_strategy if is_single_input else "result_folder"
    output_root = input_root if actual_strategy in {"same_dir", "overwrite"} else os.path.join(input_root, RESULT_FOLDER_NAME)
    checkpoint_path = _watermark_checkpoint_path(input_value, output_root)
    cleared = _watermark_checkpoint_clear(checkpoint_path)
    try:
        app.log(f"[断点续传] {'已清除' if cleared else '没有可清除的'}断点：{checkpoint_path}")
    except Exception:
        pass
    return cleared


def _watermark_make_pdf_packet(settings):
    packet_bytes = settings.get("_pdf_packet_bytes") if isinstance(settings, dict) else None
    if packet_bytes:
        return io.BytesIO(packet_bytes)
    return create_watermark_packet(
        settings["text"],
        settings["font_name"],
        settings["font_size"],
        settings["opacity"],
        settings["angle"],
        color=settings.get("color", WATERMARK_DEFAULT_COLOR),
    )


def _watermark_process_pdf(src, dst, settings):
    packet = _watermark_make_pdf_packet(settings)
    return add_watermark_to_pdf(
        str(src),
        str(dst),
        packet,
        page_range=settings["page_range"],
        check_text=settings["text"],
        force_mode=settings["force_mode"],
        copy_guard=settings.get("copy_guard_enabled", False),
        copy_guard_strength=settings.get("copy_guard_strength", "standard"),
    )


def _watermark_process_word(src, dst, settings, word_app):
    return add_watermark_to_word(
        word_app,
        str(src),
        str(dst),
        settings["text"],
        settings["font_name"],
        settings["font_size"],
        settings["opacity"],
        settings["angle"],
        page_range=settings["page_range"],
        force_mode=settings["force_mode"],
        color=settings.get("color", WATERMARK_DEFAULT_COLOR),
        copy_guard=settings.get("copy_guard_enabled", False),
        copy_guard_strength=settings.get("copy_guard_strength", "standard"),
    )


def _watermark_convert_doc_to_pdf(src, raw_pdf, word_app):
    return _convert_doc_to_pdf_safely(word_app, str(src), str(raw_pdf))


def _watermark_convert_ppt_to_pdf(src, raw_pdf, ppt_app):
    return _convert_ppt_to_pdf_safely(ppt_app, str(src), str(raw_pdf))


def _watermark_format_path_error(path_value, exc):
    message = str(exc)
    try:
        parts = Path(path_value).parts
        if any(part != part.rstrip(" .") for part in parts):
            message += " | path contains a folder/file name ending with space or dot; Windows may reject this path"
    except Exception:
        pass
    return message


def _watermark_update_progress(app, current_file="", stage="添加水印", completed=0, total=1, fraction=None):
    progress_fraction = _clamp_progress_value(
        completed / max(1, total) if fraction is None else fraction
    )
    progress_bar = getattr(app, "progress_bar", None)
    if progress_bar is not None:
        try:
            progress_bar.set(progress_fraction)
        except Exception:
            pass
    return _set_progress_status(
        app,
        current_file=current_file,
        stage=stage,
        completed=completed,
        total=total,
        fraction=progress_fraction,
    )


def _copy_watermark_skipped_files(app, skipped_files, input_root, output_root, actual_strategy, logs, result):
    if not skipped_files or not _get_watermark_copy_skipped_to_output(app):
        return [], []

    copy_root = Path(output_root)
    if actual_strategy in {"same_dir", "overwrite"}:
        copy_root = Path(input_root) / RESULT_FOLDER_NAME
    try:
        copy_root.mkdir(parents=True, exist_ok=True)
    except Exception as exc:
        failed = [str(Path(_normalize_input_path_value(item))) for item in skipped_files]
        _watermark_log(app, logs, f"[批量水印] 复制跳过文件失败: {copy_root} | {_watermark_format_path_error(copy_root, exc)}")
        return [], failed

    copied = []
    failed = []
    for item in skipped_files:
        src = Path(_normalize_input_path_value(item))
        try:
            if not src.exists() or not src.is_file():
                raise FileNotFoundError(str(src))
            try:
                rel = src.relative_to(Path(input_root))
            except Exception:
                rel = Path(src.name)
            rel = Path(*[str(part).strip().rstrip(" .") or _sanitize_filename_component(part, fallback="folder") for part in rel.parts])
            dst = copy_root / rel
            if dst.resolve() == src.resolve():
                _watermark_log(app, logs, f"[批量水印] 跳过文件已在输出位置，无需复制: {src.name}")
                copied.append(str(dst))
                _add_task_result_output(result, str(dst))
                continue
            dst.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(str(src), str(dst))
            copied.append(str(dst))
            _add_task_result_output(result, str(dst))
            _watermark_log(app, logs, f"[批量水印] 已复制跳过文件: {src.name} -> {dst}")
        except Exception as exc:
            failed.append(str(src))
            _watermark_log(app, logs, f"[批量水印] 复制跳过文件失败: {src.name} | {exc}")
    return copied, failed


def _normalize_watermark_file_key(path_value):
    normalized = _normalize_input_path_value(path_value)
    if not normalized:
        return ""
    try:
        return os.path.normcase(str(Path(normalized).resolve()))
    except Exception:
        try:
            return os.path.normcase(os.path.abspath(normalized))
        except Exception:
            return os.path.normcase(str(normalized))


def _collect_watermark_input_files(input_value):
    normalized_input = _normalize_input_path_value(input_value)
    if not normalized_input:
        return []
    input_path = Path(normalized_input)
    if input_path.is_file():
        return [str(input_path)]
    if not input_path.is_dir():
        return []

    collected = []
    result_folder_name = str(RESULT_FOLDER_NAME).strip()
    for root, dirnames, filenames in os.walk(str(input_path)):
        current_root = Path(root)
        if current_root == input_path:
            dirnames[:] = [name for name in dirnames if str(name).strip() != result_folder_name]
        for filename in filenames:
            collected.append(str(Path(root) / filename))
    return collected


def _watermark_pdf_parallel_enabled(app, all_files, settings, is_single_input):
    if is_single_input or not _is_parallel_enabled(app):
        return False
    if len(all_files) <= 1:
        return False
    if settings.get("convert_pdf"):
        return False
    for file_path in all_files:
        try:
            if Path(_normalize_input_path_value(file_path)).suffix.lower() not in WATERMARK_PDF_EXTS:
                return False
        except Exception:
            return False
    return _get_watermark_pdf_worker_count(len(all_files)) > 1


def _write_watermark_failed_report(failed_items, input_root, output_root, actual_strategy, result):
    if not failed_items:
        return
    report_root = Path(output_root if actual_strategy == "result_folder" else input_root)
    try:
        report_root.mkdir(parents=True, exist_ok=True)
    except Exception:
        report_root = Path(input_root)
        report_root.mkdir(parents=True, exist_ok=True)
    report_path = report_root / "!失败文件清单.txt"
    lines = ["以下文件处理失败："]
    for item in failed_items:
        try:
            lines.append(os.path.relpath(item, input_root))
        except Exception:
            lines.append(str(item))
    report_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    _add_task_result_output(result, str(report_path))


def _run_watermark_pdf_parallel_task_impl(
    app,
    all_files,
    input_root,
    output_root,
    actual_strategy,
    settings,
    logs,
    result,
    overall_total,
    rule_skipped_files,
    type_skipped_files,
    unsupported_skipped_files,
    checkpoint=None,
    resumed_entries=None,
):
    pdf_settings = dict(settings)
    try:
        packet = _watermark_make_pdf_packet(settings)
        packet.seek(0)
        pdf_settings["_pdf_packet_bytes"] = packet.getvalue()
    except Exception:
        pdf_settings.pop("_pdf_packet_bytes", None)

    workers = _get_watermark_pdf_worker_count(len(all_files))
    _watermark_log(app, logs, f"[批量并行] PDF 水印启用 {workers} 个线程。")

    resumed_entries = list(resumed_entries or [])
    success_outputs = [str(item.get("output")) for item in resumed_entries if item.get("output")]
    failed_items = []
    skipped_items = [str(item.get("source")) for item in resumed_entries if item.get("source")]
    copied_skipped_outputs = []
    deferred_skipped_copy_items = []
    processed_count = len(resumed_entries)
    jobs = []

    if resumed_entries:
        _watermark_log(app, logs, f"[断点续传] 已确认 {len(resumed_entries)} 个已完成文件，跳过逐文件处理。")
        for item in resumed_entries:
            _add_task_result_output(result, str(item.get("output") or ""))

    for file_path in all_files:
        if getattr(app, "stop_event", False):
            break
        src = Path(_normalize_input_path_value(file_path))
        resolved_src = _resolve_watermark_source_path(src, input_root)
        if resolved_src != src:
            _watermark_log(app, logs, f"[批量水印] 已修正源文件路径: {src} -> {resolved_src}")
            src = resolved_src
        if not src.exists():
            failed_items.append(str(src))
            processed_count += 1
            _watermark_log(app, logs, f"[批量水印] 失败: {src.name} | 源文件不存在")
            _watermark_update_progress(app, current_file=str(src), stage="添加水印", completed=processed_count, total=max(overall_total, 1))
            continue
        target_path = Path(
            _build_watermark_output_path(
                src,
                input_root,
                output_root,
                actual_strategy,
                convert_to_pdf=False,
                single_input=False,
            )
        )
        try:
            target_path.parent.mkdir(parents=True, exist_ok=True)
        except Exception as exc:
            failed_items.append(str(src))
            processed_count += 1
            _watermark_log(app, logs, f"[批量水印] 失败: {src.name} | 输出路径准备失败: {_watermark_format_path_error(target_path.parent, exc)}")
            _watermark_update_progress(app, current_file=str(src), stage="添加水印", completed=processed_count, total=max(overall_total, 1))
            continue
        if is_nonempty_file(target_path):
            success_outputs.append(str(target_path))
            skipped_items.append(str(src))
            _watermark_checkpoint_mark(checkpoint, src, target_path, "skipped")
            _add_task_result_output(result, str(target_path))
            _watermark_log(app, logs, f"[批量水印] 断点续跑：已存在结果，跳过处理: {src.name} -> {target_path}")
            processed_count += 1
            _watermark_update_progress(app, current_file=str(src), stage="断点续跑", completed=processed_count, total=max(overall_total, 1))
            continue
        jobs.append((src, target_path))

    def process_pdf_job(job):
        src, output_path = job
        status = _watermark_process_pdf(src, output_path, pdf_settings)
        kind = _watermark_status_kind(status)
        item = {
            "src": str(src),
            "name": src.name,
            "output": str(output_path),
            "status": status,
            "kind": kind,
            "copied_output": "",
        }
        if kind == "skipped" and _watermark_should_preserve_original(status) and output_path.resolve() != src.resolve():
            try:
                output_path.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(str(src), str(output_path))
                item["copied_output"] = str(output_path)
            except Exception as exc:
                item["kind"] = "failed"
                item["status"] = f"ERROR:preserve original failed: {exc}"
        return item

    with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as executor:
        # Keep only a small number of PDF rewrites queued.  Large folders used
        # to create one Future per file, increasing memory use and making Stop
        # wait for the entire already-submitted batch.
        job_iterator = iter(jobs)
        future_map = {}

        def submit_next_job():
            try:
                job = next(job_iterator)
            except StopIteration:
                return False
            future_map[executor.submit(process_pdf_job, job)] = job
            return True

        for _ in range(min(len(jobs), workers * WATERMARK_PDF_QUEUE_FACTOR)):
            if not submit_next_job():
                break

        while future_map:
            done, _ = concurrent.futures.wait(
                future_map,
                return_when=concurrent.futures.FIRST_COMPLETED,
            )
            for future in done:
                src, output_path = future_map.pop(future)
                try:
                    item = future.result()
                except Exception as exc:
                    item = {
                        "src": str(src),
                        "name": src.name,
                        "output": str(output_path),
                        "status": f"ERROR:{exc}",
                        "kind": "failed",
                        "copied_output": "",
                    }
                kind = item.get("kind")
                src_text = item.get("src") or str(src)
                output_text = item.get("output") or str(output_path)
                if kind == "success":
                    success_outputs.append(output_text)
                    _watermark_checkpoint_mark(checkpoint, src_text, output_text, "success")
                    _add_task_result_output(result, output_text)
                    _watermark_log(app, logs, f"[批量水印] 成功: {item.get('name') or Path(src_text).name} -> {output_text}")
                    if settings.get("delete_source"):
                        try:
                            Path(src_text).unlink()
                        except Exception as delete_exc:
                            _watermark_log(app, logs, f"[批量水印] 删除源文件失败: {Path(src_text).name} | {delete_exc}")
                elif kind == "skipped":
                    _watermark_checkpoint_mark(checkpoint, src_text, output_text, "skipped")
                    copied_output = item.get("copied_output") or ""
                    if copied_output:
                        copied_skipped_outputs.append(copied_output)
                        _add_task_result_output(result, copied_output)
                        _watermark_log(app, logs, f"[批量水印] 受保护文件无法加水印，已保留原文件: {Path(src_text).name} -> {copied_output}")
                    elif _get_watermark_copy_skipped_to_output(app):
                        deferred_skipped_copy_items.append(src_text)
                    skipped_items.append(src_text)
                    _watermark_log(app, logs, f"[批量水印] 跳过: {Path(src_text).name} | {item.get('status')}")
                else:
                    failed_items.append(src_text)
                    _watermark_log(app, logs, f"[批量水印] 失败: {Path(src_text).name} | {item.get('status') or 'unknown'}")
                processed_count += 1
                _watermark_update_progress(app, current_file=src_text, stage="添加水印", completed=processed_count, total=max(overall_total, 1))
                if not getattr(app, "stop_event", False):
                    submit_next_job()

    if rule_skipped_files:
        copied_skipped, failed_copies = _copy_watermark_skipped_files(app, rule_skipped_files, input_root, output_root, actual_strategy, logs, result)
        copied_skipped_outputs.extend(copied_skipped)
        skipped_items.extend(str(Path(_normalize_input_path_value(item))) for item in rule_skipped_files)
        failed_items.extend(failed_copies)
    if type_skipped_files:
        copied_skipped, failed_copies = _copy_watermark_skipped_files(app, type_skipped_files, input_root, output_root, actual_strategy, logs, result)
        copied_skipped_outputs.extend(copied_skipped)
        skipped_items.extend(str(Path(_normalize_input_path_value(item))) for item in type_skipped_files)
        failed_items.extend(failed_copies)
    if deferred_skipped_copy_items:
        copied_skipped, failed_copies = _copy_watermark_skipped_files(app, deferred_skipped_copy_items, input_root, output_root, actual_strategy, logs, result)
        copied_skipped_outputs.extend(copied_skipped)
        failed_items.extend(failed_copies)
    if unsupported_skipped_files:
        copied_skipped, failed_copies = _copy_watermark_skipped_files(app, unsupported_skipped_files, input_root, output_root, actual_strategy, logs, result)
        copied_skipped_outputs.extend(copied_skipped)
        skipped_items.extend(str(Path(_normalize_input_path_value(item))) for item in unsupported_skipped_files)
        failed_items.extend(failed_copies)

    _write_watermark_failed_report(failed_items, input_root, output_root, actual_strategy, result)
    result["failed_items"] = list(failed_items)
    _set_task_result_counts(
        result,
        processed=processed_count + len(rule_skipped_files) + len(type_skipped_files) + len(unsupported_skipped_files),
        success=len(success_outputs),
        failed=len(failed_items),
        skipped=len(skipped_items),
    )
    if getattr(app, "stop_event", False):
        _set_task_result_finished(result, "stopped", message="用户停止批量水印任务", detail="用户停止批量水印任务", stopped=True)
    elif failed_items:
        _set_task_result_finished(result, "failed", message=f"批量水印失败 {len(failed_items)} 个文件", detail=f"批量水印失败 {len(failed_items)} 个文件", error=f"failed_items={len(failed_items)}")
    elif success_outputs or copied_skipped_outputs:
        _set_task_result_finished(result, "success", message="批量水印完成", detail="批量水印完成")
    else:
        _set_task_result_finished(result, "skipped", message="没有生成新的水印文件", detail="没有生成新的水印文件", skipped=True)
    _flush_watermark_ui_logs(app, force=True)
    return result


def _run_watermark_pdf_parallel_task(
    app,
    all_files,
    input_root,
    output_root,
    actual_strategy,
    settings,
    logs,
    result,
    overall_total,
    rule_skipped_files,
    type_skipped_files,
    unsupported_skipped_files,
    checkpoint=None,
    resumed_entries=None,
):
    try:
        task_result = _run_watermark_pdf_parallel_task_impl(
            app,
            all_files,
            input_root,
            output_root,
            actual_strategy,
            settings,
            logs,
            result,
            overall_total,
            rule_skipped_files,
            type_skipped_files,
            unsupported_skipped_files,
            checkpoint=checkpoint,
            resumed_entries=resumed_entries,
        )
    except Exception as exc:
        _watermark_checkpoint_finish(
            checkpoint,
            "interrupted",
            int(result.get("processed_count") or 0),
            overall_total,
            f"exception:{type(exc).__name__}",
        )
        raise

    if getattr(app, "stop_event", False):
        checkpoint_state = "paused"
        checkpoint_message = "user_stop"
    elif result.get("failed_items"):
        checkpoint_state = "failed"
        checkpoint_message = "failed_items"
    else:
        checkpoint_state = "completed"
        checkpoint_message = "completed"
    _watermark_checkpoint_finish(
        checkpoint,
        checkpoint_state,
        int(result.get("processed_count") or 0),
        overall_total,
        checkpoint_message,
    )
    return task_result


def _run_watermark_task(app, input_value):
    normalized_input = _normalize_input_path_value(input_value)
    is_single_input = bool(normalized_input and os.path.isfile(normalized_input))
    input_root = os.path.dirname(normalized_input) if is_single_input else normalized_input
    requested_strategy = _get_task_output_strategy(app, "watermark")
    actual_strategy = requested_strategy if is_single_input else "result_folder"
    output_root = input_root if actual_strategy in {"same_dir", "overwrite"} else os.path.join(input_root, RESULT_FOLDER_NAME)
    settings = _get_watermark_settings(app)
    logs = []
    _flush_watermark_ui_logs(app, force=True)
    app._fx_watermark_ui_log_pending = []

    result = _get_last_task_result(app)
    if result is None:
        result = _start_task_result(app, normalized_input, "watermark")
    _set_task_result_output_strategy(result, "watermark", actual_strategy)
    result["output_strategy"] = actual_strategy
    result["output_strategy_label"] = _get_output_strategy_label(actual_strategy)
    _set_task_result_output_root(result, output_root)
    try:
        app._fx_last_task_logs = logs
    except Exception:
        pass

    filename_rule = _get_watermark_filename_rule(app)
    previous_runtime_rule = getattr(app, "_fx_wm_filename_rule_runtime", None)
    previous_skipped_files = getattr(app, "_fx_wm_filename_rule_skipped_files", None)
    skip_var = getattr(app, "wm_skip_hyphen_var", None)
    previous_skip_value = None
    restore_skip_value = False
    previous_rule_loading = getattr(app, "_fx_wm_filename_rule_loading", False)
    try:
        app._fx_wm_filename_rule_runtime = filename_rule
        app._fx_wm_filename_rule_skipped_files = []
        if filename_rule and skip_var is not None:
            try:
                app._fx_wm_filename_rule_loading = True
                previous_skip_value = skip_var.get()
                skip_var.set(False)
                restore_skip_value = True
            except Exception as exc:
                _debug(f"watermark_task:disable_builtin_skip_error:{exc}")
        all_files = list(app.collect_input_files(normalized_input, "watermark"))
        rule_skipped_files = list(getattr(app, "_fx_wm_filename_rule_skipped_files", []) or [])
    finally:
        app._fx_wm_filename_rule_runtime = previous_runtime_rule
        if previous_skipped_files is not None and not getattr(app, "_fx_wm_filename_rule_skipped_files", None):
            app._fx_wm_filename_rule_skipped_files = previous_skipped_files
        if restore_skip_value and skip_var is not None:
            try:
                skip_var.set(previous_skip_value)
            except Exception:
                pass
        app._fx_wm_filename_rule_loading = previous_rule_loading
    all_files, type_skipped_files = _filter_watermark_files_by_type_skip(app, all_files)
    all_files, unsupported_skipped_files = _split_watermark_supported_files(all_files)
    known_file_keys = set()
    for item in list(all_files) + list(rule_skipped_files) + list(type_skipped_files) + list(unsupported_skipped_files):
        item_key = _normalize_watermark_file_key(item)
        if item_key:
            known_file_keys.add(item_key)
    # A second full directory walk is only necessary when skipped unsupported
    # files must be copied to the result folder.  Most watermark jobs do not
    # enable that option, so avoid scanning a large tree twice before work starts.
    if _get_watermark_copy_skipped_to_output(app):
        for item in _collect_watermark_input_files(normalized_input):
            item_key = _normalize_watermark_file_key(item)
            suffix = Path(_normalize_input_path_value(item)).suffix.lower()
            if item_key and item_key not in known_file_keys and suffix not in WATERMARK_SUPPORTED_EXTS:
                unsupported_skipped_files.append(str(Path(_normalize_input_path_value(item))))
    checkpoint = _watermark_checkpoint_context(
        normalized_input,
        input_root,
        output_root,
        actual_strategy,
        settings,
        len(all_files),
    )
    try:
        app._fx_watermark_checkpoint_active = checkpoint
    except Exception:
        pass
    previous_checkpoint_state = str((checkpoint or {}).get("data", {}).get("previous_state") or "")
    if previous_checkpoint_state in {"running", "interrupted"}:
        _watermark_log(
            app,
            logs,
            f"[断点续传] 上次任务未正常结束（{previous_checkpoint_state}），将从已记录的完成项继续。",
        )
    all_files, resumed_entries = _watermark_checkpoint_split(
        checkpoint,
        all_files,
        input_root,
        output_root,
        actual_strategy,
        settings,
        is_single_input,
    )
    resumed_count = len(resumed_entries)
    result["resumed_count"] = resumed_count
    total = len(all_files)
    overall_total = total + resumed_count + len(rule_skipped_files) + len(type_skipped_files) + len(unsupported_skipped_files)
    if total <= 0 and resumed_count <= 0 and not rule_skipped_files and not type_skipped_files and not unsupported_skipped_files:
        _watermark_log(app, logs, "[批量水印] 未找到可处理文件")
        _set_task_result_counts(result, processed=0, success=0, failed=0, skipped=1)
        _set_task_result_finished(result, "skipped", message="未找到可处理文件", detail="未找到可处理文件", skipped=True)
        _watermark_checkpoint_finish(checkpoint, "completed", 0, overall_total, "no_files")
        return result
    if actual_strategy == "result_folder" or ((rule_skipped_files or type_skipped_files or unsupported_skipped_files) and _get_watermark_copy_skipped_to_output(app)):
        os.makedirs(output_root, exist_ok=True)

    _watermark_log(app, logs, f"[批量水印] 将处理 {total} 个文件 | 输出策略：{_get_output_strategy_label(actual_strategy)}")
    if resumed_count:
        _watermark_log(app, logs, f"[断点续传] 已确认 {resumed_count} 个已完成文件，跳过逐文件处理。")
    if settings.get("copy_guard_enabled"):
        strength_label = WATERMARK_COPY_GUARD_STRENGTH_LABELS.get(
            settings.get("copy_guard_strength", "standard"),
            "标准",
        )
        _watermark_log(
            app,
            logs,
            f"[复制干扰层] 已开启（{strength_label}）| 应用于 PDF、直接输出的 Word，以及转成 PDF 的 Word/PPT；默认显示和打印不变。",
        )
    if rule_skipped_files:
        copy_hint = "会复制到输出文件夹" if _get_watermark_copy_skipped_to_output(app) else "仅跳过不复制"
        _watermark_log(app, logs, f"[批量水印] 文件名规则跳过 {len(rule_skipped_files)} 个文件 | {copy_hint}")
    if type_skipped_files:
        skipped_groups = sorted(_get_watermark_type_skip_groups(app))
        copy_hint = "会复制到输出文件夹" if _get_watermark_copy_skipped_to_output(app) else "仅跳过不复制"
        _watermark_log(app, logs, f"[批量水印] 按文件类型跳过 {len(type_skipped_files)} 个文件 | 类型：{', '.join(skipped_groups)} | {copy_hint}")
    if unsupported_skipped_files:
        copy_hint = "会复制到输出文件夹" if _get_watermark_copy_skipped_to_output(app) else "仅跳过不复制"
        _watermark_log(app, logs, f"[批量水印] 非水印可处理文件跳过 {len(unsupported_skipped_files)} 个文件 | {copy_hint}")
    if _watermark_pdf_parallel_enabled(app, all_files, settings, is_single_input):
        return _run_watermark_pdf_parallel_task(
            app,
            all_files,
            input_root,
            output_root,
            actual_strategy,
            settings,
            logs,
            result,
            overall_total,
            rule_skipped_files,
            type_skipped_files,
            unsupported_skipped_files,
            checkpoint=checkpoint,
            resumed_entries=resumed_entries,
        )

    success_outputs = [str(item.get("output")) for item in resumed_entries if item.get("output")]
    failed_items = []
    skipped_items = [str(item.get("source")) for item in resumed_entries if item.get("source")]
    copied_skipped_outputs = []
    deferred_skipped_copy_items = []
    processed_count = resumed_count
    for item in resumed_entries:
        _add_task_result_output(result, str(item.get("output") or ""))
    word_app = None
    ppt_app = None
    pythoncom_initialized = False
    parallel_executor = None
    parallel_future_map = {}
    parallel_pdf_jobs = []
    parallel_pdf_job_index = 0
    serial_files = list(all_files)

    def get_word_app():
        nonlocal word_app, pythoncom_initialized
        if word_app is None:
            if not pythoncom_initialized:
                pythoncom.CoInitialize()
                pythoncom_initialized = True
            word_app = _create_hidden_word_app()
        return word_app

    def get_ppt_app():
        nonlocal ppt_app, pythoncom_initialized
        if ppt_app is None:
            if not pythoncom_initialized:
                pythoncom.CoInitialize()
                pythoncom_initialized = True
            ppt_app = _safe_office_dispatch_ex("PowerPoint.Application")
            try:
                ppt_app.Visible = False
            except Exception:
                pass
        return ppt_app

    def process_parallel_pdf_job(job):
        src, output_path, pdf_settings = job
        status = _watermark_process_pdf(src, output_path, pdf_settings)
        kind = _watermark_status_kind(status)
        item = {
            "src": str(src),
            "name": src.name,
            "output": str(output_path),
            "status": status,
            "kind": kind,
            "copied_output": "",
        }
        if kind == "skipped" and _watermark_should_preserve_original(status) and output_path.resolve() != src.resolve():
            try:
                output_path.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(str(src), str(output_path))
                item["copied_output"] = str(output_path)
            except Exception as exc:
                item["kind"] = "failed"
                item["status"] = f"ERROR:preserve original failed: {exc}"
        return item

    def record_parallel_pdf_result(future):
        nonlocal processed_count
        src, output_path, _pdf_settings = parallel_future_map.pop(future)
        try:
            item = future.result()
        except Exception as exc:
            item = {
                "src": str(src),
                "name": src.name,
                "output": str(output_path),
                "status": f"ERROR:{exc}",
                "kind": "failed",
                "copied_output": "",
            }
        kind = item.get("kind")
        src_text = item.get("src") or str(src)
        output_text = item.get("output") or str(output_path)
        if kind == "success":
            success_outputs.append(output_text)
            _watermark_checkpoint_mark(checkpoint, src_text, output_text, "success")
            _add_task_result_output(result, output_text)
            _watermark_log(app, logs, f"[批量水印] 成功: {item.get('name') or Path(src_text).name} -> {output_text}")
            if settings.get("delete_source"):
                try:
                    Path(src_text).unlink()
                except Exception as delete_exc:
                    _watermark_log(app, logs, f"[批量水印] 删除源文件失败: {Path(src_text).name} | {delete_exc}")
        elif kind == "skipped":
            _watermark_checkpoint_mark(checkpoint, src_text, output_text, "skipped")
            copied_output = item.get("copied_output") or ""
            if copied_output:
                copied_skipped_outputs.append(copied_output)
                _add_task_result_output(result, copied_output)
                _watermark_log(app, logs, f"[批量水印] 受保护文件无法加水印，已保留原文件: {Path(src_text).name} -> {copied_output}")
            elif _get_watermark_copy_skipped_to_output(app):
                deferred_skipped_copy_items.append(src_text)
            skipped_items.append(src_text)
            _watermark_log(app, logs, f"[批量水印] 跳过: {Path(src_text).name} | {item.get('status')}")
        else:
            failed_items.append(src_text)
            _watermark_log(app, logs, f"[批量水印] 失败: {Path(src_text).name} | {item.get('status') or 'unknown'}")
        processed_count += 1
        _watermark_update_progress(app, current_file=src_text, stage="添加水印", completed=processed_count, total=max(overall_total, 1))

    pdf_files = []
    office_files = []
    for file_path in all_files:
        try:
            is_pdf = Path(_normalize_input_path_value(file_path)).suffix.lower() in WATERMARK_PDF_EXTS
        except Exception:
            is_pdf = False
        (pdf_files if is_pdf else office_files).append(file_path)
    use_mixed_pdf_parallel = (
        not is_single_input
        and _is_parallel_enabled(app)
        and not settings.get("convert_pdf")
        and len(pdf_files) > 1
        and bool(office_files)
        and _get_watermark_pdf_worker_count(len(pdf_files)) > 1
    )
    if use_mixed_pdf_parallel:
        pdf_settings = dict(settings)
        try:
            packet = _watermark_make_pdf_packet(settings)
            packet.seek(0)
            pdf_settings["_pdf_packet_bytes"] = packet.getvalue()
        except Exception:
            pdf_settings.pop("_pdf_packet_bytes", None)
        serial_files = office_files
        for file_path in pdf_files:
            if getattr(app, "stop_event", False):
                break
            src = Path(_normalize_input_path_value(file_path))
            src = _resolve_watermark_source_path(src, input_root)
            target_path = Path(
                _build_watermark_output_path(
                    src,
                    input_root,
                    output_root,
                    actual_strategy,
                    convert_to_pdf=False,
                    single_input=False,
                )
            )
            try:
                target_path.parent.mkdir(parents=True, exist_ok=True)
            except Exception as exc:
                failed_items.append(str(src))
                processed_count += 1
                _watermark_log(app, logs, f"[批量水印] 失败: {src.name} | 输出路径准备失败: {_watermark_format_path_error(target_path.parent, exc)}")
                continue
            if not src.exists():
                failed_items.append(str(src))
                processed_count += 1
                _watermark_log(app, logs, f"[批量水印] 失败: {src.name} | 源文件不存在")
                continue
            if is_nonempty_file(target_path):
                success_outputs.append(str(target_path))
                skipped_items.append(str(src))
                _watermark_checkpoint_mark(checkpoint, src, target_path, "skipped")
                _add_task_result_output(result, str(target_path))
                processed_count += 1
                _watermark_log(app, logs, f"[批量水印] 断点续跑：已存在结果，跳过处理: {src.name} -> {target_path}")
                continue
            parallel_pdf_jobs.append((src, target_path, pdf_settings))
        if parallel_pdf_jobs:
            workers = _get_watermark_pdf_worker_count(len(parallel_pdf_jobs))
            _watermark_log(app, logs, f"[批量并行] PDF 水印启用 {workers} 个线程，Word/PPT 保持单线程。")
            parallel_executor = concurrent.futures.ThreadPoolExecutor(max_workers=workers)
            for job in parallel_pdf_jobs[: workers * WATERMARK_PDF_QUEUE_FACTOR]:
                parallel_future_map[parallel_executor.submit(process_parallel_pdf_job, job)] = job
            parallel_pdf_job_index = min(len(parallel_pdf_jobs), workers * WATERMARK_PDF_QUEUE_FACTOR)

    try:
        for index, file_path in enumerate(serial_files, start=1):
            if getattr(app, "stop_event", False):
                break
            src = Path(_normalize_input_path_value(file_path))
            resolved_src = _resolve_watermark_source_path(src, input_root)
            if resolved_src != src:
                _watermark_log(app, logs, f"[批量水印] 已修正源文件路径: {src} -> {resolved_src}")
                src = resolved_src
            if not src.exists():
                failed_items.append(str(src))
                processed_count += 1
                _watermark_log(app, logs, f"[批量水印] 失败: {src.name} | 源文件不存在")
                _watermark_update_progress(
                    app,
                    current_file=str(src),
                    stage="添加水印",
                    completed=processed_count,
                    total=max(overall_total, 1),
                )
                continue
            suffix = src.suffix.lower()
            convert_to_pdf = settings["convert_pdf"] and suffix in WATERMARK_WORD_EXTS | WATERMARK_PPT_EXTS
            target_path = Path(
                _build_watermark_output_path(
                    src,
                    input_root,
                    output_root,
                    actual_strategy,
                    convert_to_pdf=convert_to_pdf,
                    single_input=is_single_input,
                )
            )
            try:
                target_path.parent.mkdir(parents=True, exist_ok=True)
            except Exception as exc:
                failed_items.append(str(src))
                processed_count += 1
                _watermark_log(app, logs, f"[批量水印] 失败: {src.name} | 输出路径准备失败: {_watermark_format_path_error(target_path.parent, exc)}")
                _watermark_update_progress(
                    app,
                    current_file=str(src),
                    stage="添加水印",
                    completed=processed_count,
                    total=max(overall_total, 1),
                )
                continue
            output_path = target_path
            stage_dir = None
            stage_path = target_path

            if actual_strategy == "overwrite" and target_path.resolve() == src.resolve():
                stage_dir = Path(tempfile.mkdtemp(prefix="fx_wm_stage_", dir=str(src.parent)))
                stage_path = stage_dir / src.name
            elif actual_strategy == "overwrite" and target_path.suffix.lower() != src.suffix.lower():
                actual_strategy = "same_dir"
                result["output_strategy"] = actual_strategy
                result["output_strategy_label"] = _get_output_strategy_label(actual_strategy)
                output_root = input_root
                _set_task_result_output_root(result, output_root)
                target_path = Path(
                    _build_watermark_output_path(
                        src,
                        input_root,
                        output_root,
                        "same_dir",
                        convert_to_pdf=convert_to_pdf,
                        single_input=True,
                    )
                )
                try:
                    target_path.parent.mkdir(parents=True, exist_ok=True)
                except Exception as exc:
                    failed_items.append(str(src))
                    processed_count += 1
                    _watermark_log(app, logs, f"[批量水印] 失败: {src.name} | 输出路径准备失败: {_watermark_format_path_error(target_path.parent, exc)}")
                    _watermark_update_progress(
                        app,
                        current_file=str(src),
                        stage="添加水印",
                        completed=processed_count,
                        total=max(overall_total, 1),
                    )
                    continue
                output_path = target_path
                stage_path = target_path
                _watermark_log(app, logs, f"[批量水印] {src.name} 转 PDF 后扩展名变化，已改为同目录新文件输出")

            _watermark_update_progress(
                app,
                current_file=str(src),
                stage="添加水印",
                completed=processed_count,
                total=max(overall_total, 1),
            )
            try:
                same_file_output = output_path.resolve() == src.resolve()
            except Exception:
                same_file_output = False
            if not same_file_output and is_nonempty_file(output_path):
                success_outputs.append(str(output_path))
                skipped_items.append(str(src))
                _watermark_checkpoint_mark(checkpoint, src, output_path, "skipped")
                _add_task_result_output(result, str(output_path))
                _watermark_log(app, logs, f"[批量水印] 断点续跑：已存在结果，跳过处理: {src.name} -> {output_path}")
                processed_count += 1
                _watermark_update_progress(
                    app,
                    current_file=str(src),
                    stage="断点续跑",
                    completed=processed_count,
                    total=max(overall_total, 1),
                )
                continue
            try:
                if suffix in WATERMARK_PDF_EXTS:
                    status = _watermark_process_pdf(src, stage_path, settings)
                elif suffix in WATERMARK_WORD_EXTS:
                    if convert_to_pdf:
                        raw_fd, raw_pdf_name = tempfile.mkstemp(
                            prefix="fx_wm_raw_",
                            suffix=".pdf",
                            dir=str(target_path.parent),
                        )
                        os.close(raw_fd)
                        raw_pdf = Path(raw_pdf_name)
                        try:
                            convert_status = _watermark_convert_doc_to_pdf(src, raw_pdf, get_word_app())
                            if str(convert_status).strip() != "SUCCESS" or not raw_pdf.exists():
                                status = f"ERROR:Word 转 PDF 失败: {convert_status}"
                            else:
                                status = _watermark_process_pdf(raw_pdf, stage_path, settings)
                        finally:
                            try:
                                if raw_pdf.exists():
                                    raw_pdf.unlink()
                            except Exception:
                                pass
                    else:
                        status = _watermark_process_word(src, stage_path, settings, get_word_app())
                elif suffix in WATERMARK_PPT_EXTS:
                    raw_fd, raw_pdf_name = tempfile.mkstemp(
                        prefix="fx_wm_raw_",
                        suffix=".pdf",
                        dir=str(target_path.parent),
                    )
                    os.close(raw_fd)
                    raw_pdf = Path(raw_pdf_name)
                    try:
                        convert_status = _watermark_convert_ppt_to_pdf(src, raw_pdf, get_ppt_app())
                        if str(convert_status).strip() != "SUCCESS" or not raw_pdf.exists():
                            status = f"ERROR:PPT 转 PDF 失败: {convert_status}"
                        else:
                            status = _watermark_process_pdf(raw_pdf, stage_path, settings)
                    finally:
                        try:
                            if raw_pdf.exists():
                                raw_pdf.unlink()
                        except Exception:
                            pass
                else:
                    status = f"SKIP:unsupported file type: {src.suffix}"

                kind = _watermark_status_kind(status)
                if kind == "success":
                    if stage_dir is not None:
                        _watermark_replace_original(stage_path, src)
                        output_path = src
                    if settings["delete_source"] and output_path.resolve() != src.resolve():
                        try:
                            src.unlink()
                        except Exception as delete_exc:
                            _watermark_log(app, logs, f"[批量水印] 删除源文件失败: {src.name} | {delete_exc}")
                    success_outputs.append(str(output_path))
                    _watermark_checkpoint_mark(checkpoint, src, output_path, "success")
                    _add_task_result_output(result, str(output_path))
                    _watermark_log(app, logs, f"[批量水印] 成功: {src.name} -> {output_path}")
                elif kind == "skipped":
                    if _watermark_should_preserve_original(status) and not same_file_output:
                        try:
                            output_path.parent.mkdir(parents=True, exist_ok=True)
                            shutil.copy2(str(src), str(output_path))
                            copied_skipped_outputs.append(str(output_path))
                            _add_task_result_output(result, str(output_path))
                            _watermark_log(app, logs, f"[批量水印] 受保护文件无法加水印，已保留原文件: {src.name} -> {output_path}")
                        except Exception as copy_exc:
                            failed_items.append(str(src))
                            _watermark_log(app, logs, f"[批量水印] 保留受保护文件失败: {src.name} | {copy_exc}")
                    elif _get_watermark_copy_skipped_to_output(app):
                        deferred_skipped_copy_items.append(str(src))
                    skipped_items.append(str(src))
                    _watermark_log(app, logs, f"[批量水印] 跳过: {src.name} | {status}")
                else:
                    failed_items.append(str(src))
                    _watermark_log(app, logs, f"[批量水印] 失败: {src.name} | {status}")
            except Exception as exc:
                failed_items.append(str(src))
                _watermark_log(app, logs, f"[批量水印] 失败: {src.name} | {exc}")
            finally:
                processed_count += 1
                if stage_dir is not None:
                    try:
                        shutil.rmtree(stage_dir, ignore_errors=True)
                    except Exception:
                        pass
                _watermark_update_progress(
                    app,
                    current_file=str(src),
                    stage="添加水印",
                    completed=processed_count,
                    total=max(overall_total, 1),
                )

        if parallel_executor is not None:
            workers = _get_watermark_pdf_worker_count(len(pdf_files))
            while parallel_future_map:
                done, _ = concurrent.futures.wait(
                    parallel_future_map,
                    return_when=concurrent.futures.FIRST_COMPLETED,
                )
                for future in done:
                    record_parallel_pdf_result(future)
                    if parallel_pdf_job_index < len(parallel_pdf_jobs) and not getattr(app, "stop_event", False):
                        job = parallel_pdf_jobs[parallel_pdf_job_index]
                        parallel_pdf_job_index += 1
                        parallel_future_map[parallel_executor.submit(process_parallel_pdf_job, job)] = job
            parallel_executor.shutdown(wait=True)
            parallel_executor = None

        if rule_skipped_files:
            copied_skipped, failed_copies = _copy_watermark_skipped_files(
                app,
                rule_skipped_files,
                input_root,
                output_root,
                actual_strategy,
                logs,
                result,
            )
            copied_skipped_outputs.extend(copied_skipped)
            skipped_items.extend(str(Path(_normalize_input_path_value(item))) for item in rule_skipped_files)
            failed_items.extend(failed_copies)
        if type_skipped_files:
            copied_skipped, failed_copies = _copy_watermark_skipped_files(
                app,
                type_skipped_files,
                input_root,
                output_root,
                actual_strategy,
                logs,
                result,
            )
            copied_skipped_outputs.extend(copied_skipped)
            skipped_items.extend(str(Path(_normalize_input_path_value(item))) for item in type_skipped_files)
            failed_items.extend(failed_copies)
        if deferred_skipped_copy_items:
            copied_skipped, failed_copies = _copy_watermark_skipped_files(
                app,
                deferred_skipped_copy_items,
                input_root,
                output_root,
                actual_strategy,
                logs,
                result,
            )
            copied_skipped_outputs.extend(copied_skipped)
            failed_items.extend(failed_copies)
        if unsupported_skipped_files:
            copied_skipped, failed_copies = _copy_watermark_skipped_files(
                app,
                unsupported_skipped_files,
                input_root,
                output_root,
                actual_strategy,
                logs,
                result,
            )
            copied_skipped_outputs.extend(copied_skipped)
            skipped_items.extend(str(Path(_normalize_input_path_value(item))) for item in unsupported_skipped_files)
            failed_items.extend(failed_copies)

        if failed_items:
            report_root = Path(output_root if actual_strategy == "result_folder" else input_root)
            try:
                report_root.mkdir(parents=True, exist_ok=True)
            except Exception as exc:
                _watermark_log(app, logs, f"[批量水印] 失败清单写入失败: {report_root} | {_watermark_format_path_error(report_root, exc)}")
                report_root = Path(input_root)
                report_root.mkdir(parents=True, exist_ok=True)
            report_path = report_root / "!失败文件清单.txt"
            lines = ["以下文件处理失败："]
            for item in failed_items:
                try:
                    lines.append(os.path.relpath(item, input_root))
                except Exception:
                    lines.append(str(item))
            report_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
            _add_task_result_output(result, str(report_path))

        result["failed_items"] = list(failed_items)
        _set_task_result_counts(
            result,
            processed=processed_count + len(rule_skipped_files) + len(type_skipped_files) + len(unsupported_skipped_files),
            success=len(success_outputs),
            failed=len(failed_items),
            skipped=len(skipped_items),
        )
        if getattr(app, "stop_event", False):
            _watermark_checkpoint_finish(checkpoint, "paused", processed_count, overall_total, "user_stop")
            _set_task_result_finished(result, "stopped", message="用户停止批量水印任务", detail="用户停止批量水印任务", stopped=True)
        elif failed_items:
            _watermark_checkpoint_finish(checkpoint, "failed", processed_count, overall_total, "failed_items")
            _set_task_result_finished(
                    result,
                    "failed",
                    message=f"批量水印失败 {len(failed_items)} 个文件",
                    detail=f"批量水印失败 {len(failed_items)} 个文件",
                    error=f"failed_items={len(failed_items)}",
                )
        elif success_outputs or copied_skipped_outputs:
            _watermark_checkpoint_finish(checkpoint, "completed", processed_count, overall_total, "completed")
            _set_task_result_finished(result, "success", message="批量水印完成", detail="批量水印完成")
        else:
            _watermark_checkpoint_finish(checkpoint, "completed", processed_count, overall_total, "completed")
            _set_task_result_finished(result, "skipped", message="没有生成新的水印文件", detail="没有生成新的水印文件", skipped=True)
        _flush_watermark_ui_logs(app, force=True)
        return result
    except Exception as exc:
        _watermark_checkpoint_finish(
            checkpoint,
            "interrupted",
            processed_count,
            overall_total,
            f"exception:{type(exc).__name__}",
        )
        raise
    finally:
        if parallel_executor is not None:
            try:
                parallel_executor.shutdown(wait=True, cancel_futures=True)
            except TypeError:
                parallel_executor.shutdown(wait=True)
            except Exception:
                pass
        try:
            if word_app is not None:
                word_app.Quit()
        except Exception:
            pass
        try:
            if ppt_app is not None:
                ppt_app.Quit()
        except Exception:
            pass
        if pythoncom_initialized:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass


def _patch_watermark_task():
    try:
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_watermark_task:missing:{exc}")
        return
    if getattr(original_run_process, "__fx_watermark_task_patch__", False):
        return

    def patched_run_process(self, input_folder, task_type):
        if task_type == "watermark":
            try:
                return _run_watermark_task(self, input_folder)
            except Exception as exc:
                try:
                    self.log(f"[批量水印] 严重错误: {exc}")
                except Exception:
                    pass
                checkpoint = getattr(self, "_fx_watermark_checkpoint_active", None)
                _watermark_checkpoint_finish(
                    checkpoint,
                    "interrupted",
                    int((_get_last_task_result(self) or {}).get("processed_count") or 0),
                    int((_get_last_task_result(self) or {}).get("total_count") or 0),
                    f"exception:{type(exc).__name__}",
                )
                _finalize_current_task_result(self, "failed", message=str(exc), detail=str(exc), error=str(exc))
            finally:
                try:
                    self._fx_watermark_checkpoint_active = None
                except Exception:
                    pass
                _set_watermark_checkpoint_action_state(self, "normal")
                try:
                    self.reset_ui()
                except Exception:
                    pass
            return None
        return original_run_process(self, input_folder, task_type)

    patched_run_process.__fx_watermark_task_patch__ = True
    patched_run_process.__wrapped__ = original_run_process
    FengxiToolboxApp.run_process = patched_run_process
    _debug("patch_watermark_task:installed")


_patch_watermark_task()


def _get_parallel_detail_key(app, task_type=None):
    task_type = str(task_type or getattr(app, "current_task", "") or "")
    detail = ""
    try:
        if task_type == "pdf" and getattr(app, "pdf_mode_var", None) is not None:
            detail = str(app.pdf_mode_var.get() or "")
        elif task_type == "image" and getattr(app, "img_mode_var", None) is not None:
            detail = str(app.img_mode_var.get() or "")
        elif task_type == "file" and getattr(app, "file_mode_var", None) is not None:
            detail = str(app.file_mode_var.get() or "")
        elif task_type == "meta" and getattr(app, "meta_mode_var", None) is not None:
            detail = str(app.meta_mode_var.get() or "")
    except Exception:
        detail = ""
    return task_type, detail


def _get_parallel_mode_message(app, task_type=None):
    task_type, detail = _get_parallel_detail_key(app, task_type)
    if not task_type:
        return "并行状态：请选择具体功能。"
    if getattr(app, "_fx_single_input_target", None):
        return "并行状态：单文件输入会自动单线程，避免 UI 日志线程冲突。"
    spec = _get_feature_spec(task_type)
    parallel = spec.get("parallel") if isinstance(spec.get("parallel"), dict) else {}
    details = parallel.get("detail") if isinstance(parallel.get("detail"), dict) else {}
    detail_value = details.get(detail)
    if detail_value is None:
        detail_value = details.get("")
    if detail_value is not None:
        detail_mode = parallel.get("mode", "")
        detail_message = detail_value
        if isinstance(detail_value, (tuple, list)) and len(detail_value) >= 2:
            detail_mode, detail_message = detail_value[0], detail_value[1]
        if detail_mode == "forced_single":
            return "并行状态：稳定单线程 · " + str(detail_message or "当前功能使用专用流程。")
    if parallel.get("mode") == "forced_single":
        return "并行状态：稳定单线程 · 当前功能使用专用流程。"
    if parallel.get("mode") == "safe":
        return "并行状态：可提速 · " + str(parallel.get("hint") or "多文件批处理可尝试并行。")
    return "并行状态：按当前功能自动选择。"


def _refresh_parallel_mode_hint(app):
    switch = getattr(app, "chk_multithread", None)
    if switch is not None:
        try:
            switch.configure(text=PARALLEL_SWITCH_TEXT)
        except Exception:
            pass
    label = getattr(app, "_fx_parallel_hint_label", None)
    if label is not None:
        try:
            if label.winfo_exists():
                label.destroy()
        except Exception:
            try:
                label.destroy()
            except Exception:
                pass
        try:
            app._fx_parallel_hint_label = None
        except Exception:
            pass
    var = getattr(app, "_fx_parallel_hint_var", None)
    if var is None:
        return ""
    try:
        var.set("")
    except Exception:
        pass
    return ""


def _install_parallel_mode_hint(app):
    if getattr(app, "_fx_parallel_hint_ready", False):
        _refresh_parallel_mode_hint(app)
        return
    switch = getattr(app, "chk_multithread", None)
    if switch is None:
        return
    try:
        switch.configure(text=PARALLEL_SWITCH_TEXT)
    except Exception:
        pass
    try:
        app._fx_parallel_hint_var = tkinter.StringVar(master=app, value="")
    except Exception:
        app._fx_parallel_hint_var = None

    def on_parallel_toggle(*_args, target=app):
        _refresh_parallel_mode_hint(target)

    try:
        getattr(app, "enable_multithread", None).trace_add("write", on_parallel_toggle)
    except Exception:
        pass

    app._fx_parallel_hint_ready = True
    _refresh_parallel_mode_hint(app)


def _normalize_preset_category(value):
    normalized = str(value or "").strip()
    if normalized in PRESET_CATEGORY_LABELS:
        return normalized
    mapped = PRESET_LABEL_TO_CATEGORY.get(normalized)
    if mapped:
        return mapped
    return "watermark"


def _get_current_preset_category(app):
    task_type = str(getattr(app, "current_task", "") or "")
    if task_type == "pdf":
        mode = str(_safe_var_get(app, "pdf_mode_var", "") or "")
        if mode == "ocr":
            return "ocr"
        if mode == "compress":
            return "pdf_compress"
    if task_type == "file":
        return "rename"
    if task_type == "watermark":
        return "watermark"
    if task_type == "zip":
        return "zip"
    return "watermark"


def _make_preset_id():
    return _prefs_make_preset_id()


def _load_presets():
    return _prefs_load_presets(_user_prefs_context())


def _save_presets(presets):
    return _prefs_save_presets(presets, _user_prefs_context())


def _save_preset_entry(name, category, settings):
    return _prefs_save_preset_entry(
        name,
        category,
        settings,
        _user_prefs_context(),
        default_name_suffix=_format_queue_time(),
    )


def _delete_preset_entry(preset_id):
    return _prefs_delete_preset_entry(preset_id, _user_prefs_context())


def _find_preset_entry(preset_id):
    return _prefs_find_preset_entry(preset_id, _user_prefs_context())


def _preset_pick_display_from_key(display_value, key_value, mapping):
    if isinstance(mapping, dict):
        if display_value in mapping:
            return display_value
        for display, key in mapping.items():
            if key == key_value:
                return display
    return display_value


def _capture_preset_settings(app, category=None):
    category = _normalize_preset_category(category or _get_current_preset_category(app))
    task_type = PRESET_CATEGORY_TO_TASK.get(category)
    if task_type:
        try:
            _ensure_lazy_tab_initialized(app, task_type)
        except Exception as exc:
            _debug(f"preset:capture_lazy_error:{category}:{exc}")

    output_strategy = _safe_var_get(app, "output_strategy_var", "")
    settings = {"category": category, "output_strategy": output_strategy}

    if category == "watermark":
        settings.update(
            {
                "wm_text": _read_watermark_text_widget(app),
                "selected_font": _safe_var_get(app, "selected_font", ""),
                "wm_range_var": _safe_var_get(app, "wm_range_var", "all"),
                "wm_overwrite_var": _safe_var_get(app, "wm_overwrite_var", "smart"),
                "allow_simsun": bool(_safe_var_get(app, "allow_simsun", False)),
                "wm_delete_var": bool(_safe_var_get(app, "wm_delete_var", False)),
                "wm_convert_pdf": bool(_safe_var_get(app, "wm_convert_pdf", False)),
                "wm_skip_hyphen_var": bool(_safe_var_get(app, "wm_skip_hyphen_var", False)),
                "wm_skip_name_position_var": _safe_var_get(app, "wm_skip_name_position_var", "结尾"),
                "wm_skip_name_text_var": _safe_var_get(app, "wm_skip_name_text_var", "-"),
                "wm_copy_skipped_var": bool(_safe_var_get(app, "wm_copy_skipped_var", False)),
                "wm_skip_pdf_type_var": bool(_safe_var_get(app, "wm_skip_pdf_type_var", False)),
                "wm_skip_word_type_var": bool(_safe_var_get(app, "wm_skip_word_type_var", False)),
                "wm_skip_ppt_type_var": bool(_safe_var_get(app, "wm_skip_ppt_type_var", False)),
                "wm_color_var": _get_watermark_color(app),
                "wm_copy_guard_enabled_var": bool(_safe_var_get(app, "wm_copy_guard_enabled_var", False)),
                "wm_copy_guard_strength_var": WATERMARK_COPY_GUARD_STRENGTH_LABELS.get(
                    _get_watermark_copy_guard_strength(app),
                    "标准",
                ),
                "slider_size": _safe_named_widget_get(app, "slider_size", 60),
                "slider_opacity": _safe_named_widget_get(app, "slider_opacity", 0.08),
                "slider_angle": _safe_named_widget_get(app, "slider_angle", 45),
            }
        )
    elif category == "ocr":
        backend_display = _safe_var_get(app, "pdf_ocr_backend", "")
        language_display = _safe_var_get(app, "pdf_ocr_language", "")
        mode_display = _safe_var_get(app, "pdf_ocr_mode", "")
        preprocess_display = _safe_var_get(app, "pdf_ocr_preprocess", "")
        settings.update(
            {
                "pdf_mode_var": "ocr",
                "pdf_pwd_entry": _safe_named_widget_get(app, "pdf_pwd_entry", ""),
                "pdf_delete_var": bool(_safe_var_get(app, "pdf_delete_var", False)),
                "pdf_ocr_model_root": _safe_var_get(app, "pdf_ocr_model_root", ""),
                "pdf_ocr_backend": backend_display,
                "pdf_ocr_backend_key": getattr(app, "_fx_pdf_ocr_backend_map", {}).get(backend_display, ""),
                "pdf_ocr_language": language_display,
                "pdf_ocr_language_key": getattr(app, "_fx_pdf_ocr_lang_map", {}).get(language_display, ""),
                "pdf_ocr_mode": mode_display,
                "pdf_ocr_mode_key": getattr(app, "_fx_pdf_ocr_mode_map", {}).get(mode_display, ""),
                "pdf_ocr_preprocess": preprocess_display,
                "pdf_ocr_preprocess_key": getattr(app, "_fx_pdf_ocr_preprocess_map", {}).get(preprocess_display, ""),
                "pdf_ocr_cls": bool(_safe_var_get(app, "pdf_ocr_cls", False)),
                "pdf_ocr_compare_report": bool(_safe_var_get(app, "pdf_ocr_compare_report", False)),
                "pdf_ocr_normalize_rotation": bool(_safe_var_get(app, "pdf_ocr_normalize_rotation", True)),
            }
        )
    elif category == "pdf_compress":
        settings.update(
            {
                "pdf_mode_var": "compress",
                "pdf_pwd_entry": _safe_named_widget_get(app, "pdf_pwd_entry", ""),
                "pdf_delete_var": bool(_safe_var_get(app, "pdf_delete_var", False)),
                "pdf_compress_level_var": _safe_var_get(app, "pdf_compress_level_var", "标准"),
                "pdf_image_compress_level_var": _safe_var_get(app, "pdf_image_compress_level_var", "标准"),
            }
        )
    elif category == "audio":
        transcribe_args = _get_audio_transcribe_args(app)
        settings.update(
            {
                "audio_mode_var": _safe_var_get(app, "audio_mode_var", "video2mp3"),
                "audio_target_fmt": _safe_var_get(app, "audio_target_fmt", "mp3"),
                "audio_bitrate": _safe_var_get(app, "audio_bitrate", "192k"),
                "audio_delete_var": bool(_safe_var_get(app, "audio_delete_var", False)),
                "audio_transcribe_model": transcribe_args.get("model_name", "base"),
                "audio_transcribe_language": transcribe_args.get("language", "自动识别"),
                "audio_transcribe_format": transcribe_args.get("output_format", "txt"),
            }
        )
    elif category == "rename":
        settings.update(
            {
                "file_mode_var": "rename",
                "rename_type_var": _safe_var_get(app, "rename_type_var", "add"),
                "rename_prefix": _safe_named_widget_get(app, "rename_prefix", ""),
                "rename_suffix": _safe_named_widget_get(app, "rename_suffix", ""),
                "rename_find": _safe_named_widget_get(app, "rename_find", ""),
                "rename_rep": _safe_named_widget_get(app, "rename_rep", ""),
                "rename_cut_head": _safe_named_widget_get(app, "rename_cut_head", ""),
                "rename_cut_tail": _safe_named_widget_get(app, "rename_cut_tail", ""),
            }
        )
    elif category == "zip":
        settings.update(
            {
                "zip_mode_var": _safe_var_get(app, "zip_mode_var", "total"),
                "zip_min_depth_var": _safe_var_get(app, "zip_min_depth_var", ""),
                "zip_max_depth_var": _safe_var_get(app, "zip_max_depth_var", ""),
                "zip_depth_range_var": _get_zip_max_depth(app),
                "zip_archive_policy_var": _get_zip_archive_policy(app),
            }
        )
    return settings


def _switch_to_preset_task(app, category):
    task_type = PRESET_CATEGORY_TO_TASK.get(_normalize_preset_category(category))
    if not task_type:
        return
    try:
        if callable(getattr(app, "switch_tab", None)):
            app.switch_tab(task_type)
        else:
            _ensure_lazy_tab_initialized(app, task_type)
            app.current_task = task_type
    except Exception as exc:
        _debug(f"preset:switch_task_error:{category}:{exc}")
        try:
            _ensure_lazy_tab_initialized(app, task_type)
            app.current_task = task_type
        except Exception:
            pass


def _select_pdf_preset_mode(app, mode):
    _safe_var_set(app, "pdf_mode_var", mode)
    selector = getattr(app, "_fx_select_pdf_mode", None)
    if callable(selector):
        try:
            selector(mode)
            return
        except Exception:
            pass
    panels = getattr(app, "_fx_pdf_detail_panels", {})
    if isinstance(panels, dict):
        for key, panel in panels.items():
            try:
                if key == mode:
                    panel.grid(row=0, column=0, sticky="nsew", padx=10, pady=10)
                    panel.tkraise()
                else:
                    panel.grid_remove()
            except Exception:
                pass


def _apply_preset_settings(app, preset, switch_task=True):
    if not isinstance(preset, dict):
        return False, "设置不存在。"
    category = _normalize_preset_category(preset.get("category"))
    settings = preset.get("settings") if isinstance(preset.get("settings"), dict) else {}
    if switch_task:
        _switch_to_preset_task(app, category)
    else:
        try:
            _ensure_lazy_tab_initialized(app, PRESET_CATEGORY_TO_TASK.get(category, "watermark"))
        except Exception:
            pass

    output_strategy = settings.get("output_strategy")
    if output_strategy:
        _safe_var_set(app, "output_strategy_var", _get_output_strategy_label(_coerce_output_strategy_value(output_strategy)))
        _refresh_output_strategy_hint(app)

    if category == "watermark":
        _safe_named_widget_set(app, "wm_text", settings.get("wm_text", ""))
        for name in (
            "selected_font",
            "wm_range_var",
            "wm_overwrite_var",
            "wm_skip_name_position_var",
            "wm_skip_name_text_var",
            "wm_copy_guard_strength_var",
        ):
            if name in settings:
                _safe_var_set(app, name, settings.get(name))
        if "wm_color_var" in settings:
            _set_watermark_color(app, settings.get("wm_color_var"))
        for name in (
            "allow_simsun",
            "wm_delete_var",
            "wm_convert_pdf",
            "wm_skip_hyphen_var",
            "wm_copy_skipped_var",
            "wm_skip_pdf_type_var",
            "wm_skip_word_type_var",
            "wm_skip_ppt_type_var",
            "wm_copy_guard_enabled_var",
        ):
            if name in settings:
                _safe_var_set(app, name, bool(settings.get(name)))
        for name in ("slider_size", "slider_opacity", "slider_angle"):
            if name in settings:
                _safe_named_widget_set(app, name, settings.get(name))
        try:
            _refresh_watermark_preview(app)
        except Exception:
            pass
        try:
            _flush_watermark_text_persistence(app)
            _flush_watermark_filename_rule_persistence(app)
        except Exception:
            pass
    elif category == "ocr":
        _select_pdf_preset_mode(app, "ocr")
        _safe_named_widget_set(app, "pdf_pwd_entry", settings.get("pdf_pwd_entry", ""))
        _safe_var_set(app, "pdf_delete_var", bool(settings.get("pdf_delete_var", False)))
        _safe_var_set(app, "pdf_ocr_model_root", settings.get("pdf_ocr_model_root", ""))
        _safe_var_set(
            app,
            "pdf_ocr_backend",
            _preset_pick_display_from_key(
                settings.get("pdf_ocr_backend", ""),
                settings.get("pdf_ocr_backend_key", ""),
                getattr(app, "_fx_pdf_ocr_backend_map", {}),
            ),
        )
        _safe_var_set(
            app,
            "pdf_ocr_language",
            _preset_pick_display_from_key(
                settings.get("pdf_ocr_language", ""),
                settings.get("pdf_ocr_language_key", ""),
                getattr(app, "_fx_pdf_ocr_lang_map", {}),
            ),
        )
        _safe_var_set(
            app,
            "pdf_ocr_mode",
            _preset_pick_display_from_key(
                settings.get("pdf_ocr_mode", ""),
                settings.get("pdf_ocr_mode_key", ""),
                getattr(app, "_fx_pdf_ocr_mode_map", {}),
            ),
        )
        _safe_var_set(
            app,
            "pdf_ocr_preprocess",
            _preset_pick_display_from_key(
                settings.get("pdf_ocr_preprocess", ""),
                settings.get("pdf_ocr_preprocess_key", ""),
                getattr(app, "_fx_pdf_ocr_preprocess_map", {}),
            ),
        )
        _safe_var_set(app, "pdf_ocr_cls", bool(settings.get("pdf_ocr_cls", False)))
        _safe_var_set(app, "pdf_ocr_compare_report", bool(settings.get("pdf_ocr_compare_report", False)))
        _safe_var_set(app, "pdf_ocr_normalize_rotation", bool(settings.get("pdf_ocr_normalize_rotation", True)))
    elif category == "pdf_compress":
        _select_pdf_preset_mode(app, "compress")
        _safe_named_widget_set(app, "pdf_pwd_entry", settings.get("pdf_pwd_entry", ""))
        _safe_var_set(app, "pdf_delete_var", bool(settings.get("pdf_delete_var", False)))
        _safe_var_set(app, "pdf_compress_level_var", settings.get("pdf_compress_level_var", "标准"))
        _safe_var_set(app, "pdf_image_compress_level_var", settings.get("pdf_image_compress_level_var", "标准"))
    elif category == "audio":
        for name in (
            "audio_mode_var",
            "audio_target_fmt",
            "audio_bitrate",
            "audio_transcribe_model",
            "audio_transcribe_language",
            "audio_transcribe_format",
        ):
            if name in settings:
                _safe_var_set(app, name, settings.get(name))
        if "audio_delete_var" in settings:
            _safe_var_set(app, "audio_delete_var", bool(settings.get("audio_delete_var", False)))
    elif category == "rename":
        _safe_var_set(app, "file_mode_var", "rename")
        _safe_var_set(app, "rename_type_var", settings.get("rename_type_var", "add"))
        for name in ("rename_prefix", "rename_suffix", "rename_find", "rename_rep", "rename_cut_head", "rename_cut_tail"):
            _safe_named_widget_set(app, name, settings.get(name, ""))
    elif category == "zip":
        _safe_var_set(app, "zip_mode_var", settings.get("zip_mode_var", "total"))
        if "zip_min_depth_var" in settings:
            _safe_var_set(app, "zip_min_depth_var", settings.get("zip_min_depth_var", ""))
            _safe_var_set(app, "zip_max_depth_var", settings.get("zip_max_depth_var", ""))
        else:
            _set_zip_depth_fields(app, settings.get("zip_depth_range_var", settings.get("zip_max_depth_var", "")))
        _set_zip_archive_policy(app, settings.get("zip_archive_policy_var", settings.get("zip_existing_archive_policy", "reuse_existing")))
    else:
        return False, "暂不支持该设置类型。"

    return True, f"已恢复设置：{preset.get('name', '')}"


def _load_last_settings():
    return _prefs_load_last_settings(_user_prefs_context())


def _save_last_settings_entry(category, settings, update_active=True):
    return _prefs_save_last_settings_entry(
        category,
        settings,
        _user_prefs_context(),
        update_active=update_active,
    )


def _last_settings_category_ready(app, category):
    category = _normalize_preset_category(category)
    if category == "watermark":
        return getattr(app, "wm_text", None) is not None and getattr(app, "selected_font", None) is not None
    if category in {"ocr", "pdf_compress"}:
        return getattr(app, "pdf_mode_var", None) is not None
    if category == "audio":
        return getattr(app, "audio_mode_var", None) is not None
    if category == "rename":
        return getattr(app, "rename_type_var", None) is not None and getattr(app, "rename_prefix", None) is not None
    if category == "zip":
        return (
            getattr(app, "zip_mode_var", None) is not None
            and getattr(app, "zip_min_depth_var", None) is not None
            and getattr(app, "zip_max_depth_var", None) is not None
            and getattr(app, "zip_archive_policy_var", None) is not None
        )
    return False


def _get_current_last_settings_category(app):
    task_type = str(getattr(app, "current_task", "") or "")
    if task_type == "watermark":
        return "watermark"
    if task_type == "pdf":
        mode = str(_safe_var_get(app, "pdf_mode_var", "") or "")
        if mode == "ocr":
            return "ocr"
        if mode == "compress":
            return "pdf_compress"
        return None
    if task_type == "file":
        mode = str(_safe_var_get(app, "file_mode_var", "rename") or "")
        if mode == "rename":
            return "rename"
    if task_type == "audio":
        return "audio"
    if task_type == "zip":
        return "zip"
    return None


def _save_last_settings_category(app, category=None, update_active=True):
    raw_category = category if category is not None else _get_current_last_settings_category(app)
    if not raw_category:
        return None
    category = _normalize_preset_category(raw_category)
    if not _last_settings_category_ready(app, category):
        return None
    settings = _capture_preset_settings(app, category)
    return _save_last_settings_entry(category, settings, update_active=update_active)


def _save_current_last_settings(app):
    category = _get_current_last_settings_category(app)
    if not category:
        return None
    return _save_last_settings_category(app, category)


def _flush_watermark_last_settings_persistence(app):
    if getattr(app, "_fx_last_settings_loading", False):
        return None
    if not _last_settings_category_ready(app, "watermark"):
        return None
    return _save_last_settings_category(app, "watermark", update_active=True)


def _schedule_watermark_last_settings_persistence(app, delay_ms=350):
    if getattr(app, "_fx_last_settings_loading", False):
        return
    after_id = getattr(app, "_fx_wm_last_settings_after_id", None)
    if after_id:
        try:
            app.after_cancel(after_id)
        except Exception:
            pass

    def persist(target=app):
        try:
            target._fx_wm_last_settings_after_id = None
        except Exception:
            pass
        try:
            _flush_watermark_last_settings_persistence(target)
        except Exception as exc:
            _debug(f"last_settings:watermark_auto_save_error:{exc}")

    try:
        app._fx_wm_last_settings_after_id = app.after(delay_ms, persist)
    except Exception:
        persist()


def _install_watermark_last_settings_memory(app):
    if getattr(app, "_fx_wm_last_settings_memory_ready", False):
        return

    def changed(*_args, target=app):
        _schedule_watermark_last_settings_persistence(target)

    variable_names = (
        "selected_font",
        "wm_range_var",
        "wm_overwrite_var",
        "allow_simsun",
        "wm_delete_var",
        "wm_convert_pdf",
        "wm_skip_hyphen_var",
        "wm_skip_name_position_var",
        "wm_skip_name_text_var",
        "wm_copy_skipped_var",
        "wm_skip_pdf_type_var",
        "wm_skip_word_type_var",
        "wm_skip_ppt_type_var",
        "wm_color_var",
        "wm_copy_guard_enabled_var",
        "wm_copy_guard_strength_var",
        "output_strategy_var",
    )
    trace_ids = []
    for name in variable_names:
        var = getattr(app, name, None)
        if not isinstance(var, tkinter.Variable):
            continue
        try:
            trace_ids.append((name, var.trace_add("write", changed)))
        except Exception:
            pass
    app._fx_wm_last_settings_trace_ids = trace_ids

    for name in ("slider_size", "slider_opacity", "slider_angle"):
        widget = getattr(app, name, None)
        if widget is None:
            continue
        try:
            original_command = getattr(widget, "_command", None)
            if not getattr(original_command, "__fx_wm_last_settings_slider_command__", False):
                def slider_command(value=None, target=app, original=original_command):
                    if callable(original):
                        try:
                            original(value)
                        except TypeError:
                            original()
                    _schedule_watermark_last_settings_persistence(target)

                slider_command.__fx_wm_last_settings_slider_command__ = True
                widget.configure(command=slider_command)
        except Exception:
            pass
        for event_name in ("<ButtonRelease-1>", "<B1-Motion>"):
            try:
                widget.bind(event_name, lambda _event=None, target=app: _schedule_watermark_last_settings_persistence(target), add="+")
            except Exception:
                pass

    app._fx_wm_last_settings_memory_ready = True


def _save_initialized_last_settings(app):
    saved = {}
    current_category = _get_current_last_settings_category(app)
    for category in ("watermark", "ocr", "pdf_compress", "audio", "rename"):
        try:
            if _last_settings_category_ready(app, category):
                entry = _save_last_settings_category(app, category, update_active=(category == current_category))
                if entry:
                    saved[category] = entry
        except Exception as exc:
            _debug(f"last_settings:save_initialized_error:{category}:{exc}")
    return saved


def _get_active_last_settings_category(task_name):
    return _prefs_get_active_last_settings_category(task_name, _user_prefs_context())


def _restore_last_settings_category(app, category):
    category = _normalize_preset_category(category)
    entry = _load_last_settings().get(category)
    if not isinstance(entry, dict):
        return False, "没有可恢复的上次设置。"
    if not isinstance(entry.get("settings"), dict):
        return False, "上次设置为空。"
    previous_loading = getattr(app, "_fx_last_settings_loading", False)
    app._fx_last_settings_loading = True
    try:
        ok, message = _apply_preset_settings(
            app,
            {
                "name": "上次设置",
                "category": category,
                "settings": entry.get("settings") or {},
            },
            switch_task=False,
        )
        return ok, message
    finally:
        app._fx_last_settings_loading = previous_loading


def _restore_last_settings_for_task(app, task_name):
    restored = getattr(app, "_fx_last_settings_restored_tasks", None)
    if not isinstance(restored, set):
        restored = set()
        app._fx_last_settings_restored_tasks = restored
    if task_name in restored:
        return False
    category = _get_active_last_settings_category(task_name)
    if not category:
        restored.add(task_name)
        return False
    ok, _message = _restore_last_settings_category(app, category)
    restored.add(task_name)
    return ok


def _queue_snapshot_app_state(app, task_type):
    variables = {}
    widgets = {}
    try:
        _ensure_lazy_tab_initialized(app, task_type)
    except Exception as exc:
        _debug(f"queue:snapshot_lazy_error:{task_type}:{exc}")

    for name, value in list(vars(app).items()):
        try:
            if isinstance(value, tkinter.Variable):
                variables[name] = value.get()
        except Exception:
            pass
        try:
            if isinstance(
                value,
                (
                    customtkinter.CTkEntry,
                    customtkinter.CTkTextbox,
                    customtkinter.CTkComboBox,
                    customtkinter.CTkOptionMenu,
                ),
            ):
                widget_value = _safe_widget_get(value)
                if widget_value is not None:
                    widgets[name] = widget_value
        except Exception:
            pass
    return {"variables": variables, "widgets": widgets}


def _queue_restore_app_state(app, task):
    task_type = task.get("task_type")
    if task_type:
        try:
            _ensure_lazy_tab_initialized(app, task_type)
        except Exception as exc:
            _debug(f"queue:restore_lazy_error:{task_type}:{exc}")
    snapshot = task.get("snapshot") if isinstance(task, dict) else {}
    if not isinstance(snapshot, dict):
        return
    variables = snapshot.get("variables") if isinstance(snapshot.get("variables"), dict) else {}
    widgets = snapshot.get("widgets") if isinstance(snapshot.get("widgets"), dict) else {}
    for name, value in variables.items():
        try:
            var = getattr(app, name, None)
            if isinstance(var, tkinter.Variable):
                var.set(value)
        except Exception:
            pass
    for name, value in widgets.items():
        try:
            widget = getattr(app, name, None)
            _safe_widget_set(widget, value)
        except Exception:
            pass
    try:
        if task.get("input"):
            app.input_path.set(task.get("input"))
    except Exception:
        pass
    if task_type:
        try:
            app.current_task = task_type
        except Exception:
            pass


def _queue_describe_task(app, task_type, input_path):
    label = _get_feature_label(task_type, fallback="未知任务")
    name = os.path.basename(str(input_path or "").rstrip("\\/")) or str(input_path or "未选择路径")
    detail = ""
    try:
        if task_type == "pdf" and getattr(app, "pdf_mode_var", None) is not None:
            detail = app.pdf_mode_var.get()
        elif task_type == "image" and getattr(app, "img_mode_var", None) is not None:
            detail = app.img_mode_var.get()
        elif task_type == "zip" and getattr(app, "zip_mode_var", None) is not None:
            detail = app.zip_mode_var.get()
        elif task_type == "convert":
            detail = _get_convert_preview_detail(app)
        elif task_type == "file" and getattr(app, "file_mode_var", None) is not None:
            detail = app.file_mode_var.get()
        elif task_type == "watermark":
            detail = "add"
        elif task_type == "remove_wm":
            detail = "remove"
    except Exception:
        detail = ""
    suffix = f" · {detail}" if detail else ""
    return f"{label}{suffix} · {name}"


def _get_queue_history_context():
    return QueueHistoryContext(
        history_file=_get_queue_history_file,
        retention_days=QUEUE_HISTORY_RETENTION_DAYS,
        history_limit=QUEUE_HISTORY_LIMIT,
        status_labels=QUEUE_STATUS_LABELS,
        status_label_to_value=QUEUE_HISTORY_STATUS_LABEL_TO_VALUE,
        task_label_to_value=QUEUE_HISTORY_TASK_LABEL_TO_VALUE,
        failure_label_to_value=QUEUE_HISTORY_FAILURE_LABEL_TO_VALUE,
        classify_failure_reason=_classify_failure_reason,
        task_result_snapshot=_task_result_snapshot,
        debug=_debug,
    )


def _load_queue_history():
    return load_queue_history(_get_queue_history_context())


def _save_queue_history(entries):
    return save_queue_history(entries, _get_queue_history_context())


def _queue_history_entry_timestamp(entry):
    return queue_history_entry_timestamp(entry)


def _prune_queue_history_entries(entries, now=None):
    return prune_queue_history_entries(entries, _get_queue_history_context(), now=now)


def _normalize_queue_history_entry(task):
    return normalize_queue_history_entry(task, _get_queue_history_context())


def _append_queue_history(app, task):
    history = getattr(app, "_fx_task_history", None)
    if history is None:
        history = _load_queue_history()
        app._fx_task_history = history
    history.append(_normalize_queue_history_entry(task))
    history[:] = _prune_queue_history_entries(history)
    _save_queue_history(history)


def _queue_status_text(status):
    return queue_status_text(status, _get_queue_history_context())


def _queue_status_color(status):
    if status == "success":
        return "#A7D39B"
    if status == "failed":
        return "#F0A6A6"
    if status == "skipped":
        return "#B5C0CC"
    if status == "running":
        return "#F6D28B"
    if status == "stopped":
        return "#B5C0CC"
    return globals().get("COLOR_TEXT_SOFT", "#B2C0C8")


def _build_queue_history_search_blob(entry):
    return build_queue_history_search_blob(entry, _get_queue_history_context())


def _classify_failure_reason(entry):
    item = dict(entry or {})
    task_result = item.get("task_result") if isinstance(item.get("task_result"), dict) else {}
    error_text = str(task_result.get("error") or item.get("error") or "").strip()
    detail_text = str(task_result.get("detail") or item.get("detail") or "").strip()
    logs = [str(text).strip() for text in (item.get("logs") or []) if str(text).strip()]
    failed_items = list(task_result.get("failed_items") or item.get("failed_items") or [])
    reason_text = error_text or detail_text
    if not reason_text and logs:
        reason_text = logs[-1]
    reason_lower = reason_text.lower()
    if any(marker in reason_lower for marker in ("路径不存在", "not found", "no such file", "找不到", "不存在")):
        return "path_missing", reason_text
    if any(marker in reason_lower for marker in ("权限", "access denied", "permission", "拒绝访问")):
        return "permission", reason_text
    if any(marker in reason_lower for marker in ("超时", "timeout", "timed out")):
        return "timeout", reason_text
    if any(marker in reason_lower for marker in ("依赖", "dll", "库", "module", "import", "onnxruntime", "com")):
        return "dependency", reason_text
    if failed_items:
        return "partial_failure", reason_text or f"{len(failed_items)} 个失败项"
    if any(marker in " ".join(logs).lower() for marker in ("❌", "🔥", "错误", "失败", "error", "failed", "traceback")):
        return "log_failure", reason_text or "日志包含失败标记"
    if reason_text:
        return "generic_failure", reason_text
    return "unknown", ""


def _filter_queue_history_entries(
    entries,
    status_filter=QUEUE_HISTORY_STATUS_DEFAULT,
    task_filter=QUEUE_HISTORY_TASK_DEFAULT,
    failure_filter=QUEUE_HISTORY_FAILURE_DEFAULT,
    keyword="",
):
    return filter_queue_history_entries(
        entries,
        _get_queue_history_context(),
        status_filter=status_filter,
        task_filter=task_filter,
        failure_filter=failure_filter,
        keyword=keyword,
    )


def _get_queue_history_filters(app):
    try:
        status_filter = getattr(app, "_fx_history_filter_status_var", None)
        status_filter = status_filter.get() if status_filter is not None else QUEUE_HISTORY_STATUS_DEFAULT
    except Exception:
        status_filter = QUEUE_HISTORY_STATUS_DEFAULT
    try:
        task_filter = getattr(app, "_fx_history_filter_task_var", None)
        task_filter = task_filter.get() if task_filter is not None else QUEUE_HISTORY_TASK_DEFAULT
    except Exception:
        task_filter = QUEUE_HISTORY_TASK_DEFAULT
    try:
        failure_filter = getattr(app, "_fx_history_filter_failure_var", None)
        failure_filter = failure_filter.get() if failure_filter is not None else QUEUE_HISTORY_FAILURE_DEFAULT
    except Exception:
        failure_filter = QUEUE_HISTORY_FAILURE_DEFAULT
    try:
        keyword = getattr(app, "_fx_history_search_var", None)
        keyword = keyword.get() if keyword is not None else ""
    except Exception:
        keyword = ""
    return (
        str(status_filter or QUEUE_HISTORY_STATUS_DEFAULT),
        str(task_filter or QUEUE_HISTORY_TASK_DEFAULT),
        str(failure_filter or QUEUE_HISTORY_FAILURE_DEFAULT),
        str(keyword or ""),
    )


def _get_filtered_queue_history(app):
    status_filter, task_filter, failure_filter, keyword = _get_queue_history_filters(app)
    history = list(getattr(app, "_fx_task_history", []) or [])
    return _filter_queue_history_entries(
        history,
        status_filter=status_filter,
        task_filter=task_filter,
        failure_filter=failure_filter,
        keyword=keyword,
    )


def _refresh_history_summary(app, filtered_history=None, full_history=None):
    summary_var = getattr(app, "_fx_history_summary_var", None)
    if summary_var is None:
        return
    if full_history is None:
        full_history = list(getattr(app, "_fx_task_history", []) or [])
    if filtered_history is None:
        filtered_history = _get_filtered_queue_history(app)
    shown = len(filtered_history)
    total = len(full_history)
    failed = sum(1 for item in filtered_history if item.get("status") == "failed")
    try:
        summary_var.set(f"显示 {shown}/{total} 条 · 当前失败 {failed}")
    except Exception:
        pass


def _build_task_history_detail_text(entry):
    item = dict(entry or {})
    task_result = item.get("task_result") if isinstance(item.get("task_result"), dict) else {}
    lines = [
        f"标题：{item.get('title', '')}",
        f"功能：{_get_feature_label(item.get('task_type'), fallback='未知任务')}",
        f"状态：{_queue_status_text(item.get('status'))}",
        f"输入：{item.get('input', '')}",
        f"创建时间：{_format_queue_time(item.get('created_at'))}",
        f"结束时间：{_format_queue_time(item.get('finished_at') or item.get('created_at'))}",
    ]
    duration_seconds = task_result.get("duration_seconds", item.get("duration_seconds", 0.0))
    try:
        duration_seconds = float(duration_seconds or 0.0)
    except Exception:
        duration_seconds = 0.0
    if duration_seconds > 0:
        lines.append(f"耗时：{duration_seconds:.3f}s")
    detail = task_result.get("detail") or item.get("detail", "")
    if detail:
        lines.append(f"详情：{detail}")
    error = task_result.get("error") or item.get("error", "")
    if error:
        lines.append(f"错误：{error}")
    output_root = task_result.get("output_root") or item.get("output_root", "")
    output_strategy_label = task_result.get("output_strategy_label") or item.get("output_strategy_label", "")
    output_strategy = task_result.get("output_strategy") or item.get("output_strategy", "")
    if output_strategy_label:
        lines.append(f"输出策略：{output_strategy_label}")
    elif output_strategy:
        lines.append(f"输出策略：{output_strategy}")
    if output_root:
        lines.append(f"输出目录：{output_root}")
    outputs = list(task_result.get("outputs") or item.get("outputs") or [])
    if outputs:
        lines.append("")
        lines.append("输出文件：")
        lines.extend(f"- {path}" for path in outputs)
    failed_items = list(task_result.get("failed_items") or item.get("failed_items") or [])
    logs = [str(text).strip() for text in (item.get("logs") or []) if str(text).strip()]
    if item.get("status") == "failed" or error or failed_items:
        lines.append("")
        lines.append("失败概览：")
        lines.append(f"- 状态：{_queue_status_text(item.get('status'))}")
        if error:
            lines.append(f"- 失败原因：{error}")
        elif detail:
            lines.append(f"- 失败原因：{detail}")
        if failed_items:
            lines.append(f"- 失败项数量：{len(failed_items)}")
        failed_log_lines = [
            text for text in logs if any(marker in text.lower() for marker in ("❌", "🔥", "错误", "失败", "error", "failed", "traceback"))
        ]
        if failed_log_lines:
            lines.append(f"- 关键失败日志：{len(failed_log_lines)} 条")
    if error:
        lines.append("")
        lines.append("失败原因：")
        lines.append(f"- {error}")
    elif item.get("status") == "failed" and detail:
        lines.append("")
        lines.append("失败原因：")
        lines.append(f"- {detail}")
    if failed_items:
        lines.append("")
        lines.append("失败项：")
        lines.extend(f"- {path}" for path in failed_items)
    failed_log_lines = [
        text for text in logs if any(marker in text.lower() for marker in ("❌", "🔥", "错误", "失败", "error", "failed", "traceback"))
    ]
    if failed_log_lines:
        lines.append("")
        lines.append("关键日志：")
        lines.extend(f"- {text}" for text in failed_log_lines[-8:])
    if logs:
        lines.append("")
        lines.append("日志片段：")
        lines.extend(f"- {text}" for text in logs[-12:])
    lines.append("")
    lines.append("结构化结果 JSON：")
    lines.append(json.dumps(task_result, ensure_ascii=False, indent=2) if task_result else "{}")
    return "\n".join(lines)


def _apply_task_history_detail_highlights(detail_box, detail_text, entry):
    if detail_box is None:
        return
    try:
        detail_box.tag_remove("fx_history_fail_header", "1.0", "end")
        detail_box.tag_remove("fx_history_fail_text", "1.0", "end")
        detail_box.tag_remove("fx_history_fail_item", "1.0", "end")
        detail_box.tag_remove("fx_history_log_error", "1.0", "end")
        detail_box.tag_config("fx_history_fail_header", foreground="#F6C66A")
        detail_box.tag_config("fx_history_fail_text", foreground="#FF8A80")
        detail_box.tag_config("fx_history_fail_item", foreground="#FFD6A5")
        detail_box.tag_config("fx_history_log_error", foreground="#FFB4AB")
    except Exception:
        return

    item = dict(entry or {})
    task_result = item.get("task_result") if isinstance(item.get("task_result"), dict) else {}
    failed_items = list(task_result.get("failed_items") or item.get("failed_items") or [])
    status = str(item.get("status") or "")
    error = str(task_result.get("error") or item.get("error") or "").strip()
    detail = str(task_result.get("detail") or item.get("detail") or "").strip()
    logs = [str(text).strip() for text in (item.get("logs") or []) if str(text).strip()]
    highlight_lines = []
    if status == "failed":
        highlight_lines.append(("状态：失败", "fx_history_fail_text"))
    if error:
        highlight_lines.append((f"错误：{error}", "fx_history_fail_text"))
        highlight_lines.append((f"- {error}", "fx_history_fail_text"))
    elif status == "failed" and detail:
        highlight_lines.append((f"- {detail}", "fx_history_fail_text"))
    for header in ("失败概览：", "失败原因：", "失败项：", "关键日志："):
        highlight_lines.append((header, "fx_history_fail_header"))
    for item_path in failed_items:
        highlight_lines.append((f"- {item_path}", "fx_history_fail_item"))
    for log_line in logs:
        lowered = log_line.lower()
        if any(marker in lowered for marker in ("❌", "🔥", "错误", "失败", "error", "failed", "traceback")):
            highlight_lines.append((f"- {log_line}", "fx_history_log_error"))

    for text, tag in highlight_lines:
        if not text:
            continue
        start_index = "1.0"
        while True:
            try:
                match_index = detail_box.search(text, start_index, stopindex="end")
            except Exception:
                match_index = ""
            if not match_index:
                break
            end_index = f"{match_index}+{len(text)}c"
            try:
                detail_box.tag_add(tag, match_index, end_index)
            except Exception:
                break
            start_index = end_index


def _get_task_history_export_context():
    return TaskHistoryExportContext(
        normalize_path=_normalize_input_path_value,
        export_task_result=_export_task_result,
        sanitize_filename_component=_sanitize_filename_component,
        format_queue_time=_format_queue_time,
        get_feature_label=_get_feature_label,
        queue_status_text=_queue_status_text,
        classify_failure_reason=_classify_failure_reason,
        failure_value_to_label=QUEUE_HISTORY_FAILURE_VALUE_TO_LABEL,
        project_root=Path(__file__).resolve().parent,
        user_home=Path.home(),
        probe_environment=_probe_diagnostic_environment,
        load_queue_history=_load_queue_history,
        debug=_debug,
    )


def _build_task_history_export_filename(entry):
    return build_task_history_export_filename(entry, _get_task_history_export_context())


def _export_task_history_entry(entry, output_path):
    return export_task_history_entry(entry, output_path, _get_task_history_export_context())


def _prompt_export_task_history_entry(app, entry, output_path=None):
    item = dict(entry or {})
    task_result = item.get("task_result") if isinstance(item.get("task_result"), dict) else {}
    if not task_result:
        try:
            tkinter.messagebox.showwarning("导出结果", "当前历史记录没有可导出的结构化结果。", parent=app)
        except Exception:
            pass
        return False
    selected_path = output_path
    if not selected_path:
        initial_dir = task_result.get("output_root") or os.path.dirname(_normalize_input_path_value(item.get("input"))) or str(Path.cwd())
        try:
            selected_path = tkinter.filedialog.asksaveasfilename(
                title="导出任务结果",
                parent=app,
                defaultextension=".json",
                initialdir=_normalize_input_path_value(initial_dir) or str(Path.cwd()),
                initialfile=_build_task_history_export_filename(item),
                filetypes=[("JSON 文件", "*.json"), ("所有文件", "*.*")],
            )
        except Exception as exc:
            _debug(f"queue:history_export_dialog_error:{exc}")
            selected_path = ""
    ok, payload = _export_task_history_entry(item, selected_path)
    try:
        if ok:
            tkinter.messagebox.showinfo("导出结果", f"任务结果已导出到：\n{payload}", parent=app)
        elif selected_path:
            tkinter.messagebox.showerror("导出结果", payload, parent=app)
    except Exception:
        pass
    if ok:
        try:
            app.log(f"[任务历史] 已导出结果：{payload}")
        except Exception:
            pass
    return ok


def _build_task_history_log_export_text(entry):
    return build_task_history_log_export_text(entry, _get_task_history_export_context())


def _build_task_history_log_export_filename(entry):
    return build_task_history_log_export_filename(entry, _get_task_history_export_context())


def _export_task_history_log(entry, output_path):
    return export_task_history_log(entry, output_path, _get_task_history_export_context())


def _prompt_export_task_history_log(app, entry, output_path=None):
    item = dict(entry or {})
    selected_path = output_path
    if not selected_path:
        initial_dir = _normalize_input_path_value(item.get("input")) or str(Path.cwd())
        try:
            selected_path = tkinter.filedialog.asksaveasfilename(
                title="导出任务日志",
                parent=app,
                defaultextension=".txt",
                initialdir=initial_dir,
                initialfile=_build_task_history_log_export_filename(item),
                filetypes=[("文本文件", "*.txt"), ("所有文件", "*.*")],
            )
        except Exception as exc:
            _debug(f"queue:history_log_export_dialog_error:{exc}")
            selected_path = ""
    ok, payload = _export_task_history_log(item, selected_path)
    try:
        if ok:
            tkinter.messagebox.showinfo("导出日志", f"任务日志已导出到：\n{payload}", parent=app)
        elif selected_path:
            tkinter.messagebox.showerror("导出日志", payload, parent=app)
    except Exception:
        pass
    if ok:
        try:
            app.log(f"[任务历史] 已导出日志：{payload}")
        except Exception:
            pass
    return ok


def _build_task_history_report_export_filename(entry):
    return build_task_history_report_export_filename(entry, _get_task_history_export_context())


def _build_task_history_report_text(entry):
    return build_task_history_report_text(entry, _get_task_history_export_context())


def _export_task_history_report(entry, output_path):
    return export_task_history_report(entry, output_path, _get_task_history_export_context())


def _build_task_history_diagnostic_filename(entry):
    return build_task_history_diagnostic_filename(entry, _get_task_history_export_context())


def _diagnostic_path_replacements():
    return diagnostic_path_replacements(_get_task_history_export_context())


def _redact_diagnostic_text(value):
    return redact_diagnostic_text(value, _get_task_history_export_context())


def _redact_diagnostic_payload(value):
    return redact_diagnostic_payload(value, _get_task_history_export_context())


def _diagnostic_write_json(zip_file, arcname, payload):
    return diagnostic_write_json(zip_file, arcname, payload, _get_task_history_export_context())


def _diagnostic_write_text(zip_file, arcname, text):
    return diagnostic_write_text(zip_file, arcname, text, _get_task_history_export_context())


def _probe_diagnostic_environment():
    ffmpeg_path = _locate_ffmpeg()
    ocr_status = {}
    try:
        from tools.fx_pdf_ocr import discover_backend_status

        ocr_status = discover_backend_status(detailed=True)
    except Exception as exc:
        ocr_status = {"error": str(exc)}

    office_status = {}
    for label, progid in (("word", "Word.Application"), ("powerpoint", "PowerPoint.Application")):
        status = {"available": False, "version": "", "error": ""}
        try:
            pythoncom.CoInitialize()
            try:
                app = win32com.client.DispatchEx(progid)
                status["available"] = True
                status["version"] = str(getattr(app, "Version", "") or "")
                try:
                    app.Quit()
                except Exception:
                    pass
            finally:
                pythoncom.CoUninitialize()
        except Exception as exc:
            status["error"] = str(exc)
        office_status[label] = status

    return {
        "app": {
            "release_version": APP_RELEASE_VERSION,
            "display_version": APP_DISPLAY_VERSION,
            "frozen": bool(getattr(sys, "frozen", False)),
        },
        "system": {
            "platform": platform.platform(),
            "python": sys.version,
            "executable": sys.executable,
            "cwd": str(Path.cwd()),
        },
        "dependencies": {
            "ffmpeg": {"available": bool(ffmpeg_path), "path": ffmpeg_path or ""},
            "ocr": ocr_status,
            "office": office_status,
        },
        "performance": {
            "log_path": str(_get_performance_log_file()),
            "recent": _load_recent_performance_entries(limit=40),
        },
    }


def _build_diagnostic_summary(entry, environment):
    return build_diagnostic_summary(entry, environment, _get_task_history_export_context())


def _build_recent_history_diagnostic_snapshot(entry, limit=12):
    return build_recent_history_diagnostic_snapshot(entry, _get_task_history_export_context(), limit=limit)


def _export_task_history_diagnostic_package(entry, output_path):
    return export_task_history_diagnostic_package(entry, output_path, _get_task_history_export_context())


def _prompt_export_task_history_diagnostic_package(app, entry, output_path=None):
    item = dict(entry or {})
    selected_path = output_path
    if not selected_path:
        initial_dir = _resolve_task_history_open_target(item) or str(Path.cwd())
        try:
            selected_path = tkinter.filedialog.asksaveasfilename(
                title="导出诊断包",
                parent=app,
                defaultextension=".zip",
                initialdir=_normalize_input_path_value(initial_dir) or str(Path.cwd()),
                initialfile=_build_task_history_diagnostic_filename(item),
                filetypes=[("ZIP 诊断包", "*.zip"), ("所有文件", "*.*")],
            )
        except Exception as exc:
            _debug(f"queue:diagnostic_export_dialog_error:{exc}")
            selected_path = ""
    ok, payload = _export_task_history_diagnostic_package(item, selected_path)
    try:
        if ok:
            tkinter.messagebox.showinfo("导出诊断包", f"诊断包已导出到：\n{payload}", parent=app)
        elif selected_path:
            tkinter.messagebox.showerror("导出诊断包", payload, parent=app)
    except Exception:
        pass
    if ok:
        try:
            app.log(f"[任务历史] 已导出诊断包：{payload}")
        except Exception:
            pass
    return ok


def _prompt_export_task_history_report(app, entry, output_path=None):
    item = dict(entry or {})
    selected_path = output_path
    if not selected_path:
        initial_dir = _resolve_task_history_open_target(item) or str(Path.cwd())
        try:
            selected_path = tkinter.filedialog.asksaveasfilename(
                title="导出任务报告",
                parent=app,
                defaultextension=".md",
                initialdir=_normalize_input_path_value(initial_dir) or str(Path.cwd()),
                initialfile=_build_task_history_report_export_filename(item),
                filetypes=[("Markdown 文件", "*.md"), ("文本文件", "*.txt"), ("所有文件", "*.*")],
            )
        except Exception as exc:
            _debug(f"queue:history_report_export_dialog_error:{exc}")
            selected_path = ""
    ok, payload = _export_task_history_report(item, selected_path)
    try:
        if ok:
            tkinter.messagebox.showinfo("导出报告", f"任务报告已导出到：\n{payload}", parent=app)
        elif selected_path:
            tkinter.messagebox.showerror("导出报告", payload, parent=app)
    except Exception:
        pass
    if ok:
        try:
            app.log(f"[任务历史] 已导出报告：{payload}")
        except Exception:
            pass
    return ok


def _resolve_task_history_open_target(entry):
    item = dict(entry or {})
    task_result = item.get("task_result") if isinstance(item.get("task_result"), dict) else {}
    candidates = [
        task_result.get("output_root"),
        item.get("output_root"),
    ]
    outputs = list(task_result.get("outputs") or item.get("outputs") or [])
    if outputs:
        candidates.extend(outputs)
    candidates.append(item.get("input"))
    for value in candidates:
        normalized = _normalize_input_path_value(value)
        if not normalized:
            continue
        if os.path.isdir(normalized):
            return normalized
        if os.path.isfile(normalized):
            return os.path.dirname(normalized)
    return ""


def _open_task_history_output(entry):
    target = _resolve_task_history_open_target(entry)
    if not target:
        return False, "当前历史记录没有可打开的输出位置。"
    try:
        os.startfile(target)
        return True, target
    except Exception as exc:
        _debug(f"queue:history_open_output_error:{exc}")
        return False, f"打开失败：{target}"


def _prompt_open_task_history_output(app, entry):
    ok, payload = _open_task_history_output(entry)
    try:
        if (not ok) and payload:
            tkinter.messagebox.showwarning("打开输出位置", payload, parent=app)
    except Exception:
        pass
    if ok:
        try:
            app.log(f"[任务历史] 已打开输出位置：{payload}")
        except Exception:
            pass
    return ok


def _show_task_history_detail(app, entry):
    _ensure_queue_state(app)
    detail_text = _build_task_history_detail_text(entry)
    detail_window = getattr(app, "_fx_history_detail_window", None)
    try:
        if detail_window is not None and detail_window.winfo_exists():
            detail_window.deiconify()
            detail_window.lift()
        else:
            detail_window = customtkinter.CTkToplevel(app)
            _apply_window_icon(detail_window)
            app._fx_history_detail_window = detail_window
            detail_window.title("任务历史详情")
            detail_window.geometry("820x620")
            detail_window.minsize(680, 480)
            try:
                detail_window.configure(fg_color=globals().get("COLOR_CARD_ALT", "#303030"))
            except Exception:
                pass
            detail_window.grid_columnconfigure(0, weight=1)
            detail_window.grid_rowconfigure(1, weight=1)
            customtkinter.CTkLabel(
                detail_window,
                text="任务历史详情",
                text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
                font=customtkinter.CTkFont(size=20, weight="bold"),
                anchor="w",
            ).grid(row=0, column=0, sticky="ew", padx=18, pady=(16, 8))
            detail_box = customtkinter.CTkTextbox(
                detail_window,
                wrap="word",
                fg_color=globals().get("COLOR_CARD", "#2B2B2B"),
                text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
                border_width=1,
                border_color=globals().get("COLOR_BORDER", "#3A3A3A"),
                corner_radius=12,
            )
            detail_box.grid(row=1, column=0, sticky="nsew", padx=18, pady=(0, 12))
            detail_window._fx_detail_box = detail_box
            actions = customtkinter.CTkFrame(detail_window, fg_color="transparent")
            actions.grid(row=2, column=0, sticky="ew", padx=18, pady=(0, 16))
            actions.grid_columnconfigure(0, weight=1)

            def copy_detail():
                try:
                    app.clipboard_clear()
                    app.clipboard_append(detail_window._fx_detail_box.get("1.0", "end-1c"))
                    app.update()
                except Exception as exc:
                    _debug(f"queue:history_detail_copy_error:{exc}")

            def export_detail_result():
                current_entry = getattr(detail_window, "_fx_entry", None)
                _prompt_export_task_history_entry(app, current_entry)

            def export_detail_log():
                current_entry = getattr(detail_window, "_fx_entry", None)
                _prompt_export_task_history_log(app, current_entry)

            def export_detail_report():
                current_entry = getattr(detail_window, "_fx_entry", None)
                _prompt_export_task_history_report(app, current_entry)

            def export_detail_diagnostic():
                current_entry = getattr(detail_window, "_fx_entry", None)
                _prompt_export_task_history_diagnostic_package(app, current_entry)

            def open_output_location():
                current_entry = getattr(detail_window, "_fx_entry", None)
                _prompt_open_task_history_output(app, current_entry)

            customtkinter.CTkButton(
                actions,
                text="导出结果",
                command=export_detail_result,
                height=34,
                width=92,
                corner_radius=10,
                fg_color="#56603D",
                hover_color="#68714A",
                text_color="#F5F8EA",
            ).grid(row=0, column=1, sticky="e", padx=(0, 8))
            customtkinter.CTkButton(
                actions,
                text="导出报告",
                command=export_detail_report,
                height=34,
                width=92,
                corner_radius=10,
                fg_color="#5A4768",
                hover_color="#69557B",
                text_color="#F4EEFF",
            ).grid(row=0, column=2, sticky="e", padx=(0, 8))
            customtkinter.CTkButton(
                actions,
                text="诊断包",
                command=export_detail_diagnostic,
                height=34,
                width=92,
                corner_radius=10,
                fg_color="#68474A",
                hover_color="#7A5558",
                text_color="#FFF0F1",
            ).grid(row=0, column=3, sticky="e", padx=(0, 8))
            customtkinter.CTkButton(
                actions,
                text="打开位置",
                command=open_output_location,
                height=34,
                width=92,
                corner_radius=10,
                fg_color="#3F5B57",
                hover_color="#4D6C67",
                text_color="#EAF6F3",
            ).grid(row=0, column=4, sticky="e", padx=(0, 8))
            customtkinter.CTkButton(
                actions,
                text="导出日志",
                command=export_detail_log,
                height=34,
                width=92,
                corner_radius=10,
                fg_color="#5A4E3D",
                hover_color="#6A5B49",
                text_color="#F8F1E6",
            ).grid(row=0, column=5, sticky="e", padx=(0, 8))

            customtkinter.CTkButton(
                actions,
                text="复制详情",
                command=copy_detail,
                height=34,
                width=92,
                corner_radius=10,
                fg_color="#44566C",
                hover_color="#51657D",
                text_color="#EEF5FF",
            ).grid(row=0, column=6, sticky="e")
        detail_box = getattr(detail_window, "_fx_detail_box", None)
        detail_window._fx_entry = dict(entry or {})
        if detail_box is not None:
            detail_box.delete("1.0", "end")
            detail_box.insert("1.0", detail_text)
            _apply_task_history_detail_highlights(detail_box, detail_text, detail_window._fx_entry)
            detail_box.see("1.0")
    except Exception as exc:
        _debug(f"queue:history_detail_error:{exc}")


def _reset_queue_history_filters(app):
    try:
        getattr(app, "_fx_history_filter_status_var", None).set(QUEUE_HISTORY_STATUS_DEFAULT)
    except Exception:
        pass
    try:
        getattr(app, "_fx_history_filter_task_var", None).set(QUEUE_HISTORY_TASK_DEFAULT)
    except Exception:
        pass
    try:
        getattr(app, "_fx_history_filter_failure_var", None).set(QUEUE_HISTORY_FAILURE_DEFAULT)
    except Exception:
        pass
    try:
        getattr(app, "_fx_history_search_var", None).set("")
    except Exception:
        pass
    _refresh_history_panel(app)


def _install_queue_history_filter_hooks(app):
    if getattr(app, "_fx_history_filter_hooks_ready", False):
        return

    def on_change(*_args):
        _refresh_history_panel(app)

    app._fx_history_filter_trace_callback = on_change
    for attr in ("_fx_history_filter_status_var", "_fx_history_filter_task_var", "_fx_history_filter_failure_var", "_fx_history_search_var"):
        var = getattr(app, attr, None)
        if isinstance(var, tkinter.Variable):
            try:
                var.trace_add("write", on_change)
            except Exception:
                pass
    app._fx_history_filter_hooks_ready = True


def _task_result_matches_task(task_result, task):
    if not isinstance(task_result, dict) or not isinstance(task, dict):
        return False
    task_input = _normalize_input_path_value(task.get("input"))
    result_input = _normalize_input_path_value(task_result.get("input"))
    task_type = str(task.get("task_type") or "")
    result_type = str(task_result.get("task_type") or "")
    if task_input and result_input and os.path.normcase(task_input) != os.path.normcase(result_input):
        return False
    if task_type and result_type and task_type != result_type:
        return False
    return True


def _queue_set_task_status(app, task, status, detail=""):
    task_result = task.get("task_result")
    if not _task_result_matches_task(task_result, task):
        task_result = _get_last_task_result(app)
    if not _task_result_matches_task(task_result, task):
        task_result = None
    if task_result is None and str(status or "") in {"success", "failed", "stopped", "skipped"}:
        task_result = _new_task_result(task.get("input"), task.get("task_type"))
        _set_task_result_finished(
            task_result,
            status=status,
            message=str(detail or status),
            detail=str(detail or status),
            error=str(detail or "") if str(status or "") == "failed" else "",
            stopped=str(status or "") == "stopped",
            skipped=str(status or "") == "skipped",
        )
    if isinstance(task_result, dict):
        task["task_result"] = dict(task_result)
        detail = task_result.get("detail") or detail
        task["output_strategy"] = task_result.get("output_strategy", "")
        task["output_strategy_label"] = task_result.get("output_strategy_label", "")
        task["output_root"] = task_result.get("output_root", "")
        task["outputs"] = list(task_result.get("outputs") or [])
        task["error"] = task_result.get("error", "")
        task["duration_seconds"] = task_result.get("duration_seconds", 0.0)
        task["skipped"] = bool(task_result.get("skipped", False))
        status = task_result.get("status") or status
    task["status"] = status
    task["finished_at"] = time.time() if status in {"success", "failed", "stopped", "skipped"} else task.get("finished_at")
    if detail:
        task["detail"] = str(detail)
    status_var = task.get("status_var")
    if status_var is not None:
        try:
            status_var.set(_queue_status_text(status))
        except Exception:
            pass
    row = task.get("row")
    if row is not None:
        try:
            for child in row.winfo_children():
                if isinstance(child, customtkinter.CTkLabel) and getattr(child, "_fx_queue_status_label", False):
                    child.configure(text_color=_queue_status_color(status))
        except Exception:
            pass
    _refresh_queue_panel(app)


def _queue_task_had_errors(task):
    task_result = task.get("task_result")
    if isinstance(task_result, dict):
        if task_result.get("status") in {"failed"} or task_result.get("error"):
            return True
        if task_result.get("failed_items"):
            return True
    if task.get("exception"):
        return True
    logs = task.get("logs") or []
    if getattr(task.get("app", None), "stop_event", False):
        return True
    joined = "\n".join(str(item) for item in logs).lower()
    return any(marker.lower() in joined for marker in QUEUE_ERROR_MARKERS)


def _queue_capture_task_logs(app, task):
    original_log = getattr(app, "log", None)
    if not callable(original_log):
        return None

    def patched_log(message, *args, **kwargs):
        text = str(message)
        try:
            task.setdefault("logs", []).append(text)
        except Exception:
            pass
        task_result = _get_last_task_result(app)
        if _task_result_matches_task(task_result, task):
            try:
                task_result.setdefault("logs", []).append(text)
            except Exception:
                pass
            try:
                task_result.setdefault("log_lines", []).append(text)
            except Exception:
                pass
        return original_log(message, *args, **kwargs)

    app.log = patched_log
    return original_log


def _queue_restore_task_logs(app, original_log):
    if original_log is not None:
        try:
            app.log = original_log
        except Exception:
            pass


def _split_failed_item_path(item):
    text = str(item or "").strip()
    if not text:
        return ""
    if ": " in text:
        head, _tail = text.split(": ", 1)
        if head and os.path.exists(head):
            return head
    return text


def _resolve_retry_failed_item_paths(source):
    if not isinstance(source, dict):
        return []
    source_input = _normalize_input_path_value(source.get("input"))
    if not source_input:
        return []
    input_root = os.path.dirname(source_input) if os.path.isfile(source_input) else source_input
    task_result = source.get("task_result") if isinstance(source.get("task_result"), dict) else {}
    failed_items = list(task_result.get("failed_items") or source.get("failed_items") or [])
    resolved = []
    seen = set()
    for item in failed_items:
        path_value = _split_failed_item_path(item)
        if not path_value:
            continue
        normalized = _normalize_input_path_value(path_value)
        if normalized and os.path.exists(normalized):
            candidate = normalized
        else:
            candidate = _normalize_input_path_value(os.path.join(input_root, path_value)) if input_root else ""
        if not candidate or not os.path.exists(candidate):
            continue
        key = os.path.normcase(candidate)
        if key in seen:
            continue
        seen.add(key)
        resolved.append(candidate)
    return resolved


def _build_retry_subset_input(app, source):
    if not isinstance(source, dict):
        return None
    failed_paths = _resolve_retry_failed_item_paths(source)
    if not failed_paths:
        return None
    source_input = _normalize_input_path_value(source.get("input"))
    if len(failed_paths) == 1 and os.path.isfile(failed_paths[0]):
        return {
            "mode": "single_file",
            "input": failed_paths[0],
            "failed_items": failed_paths,
            "cleanup": None,
        }
    task_type = str(source.get("task_type") or "")
    if task_type not in {"pdf", "image", "remove_wm"}:
        return None
    try:
        staging_root = Path(tempfile.mkdtemp(prefix="fx_retry_subset_"))
        for src in failed_paths:
            if os.path.isfile(source_input):
                relative = os.path.basename(src)
            else:
                try:
                    relative = os.path.relpath(src, source_input)
                except Exception:
                    relative = os.path.basename(src)
            dst = staging_root / relative
            dst.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(src, dst)
        return {
            "mode": "staging_dir",
            "input": str(staging_root),
            "failed_items": failed_paths,
            "cleanup": str(staging_root),
        }
    except Exception as exc:
        _debug(f"queue:retry_subset_build_error:{exc}")
        try:
            shutil.rmtree(staging_root, ignore_errors=True)
        except Exception:
            pass
        return None


def _queue_build_task(app, input_path=None, task_type=None, retry_source=None):
    normalized_input = _normalize_input_path_value(input_path if input_path is not None else app.input_path.get())
    task_type = task_type or getattr(app, "current_task", "")
    if not normalized_input:
        raise ValueError("请先选择文件或文件夹")
    if not os.path.exists(normalized_input):
        raise ValueError(f"路径不存在: {normalized_input}")
    if task_type == "help":
        raise ValueError("使用教程不能加入任务队列")
    snapshot = retry_source.get("snapshot") if isinstance(retry_source, dict) else None
    if not isinstance(snapshot, dict):
        snapshot = _queue_snapshot_app_state(app, task_type)
    task = {
        "id": _make_queue_task_id(),
        "input": normalized_input,
        "task_type": task_type,
        "title": retry_source.get("title") if isinstance(retry_source, dict) else _queue_describe_task(app, task_type, normalized_input),
        "snapshot": snapshot,
        "status": "queued",
        "created_at": time.time(),
        "started_at": None,
        "finished_at": None,
        "detail": "",
        "logs": [],
    }
    if isinstance(retry_source, dict):
        task["retry_source_id"] = retry_source.get("id")
    return task


def _ensure_queue_state(app):
    if not hasattr(app, "_fx_task_queue"):
        app._fx_task_queue = []
    if not hasattr(app, "_fx_task_history"):
        app._fx_task_history = _load_queue_history()
    if not hasattr(app, "_fx_task_queue_running"):
        app._fx_task_queue_running = False
    if not hasattr(app, "_fx_queue_window"):
        app._fx_queue_window = None
    if not hasattr(app, "_fx_history_filter_status_var"):
        app._fx_history_filter_status_var = tkinter.StringVar(master=app, value=QUEUE_HISTORY_STATUS_DEFAULT)
    if not hasattr(app, "_fx_history_filter_task_var"):
        app._fx_history_filter_task_var = tkinter.StringVar(master=app, value=QUEUE_HISTORY_TASK_DEFAULT)
    if not hasattr(app, "_fx_history_filter_failure_var"):
        app._fx_history_filter_failure_var = tkinter.StringVar(master=app, value=QUEUE_HISTORY_FAILURE_DEFAULT)
    if not hasattr(app, "_fx_history_search_var"):
        app._fx_history_search_var = tkinter.StringVar(master=app, value="")
    if not hasattr(app, "_fx_history_summary_var"):
        app._fx_history_summary_var = tkinter.StringVar(master=app, value="")


def _refresh_queue_status_summary(app):
    summary_var = getattr(app, "_fx_queue_summary_var", None)
    if summary_var is None:
        return
    try:
        queue = getattr(app, "_fx_task_queue", [])
        running = sum(1 for item in queue if item.get("status") == "running")
        queued = sum(1 for item in queue if item.get("status") == "queued")
        failed = sum(1 for item in getattr(app, "_fx_task_history", []) if item.get("status") == "failed")
        if running:
            summary = f"队列执行中 · 等待 {queued} · 可重试失败 {failed}"
        elif queued:
            summary = f"队列待执行 {queued} 项 · 可重试失败 {failed}"
        else:
            summary = f"队列空闲 · 可重试失败 {failed}"
        summary_var.set(summary)
    except Exception:
        pass


def _refresh_queue_panel(app):
    _refresh_queue_status_summary(app)
    panel = getattr(app, "_fx_queue_list_frame", None)
    if panel is None:
        return
    try:
        for child in panel.winfo_children():
            child.destroy()
        queue = getattr(app, "_fx_task_queue", [])
        if not queue:
            customtkinter.CTkLabel(
                panel,
                text="当前队列为空。配置好功能后点“加入队列”，就可以批量排队执行。",
                text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
                font=customtkinter.CTkFont(size=12),
                justify="left",
            ).pack(anchor="w", padx=16, pady=14)
            return
        for index, task in enumerate(queue, 1):
            row = customtkinter.CTkFrame(panel, fg_color=globals().get("COLOR_CARD_ALT", "#303030"), corner_radius=10)
            row.pack(fill="x", padx=10, pady=(8 if index == 1 else 4, 4))
            task["row"] = row
            row.grid_columnconfigure(1, weight=1)
            customtkinter.CTkLabel(
                row,
                text=f"{index:02d}",
                text_color="#F0D39A",
                font=customtkinter.CTkFont(size=12, weight="bold"),
                width=34,
            ).grid(row=0, column=0, rowspan=2, padx=(10, 6), pady=8, sticky="n")
            customtkinter.CTkLabel(
                row,
                text=task.get("title", ""),
                text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
                font=customtkinter.CTkFont(size=12, weight="bold"),
                anchor="w",
            ).grid(row=0, column=1, padx=4, pady=(8, 1), sticky="ew")
            customtkinter.CTkLabel(
                row,
                text=task.get("input", ""),
                text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
                font=customtkinter.CTkFont(size=10),
                anchor="w",
            ).grid(row=1, column=1, padx=4, pady=(0, 8), sticky="ew")
            status_label = customtkinter.CTkLabel(
                row,
                text=_queue_status_text(task.get("status")),
                text_color=_queue_status_color(task.get("status")),
                font=customtkinter.CTkFont(size=12, weight="bold"),
                width=64,
            )
            status_label._fx_queue_status_label = True
            status_label.grid(row=0, column=2, rowspan=2, padx=(6, 10), pady=8, sticky="e")
    except Exception as exc:
        _debug(f"queue:refresh_panel_error:{exc}")


def _refresh_history_panel(app):
    panel = getattr(app, "_fx_history_list_frame", None)
    history = list(getattr(app, "_fx_task_history", []) or [])
    filtered_history = _get_filtered_queue_history(app)
    _refresh_history_summary(app, filtered_history=filtered_history, full_history=history)
    if panel is None:
        return
    try:
        for child in panel.winfo_children():
            child.destroy()
        if not history:
            customtkinter.CTkLabel(
                panel,
                text="暂无历史记录。完成队列任务后会自动记录到这里。",
                text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
                font=customtkinter.CTkFont(size=12),
                justify="left",
            ).pack(anchor="w", padx=16, pady=14)
            return
        if not filtered_history:
            customtkinter.CTkLabel(
                panel,
                text="没有符合筛选条件的历史记录。可以调整状态、功能类型或搜索关键词。",
                text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
                font=customtkinter.CTkFont(size=12),
                justify="left",
            ).pack(anchor="w", padx=16, pady=14)
            return
        for item in reversed(filtered_history[-QUEUE_HISTORY_LIMIT:]):
            row = customtkinter.CTkFrame(panel, fg_color=globals().get("COLOR_CARD_ALT", "#303030"), corner_radius=10)
            row.pack(fill="x", padx=10, pady=(6, 4))
            row.grid_columnconfigure(0, weight=1)
            top = customtkinter.CTkFrame(row, fg_color="transparent")
            top.grid(row=0, column=0, sticky="ew", padx=10, pady=(8, 2))
            top.grid_columnconfigure(0, weight=1)
            customtkinter.CTkLabel(
                top,
                text=item.get("title", ""),
                text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
                font=customtkinter.CTkFont(size=12, weight="bold"),
                anchor="w",
            ).grid(row=0, column=0, sticky="ew")
            customtkinter.CTkLabel(
                top,
                text=_queue_status_text(item.get("status")),
                text_color=_queue_status_color(item.get("status")),
                font=customtkinter.CTkFont(size=12, weight="bold"),
                width=58,
            ).grid(row=0, column=1, padx=(8, 0), sticky="e")
            task_result = item.get("task_result") if isinstance(item.get("task_result"), dict) else {}
            stats_parts = []
            processed_count = int(task_result.get("processed_count") or 0)
            success_count = int(task_result.get("success_count") or 0)
            failed_count = int(task_result.get("failed_count") or 0)
            skipped_count = int(task_result.get("skipped_count") or 0)
            duration_seconds = float(task_result.get("duration_seconds") or item.get("duration_seconds") or 0.0)
            if processed_count:
                stats_parts.append(f"结果 {success_count}/{processed_count}")
            if failed_count:
                stats_parts.append(f"失败 {failed_count}")
            if skipped_count:
                stats_parts.append(f"跳过 {skipped_count}")
            if duration_seconds > 0:
                stats_parts.append(f"耗时 {duration_seconds:.2f}s")
            detail_lines = [
                item.get("input", ""),
                f"功能：{_get_feature_label(item.get('task_type'), fallback='未知任务')} · 时间：{_format_queue_time(item.get('finished_at') or item.get('created_at'))}",
            ]
            if stats_parts:
                detail_lines.append(" · ".join(stats_parts))
            output_root = task_result.get("output_root") or item.get("output_root", "")
            if output_root:
                detail_lines.append(f"输出：{output_root}")
            error_text = task_result.get("error") or item.get("error", "")
            if error_text:
                detail_lines.append(str(error_text)[:160])
            failure_kind, failure_reason = _classify_failure_reason(item)
            if item.get("status") == "failed" and failure_kind and failure_kind != "unknown":
                failure_labels = {
                    "path_missing": "路径缺失",
                    "permission": "权限问题",
                    "timeout": "超时",
                    "dependency": "依赖问题",
                    "partial_failure": "部分失败",
                    "log_failure": "日志失败",
                    "generic_failure": "普通失败",
                }
                detail_lines.append(f"失败分类：{failure_labels.get(failure_kind, failure_kind)}")
                if failure_reason:
                    detail_lines.append(failure_reason[:120])
            if item.get("detail"):
                detail_lines.append(str(item.get("detail"))[:160])
            customtkinter.CTkLabel(
                row,
                text="\n".join(line for line in detail_lines if line),
                text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
                font=customtkinter.CTkFont(size=10),
                justify="left",
                anchor="w",
            ).grid(row=1, column=0, sticky="ew", padx=10, pady=(0, 8))
            is_failed = item.get("status") == "failed"
            action_col = customtkinter.CTkFrame(row, fg_color="transparent")
            action_col.grid(row=0, column=1, rowspan=2, padx=(0, 10), pady=8, sticky="ne")
            customtkinter.CTkButton(
                action_col,
                text="详情",
                command=lambda source=item: _show_task_history_detail(app, source),
                height=30,
                width=72,
                corner_radius=8,
                fg_color="transparent",
                hover_color="#303030",
                border_width=1,
                border_color="#566274",
                text_color="#E8EDF5",
            ).pack(anchor="e", pady=(0, 6))
            customtkinter.CTkButton(
                action_col,
                text="重试" if is_failed else "回放",
                command=(lambda source=item: _queue_retry_history_task(app, source))
                if is_failed
                else (lambda source=item: _queue_replay_history_task(app, source)),
                height=30,
                width=72,
                corner_radius=8,
                fg_color="#6D4F3B" if is_failed else "#44566C",
                hover_color="#7C5B43" if is_failed else "#51657D",
                text_color="#FFF6E6" if is_failed else "#EEF5FF",
            ).pack(anchor="e")
    except Exception as exc:
        _debug(f"queue:refresh_history_error:{exc}")


def _show_task_queue_window(app):
    _ensure_queue_state(app)
    window = getattr(app, "_fx_queue_window", None)
    try:
        if window is not None and window.winfo_exists():
            window.deiconify()
            window.lift()
            _refresh_queue_panel(app)
            _refresh_history_panel(app)
            return window
    except Exception:
        pass

    window = customtkinter.CTkToplevel(app)
    _apply_window_icon(window)
    app._fx_queue_window = window
    window.title("任务队列与历史记录")
    window.geometry("860x620")
    window.minsize(760, 520)
    try:
        window.configure(fg_color=globals().get("COLOR_CARD_ALT", "#303030"))
    except Exception:
        pass
    window.grid_columnconfigure(0, weight=1)
    window.grid_columnconfigure(1, weight=1)
    window.grid_rowconfigure(1, weight=1)

    header = customtkinter.CTkFrame(window, fg_color="transparent")
    header.grid(row=0, column=0, columnspan=2, sticky="ew", padx=18, pady=(16, 8))
    header.grid_columnconfigure(0, weight=1)
    customtkinter.CTkLabel(
        header,
        text="任务队列",
        text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
        font=customtkinter.CTkFont(size=20, weight="bold"),
        anchor="w",
    ).grid(row=0, column=0, sticky="w")
    customtkinter.CTkLabel(
        header,
        textvariable=app._fx_queue_summary_var,
        text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
        font=customtkinter.CTkFont(size=12),
        anchor="e",
    ).grid(row=0, column=1, sticky="e", padx=(16, 0))

    queue_card = customtkinter.CTkFrame(window, fg_color=globals().get("COLOR_CARD", "#2B2B2B"), corner_radius=14, border_width=1, border_color=globals().get("COLOR_BORDER", "#3A3A3A"))
    queue_card.grid(row=1, column=0, sticky="nsew", padx=(18, 8), pady=(0, 16))
    queue_card.grid_rowconfigure(1, weight=1)
    queue_card.grid_columnconfigure(0, weight=1)
    customtkinter.CTkLabel(
        queue_card,
        text="等待执行",
        text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
        font=customtkinter.CTkFont(size=14, weight="bold"),
        anchor="w",
    ).grid(row=0, column=0, sticky="ew", padx=14, pady=(12, 4))
    app._fx_queue_list_frame = customtkinter.CTkScrollableFrame(queue_card, fg_color="transparent")
    app._fx_queue_list_frame.grid(row=1, column=0, sticky="nsew", padx=8, pady=(0, 8))
    queue_actions = customtkinter.CTkFrame(queue_card, fg_color="transparent")
    queue_actions.grid(row=2, column=0, sticky="ew", padx=12, pady=(0, 12))
    queue_actions.grid_columnconfigure(0, weight=1)
    customtkinter.CTkButton(
        queue_actions,
        text="开始队列",
        command=lambda target=app: _start_task_queue(target),
        height=34,
        corner_radius=10,
        fg_color="#7A695B",
        hover_color="#8B7867",
        text_color="#FFFFFF",
    ).grid(row=0, column=0, sticky="ew", padx=(0, 8))
    customtkinter.CTkButton(
        queue_actions,
        text="清空等待",
        command=lambda target=app: _clear_queued_tasks(target),
        height=34,
        width=92,
        corner_radius=10,
        fg_color="transparent",
        hover_color="#303030",
        border_width=1,
        border_color="#566274",
        text_color="#E8EDF5",
    ).grid(row=0, column=1, sticky="e")

    history_card = customtkinter.CTkFrame(window, fg_color=globals().get("COLOR_CARD", "#2B2B2B"), corner_radius=14, border_width=1, border_color=globals().get("COLOR_BORDER", "#3A3A3A"))
    history_card.grid(row=1, column=1, sticky="nsew", padx=(8, 18), pady=(0, 16))
    history_card.grid_rowconfigure(2, weight=1)
    history_card.grid_columnconfigure(0, weight=1)
    history_header = customtkinter.CTkFrame(history_card, fg_color="transparent")
    history_header.grid(row=0, column=0, sticky="ew", padx=14, pady=(12, 4))
    history_header.grid_columnconfigure(0, weight=1)
    customtkinter.CTkLabel(
        history_header,
        text="历史与失败重试",
        text_color=globals().get("COLOR_TEXT", "#E6EEF2"),
        font=customtkinter.CTkFont(size=14, weight="bold"),
        anchor="w",
    ).grid(row=0, column=0, sticky="w")
    customtkinter.CTkLabel(
        history_header,
        textvariable=app._fx_history_summary_var,
        text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
        font=customtkinter.CTkFont(size=11),
        anchor="e",
    ).grid(row=0, column=1, sticky="e", padx=(12, 0))
    history_filters = customtkinter.CTkFrame(history_card, fg_color="transparent")
    history_filters.grid(row=1, column=0, sticky="ew", padx=12, pady=(0, 8))
    history_filters.grid_columnconfigure(3, weight=1)
    customtkinter.CTkOptionMenu(
        history_filters,
        variable=app._fx_history_filter_status_var,
        values=[label for label, _value in QUEUE_HISTORY_STATUS_OPTIONS],
        command=lambda _value, target=app: _refresh_history_panel(target),
        height=32,
        width=110,
    ).grid(row=0, column=0, padx=(0, 8), sticky="w")
    customtkinter.CTkOptionMenu(
        history_filters,
        variable=app._fx_history_filter_task_var,
        values=[label for label, _value in QUEUE_HISTORY_TASK_OPTIONS],
        command=lambda _value, target=app: _refresh_history_panel(target),
        height=32,
        width=124,
    ).grid(row=0, column=1, padx=(0, 8), sticky="w")
    customtkinter.CTkOptionMenu(
        history_filters,
        variable=app._fx_history_filter_failure_var,
        values=[label for label, _value in QUEUE_HISTORY_FAILURE_OPTIONS],
        command=lambda _value, target=app: _refresh_history_panel(target),
        height=32,
        width=118,
    ).grid(row=0, column=2, padx=(0, 8), sticky="w")
    app._fx_history_search_entry = customtkinter.CTkEntry(
        history_filters,
        textvariable=app._fx_history_search_var,
        height=32,
        placeholder_text="搜索路径、错误、失败原因、输出位置...",
    )
    app._fx_history_search_entry.grid(row=0, column=3, sticky="ew")
    app._fx_history_list_frame = customtkinter.CTkScrollableFrame(history_card, fg_color="transparent")
    app._fx_history_list_frame.grid(row=2, column=0, sticky="nsew", padx=8, pady=(0, 8))
    history_actions = customtkinter.CTkFrame(history_card, fg_color="transparent")
    history_actions.grid(row=3, column=0, sticky="ew", padx=12, pady=(0, 12))
    history_actions.grid_columnconfigure(0, weight=1)
    customtkinter.CTkButton(
        history_actions,
        text="只重试失败",
        command=lambda target=app: _queue_retry_failed_history(target),
        height=34,
        corner_radius=10,
        fg_color="#6D4F3B",
        hover_color="#7C5B43",
        text_color="#FFF6E6",
    ).grid(row=0, column=0, sticky="ew", padx=(0, 8))
    customtkinter.CTkButton(
        history_actions,
        text="重置筛选",
        command=lambda target=app: _reset_queue_history_filters(target),
        height=34,
        width=92,
        corner_radius=10,
        fg_color="transparent",
        hover_color="#303030",
        border_width=1,
        border_color="#566274",
        text_color="#E8EDF5",
    ).grid(row=0, column=1, sticky="e", padx=(0, 8))
    customtkinter.CTkButton(
        history_actions,
        text="刷新",
        command=lambda target=app: (_refresh_queue_panel(target), _refresh_history_panel(target)),
        height=34,
        width=82,
        corner_radius=10,
        fg_color="transparent",
        hover_color="#303030",
        border_width=1,
        border_color="#566274",
        text_color="#E8EDF5",
    ).grid(row=0, column=2, sticky="e")

    def on_close():
        try:
            window.withdraw()
        except Exception:
            pass

    window.protocol("WM_DELETE_WINDOW", on_close)
    _install_queue_history_filter_hooks(app)
    _refresh_queue_panel(app)
    _refresh_history_panel(app)
    return window


def _queue_add_current_task(app, start_after_add=False):
    _ensure_queue_state(app)
    try:
        task = _queue_build_task(app)
    except Exception as exc:
        try:
            app.log(f"❌ [队列] 无法加入队列：{exc}")
            tkinter.messagebox.showwarning("任务队列", str(exc))
        except Exception:
            pass
        return None
    app._fx_task_queue.append(task)
    try:
        app.log(f"📌 [队列] 已加入：{task['title']}")
    except Exception:
        pass
    _refresh_queue_panel(app)
    if start_after_add:
        _start_task_queue(app)
    return task


def _clear_queued_tasks(app):
    _ensure_queue_state(app)
    if getattr(app, "_fx_task_queue_running", False):
        try:
            app.log("ℹ️ [队列] 正在执行时不能清空等待项。")
        except Exception:
            pass
        return
    app._fx_task_queue = [item for item in app._fx_task_queue if item.get("status") == "running"]
    _refresh_queue_panel(app)


def _queue_retry_history_task(app, source):
    return _queue_replay_history_task(app, source, action_name="重试", prefer_failed_subset=True)


def _queue_replay_history_task(app, source, action_name="回放", prefer_failed_subset=False):
    _ensure_queue_state(app)
    try:
        retry_input = source.get("input")
        retry_mode = "full"
        retry_failed_items = []
        retry_cleanup = None
        if prefer_failed_subset:
            subset = _build_retry_subset_input(app, source)
            if isinstance(subset, dict) and subset.get("input"):
                retry_input = subset.get("input")
                retry_mode = str(subset.get("mode") or "subset")
                retry_failed_items = list(subset.get("failed_items") or [])
                retry_cleanup = subset.get("cleanup")
        task = _queue_build_task(app, retry_input, source.get("task_type"), retry_source=source)
        if isinstance(source.get("task_result"), dict):
            task["snapshot"]["task_result"] = dict(source.get("task_result"))
        task["history_source_id"] = source.get("id")
        task["retry_mode"] = retry_mode
        if retry_failed_items:
            task["retry_failed_items"] = list(retry_failed_items)
        if retry_cleanup:
            task["retry_staging_root"] = retry_cleanup
        app._fx_task_queue.append(task)
        if prefer_failed_subset and retry_mode != "full":
            app.log(f"🔁 [队列] 已加入{action_name}：{task['title']}（仅失败项）")
        else:
            app.log(f"🔁 [队列] 已加入{action_name}：{task['title']}")
        _refresh_queue_panel(app)
        return task
    except Exception as exc:
        try:
            app.log(f"❌ [队列] {action_name}失败：{exc}")
        except Exception:
            pass
    _refresh_queue_panel(app)
    return None


def _queue_retry_failed_history(app):
    _ensure_queue_state(app)
    failed = [item for item in getattr(app, "_fx_task_history", []) if item.get("status") == "failed"]
    if not failed:
        try:
            app.log("ℹ️ [队列] 暂无失败历史可重试。")
        except Exception:
            pass
        return
    for item in failed:
        _queue_retry_history_task(app, item)


def _run_task_queue_worker(app):
    _ensure_queue_state(app)
    app._fx_task_queue_running = True
    try:
        while True:
            task = next((item for item in app._fx_task_queue if item.get("status") == "queued"), None)
            if task is None:
                break
            task["started_at"] = time.time()
            task["logs"] = []
            _queue_set_task_status(app, task, "running")
            try:
                app.after(0, lambda target=app, current=task: _queue_restore_app_state(target, current))
            except Exception:
                _queue_restore_app_state(app, task)
            time.sleep(0.05)
            original_log = _queue_capture_task_logs(app, task)
            try:
                _clear_last_task_result(app)
                _start_task_result(app, task.get("input"), task.get("task_type"))
                app.stop_event = False
                app.is_running = True
                try:
                    app.progress_bar.set(0)
                except Exception:
                    pass
                _set_progress_status(app, current_file=task.get("input"), stage="队列任务准备", fraction=0.0)
                app.run_process(task.get("input"), task.get("task_type"))
                task_result = _infer_task_result_from_context(
                    app,
                    task.get("input"),
                    task.get("task_type"),
                    return_value=None,
                    logs=task.get("logs"),
                )
                if isinstance(task_result, dict):
                    task["task_result"] = dict(task_result)
                    _queue_set_task_status(app, task, task_result.get("status", "unknown"), task_result.get("detail", ""))
                elif getattr(app, "stop_event", False):
                    _queue_set_task_status(app, task, "stopped", "用户停止")
                elif _queue_task_had_errors(task):
                    _queue_set_task_status(app, task, "failed", "执行日志中包含失败或错误信息")
                else:
                    _queue_set_task_status(app, task, "success", "执行完成")
            except Exception as exc:
                task["exception"] = str(exc)
                _queue_set_task_status(app, task, "failed", exc)
                try:
                    app.log(f"🔥 [队列] {task.get('title', '')} 执行异常：{exc}")
                except Exception:
                    pass
            finally:
                _queue_restore_task_logs(app, original_log)
                try:
                    app.is_running = False
                except Exception:
                    pass
                _append_queue_history(app, task)
                _clear_last_task_result(app)
                retry_cleanup = task.get("retry_staging_root")
                if retry_cleanup:
                    try:
                        shutil.rmtree(retry_cleanup, ignore_errors=True)
                    except Exception:
                        pass
                try:
                    app._fx_task_queue = [item for item in app._fx_task_queue if item.get("id") != task.get("id")]
                except Exception:
                    pass
                _refresh_queue_panel(app)
                _refresh_history_panel(app)
        try:
            app.log("✅ [队列] 队列执行完成。")
        except Exception:
            pass
    finally:
        app._fx_task_queue_running = False
        try:
            app.reset_ui()
        except Exception:
            pass
        _refresh_queue_panel(app)
        _refresh_history_panel(app)


def _start_task_queue(app):
    _ensure_queue_state(app)
    if getattr(app, "_fx_task_queue_running", False):
        try:
            app.log("ℹ️ [队列] 队列已经在执行中。")
        except Exception:
            pass
        return
    if not any(item.get("status") == "queued" for item in app._fx_task_queue):
        try:
            app.log("ℹ️ [队列] 当前没有等待执行的任务。")
        except Exception:
            pass
        return
    try:
        app.log("🚦 [队列] 开始顺序执行任务。")
    except Exception:
        pass
    threading.Thread(target=_run_task_queue_worker, args=(app,), daemon=True).start()


def _install_progress_status_label(app):
    if getattr(app, "_fx_progress_status_ready", False):
        return
    try:
        app._fx_progress_status_var = tkinter.StringVar(master=app, value=PROGRESS_STATUS_IDLE_TEXT)
    except Exception:
        app._fx_progress_status_var = None
    bottom_bar = getattr(app, "bottom_bar", None)
    if bottom_bar is None:
        _set_progress_status(app)
        app._fx_progress_status_ready = True
        return
    try:
        progress_bar = getattr(app, "progress_bar", None)
        if progress_bar is not None:
            progress_bar.grid_configure(row=0, column=0, columnspan=1, padx=(24, 12), pady=(10, 8), sticky="ew")
        bottom_bar.grid_columnconfigure(0, weight=1)
        bottom_bar.grid_columnconfigure(1, weight=1)
    except Exception:
        pass
    try:
        stale_label = getattr(app, "_fx_progress_status_label", None)
        if stale_label is not None and stale_label.winfo_exists():
            stale_label.destroy()
    except Exception:
        pass
    try:
        label = customtkinter.CTkLabel(
            bottom_bar,
            textvariable=app._fx_progress_status_var,
            text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
            font=customtkinter.CTkFont(family="Microsoft YaHei UI", size=12),
            anchor="w",
            justify="left",
            height=30,
        )
        label.grid(row=0, column=1, padx=(0, 24), pady=(8, 8), sticky="ew")
        app._fx_progress_status_label = label
    except Exception as exc:
        _debug(f"progress_status:label_error:{exc}")
    _set_progress_status(app)
    app._fx_progress_status_ready = True


def _fit_queue_bottom_action_widths(app, action_row=None):
    row = action_row
    if row is None:
        try:
            for child in app.bottom_bar.winfo_children():
                if isinstance(child, customtkinter.CTkFrame) and child.winfo_children():
                    row = child
                    break
        except Exception:
            row = None
    if row is None:
        return

    try:
        row.pack_configure(fill="x", expand=True)
    except Exception:
        pass
    try:
        row.update_idletasks()
        available = int(row.winfo_width())
    except Exception:
        available = 0

    if available >= 1160:
        widths = {
            "btn_run": 220,
            "btn_stop": 132,
            "wm_clear_checkpoint_button": 116,
            "btn_queue_add": 96,
            "btn_queue_panel": 124,
        }
    elif available >= 980:
        widths = {
            "btn_run": 198,
            "btn_stop": 122,
            "wm_clear_checkpoint_button": 106,
            "btn_queue_add": 92,
            "btn_queue_panel": 120,
        }
    else:
        widths = {
            "btn_run": 178,
            "btn_stop": 114,
            "wm_clear_checkpoint_button": 98,
            "btn_queue_add": 88,
            "btn_queue_panel": 116,
        }
    for attr, width in widths.items():
        widget = getattr(app, attr, None)
        if widget is None:
            continue
        try:
            widget.configure(width=width)
        except Exception:
            pass


def _install_queue_bottom_actions(app):
    _ensure_queue_state(app)
    if getattr(app, "_fx_queue_actions_ready", False):
        return
    action_row = None
    try:
        for child in app.bottom_bar.winfo_children():
            if isinstance(child, customtkinter.CTkFrame) and child.winfo_children():
                action_row = child
                break
    except Exception:
        action_row = None
    if action_row is None:
        return
    app._fx_queue_summary_var = tkinter.StringVar(value="")
    try:
        app.btn_queue_add = customtkinter.CTkButton(
            action_row,
            text="加入队列",
            command=lambda target=app: _queue_add_current_task(target),
            height=40,
            width=104,
            corner_radius=10,
            fg_color="transparent",
            hover_color="#303030",
            border_width=1,
            border_color="#B89352",
            text_color="#F6E2B2",
        )
        app.btn_queue_add.pack(side="left", padx=(0, 12))
        app.btn_queue_panel = customtkinter.CTkButton(
            action_row,
            text="队列历史",
            command=lambda target=app: _show_task_queue_window(target),
            height=40,
            width=104,
            corner_radius=10,
            fg_color="transparent",
            hover_color="#303030",
            border_width=1,
            border_color="#566274",
            text_color="#E8EDF5",
        )
        app.btn_queue_panel.pack(side="left", padx=(0, 18))
    except Exception as exc:
        _debug(f"queue:bottom_actions_error:{exc}")
        return
    app._fx_queue_actions_ready = True
    _fit_queue_bottom_action_widths(app, action_row)
    try:
        app.after(250, lambda target=app: _fit_queue_bottom_action_widths(target, action_row))
    except Exception:
        pass
    _refresh_queue_status_summary(app)


def _install_output_strategy_controls(app):
    if getattr(app, "_fx_output_strategy_controls_ready", False):
        return
    top_bar = getattr(app, "top_bar", None)
    if top_bar is None:
        return
    try:
        app.output_strategy_var = tkinter.StringVar(master=app, value=_get_saved_output_strategy())
        app.output_strategy_hint_var = tkinter.StringVar(master=app, value="")
    except Exception as exc:
        _debug(f"output_strategy:vars_error:{exc}")
        return

    try:
        controls = customtkinter.CTkFrame(top_bar, fg_color="transparent")
        controls.grid(row=0, column=2, rowspan=2, sticky="e", padx=(0, 24), pady=(10, 12))
        controls.grid_columnconfigure(0, weight=0)
        controls.grid_columnconfigure(1, weight=1)
        app._fx_output_strategy_controls = controls

        title = customtkinter.CTkLabel(
            controls,
            text="输出策略",
            text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
            font=customtkinter.CTkFont(size=11),
            anchor="w",
        )
        title.grid(row=0, column=0, sticky="w", padx=(0, 10), pady=(0, 4))

        combo = customtkinter.CTkComboBox(
            controls,
            values=[OUTPUT_STRATEGY_VALUE_TO_LABEL[key] for key in OUTPUT_STRATEGY_VALUES],
            variable=app.output_strategy_var,
            width=260,
            height=34,
            command=lambda _value=None, target=app: _refresh_output_strategy_hint(target),
        )
        combo.grid(row=0, column=1, sticky="ew", pady=(0, 4))
        app.output_strategy_combo = combo

        hint = customtkinter.CTkLabel(
            controls,
            textvariable=app.output_strategy_hint_var,
            text_color=globals().get("COLOR_TEXT_SOFT", "#B2C0C8"),
            font=customtkinter.CTkFont(size=11),
            justify="left",
            anchor="w",
            wraplength=340,
        )
        hint.grid(row=1, column=0, columnspan=2, sticky="w")
        app.output_strategy_hint_label = hint
    except Exception as exc:
        _debug(f"output_strategy:controls_error:{exc}")
        return

    _install_output_strategy_memory(app)
    app._fx_output_strategy_controls_ready = True


def _patch_task_queue_history():
    try:
        original_setup_main_area = FengxiToolboxApp.setup_main_area
        original_on_start_click = FengxiToolboxApp.on_start_click
    except Exception as exc:
        _debug(f"queue:patch_missing:{exc}")
        return
    if getattr(original_setup_main_area, "__fx_queue_history_patch__", False):
        return

    def patched_setup_main_area(self, *args, **kwargs):
        result = original_setup_main_area(self, *args, **kwargs)
        try:
            _install_output_strategy_controls(self)
        except Exception as exc:
            _debug(f"output_strategy:install_after_main_area_error:{exc}")
        try:
            _install_progress_status_label(self)
        except Exception as exc:
            _debug(f"progress_status:install_after_main_area_error:{exc}")
        try:
            _install_parallel_mode_hint(self)
        except Exception as exc:
            _debug(f"parallel_hint:install_after_main_area_error:{exc}")
        try:
            _install_queue_bottom_actions(self)
        except Exception as exc:
            _debug(f"queue:install_after_main_area_error:{exc}")
        try:
            _restore_last_settings_for_task(self, DEFAULT_STARTUP_TAB)
        except Exception as exc:
            _debug(f"last_settings:restore_startup_error:{exc}")
        try:
            _install_watermark_color_preview_ui(self)
        except Exception as exc:
            _debug(f"watermark_color:repair_after_main_area_error:{exc}")
        try:
            _install_watermark_last_settings_memory(self)
        except Exception as exc:
            _debug(f"last_settings:watermark_memory_after_main_area_error:{exc}")
        return result

    def patched_on_start_click(self):
        if getattr(self, "_fx_start_via_queue", False):
            return original_on_start_click(self)
        try:
            _save_current_last_settings(self)
        except Exception as exc:
            _debug(f"last_settings:on_start_save_error:{exc}")
        return original_on_start_click(self)

    patched_setup_main_area.__fx_queue_history_patch__ = True
    patched_on_start_click.__fx_queue_history_patch__ = True
    FengxiToolboxApp.setup_main_area = patched_setup_main_area
    FengxiToolboxApp.on_start_click = patched_on_start_click
    _debug("queue:patch_installed")


_patch_task_queue_history()


def _patch_runtime_progress_reporting():
    try:
        original_run_process = FengxiToolboxApp.run_process
    except Exception as exc:
        _debug(f"patch_runtime_progress:missing:{exc}")
        return

    if getattr(original_run_process, "__fx_runtime_progress_patch__", False):
        return

    runtime_run_process = _unwrap_runtime_run_process(original_run_process)
    if runtime_run_process is None:
        _debug("patch_runtime_progress:no_runtime_run_process")
        return
    site_map = _build_runtime_progress_site_map(runtime_run_process)

    def patched_run_process(self, input_folder, task_type):
        _start_task_result(self, input_folder, task_type)
        tracker = None
        result = None
        guard_reason = f"run_process:{task_type}"
        _fx_background_guard_begin(self, guard_reason)
        try:
            if task_type == "convert" and _get_convert_mode(self) in CONVERT_EXTENDED_MODES:
                try:
                    result = _run_convert_extended_task(self, input_folder)
                    return result
                finally:
                    try:
                        self.reset_ui()
                    except Exception:
                        pass
            tracker = _install_run_progress_tracker(
                self,
                input_folder,
                task_type,
                runtime_run_process=runtime_run_process,
                site_map=site_map,
            )
            result = original_run_process(self, input_folder, task_type)
            return result
        except Exception as exc:
            _finalize_current_task_result(
                self,
                "failed",
                message=str(exc),
                detail=str(exc),
                error=str(exc),
            )
            raise
        finally:
            if tracker is not None:
                try:
                    tracker.finalize_pending()
                except Exception:
                    pass
                try:
                    tracker.restore()
                except Exception as exc:
                    _debug(f"patch_runtime_progress:restore_error:{exc}")
            try:
                finalized_result = _infer_task_result_from_context(
                    self,
                    input_folder,
                    task_type,
                    return_value=result,
                    logs=getattr(self, "_fx_last_task_logs", None),
                )
                _schedule_task_completion_notice(self, finalized_result)
            except Exception as exc:
                _debug(f"patch_runtime_progress:finalize_result_error:{exc}")
            _fx_background_guard_end(self, guard_reason)

    patched_run_process.__fx_runtime_progress_patch__ = True
    FengxiToolboxApp.run_process = patched_run_process
    _debug(
        "patch_runtime_progress:installed:"
        f"process={sorted(site_map.get('before_process_single', set()))}:"
        f"direct={sorted(site_map.get('direct_pre', set()))}:"
        f"pass={sorted(site_map.get('pass_through', set()))}"
    )


_patch_runtime_progress_reporting()


def _get_internal_attr(obj, name, default=None):
    try:
        return object.__getattribute__(obj, name)
    except AttributeError:
        return default


def _guess_lazy_tab_for_attr(name):
    for prefix, task_name in LAZY_ATTR_PREFIXES:
        if name.startswith(prefix):
            return task_name
    return None


def _ensure_lazy_tab_initialized(app, task_name):
    spec = LAZY_TAB_SPECS.get(task_name)
    if spec is None:
        return False

    lazy_state = _get_internal_attr(app, "_fx_lazy_tabs_state", None)
    if not lazy_state:
        return False
    if lazy_state.get(task_name):
        return True
    lazy_initializing = _get_internal_attr(app, "_fx_lazy_tabs_initializing", None)
    if lazy_initializing is None:
        lazy_initializing = set()
        try:
            app._fx_lazy_tabs_initializing = lazy_initializing
        except Exception:
            pass
    if task_name in lazy_initializing:
        _debug(f"lazy_tab:init:{task_name}:reentrant_skip")
        return False

    initializers = _get_internal_attr(app, "_fx_lazy_tab_initializers", {})
    initializer = initializers.get(task_name)
    if not callable(initializer):
        initializer = getattr(app, spec["init"], None)
    if not callable(initializer):
        return False

    started_at = time.perf_counter()
    _debug(f"lazy_tab:init:{task_name}:start")
    lazy_initializing.add(task_name)
    try:
        initializer()
        lazy_state[task_name] = True
        try:
            if not _get_internal_attr(app, "_fx_startup_visible_pending", False):
                _tighten_layout(app, task_name=task_name)
        except Exception as exc:
            _debug(f"lazy_tab:layout_refresh_error:{task_name}:{exc}")
        try:
            _restore_last_settings_for_task(app, task_name)
        except Exception as exc:
            _debug(f"last_settings:lazy_restore_error:{task_name}:{exc}")
        try:
            app.update_idletasks()
        except Exception:
            pass
    except Exception:
        _record_performance(
            "lazy_tab_init",
            started_at=started_at,
            task_name=task_name,
            details={"status": "error"},
        )
        raise
    finally:
        try:
            lazy_initializing.discard(task_name)
        except Exception:
            pass
    _record_performance(
        "lazy_tab_init",
        started_at=started_at,
        task_name=task_name,
        details={"status": "success"},
    )
    _debug(f"lazy_tab:init:{task_name}:done")
    return True


def _show_ready_window(app):
    started_at = time.perf_counter()
    _install_fast_close_protocol(app)
    try:
        app.update_idletasks()
    except Exception as exc:
        _debug(f"startup:update_idletasks_error:{exc}")
    try:
        app.deiconify()
        _ensure_startup_window_visible(app)
        app.lift()
        try:
            app.focus_force()
        except Exception:
            pass
        try:
            app._fx_startup_visible_pending = False
        except Exception:
            pass
        _debug("startup:window_shown")
    except Exception as exc:
        _debug(f"startup:window_show_error:{exc}")
    try:
        app.after(100, lambda target=app: _ensure_startup_window_visible(target))
    except Exception:
        pass
    try:
        app.after(STARTUP_DEFAULT_TAB_INIT_DELAY_MS, lambda target=app: _ensure_lazy_tab_initialized(target, DEFAULT_STARTUP_TAB))
    except Exception:
        _ensure_lazy_tab_initialized(app, DEFAULT_STARTUP_TAB)
    try:
        app.after(160, lambda target=app: _run_startup_layout_refresh(target))
    except Exception:
        _run_startup_layout_refresh(app)
    _record_performance("startup_show_ready", started_at=started_at)
    _record_performance(
        "startup_total",
        started_at=BOOTSTRAP_STARTED_AT,
        details={"default_tab": DEFAULT_STARTUP_TAB},
    )


def _get_startup_visible_geometry(window_x, window_y, window_width, window_height, viewport_x, viewport_y, viewport_width, viewport_height, margin=36):
    """Return a centered geometry only when the startup window is outside its viewport."""
    try:
        window_x = int(window_x)
        window_y = int(window_y)
        window_width = max(1, int(window_width))
        window_height = max(1, int(window_height))
        viewport_x = int(viewport_x)
        viewport_y = int(viewport_y)
        viewport_width = max(1, int(viewport_width))
        viewport_height = max(1, int(viewport_height))
        margin = max(0, int(margin))
    except Exception:
        return None

    viewport_right = viewport_x + viewport_width
    viewport_bottom = viewport_y + viewport_height
    window_right = window_x + window_width
    window_bottom = window_y + window_height
    if (
        window_x >= viewport_x + margin
        and window_y >= viewport_y + margin
        and window_right <= viewport_right - margin
        and window_bottom <= viewport_bottom - margin
    ):
        return None

    usable_width = max(1, viewport_width - margin * 2)
    usable_height = max(1, viewport_height - margin * 2)
    target_width = min(window_width, usable_width)
    target_height = min(window_height, usable_height)
    target_x = viewport_x + max(margin, (viewport_width - target_width) // 2)
    target_y = viewport_y + max(margin, (viewport_height - target_height) // 2)
    return f"{target_width}x{target_height}+{target_x}+{target_y}"


def _ensure_startup_window_visible(app):
    """Recover a window that was restored beyond the current desktop bounds."""
    try:
        app.update_idletasks()
        viewport_x = app.winfo_vrootx()
        viewport_y = app.winfo_vrooty()
        viewport_width = app.winfo_vrootwidth()
        viewport_height = app.winfo_vrootheight()
        if int(viewport_width or 0) <= 1 or int(viewport_height or 0) <= 1:
            viewport_x = 0
            viewport_y = 0
            viewport_width = app.winfo_screenwidth()
            viewport_height = app.winfo_screenheight()
        geometry = _get_startup_visible_geometry(
            app.winfo_rootx(),
            app.winfo_rooty(),
            app.winfo_width(),
            app.winfo_height(),
            viewport_x,
            viewport_y,
            viewport_width,
            viewport_height,
        )
        if geometry:
            app.geometry(geometry)
            app.update_idletasks()
            _debug(f"startup:window_repositioned:{geometry}")
        return geometry
    except Exception as exc:
        _debug(f"startup:window_visibility_check_error:{exc}")
        return None


def _run_startup_layout_refresh(app):
    if getattr(app, "_fx_startup_layout_refreshed", False):
        return
    try:
        app._fx_startup_layout_refreshed = True
    except Exception:
        pass
    started_at = time.perf_counter()
    try:
        task_name = getattr(app, "current_task", DEFAULT_STARTUP_TAB)
        if task_name not in TAB_LAYOUT_ATTRS:
            task_name = DEFAULT_STARTUP_TAB
        _schedule_startup_layout_stage(app, task_name, "shell", STARTUP_LAYOUT_FIRST_DELAY_MS)
        _debug("startup:layout_refresh_scheduled")
        _record_performance(
            "startup_layout_refresh",
            started_at=started_at,
            task_name=task_name,
            details={"status": "scheduled"},
        )
    except Exception as exc:
        _debug(f"startup:layout_refresh_error:{exc}")
        _record_performance(
            "startup_layout_refresh",
            started_at=started_at,
            details={"status": "error"},
        )


def _schedule_startup_layout_stage(app, task_name, stage, delay_ms=STARTUP_LAYOUT_STAGE_DELAY_MS):
    try:
        app.after(int(delay_ms), lambda target=app, current=task_name, current_stage=stage: _run_startup_layout_stage(target, current, current_stage))
    except Exception:
        _run_startup_layout_stage(app, task_name, stage)


def _run_startup_layout_stage(app, task_name, stage):
    stage = str(stage or "")
    event_name = f"startup_layout_{stage}" if stage else "startup_layout_stage"
    started_at = time.perf_counter()
    status = "success"
    try:
        if stage == "shell":
            _apply_shell_layout_tightening(app)
            next_stage = "tab"
        elif stage == "tab":
            _tighten_single_tab_layout(app, task_name)
            event_name = "startup_layout_tighten_visible"
            next_stage = "refresh"
        elif stage == "refresh":
            if task_name in TAB_LAYOUT_ATTRS:
                _refresh_visible_tab_layout(app, task_name)
            event_name = "startup_layout_refresh_visible"
            next_stage = ""
            _debug("startup:layout_refreshed")
        else:
            return
    except Exception as exc:
        status = "error"
        _debug(f"startup:layout_stage_error:{stage}:{exc}")
        next_stage = ""
    finally:
        _record_performance(
            event_name,
            started_at=started_at,
            task_name=task_name,
            details={"status": status},
        )
    if next_stage:
        _schedule_startup_layout_stage(app, task_name, next_stage, STARTUP_LAYOUT_STAGE_DELAY_MS)


def _request_fast_close(app):
    if getattr(app, "_fx_fast_close_started", False):
        return
    app._fx_fast_close_started = True
    try:
        _save_initialized_last_settings(app)
    except Exception as exc:
        _debug(f"fast_close:last_settings_save_error:{exc}")
    try:
        _flush_watermark_text_persistence(app)
    except Exception as exc:
        _debug(f"fast_close:wm_text_flush_error:{exc}")
    try:
        _flush_watermark_filename_rule_persistence(app)
    except Exception as exc:
        _debug(f"fast_close:wm_filename_rule_flush_error:{exc}")
    try:
        app.stop_event = True
    except Exception:
        pass
    try:
        app.is_running = False
    except Exception:
        pass
    try:
        app.withdraw()
        app.update_idletasks()
        _debug("fast_close:window_hidden")
    except Exception as exc:
        _debug(f"fast_close:withdraw_error:{exc}")

    def finish_destroy():
        try:
            app.quit()
        except Exception:
            pass
        try:
            app.destroy()
            _debug("fast_close:destroy_done")
        except Exception as exc:
            _debug(f"fast_close:destroy_error:{exc}")
            try:
                os._exit(0)
            except Exception:
                pass

    def force_exit_if_needed():
        try:
            if getattr(app, "_fx_fast_close_force_done", False):
                return
            app._fx_fast_close_force_done = True
            _debug("fast_close:force_exit")
            os._exit(0)
        except Exception:
            pass

    if not getattr(app, "_fx_disable_fast_close_force_exit", False):
        try:
            timer = threading.Timer(0.9, force_exit_if_needed)
            timer.daemon = True
            timer.start()
        except Exception as exc:
            _debug(f"fast_close:force_timer_error:{exc}")

    try:
        app.after_idle(finish_destroy)
    except Exception:
        finish_destroy()


def _install_fast_close_protocol(app):
    if getattr(app, "_fx_fast_close_protocol_ready", False):
        return
    try:
        app.protocol("WM_DELETE_WINDOW", lambda target=app: _request_fast_close(target))
        app._fx_fast_close_protocol_ready = True
        _debug("fast_close:protocol_installed")
    except Exception as exc:
        _debug(f"fast_close:protocol_error:{exc}")


_FX_SINGLE_INSTANCE_MUTEX_HANDLE = None


def _acquire_single_instance_lock():
    global _FX_SINGLE_INSTANCE_MUTEX_HANDLE
    if os.name != "nt":
        return True
    try:
        import ctypes

        kernel32 = ctypes.windll.kernel32
        kernel32.CreateMutexW.argtypes = (ctypes.c_void_p, ctypes.c_bool, ctypes.c_wchar_p)
        kernel32.CreateMutexW.restype = ctypes.c_void_p
        kernel32.GetLastError.argtypes = ()
        kernel32.GetLastError.restype = ctypes.c_ulong
        handle = kernel32.CreateMutexW(None, True, "Local\\FengxiToolboxSingleInstance")
        last_error = kernel32.GetLastError()
        if not handle:
            _debug(f"single_instance:create_mutex_failed:{last_error}")
            return True
        if last_error == 183:
            _debug("single_instance:already_running")
            return False
        _FX_SINGLE_INSTANCE_MUTEX_HANDLE = handle
        return True
    except Exception as exc:
        _debug(f"single_instance:error:{exc}")
        return True


def _patch_startup_performance():
    return install_startup_performance_patch(
        StartupPatchContext(
            app_class=FengxiToolboxApp,
            ctk_class=customtkinter.CTk,
            lazy_tab_specs=LAZY_TAB_SPECS,
            default_startup_tab=DEFAULT_STARTUP_TAB,
            debug=_debug,
            get_internal_attr=_get_internal_attr,
            ensure_lazy_tab_initialized=_ensure_lazy_tab_initialized,
            show_inline_help=_show_inline_help,
            show_inline_donate=_show_inline_donate,
            set_help_button_selected=_set_help_button_selected,
            set_donate_button_selected=_set_donate_button_selected,
            set_help_action_state=_set_help_action_state,
            refresh_output_strategy_hint=_refresh_output_strategy_hint,
            refresh_parallel_mode_hint=_refresh_parallel_mode_hint,
            refresh_visible_tab_layout=_refresh_visible_tab_layout,
            guess_lazy_tab_for_attr=_guess_lazy_tab_for_attr,
            record_performance=_record_performance,
        )
    )


_patch_startup_performance()


def _patch_sidebar_build_performance():
    try:
        original_setup_sidebar = FengxiToolboxApp.setup_sidebar
    except Exception as exc:
        _debug(f"patch_sidebar_build_performance:missing:{exc}")
        return

    if getattr(original_setup_sidebar, "__fx_fast_sidebar_build_patch__", False):
        return

    def patched_setup_sidebar(self, *args, **kwargs):
        return _run_with_fast_sidebar_button_construction(
            self,
            lambda: original_setup_sidebar(self, *args, **kwargs),
        )

    patched_setup_sidebar.__fx_fast_sidebar_build_patch__ = True
    FengxiToolboxApp.setup_sidebar = patched_setup_sidebar
    _debug("patch_sidebar_build_performance:installed")


_patch_sidebar_build_performance()

_patch_log_display_behavior()

try:
    _wrap_callable(customtkinter.CTk, "__init__", "ctk_init")
    _wrap_callable(pywinstyles, "apply_style", "pywinstyles.apply_style")
    _wrap_callable(pywinstyles, "change_header_color", "pywinstyles.change_header_color")
    _wrap_callable(windnd, "hook_dropfiles", "windnd.hook_dropfiles")
    _wrap_callable(Image, "open", "Image.open")
    _wrap_callable(FengxiToolboxApp, "scan_fonts", "app.scan_fonts")
    _wrap_callable(FengxiToolboxApp, "setup_sidebar", "app.setup_sidebar")
    _wrap_callable(FengxiToolboxApp, "setup_main_area", "app.setup_main_area")
    _wrap_callable(FengxiToolboxApp, "log", "app.log")
except Exception as _wrap_exc:
    _debug(f"bootstrap:wrap_error:{_wrap_exc}")


if __name__ == "__main__":
    diag_path = (os.environ.get("FX_OCR_DIAG_PATH") or "").strip()
    if diag_path:
        from tools.fx_pdf_ocr import run_packaged_ocr_diagnostics

        run_packaged_ocr_diagnostics(diag_path)
        raise SystemExit(0)
    if not _acquire_single_instance_lock():
        raise SystemExit(0)
    _debug("main:create_app")
    _record_performance("main_enter", started_at=BOOTSTRAP_STARTED_AT)
    main_step_started_at = time.perf_counter()
    app = FengxiToolboxApp()
    _record_performance("main_create_app", started_at=main_step_started_at)
    _debug("main:app_created")
    main_step_started_at = time.perf_counter()
    _apply_app_icon_bitmap(app)
    _record_performance("main_icon_bitmap_apply", started_at=main_step_started_at)
    _debug("main:icon_bitmap_applied")
    main_step_started_at = time.perf_counter()
    _apply_release_identity(app)
    _record_performance("main_release_identity", started_at=main_step_started_at)
    _debug("main:release_identity_applied")
    _debug("main:layout_tighten_deferred")
    _show_ready_window(app)
    try:
        app.after(180, lambda target=app: _apply_app_icon_after_first_paint(target))
    except Exception:
        _apply_app_icon_after_first_paint(app)
    app.mainloop()
